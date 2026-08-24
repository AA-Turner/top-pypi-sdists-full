import base64
import hashlib
from io import BytesIO
from typing import Any, List, NamedTuple, Optional, Tuple
from urllib.parse import urlparse
import os
from pydantic import BaseModel
from loguru import logger
import asyncio

from xpander_sdk.exceptions.module_exception import ModuleException
from xpander_sdk.media.caps import BEDROCK_VIDEO_FORMATS
from xpander_sdk.utils.safe_fetch import asafe_fetch, safe_fetch

# Provider per-image ceiling (Bedrock/Anthropic ~5MB); env-overridable.
_MAX_IMAGE_BYTES = int(
    os.getenv("XPANDER_MAX_INLINE_IMAGE_BYTES", str(5 * 1024 * 1024))
)
_MAX_DOC_BYTES = int(os.getenv("XPANDER_MAX_INLINE_DOC_BYTES", str(15 * 1024 * 1024)))
_FETCH_TIMEOUT = float(os.getenv("XPANDER_FILE_FETCH_TIMEOUT", "15"))

# Char ceilings for inlining human-readable file text into the prompt. A large
# CSV/JSON pasted in full blows the context window, so cap per-file and total
# and keep the URL (already listed) for the agent to fetch the rest if needed.
_MAX_INLINE_TEXT_CHARS = int(os.getenv("XPANDER_MAX_INLINE_TEXT_CHARS", "8000"))
_MAX_INLINE_TOTAL_CHARS = int(os.getenv("XPANDER_MAX_INLINE_TOTAL_CHARS", "24000"))

# Extracted text is clipped again per-file downstream; this only stops a converted
# document (a 6MB spreadsheet renders ~27M chars of Markdown) from being carried around whole.
_MAX_DOC_MARKDOWN_CHARS = max(
    1000, int(os.getenv("XPANDER_MAX_DOC_MARKDOWN_CHARS", "200000"))
)

# Concurrency for per-attachment download+convert; each one blocks on I/O then releases the GIL.
# Floored at 1 so setting it to 0 to "disable concurrency" serializes instead of raising.
_ATTACHMENT_WORKERS = max(1, int(os.getenv("XPANDER_ATTACHMENT_WORKERS", "4")))

# Concurrent downloads coexist in memory, so the worker count is also bounded by how many
# whole attachments we are willing to hold at once.
_ATTACHMENT_BYTES_BUDGET = int(
    os.getenv("XPANDER_ATTACHMENT_BYTES_BUDGET", str(64 * 1024 * 1024))
)

DOC_UNREADABLE_NOTE = (
    "No text could be extracted from this attachment, so its contents are NOT included here; "
    "fetch it via your tools/workspace if you need them"
)


def _charset_of(content_type: Optional[str]) -> str:
    """Charset from a content-type header, defaulting to utf-8 (was response.text's job)."""
    if content_type and "charset=" in content_type.lower():
        cand = (
            content_type.lower()
            .split("charset=", 1)[1]
            .split(";", 1)[0]
            .strip()
            .strip('"')
        )
        if cand:
            return cand
    return "utf-8"


def attachment_workers(item_count: int, per_item_bytes: int) -> int:
    """Worker count for a fan-out of *item_count* attachments, bounded by the in-flight byte budget."""
    by_memory = max(1, _ATTACHMENT_BYTES_BUDGET // max(1, per_item_bytes))
    return max(1, min(_ATTACHMENT_WORKERS, item_count, by_memory))


def truncate_inline_text(content: str, url: str, remaining: int) -> str:
    """Clip inlined file *content* to *remaining* (and the per-file cap), noting the URL for the rest."""
    cap = min(_MAX_INLINE_TEXT_CHARS, max(0, remaining))
    if len(content) <= cap:
        return content
    return (
        content[:cap]
        + f"\n... [truncated: {len(content):,} chars total - full file at {url}]"
    )


# Leading-byte signatures for the only raster formats vision providers accept.
_IMAGE_SIGNATURES = (
    (b"\x89PNG\r\n\x1a\n", "png", "image/png"),
    (b"\xff\xd8\xff", "jpeg", "image/jpeg"),
    (b"GIF87a", "gif", "image/gif"),
    (b"GIF89a", "gif", "image/gif"),
)

# iWork formats are proprietary, but the files are zips embedding a QuickLook preview PDF.
_IWORK_EXTS = {".pages", ".key", ".numbers"}

# What anydoc's format_from_extension() resolves, minus csv (raw injection) and pdf (native path).
_DOCUMENT_EXTS = {
    ".docx",
    ".doc",
    ".docm",
    ".xlsx",
    ".xls",
    ".xlsm",
    ".xlsb",
    ".pptx",
    ".ppt",
    ".pptm",
    ".pot",
    ".pps",
    ".ppsx",
    ".odt",
    ".ods",
    ".odp",
    ".rtf",
    ".epub",
} | _IWORK_EXTS


class FileCategorization(BaseModel):
    images: List[str]
    pdfs: List[str]
    files: List[str]
    documents: List[str] = []
    audios: List[str] = []
    videos: List[str] = []


# Route office/document mimes to the "documents" bucket so they take the anydoc text path, not the vision or raw-inject paths.
_DOCUMENT_MIMES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
    "application/vnd.oasis.opendocument.text",
    "application/vnd.oasis.opendocument.spreadsheet",
    "application/vnd.oasis.opendocument.presentation",
    "application/rtf",
    "application/epub+zip",
    "application/vnd.apple.pages",
    "application/vnd.apple.keynote",
    "application/vnd.apple.numbers",
}


# Sentinel: the ref is explicitly a kind we do not bucket (audio/video/archive), so it must NOT
# fall back to URL extension (a video named clip.png must not be classified as an image).
_UNCATEGORIZED = "__uncategorized__"


def _bucket_from_ref(ref: Any) -> Optional[str]:
    """Map an AttachmentRef to a bucket, `_UNCATEGORIZED` to force url_only, or None to fall back to extension."""
    kind = getattr(ref, "kind", None)
    kind_val = getattr(kind, "value", kind)
    if kind_val in ("image", "pdf", "document", "text", "audio", "video"):
        return {
            "image": "images",
            "pdf": "pdfs",
            "document": "documents",
            "text": "files",
            "audio": "audios",
            "video": "videos",
        }[kind_val]
    if kind_val == "archive":
        return _UNCATEGORIZED  # a known non-bucketed kind: url_only, never extension-guessed
    # kind is unknown/None -> consult mime, then (via None) the URL extension
    mime = (getattr(ref, "mime", None) or "").lower().split(";", 1)[0].strip()
    if not mime:
        return None
    # SVG is XML: read as text, not through the image pipeline.
    if mime.startswith("image/") and mime != "image/svg+xml":
        return "images"
    if mime.startswith("audio/"):
        return "audios"
    if mime.startswith("video/"):
        return "videos"
    if mime == "application/pdf":
        return "pdfs"
    if mime in _DOCUMENT_MIMES:
        return "documents"
    if mime.startswith("text/") or mime in (
        "application/json",
        "application/xml",
        "image/svg+xml",
    ):
        return "files"
    return None


def categorize_files(
    file_urls: list[str], refs: Optional[dict] = None
) -> FileCategorization:
    """
    Categorize file URLs into images, PDFs, documents, and human-readable files.

    Classification prefers an AttachmentRef's kind/mime (from *refs*, keyed by URL) so an
    extension-less or mis-named URL still lands correctly; falls back to the URL extension.
    Does not load the files.

    Args:
        file_urls (list[str]): List of file URLs
        refs (Optional[dict]): Optional {url: AttachmentRef} for mime/kind-based classification.

    Returns:
        FileCategorization: Pydantic model with categorized URLs
    """
    refs = refs or {}
    image_exts = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp"}
    pdf_exts = {".pdf"}
    human_readable_exts = {
        # SVG is XML: models read it as text, image pipelines reject it.
        ".svg",
        ".txt",
        ".csv",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".md",
        ".rst",
        ".yaml",
        ".yml",
        ".py",
        ".js",
        ".ts",
        ".java",
        ".c",
        ".cpp",
        ".h",
        ".cs",
        ".go",
        ".rb",
        ".php",
        ".sh",
        # plain-text formats that need no conversion, just injection
        ".tsv",
        ".log",
        ".ini",
        ".toml",
        ".tex",
        ".ipynb",
        ".eml",
    }

    audio_exts = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac", ".opus"}
    # Mirrors BEDROCK_VIDEO_FORMATS so an extension-only URL classifies the same way a
    # ref-carried mime would; .avi has no native tier but still transcribes.
    video_exts = {
        ".mp4",
        ".mov",
        ".mkv",
        ".webm",
        ".avi",
        ".mpeg",
        ".mpg",
        ".flv",
        ".wmv",
        ".3gp",
    }

    result = {
        "images": [],
        "pdfs": [],
        "files": [],
        "documents": [],
        "audios": [],
        "videos": [],
    }

    for url in file_urls:
        bucket = _bucket_from_ref(refs.get(url)) if refs else None
        if bucket == _UNCATEGORIZED:
            continue  # a known non-bucketed kind: leave it for the plan's url_only path
        if bucket is None:
            path = urlparse(url).path
            _, ext = os.path.splitext(path.lower())
            if ext in image_exts:
                bucket = "images"
            elif ext in pdf_exts:
                bucket = "pdfs"
            elif ext in _DOCUMENT_EXTS:
                bucket = "documents"
            elif ext in human_readable_exts:
                bucket = "files"
            elif ext in audio_exts:
                bucket = "audios"
            elif ext in video_exts:
                bucket = "videos"
        if bucket is not None:
            result[bucket].append(url)

    return FileCategorization(**result)


async def fetch_urls(
    urls: list[str], disable_attachment_injection: Optional[bool] = False
) -> list[dict[str, str]]:
    """
    Fetches the content of multiple URLs asynchronously.

    Args:
        urls (list[str]): List of URLs to fetch.

    Returns:
        list[dict[str, str]]: A list of dictionaries containing the URL and its content.
                              Example: [{"url": "...", "content": "..."}]
        disable_attachment_injection (Optional[bool]): Optional selection if to disable attachment injection to the context window.
    """

    async def fetch(url: str) -> dict[str, str]:
        try:
            if disable_attachment_injection:
                return {"url": url}
            result = await asafe_fetch(url, max_bytes=_MAX_DOC_BYTES, timeout=10.0)
            return {
                "url": url,
                "content": result.content.decode(
                    _charset_of(result.content_type), errors="replace"
                ),
            }
        except Exception:
            # The error text can name an internal host; the content is inlined into
            # the prompt, so return a neutral note instead of the raw exception.
            return {"url": url, "content": "Error: file could not be read"}

    tasks = [fetch(url) for url in urls]
    return await asyncio.gather(*tasks)


class AttachmentDecision(BaseModel):
    url: str
    category: str  # image | pdf | document | readable | other
    action: str  # inline | text_extract | url_only | skip
    reason: Optional[str] = None
    size: Optional[int] = None


class AttachmentPlan(BaseModel):
    items: List[AttachmentDecision] = []
    notes: List[str] = []

    def by_category(self, category: str) -> List[AttachmentDecision]:
        return [item for item in self.items if item.category == category]


async def estimate_sizes(urls: List[str], timeout: float = 2.0) -> dict:
    """Best-effort parallel size probe -> {url: bytes or None}; HEAD first, ranged-GET fallback for HEAD-hostile hosts."""
    if not urls:
        return {}

    async def _probe(url: str) -> Optional[int]:
        try:
            head = await asafe_fetch(url, max_bytes=0, timeout=timeout, method="HEAD")
            raw = head.headers.get("content-length")
            if head.status < 400 and raw is not None:
                return int(raw)
        except Exception:
            pass
        try:
            ranged = await asafe_fetch(
                url, max_bytes=0, timeout=timeout, extra_headers={"Range": "bytes=0-0"}
            )
            content_range = ranged.headers.get("content-range") or ""
            if "/" in content_range:
                total = content_range.rsplit("/", 1)[1].strip()
                if total.isdigit():
                    return int(total)
            # Host ignored the range and served the whole body.
            raw = ranged.headers.get("content-length")
            if raw is not None and int(raw) > 1:
                return int(raw)
        except Exception:
            pass
        return None

    results = await asyncio.gather(*[_probe(u) for u in urls], return_exceptions=True)
    return {u: (r if isinstance(r, int) else None) for u, r in zip(urls, results)}


def plan_attachments(
    urls: List[str],
    caps,
    sizes: Optional[dict] = None,
    refs: Optional[dict] = None,
) -> AttachmentPlan:
    """Decide per-URL attachment handling from (categories, capabilities, known sizes) - pure, no downloads.

    *refs* ({url: AttachmentRef}) lets classification use mime/kind and supplies known sizes so the
    size tier fires without a network probe.
    """
    plan = AttachmentPlan()
    if not urls:
        return plan
    refs = refs or {}
    sizes = dict(sizes or {})
    # A ref's size fills the size tier without a probe; an explicit `sizes` arg still wins.
    for u, ref in refs.items():
        rsize = getattr(ref, "size", None)
        if rsize is not None and u not in sizes:
            sizes[u] = rsize
    cat = categorize_files(file_urls=urls, refs=refs)
    categories = (
        [("image", u) for u in cat.images]
        + [("pdf", u) for u in cat.pdfs]
        + [("document", u) for u in cat.documents]
        + [("readable", u) for u in cat.files]
        + [("audio", u) for u in cat.audios]
        + [("video", u) for u in cat.videos]
    )
    known = {u for _, u in categories}
    categories += [("other", u) for u in urls if u not in known]

    skipped_images = 0
    inlined_images = 0
    overflow_images = 0
    unreadable: List[str] = []
    for category, url in categories:
        size = sizes.get(url)
        if category == "image":
            if not caps.supports_vision:
                plan.items.append(
                    AttachmentDecision(
                        url=url,
                        category=category,
                        action="skip",
                        reason="model has no vision",
                        size=size,
                    )
                )
                skipped_images += 1
            elif size is not None and size > caps.max_fetch_bytes:
                plan.items.append(
                    AttachmentDecision(
                        url=url,
                        category=category,
                        action="url_only",
                        reason="huge",
                        size=size,
                    )
                )
                plan.notes.append(
                    f"{url}: too large to attach inline; fetch it via your tools/workspace using the URL"
                )
            elif inlined_images >= caps.max_images:
                plan.items.append(
                    AttachmentDecision(
                        url=url,
                        category=category,
                        action="url_only",
                        reason="max_images",
                        size=size,
                    )
                )
                overflow_images += 1
            else:
                plan.items.append(
                    AttachmentDecision(
                        url=url, category=category, action="inline", size=size
                    )
                )
                inlined_images += 1
        elif category == "pdf":
            if size is not None and size > caps.max_fetch_bytes:
                plan.items.append(
                    AttachmentDecision(
                        url=url,
                        category=category,
                        action="url_only",
                        reason="huge",
                        size=size,
                    )
                )
                plan.notes.append(
                    f"{url}: too large to attach inline; fetch it via your tools/workspace using the URL"
                )
            elif caps.supports_native_pdf:
                plan.items.append(
                    AttachmentDecision(
                        url=url, category=category, action="inline", size=size
                    )
                )
            else:
                plan.items.append(
                    AttachmentDecision(
                        url=url,
                        category=category,
                        action="text_extract",
                        reason="provider rejects file blobs",
                        size=size,
                    )
                )
        elif category in ("document", "readable"):
            plan.items.append(
                AttachmentDecision(
                    url=url, category=category, action="text_extract", size=size
                )
            )
        elif category in ("audio", "video"):
            ref = refs.get(url)
            existing = (getattr(ref, "transcript", None) or "").strip()
            native_ok = (
                caps.supports_audio if category == "audio" else caps.supports_video
            )
            ceiling = (
                caps.max_audio_bytes if category == "audio" else caps.max_video_bytes
            )
            allowlist = getattr(caps, "native_video_formats", None)
            if category == "video" and native_ok and allowlist:
                # Bedrock RAISES on a container outside its allowlist, locally, before the
                # request - so an unlisted container must never reach the native tier.
                # Providers without an allowlist 400 instead, which degrade-retry catches.
                native_ok = _video_container(url, ref) in allowlist
            if existing:
                # Already converted upstream (a channel that transcribes at ingest, or an
                # earlier turn): inline the text instead of paying for it twice.
                plan.items.append(
                    AttachmentDecision(
                        url=url,
                        category=category,
                        action="transcript_cached",
                        size=size,
                    )
                )
            elif native_ok and (size is None or size <= ceiling):
                plan.items.append(
                    AttachmentDecision(
                        url=url,
                        category=category,
                        action=f"native_{category}",
                        size=size,
                    )
                )
            else:
                reason = "model cannot take this natively" if not native_ok else "huge"
                plan.items.append(
                    AttachmentDecision(
                        url=url,
                        category=category,
                        action="transcribe",
                        reason=reason,
                        size=size,
                    )
                )
        else:
            plan.items.append(
                AttachmentDecision(
                    url=url, category=category, action="url_only", size=size
                )
            )
            # No converter for this format (.msg, .zip, extension-less URLs, ...): without
            # this note the model answers about content it never saw.
            unreadable.append(url)

    if unreadable:
        plan.notes.append(
            "these attachments cannot be read directly and their content is NOT included; "
            "do not guess their contents - fetch them via your tools/workspace if you need "
            "them: " + ", ".join(unreadable)
        )

    if skipped_images:
        plural = "s" if skipped_images != 1 else ""
        plan.notes.append(
            f"{skipped_images} image{plural} listed above cannot be viewed by this model; "
            "use tools to process them if needed"
        )
    if overflow_images:
        plural = "s" if overflow_images != 1 else ""
        plan.notes.append(
            f"{overflow_images} additional image{plural} beyond the per-message limit "
            f"({caps.max_images}) listed above but not attached; fetch via your tools/workspace using the URLs"
        )
    return plan


def _download(url: str, *, max_bytes: int, timeout: float) -> Tuple[bytes, str]:
    """SSRF-guarded stream to bytes, capped at `max_bytes`; returns (content, lowercased content-type)."""
    result = safe_fetch(url, max_bytes=max_bytes, timeout=timeout)
    ctype = (result.content_type or "").split(";", 1)[0].strip().lower()
    return result.content, ctype


class FetchedBytes(NamedTuple):
    data: bytes
    content_type: str
    sha256: str


def host_allowed(url: str, allowed_hosts: set) -> bool:
    """https-only host-suffix allowlist (the bots hardening rule, engine-level)."""
    parsed = urlparse(url)
    if parsed.scheme != "https" or not parsed.hostname:
        return False
    host = parsed.hostname.lower()
    return any(host == s or host.endswith(f".{s}") for s in allowed_hosts)


def fetch_bytes(
    url: str,
    *,
    max_bytes: int,
    timeout: float,
    allowed_hosts: Optional[set] = None,
) -> FetchedBytes:
    """Content-addressed SSRF-guarded download.

    sha256 is computed here because this is the only place the bytes exist -
    no producer populates ref.sha256, and the conversion cache keys on it.
    """
    if allowed_hosts is not None and not host_allowed(url, allowed_hosts):
        raise ModuleException(400, "media host not in allowlist")
    data, ctype = _download(url, max_bytes=max_bytes, timeout=timeout)
    return FetchedBytes(data, ctype, hashlib.sha256(data).hexdigest())


def _sniff_image(data: bytes) -> Optional[Tuple[str, str]]:
    """Return (format, mime) if `data` begins with a supported image signature, else None."""
    for sig, fmt, mime in _IMAGE_SIGNATURES:
        if data.startswith(sig):
            return fmt, mime
    if len(data) >= 12 and data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return "webp", "image/webp"
    return None


def fetch_image(
    url: str, *, max_bytes: int = _MAX_IMAGE_BYTES, timeout: float = _FETCH_TIMEOUT
) -> Any:
    """Download a URL and wrap it as an inline Agno Image, only if the bytes are a genuine supported image.

    Providers reject remote image URLs and require inline bytes, and a URL's extension is not trustworthy
    (expired/redirected presigned links return HTML/XML with HTTP 200). So the content-type and leading
    magic bytes are validated and format/mime are derived from the bytes, not the URL. Raises on anything
    not a verifiable png/jpeg/gif/webp so callers can drop it instead of feeding garbage to the model.
    """
    content, ctype = _download(url, max_bytes=max_bytes, timeout=timeout)
    if ctype and not (
        ctype.startswith("image/")
        or ctype in ("application/octet-stream", "binary/octet-stream")
    ):
        raise ValueError(f"non-image content-type {ctype!r}: {url}")
    sniffed = _sniff_image(content)
    if sniffed is None:
        raise ValueError(f"bytes are not a supported image: {url}")
    fmt, mime = sniffed
    content_b64 = base64.b64encode(content).decode("utf-8")
    from agno.media import Image

    return Image.from_base64(base64_content=content_b64, format=fmt, mime_type=mime)


_CONTAINER_ALIASES = {
    "3gpp": "three_gp",
    "3gp": "three_gp",
    "quicktime": "mov",
    "x-matroska": "mkv",
    "x-msvideo": "avi",
    "x-flv": "flv",
    "x-ms-wmv": "wmv",
    "x-wav": "wav",
    "vnd.wave": "wav",
    "mp4": "mp4",
    "x-m4a": "m4a",
}


def _media_format(url: str, mime: Optional[str], fallback: str) -> str:
    """Container name agno needs; a ref mime beats the URL suffix, suffix beats nothing."""
    mime = (mime or "").lower().split(";", 1)[0].strip()
    if "/" in mime:
        container = mime.split("/", 1)[1]
    else:
        container = os.path.splitext(url.split("?", 1)[0].lower())[1].lstrip(".")
    # Provider allowlists name containers, mimes name codecs-ish: quicktime IS mov, and
    # Bedrock rejects anything it cannot match by name.
    container = _CONTAINER_ALIASES.get(container, container)
    # audio/mpeg IS mp3; video/mpeg is the mpeg container - one token, two meanings.
    if container == "mpeg" and (mime.startswith("audio/") or fallback == "mp3"):
        container = "mp3"
    return container or fallback


def _video_container(url: str, ref: Any = None) -> str:
    """Container name for the Bedrock allowlist check - same resolution prepare_video uses,
    so the plan gate never promises a native tier prepare would name differently."""
    return _media_format(url, getattr(ref, "mime", None), "")


def _budget_join(parts: Any, budget: int) -> str:
    """Accumulate up to ~budget chars then stop - the caller clips anyway, so materializing a
    hundred-MB string first (a 6MB xlsx renders ~27M chars) only burns memory."""
    out: List[str] = []
    used = 0
    for part in parts:
        out.append(part)
        used += len(part) + 1
        if used > budget:
            break
    return "\n".join(out)


def _docx_text(data: bytes) -> str:
    from docx import Document  # python-docx

    return _budget_join(
        (p.text for p in Document(BytesIO(data)).paragraphs), _MAX_DOC_MARKDOWN_CHARS
    )


def _xlsx_text(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(data), data_only=True, read_only=True)

    def rows():
        for sheet in wb.sheetnames:
            for row in wb[sheet].iter_rows(values_only=True):
                yield "\t".join("" if c is None else str(c) for c in row)

    return _budget_join(rows(), _MAX_DOC_MARKDOWN_CHARS)


def _pptx_text(data: bytes) -> str:
    from pptx import Presentation  # python-pptx

    def frames():
        for slide in Presentation(BytesIO(data)).slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    yield shape.text_frame.text

    return _budget_join(frames(), _MAX_DOC_MARKDOWN_CHARS)


def _cap_text(text: str, max_chars: Optional[int] = None) -> str:
    """Clip extracted text to *max_chars*, stating the true length up front so a second clip cannot hide it."""
    cap = max(1, _MAX_DOC_MARKDOWN_CHARS if max_chars is None else max_chars)
    if len(text) <= cap:
        return text
    return f"[showing the first {cap:,} of {len(text):,} characters]\n" + text[:cap]


_anydoc_import_warned = False


def _document_markdown(
    data: bytes, ext: str = "", max_chars: Optional[int] = None
) -> Optional[str]:
    """Convert office-document bytes to Markdown via anydoc; None means the caller falls back."""
    global _anydoc_import_warned
    try:
        import anydoc
    except Exception as e:
        # Once per process: an unimportable anydoc silently reverts every document to the
        # plain-text extractors, which is invisible at debug level.
        if not _anydoc_import_warned:
            _anydoc_import_warned = True
            logger.warning(
                f"anydoc unavailable ({e}); documents fall back to plain-text extraction"
            )
        return None

    # Resolved rather than named in an except clause, which would raise AttributeError on a
    # build that lacks it and escape this function entirely.
    convertible = getattr(anydoc, "ConvertError", None)
    try:
        fmt = anydoc.format_from_extension(ext.lstrip(".")) if ext else None
        markdown = anydoc.to_markdown_bytes(data, format=fmt)
    except Exception as e:
        # A document anydoc declines to convert is routine and per-file; anything else says
        # the library itself is behaving differently than this code expects.
        if convertible is not None and isinstance(e, convertible):
            logger.debug(f"anydoc cannot convert this document ({e}); falling back")
        else:
            logger.warning(f"anydoc conversion failed ({e}); falling back")
        return None
    return _cap_text(markdown, max_chars)


def _iwork_preview_pdf(
    data: bytes, max_bytes: Optional[int] = None
) -> Tuple[Optional[bytes], str]:
    """Return (preview_pdf_bytes, status in {"ok","too_large","missing"}) for an iWork zip archive."""
    import zipfile

    cap = _MAX_DOC_BYTES if max_bytes is None else max_bytes
    saw_oversized = False
    try:
        with zipfile.ZipFile(BytesIO(data)) as zf:
            names = {n.lower(): n for n in zf.namelist()}
            candidates = [
                names[c] for c in ("quicklook/preview.pdf", "preview.pdf") if c in names
            ]
            candidates += [
                real
                for lower, real in names.items()
                if lower.endswith(".pdf")
                and "quicklook" in lower
                and real not in candidates
            ]
            for name in candidates:
                # A member can declare a far larger uncompressed size than the archive that
                # passed the download cap; skip it and try the remaining candidates.
                if zf.getinfo(name).file_size > cap:
                    saw_oversized = True
                    logger.warning(
                        f"iWork preview {name} exceeds {cap} bytes; skipping"
                    )
                    continue
                return zf.read(name), "ok"
    except Exception as e:
        logger.warning(f"failed to open iWork archive: {e}")
    return None, ("too_large" if saw_oversized else "missing")


def extract_document_text(
    url: str, *, max_bytes: Optional[int] = None, timeout: Optional[float] = None
) -> str:
    """Download an office document and extract its text; '' when the format is not a document at all."""
    _, ext = os.path.splitext(urlparse(url).path.lower())
    cap = _MAX_DOC_BYTES if max_bytes is None else max_bytes
    try:
        content, _ = _download(
            url, max_bytes=cap, timeout=_FETCH_TIMEOUT if timeout is None else timeout
        )
        if ext in _IWORK_EXTS:
            # iWork content is proprietary; read the embedded preview PDF instead. Explicit
            # errors (not '') reach the model, so it says so instead of guessing.
            if not content.startswith(b"PK"):
                # Not a zip archive, e.g. a PEM ".key" file rather than Keynote.
                return "Error: attachment could not be read (unsupported file format)"
            preview, status = _iwork_preview_pdf(content, max_bytes=cap)
            if preview is None:
                if status == "too_large":
                    return (
                        "Error: attachment could not be read (the iWork file's embedded "
                        "preview exceeds the size limit; export the document as PDF instead)"
                    )
                return (
                    "Error: attachment could not be read (iWork file with no embedded "
                    "preview; re-save it with preview enabled, or export as PDF)"
                )
            markdown = _document_markdown(preview, ".pdf")
            if markdown:
                return markdown
            return (
                "Error: attachment could not be read (the iWork file's embedded "
                "preview could not be converted; export the document as PDF instead)"
            )
        # An empty conversion is a failure, not an empty document: fall through.
        markdown = _document_markdown(content, ext)
        if markdown:
            return markdown
        try:
            if ext in (".docx", ".doc"):
                text = _docx_text(content)
            elif ext in (".xlsx", ".xls"):
                text = _xlsx_text(content)
            elif ext == ".pptx":
                text = _pptx_text(content)
            else:
                text = ""
        except Exception as e:
            # A raising fallback is the case the note below exists for, so it must not
            # reach the outer handler and become ''.
            logger.warning(f"fallback extraction failed for {url}: {e}")
            text = ""
        if text.strip():
            return _cap_text(text)
        # Formats with no fallback extractor (.odt, .rtf, .epub, ...) and documents that
        # yield nothing end here; saying so beats the model answering about content it
        # never received.
        return DOC_UNREADABLE_NOTE if ext in _DOCUMENT_EXTS else ""
    except Exception as e:
        logger.warning(f"failed to extract text from document {url}: {e}")
        return ""


def extract_documents_text(urls: List[str]) -> List[Tuple[str, str]]:
    """Extract several documents concurrently, preserving order; each download blocks and anydoc frees the GIL."""
    if len(urls) <= 1:
        return [(url, extract_document_text(url=url)) for url in urls]
    from concurrent.futures import ThreadPoolExecutor

    workers = attachment_workers(len(urls), _MAX_DOC_BYTES)
    with ThreadPoolExecutor(max_workers=workers) as pool:
        return list(zip(urls, pool.map(lambda u: extract_document_text(url=u), urls)))


# Density gate for skipping the native page-image attachment: a text PDF has to yield both
# an absolute floor of Markdown and ~2 chars per KB of file. Scanned pages yield near-zero.
PDF_INLINE_MIN_CHARS = int(os.getenv("XPANDER_PDF_INLINE_MIN_CHARS", "1000"))
PDF_INLINE_CHARS_PER_KB = int(os.getenv("XPANDER_PDF_INLINE_CHARS_PER_KB", "2"))


def pdf_markdown_disabled() -> bool:
    """True when the env switch keeps every PDF on the native attachment path."""
    return os.getenv("XPANDER_PDF_MARKDOWN", "").strip().lower() in (
        "off",
        "0",
        "false",
    )


def _pdf_markdown_or_none(data: bytes) -> Optional[str]:
    """Markdown for a text-dense PDF, None for scanned/image PDFs so the caller keeps the native attachment."""
    if pdf_markdown_disabled():
        return None
    markdown = _document_markdown(data, ".pdf")
    floor = max(PDF_INLINE_MIN_CHARS, (len(data) // 1024) * PDF_INLINE_CHARS_PER_KB)
    if markdown is not None and len(markdown) >= floor:
        return markdown
    return None


def _looks_like_pdf(data: bytes) -> bool:
    """Whether the bytes contain the PDF magic; the spec allows leading bytes before %PDF-, so scan the header window."""
    return b"%PDF-" in data[:1024]


def fetch_file(url: str) -> Any:
    """
    Fetch a remote file from URL and wrap it as a File object.
    Automatically derives filename, name, format, and mime type.

    This is the legacy/kill-switch attachment path. It verifies the PDF magic and
    raises ValueError on non-PDF bytes (an expired/redirected link serves HTML with
    a 200); the text-based-PDF Markdown routing lives in media.prepare_pdf().

    Args:
        url (str): Remote file URL.

    Returns:
        File: Wrapped File object with base64 content.

    Raises:
        ValueError: the fetched bytes are not a PDF.
    """
    content, _ = _download(url, max_bytes=_MAX_DOC_BYTES, timeout=_FETCH_TIMEOUT)

    # A .pdf URL's extension is not trustworthy - an expired/redirected presigned link
    # returns HTML with HTTP 200. Verify the magic before wrapping it as a document.
    # Strip the query string: it can carry a signed token that must not reach logs.
    if not _looks_like_pdf(content):
        raise ValueError(
            f"content is not a PDF (expired or redirected link?): {url.split('?', 1)[0]}"
        )

    # Derive filename from URL
    filename = os.path.basename(url.split("?")[0])

    # Human-friendly name (strip extension, replace underscores)
    name = os.path.splitext(filename)[0].replace("_", " ")

    # Sniffed as a PDF above, so keep format/mime to the agno-accepted document type
    # rather than mimetypes' guess, which can be octet-stream (rejected at construction).
    ext = "pdf"
    mime = "application/pdf"

    # Encode content
    content_b64 = base64.b64encode(content).decode("utf-8")
    from agno.media import File

    return File.from_base64(
        base64_content=content_b64,
        filename=filename,
        name=name,
        format=ext,
        mime_type=mime,
    )
