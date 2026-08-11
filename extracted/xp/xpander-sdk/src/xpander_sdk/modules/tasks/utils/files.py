import base64
import mimetypes
from io import BytesIO
from typing import List, Optional, Tuple
from urllib.parse import urlparse
import os
from pydantic import BaseModel
from loguru import logger
import httpx
import asyncio

# Provider per-image ceiling (Bedrock/Anthropic ~5MB); env-overridable.
_MAX_IMAGE_BYTES = int(os.getenv("XPANDER_MAX_INLINE_IMAGE_BYTES", str(5 * 1024 * 1024)))
_MAX_DOC_BYTES = int(os.getenv("XPANDER_MAX_INLINE_DOC_BYTES", str(15 * 1024 * 1024)))
_FETCH_TIMEOUT = float(os.getenv("XPANDER_FILE_FETCH_TIMEOUT", "15"))

# Char ceilings for inlining human-readable file text into the prompt. A large
# CSV/JSON pasted in full blows the context window, so cap per-file and total
# and keep the URL (already listed) for the agent to fetch the rest if needed.
_MAX_INLINE_TEXT_CHARS = int(os.getenv("XPANDER_MAX_INLINE_TEXT_CHARS", "8000"))
_MAX_INLINE_TOTAL_CHARS = int(os.getenv("XPANDER_MAX_INLINE_TOTAL_CHARS", "24000"))

# Extracted text is clipped again per-file downstream; this only stops a converted
# document (a 6MB spreadsheet renders ~27M chars of Markdown) from being carried around whole.
_MAX_DOC_MARKDOWN_CHARS = max(1000, int(os.getenv("XPANDER_MAX_DOC_MARKDOWN_CHARS", "200000")))

# Concurrency for per-attachment download+convert; each one blocks on I/O then releases the GIL.
# Floored at 1 so setting it to 0 to "disable concurrency" serializes instead of raising.
_ATTACHMENT_WORKERS = max(1, int(os.getenv("XPANDER_ATTACHMENT_WORKERS", "4")))

# Concurrent downloads coexist in memory, so the worker count is also bounded by how many
# whole attachments we are willing to hold at once.
_ATTACHMENT_BYTES_BUDGET = int(os.getenv("XPANDER_ATTACHMENT_BYTES_BUDGET", str(64 * 1024 * 1024)))

DOC_UNREADABLE_NOTE = (
    "No text could be extracted from this attachment, so its contents are NOT included here; "
    "fetch it via your tools/workspace if you need them"
)


def attachment_workers(item_count: int, per_item_bytes: int) -> int:
    """Worker count for a fan-out of *item_count* attachments, bounded by the in-flight byte budget."""
    by_memory = max(1, _ATTACHMENT_BYTES_BUDGET // max(1, per_item_bytes))
    return max(1, min(_ATTACHMENT_WORKERS, item_count, by_memory))


def truncate_inline_text(content: str, url: str, remaining: int) -> str:
    """Clip inlined file *content* to *remaining* (and the per-file cap), noting the URL for the rest."""
    cap = min(_MAX_INLINE_TEXT_CHARS, max(0, remaining))
    if len(content) <= cap:
        return content
    return content[:cap] + f"\n... [truncated: {len(content):,} chars total - full file at {url}]"

# Leading-byte signatures for the only raster formats vision providers accept.
_IMAGE_SIGNATURES = (
    (b"\x89PNG\r\n\x1a\n", "png", "image/png"),
    (b"\xff\xd8\xff", "jpeg", "image/jpeg"),
    (b"GIF87a", "gif", "image/gif"),
    (b"GIF89a", "gif", "image/gif"),
)

# iWork formats are proprietary, but the files are zips embedding a QuickLook preview PDF.
_IWORK_EXTS = {".pages", ".key", ".numbers"}

# What anydoc's format_from_extension() resolves, minus csv (raw injection) and pdf (native path).
_DOCUMENT_EXTS = {
    ".docx", ".doc", ".docm",
    ".xlsx", ".xls", ".xlsm", ".xlsb",
    ".pptx", ".ppt", ".pptm", ".pot", ".pps", ".ppsx",
    ".odt", ".ods", ".odp",
    ".rtf", ".epub",
} | _IWORK_EXTS


class FileCategorization(BaseModel):
    images: List[str]
    pdfs: List[str]
    files: List[str]
    documents: List[str] = []


def categorize_files(file_urls: list[str]) -> FileCategorization:
    """
    Categorize a list of file URLs into images, PDFs, and human-readable files
    based on file extensions. Does not load the files.

    Args:
        file_urls (list[str]): List of file URLs

    Returns:
        FileCategorization: Pydantic model with categorized URLs
    """
    image_exts = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp"}
    pdf_exts = {".pdf"}
    human_readable_exts = {
        # SVG is XML: models read it as text, image pipelines reject it.
        ".svg",
        ".txt",
        ".csv",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".md",
        ".rst",
        ".yaml",
        ".yml",
        ".py",
        ".js",
        ".ts",
        ".java",
        ".c",
        ".cpp",
        ".h",
        ".cs",
        ".go",
        ".rb",
        ".php",
        ".sh",
        # plain-text formats that need no conversion, just injection
        ".tsv",
        ".log",
        ".ini",
        ".toml",
        ".tex",
        ".ipynb",
        ".eml",
    }

    result = {"images": [], "pdfs": [], "files": [], "documents": []}

    for url in file_urls:
        path = urlparse(url).path
        _, ext = os.path.splitext(path.lower())

        if ext in image_exts:
            result["images"].append(url)
        elif ext in pdf_exts:
            result["pdfs"].append(url)
        elif ext in _DOCUMENT_EXTS:
            result["documents"].append(url)
        elif ext in human_readable_exts:
            result["files"].append(url)

    return FileCategorization(**result)


async def fetch_urls(
    urls: list[str], disable_attachment_injection: Optional[bool] = False
) -> list[dict[str, str]]:
    """
    Fetches the content of multiple URLs asynchronously.

    Args:
        urls (list[str]): List of URLs to fetch.

    Returns:
        list[dict[str, str]]: A list of dictionaries containing the URL and its content.
                              Example: [{"url": "...", "content": "..."}]
        disable_attachment_injection (Optional[bool]): Optional selection if to disable attachment injection to the context window.
    """

    async def fetch(client: httpx.AsyncClient, url: str) -> dict[str, str]:
        try:
            if disable_attachment_injection:
                return {"url": url}
            response = await client.get(url, timeout=10.0)
            response.raise_for_status()
            return {"url": url, "content": response.text}
        except Exception as e:
            return {"url": url, "content": f"Error: {str(e)}"}

    async with httpx.AsyncClient() as client:
        tasks = [fetch(client, url) for url in urls]
        return await asyncio.gather(*tasks)


class AttachmentDecision(BaseModel):
    url: str
    category: str  # image | pdf | document | readable | other
    action: str  # inline | text_extract | url_only | skip
    reason: Optional[str] = None
    size: Optional[int] = None


class AttachmentPlan(BaseModel):
    items: List[AttachmentDecision] = []
    notes: List[str] = []

    def by_category(self, category: str) -> List[AttachmentDecision]:
        return [item for item in self.items if item.category == category]


async def estimate_sizes(urls: List[str], timeout: float = 2.0) -> dict:
    """Best-effort parallel size probe -> {url: bytes or None}; HEAD first, ranged-GET fallback for HEAD-hostile hosts."""
    if not urls:
        return {}

    async def _probe(client: httpx.AsyncClient, url: str) -> Optional[int]:
        try:
            resp = await client.head(url)
            raw = resp.headers.get("content-length")
            if resp.status_code < 400 and raw is not None:
                return int(raw)
        except Exception:
            pass
        try:
            resp = await client.get(url, headers={"Range": "bytes=0-0"})
            content_range = resp.headers.get("content-range") or ""
            if "/" in content_range:
                total = content_range.rsplit("/", 1)[1].strip()
                if total.isdigit():
                    return int(total)
            # Host ignored the range and served the whole body.
            raw = resp.headers.get("content-length")
            if raw is not None and int(raw) > 1:
                return int(raw)
        except Exception:
            pass
        return None

    async with httpx.AsyncClient(follow_redirects=True, timeout=timeout) as client:
        results = await asyncio.gather(*[_probe(client, u) for u in urls], return_exceptions=True)
    return {u: (r if isinstance(r, int) else None) for u, r in zip(urls, results)}


def plan_attachments(
    urls: List[str],
    caps,
    sizes: Optional[dict] = None,
) -> AttachmentPlan:
    """Decide per-URL attachment handling from (categories, capabilities, known sizes) - pure, no downloads."""
    plan = AttachmentPlan()
    if not urls:
        return plan
    sizes = sizes or {}
    cat = categorize_files(file_urls=urls)
    categories = (
        [("image", u) for u in cat.images]
        + [("pdf", u) for u in cat.pdfs]
        + [("document", u) for u in cat.documents]
        + [("readable", u) for u in cat.files]
    )
    known = {u for _, u in categories}
    categories += [("other", u) for u in urls if u not in known]

    skipped_images = 0
    inlined_images = 0
    unreadable: List[str] = []
    for category, url in categories:
        size = sizes.get(url)
        if category == "image":
            if not caps.supports_vision:
                plan.items.append(AttachmentDecision(url=url, category=category, action="skip", reason="model has no vision", size=size))
                skipped_images += 1
            elif size is not None and size > caps.max_fetch_bytes:
                plan.items.append(AttachmentDecision(url=url, category=category, action="url_only", reason="huge", size=size))
                plan.notes.append(f"{url}: too large to attach inline; fetch it via your tools/workspace using the URL")
            elif inlined_images >= caps.max_images:
                plan.items.append(AttachmentDecision(url=url, category=category, action="url_only", reason="max_images", size=size))
            else:
                plan.items.append(AttachmentDecision(url=url, category=category, action="inline", size=size))
                inlined_images += 1
        elif category == "pdf":
            if size is not None and size > caps.max_fetch_bytes:
                plan.items.append(AttachmentDecision(url=url, category=category, action="url_only", reason="huge", size=size))
                plan.notes.append(f"{url}: too large to attach inline; fetch it via your tools/workspace using the URL")
            elif caps.supports_native_pdf:
                plan.items.append(AttachmentDecision(url=url, category=category, action="inline", size=size))
            else:
                plan.items.append(AttachmentDecision(url=url, category=category, action="text_extract", reason="provider rejects file blobs", size=size))
        elif category in ("document", "readable"):
            plan.items.append(AttachmentDecision(url=url, category=category, action="text_extract", size=size))
        else:
            plan.items.append(AttachmentDecision(url=url, category=category, action="url_only", size=size))
            # No converter for this format (.msg, .zip, extension-less URLs, ...): without
            # this note the model answers about content it never saw.
            unreadable.append(url)

    if unreadable:
        plan.notes.append(
            "these attachments cannot be read directly and their content is NOT included; "
            "do not guess their contents - fetch them via your tools/workspace if you need "
            "them: " + ", ".join(unreadable)
        )

    if skipped_images:
        plural = "s" if skipped_images != 1 else ""
        plan.notes.append(
            f"{skipped_images} image{plural} listed above cannot be viewed by this model; "
            "use tools to process them if needed"
        )
    return plan


def _download(
    url: str, *, max_bytes: int, timeout: float
) -> Tuple[bytes, str]:
    """Stream a URL to bytes, aborting past `max_bytes`; returns (content, lowercased content-type)."""
    with httpx.Client(follow_redirects=True, timeout=timeout) as client:
        with client.stream("GET", url) as resp:
            resp.raise_for_status()
            ctype = (resp.headers.get("content-type") or "").split(";", 1)[0].strip().lower()
            buf = bytearray()
            for chunk in resp.iter_bytes():
                buf += chunk
                if len(buf) > max_bytes:
                    raise ValueError(f"file exceeds {max_bytes} bytes: {url}")
            return bytes(buf), ctype


def _sniff_image(data: bytes) -> Optional[Tuple[str, str]]:
    """Return (format, mime) if `data` begins with a supported image signature, else None."""
    for sig, fmt, mime in _IMAGE_SIGNATURES:
        if data.startswith(sig):
            return fmt, mime
    if len(data) >= 12 and data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return "webp", "image/webp"
    return None


def fetch_image(
    url: str, *, max_bytes: int = _MAX_IMAGE_BYTES, timeout: float = _FETCH_TIMEOUT
):
    """Download a URL and wrap it as an inline Agno Image, only if the bytes are a genuine supported image.

    Providers reject remote image URLs and require inline bytes, and a URL's extension is not trustworthy
    (expired/redirected presigned links return HTML/XML with HTTP 200). So the content-type and leading
    magic bytes are validated and format/mime are derived from the bytes, not the URL. Raises on anything
    not a verifiable png/jpeg/gif/webp so callers can drop it instead of feeding garbage to the model.
    """
    content, ctype = _download(url, max_bytes=max_bytes, timeout=timeout)
    if ctype and not (ctype.startswith("image/") or ctype in ("application/octet-stream", "binary/octet-stream")):
        raise ValueError(f"non-image content-type {ctype!r}: {url}")
    sniffed = _sniff_image(content)
    if sniffed is None:
        raise ValueError(f"bytes are not a supported image: {url}")
    fmt, mime = sniffed
    content_b64 = base64.b64encode(content).decode("utf-8")
    from agno.media import Image

    return Image.from_base64(base64_content=content_b64, format=fmt, mime_type=mime)


def _docx_text(data: bytes) -> str:
    from docx import Document  # python-docx

    return "\n".join(p.text for p in Document(BytesIO(data)).paragraphs)


def _xlsx_text(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(data), data_only=True, read_only=True)
    rows: List[str] = []
    for sheet in wb.sheetnames:
        for row in wb[sheet].iter_rows(values_only=True):
            rows.append("\t".join("" if c is None else str(c) for c in row))
    return "\n".join(rows)


def _pptx_text(data: bytes) -> str:
    from pptx import Presentation  # python-pptx

    out: List[str] = []
    for slide in Presentation(BytesIO(data)).slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                out.append(shape.text_frame.text)
    return "\n".join(out)


def _cap_text(text: str, max_chars: Optional[int] = None) -> str:
    """Clip extracted text to *max_chars*, stating the true length up front so a second clip cannot hide it."""
    cap = max(1, _MAX_DOC_MARKDOWN_CHARS if max_chars is None else max_chars)
    if len(text) <= cap:
        return text
    return f"[showing the first {cap:,} of {len(text):,} characters]\n" + text[:cap]


_anydoc_import_warned = False


def _document_markdown(data: bytes, ext: str = "", max_chars: Optional[int] = None) -> Optional[str]:
    """Convert office-document bytes to Markdown via anydoc; None means the caller falls back."""
    global _anydoc_import_warned
    try:
        import anydoc
    except Exception as e:
        # Once per process: an unimportable anydoc silently reverts every document to the
        # plain-text extractors, which is invisible at debug level.
        if not _anydoc_import_warned:
            _anydoc_import_warned = True
            logger.warning(f"anydoc unavailable ({e}); documents fall back to plain-text extraction")
        return None

    # Resolved rather than named in an except clause, which would raise AttributeError on a
    # build that lacks it and escape this function entirely.
    convertible = getattr(anydoc, "ConvertError", None)
    try:
        fmt = anydoc.format_from_extension(ext.lstrip(".")) if ext else None
        markdown = anydoc.to_markdown_bytes(data, format=fmt)
    except Exception as e:
        # A document anydoc declines to convert is routine and per-file; anything else says
        # the library itself is behaving differently than this code expects.
        if convertible is not None and isinstance(e, convertible):
            logger.debug(f"anydoc cannot convert this document ({e}); falling back")
        else:
            logger.warning(f"anydoc conversion failed ({e}); falling back")
        return None
    return _cap_text(markdown, max_chars)


def _iwork_preview_pdf(data: bytes, max_bytes: Optional[int] = None) -> Tuple[Optional[bytes], str]:
    """Return (preview_pdf_bytes, status in {"ok","too_large","missing"}) for an iWork zip archive."""
    import zipfile

    cap = _MAX_DOC_BYTES if max_bytes is None else max_bytes
    saw_oversized = False
    try:
        with zipfile.ZipFile(BytesIO(data)) as zf:
            names = {n.lower(): n for n in zf.namelist()}
            candidates = [names[c] for c in ("quicklook/preview.pdf", "preview.pdf") if c in names]
            candidates += [
                real for lower, real in names.items()
                if lower.endswith(".pdf") and "quicklook" in lower and real not in candidates
            ]
            for name in candidates:
                # A member can declare a far larger uncompressed size than the archive that
                # passed the download cap; skip it and try the remaining candidates.
                if zf.getinfo(name).file_size > cap:
                    saw_oversized = True
                    logger.warning(f"iWork preview {name} exceeds {cap} bytes; skipping")
                    continue
                return zf.read(name), "ok"
    except Exception as e:
        logger.warning(f"failed to open iWork archive: {e}")
    return None, ("too_large" if saw_oversized else "missing")


def extract_document_text(
    url: str, *, max_bytes: Optional[int] = None, timeout: Optional[float] = None
) -> str:
    """Download an office document and extract its text; '' when the format is not a document at all."""
    _, ext = os.path.splitext(urlparse(url).path.lower())
    cap = _MAX_DOC_BYTES if max_bytes is None else max_bytes
    try:
        content, _ = _download(
            url, max_bytes=cap, timeout=_FETCH_TIMEOUT if timeout is None else timeout
        )
        if ext in _IWORK_EXTS:
            # iWork content is proprietary; read the embedded preview PDF instead. Explicit
            # errors (not '') reach the model, so it says so instead of guessing.
            if not content.startswith(b"PK"):
                # Not a zip archive, e.g. a PEM ".key" file rather than Keynote.
                return "Error: attachment could not be read (unsupported file format)"
            preview, status = _iwork_preview_pdf(content, max_bytes=cap)
            if preview is None:
                if status == "too_large":
                    return (
                        "Error: attachment could not be read (the iWork file's embedded "
                        "preview exceeds the size limit; export the document as PDF instead)"
                    )
                return (
                    "Error: attachment could not be read (iWork file with no embedded "
                    "preview; re-save it with preview enabled, or export as PDF)"
                )
            markdown = _document_markdown(preview, ".pdf")
            if markdown:
                return markdown
            return (
                "Error: attachment could not be read (the iWork file's embedded "
                "preview could not be converted; export the document as PDF instead)"
            )
        # An empty conversion is a failure, not an empty document: fall through.
        markdown = _document_markdown(content, ext)
        if markdown:
            return markdown
        try:
            if ext in (".docx", ".doc"):
                text = _docx_text(content)
            elif ext in (".xlsx", ".xls"):
                text = _xlsx_text(content)
            elif ext == ".pptx":
                text = _pptx_text(content)
            else:
                text = ""
        except Exception as e:
            # A raising fallback is the case the note below exists for, so it must not
            # reach the outer handler and become ''.
            logger.warning(f"fallback extraction failed for {url}: {e}")
            text = ""
        if text.strip():
            return _cap_text(text)
        # Formats with no fallback extractor (.odt, .rtf, .epub, ...) and documents that
        # yield nothing end here; saying so beats the model answering about content it
        # never received.
        return DOC_UNREADABLE_NOTE if ext in _DOCUMENT_EXTS else ""
    except Exception as e:
        logger.warning(f"failed to extract text from document {url}: {e}")
        return ""


def extract_documents_text(urls: List[str]) -> List[Tuple[str, str]]:
    """Extract several documents concurrently, preserving order; each download blocks and anydoc frees the GIL."""
    if len(urls) <= 1:
        return [(url, extract_document_text(url=url)) for url in urls]
    from concurrent.futures import ThreadPoolExecutor

    workers = attachment_workers(len(urls), _MAX_DOC_BYTES)
    with ThreadPoolExecutor(max_workers=workers) as pool:
        return list(zip(urls, pool.map(lambda u: extract_document_text(url=u), urls)))


# Density gate for skipping the native page-image attachment: a text PDF has to yield both
# an absolute floor of Markdown and ~2 chars per KB of file. Scanned pages yield near-zero.
PDF_INLINE_MIN_CHARS = int(os.getenv("XPANDER_PDF_INLINE_MIN_CHARS", "1000"))
PDF_INLINE_CHARS_PER_KB = int(os.getenv("XPANDER_PDF_INLINE_CHARS_PER_KB", "2"))


def pdf_markdown_disabled() -> bool:
    """True when the env switch keeps every PDF on the native attachment path."""
    return os.getenv("XPANDER_PDF_MARKDOWN", "").strip().lower() in ("off", "0", "false")


def _pdf_markdown_or_none(data: bytes) -> Optional[str]:
    """Markdown for a text-dense PDF, None for scanned/image PDFs so the caller keeps the native attachment."""
    if pdf_markdown_disabled():
        return None
    markdown = _document_markdown(data, ".pdf")
    floor = max(PDF_INLINE_MIN_CHARS, (len(data) // 1024) * PDF_INLINE_CHARS_PER_KB)
    if markdown is not None and len(markdown) >= floor:
        return markdown
    return None


def fetch_file(url: str):
    """
    Fetch a remote file from URL and wrap it as a File object.
    Automatically derives filename, name, format, and mime type.

    This is the legacy/kill-switch attachment path and attaches bytes as-is;
    the text-based-PDF Markdown routing lives in media.prepare_pdf().

    Args:
        url (str): Remote file URL.

    Returns:
        File: Wrapped File object with base64 content.
    """
    content, _ = _download(url, max_bytes=_MAX_DOC_BYTES, timeout=_FETCH_TIMEOUT)

    # Derive filename from URL
    filename = os.path.basename(url.split("?")[0])

    # Human-friendly name (strip extension, replace underscores)
    name = os.path.splitext(filename)[0].replace("_", " ")

    # Guess format and mime type
    ext = os.path.splitext(filename)[1].lstrip(".").lower()
    mime, _ = mimetypes.guess_type(filename)
    mime = mime or "application/octet-stream"

    # Encode content
    content_b64 = base64.b64encode(content).decode("utf-8")
    from agno.media import File

    return File.from_base64(
        base64_content=content_b64,
        filename=filename,
        name=name,
        format=ext,
        mime_type=mime,
    )
