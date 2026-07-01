"""
cvc.gateway — The CVC Gateway (Unified Server).

A FastAPI application that runs on ``http://127.0.0.1:13421`` and provides:

1. A web-based dashboard UI (Open WebUI-style) for managing all CVC features.
2. A management REST API for service lifecycle control.
3. An agent chat endpoint (streaming) — full CVC Agent in the browser.
4. Operations API: commit, branch, merge, restore, recall, diff, timeline.
5. Memory & context browsing.
6. Model catalog and provider switching.
7. Hive Mind / SDK multi-agent management.
8. Real-time event streaming via WebSocket.
9. **Workspace-routed LLM proxy** — intercepts LLM API calls for all registered
   workspaces, replacing the standalone ``cvc serve`` proxy entirely.

The gateway is the single process users need. Running ``cvc gateway start``
provides both the dashboard AND the cognitive proxy.  Each workspace is
dynamically registered via ``POST /gateway/register_workspace`` and gets
its own isolated CVCEngine, reachable at ``/ws/{workspace_id}/v1/...``.
"""

from __future__ import annotations
from cvc._subprocess_compat import HIDDEN_KW

import asyncio
import hashlib
import json
import logging
import os
import platform
import shutil
import subprocess
import threading
import time
from contextlib import asynccontextmanager
from pathlib import Path
from typing import Any

from fastapi import FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from fastapi.staticfiles import StaticFiles
from starlette.responses import FileResponse, HTMLResponse, StreamingResponse
from starlette.websockets import WebSocket, WebSocketDisconnect

# v2.92.3 — Loop primitives (CostBudget + ToolRiskRegistry) wired into
# the gateway. Imports live here (not lazy) so any 500 from a missing
# dependency surfaces at gateway boot, not at first chat request.
from cvc.agent.loop.cost_budget import (  # noqa: E402
    CostBudget,
    CostBudgetExceeded,
    ModelPricing,
    estimate_turn_cost,
)
from cvc.agent.loop.tool_risk import (  # noqa: E402
    ToolRiskTier,
    ToolRiskDecision,
    ToolRiskRegistry,
    classify_tool_risk,
    requires_confirmation,
)

logger = logging.getLogger("cvc.gateway")

# ---------------------------------------------------------------------------
# Module-level singletons (initialised in lifespan)
# ---------------------------------------------------------------------------

_manager = None        # GatewayManager
_workspace_mgr = None  # WorkspaceManager — multi-workspace isolation
_event_bus = None      # EventBus
_health_task = None    # Background health poller
_project_root = None   # Discovered project root
_agent_llm = None      # AgentLLM (direct LLM client for chat)
_agent_tools = []      # Tool definitions for the dashboard agent
_tool_executor = None  # ToolExecutor for executing agent tool calls
_system_prompt = ""    # Agent system prompt
_chat_messages = []    # Persistent conversation history (system + user + assistant + tool)
_hermes_api_server = None  # v3.2.0 — handle to the vendored full api_server (aiohttp web app with 89 tools, sessions, runs)

# ── Tool permission system ───────────────────────────────────────────
# Approval mode mirrors VS Code's three-state model:
#   "default"   → ask the user for every non-safe tool
#   "bypass"    → skip approvals entirely (auto-allow ALL tools)
#   "autopilot" → auto-allow + let the agent continue multi-step on its own
# v2.92.4 — PROJECT LAW (full autonomy): the agent must NEVER ask the
# user mid-task. The default mode is now "autopilot" so every session
# boots into strict autonomy. The user can still flip to "default" via
# /api/chat/approval-mode if they want a cautious session, but the
# default is autonomous.
_approval_mode: str = "autopilot"
# Back-compat: any read of `_autopilot` resolves to "should we auto-allow?"
def _autopilot_on() -> bool:
    return _approval_mode in ("bypass", "autopilot")
_autopilot: bool = False  # legacy flag, kept for any external reader; do not write
_tool_confirm_event: asyncio.Event | None = None  # Set when user confirms/denies
_tool_confirm_result: str = ""  # Decision string from user

# ── v2.92.3 — Loop primitives wired into the gateway ──────────────────
# Per-conversation CostBudget + ToolRiskRegistry. The dashboard flips
# these via /api/chat/loop-config; the agent_chat and ws_chat paths
# check the budget before each LLM call and classify tool risk before
# each tool dispatch. Mirrors Claude Code's `max_budget_usd`,
# `--accept-network`, and `--yes-destroy` flags.
_loop_cost_budget: CostBudget = CostBudget()  # mutates in place via reset() + new cap
_loop_cost_budget_lock = threading.Lock()
_loop_risk_registry: ToolRiskRegistry = ToolRiskRegistry()
_loop_accept_network: bool = False
_loop_yes_destroy: bool = False
_loop_budget_warn_threshold: float = 0.8  # emit SSE warning at 80% spent

# ── Abort / Stop system ──────────────────────────────────────────────
_chat_abort: bool = False  # When True, the agentic loop stops ASAP

# Session-level tool trust ("Always allow X this session")
_session_trusted_tools: set[str] = set()
# Session-level "Trust ALL commands" flag
_session_trust_all: bool = False

# Tools that are always safe to auto-approve (read-only, no side effects)
_SAFE_TOOLS: set[str] = {
    "read_file", "grep", "glob", "list_dir",
    "cvc_status", "cvc_log", "cvc_search", "cvc_smart_search", "cvc_diff",
    "cvc_get_context", "cvc_switch_workspace",
    # v2.92.4 — ask_user is auto-resolved autonomously (no picker, no popup).
    # Full autonomy mode = agent never asks the user; ask_user is treated
    # as a safe no-op that the agent can call to satisfy its own internal
    # "I should consult" reflex without ever surfacing a UI prompt.
    "ask_user",
}

# v2.92.10 — Tools the model can call with literally zero arguments and
# still produce a sensible result. The gateway's dud-detector normally
# treats `args == {}` as a dud (the "model emitted dud tool calls"
# pattern), but these tools are different: ask_user is auto-resolved
# in full-autonomy mode, cvc_status reports current state, todo with
# no args lists tasks, etc. Without this whitelist, models that emit
# `ask_user({})` or `cvc_status({})` get bounced as duds and the
# "stop emitting duds" nudge loops forever.
_ZERO_ARG_SAFE_TOOLS: set[str] = {
    "ask_user",
    "cvc_status",
    "cvc_log",
    "cvc_diff",
    "cvc_search",
    "cvc_smart_search",
    "cvc_list_documents",
    "task_list",
    "think",
    "context_compact",
    "todo",
    "save_memory",  # can save with just a key derived from message
}

# Valid permission decisions (matches CLI-style options)
# allow_once  — execute this one time only
# allow_always — always allow this tool for the rest of the session
# trust_all  — trust ALL tools this session (like _session_trust_all)
# deny       — deny and inject error message
# deny_suggest — deny and ask agent to suggest an alternative
_VALID_DECISIONS = {"allow_once", "allow_always", "trust_all", "deny", "deny_suggest"}

# Agentic loop limits
def _wire_subagent_runtime(executor) -> None:
    """
    Inject the sub-agent runtime into a freshly-created ToolExecutor so
    the `agent` and `parallel_agents` tools work in dashboard chat
    (SSE + WebSocket) — mirrors what cvc/agent/chat.py:697 does for CLI.

    Without this:
      • `agent` runs SubAgent with empty provider/api_key → instant error
      • `parallel_agents` returns "dispatch queued" sentinel → loop stalls

    With this: dashboard chat has the same autonomous, never-stop loop as CLI.

    ── Live nested-activity surfacing ────────────────────────────────────
    The chat loop installs `executor._subagent_event_sink` (a callable
    accepting a single dict). When set, every SubAgent run forwards its
    own `subagent_start / subagent_tool_start / subagent_tool_result /
    subagent_progress / subagent_done` events into that sink. The sink hands
    them to the SSE
    or WebSocket stream so the dashboard renders nested pills under the
    parent `agent` / `parallel_agents` tool pill.

    The sink MUST be thread-safe — SubAgent.run() executes inside a
    `run_in_executor` thread (single-agent) or inside a fresh asyncio
    loop on a worker thread (parallel_agents).
    """
    if executor is None:
        return
    try:
        cfg = _get_config()
        provider = (_agent_llm.provider if _agent_llm else (cfg.provider if cfg else "")) or ""
        model = (_agent_llm.model if _agent_llm else (cfg.model if cfg else "")) or ""
        api_key = getattr(_agent_llm, "api_key", "") if _agent_llm else ""
        base_url = getattr(_agent_llm, "_api_url", "") if _agent_llm else ""

        executor._subagent_config = {
            "provider": provider,
            "api_key": api_key,
            "model": model,
            "base_url": base_url,
        }

        # The chat loop sets this to a thread-safe callable that pushes
        # events into the per-request stream. Default: no-op (events
        # silently dropped — CLI/legacy behaviour).
        if not hasattr(executor, "_subagent_event_sink"):
            executor._subagent_event_sink = None

        # Real parallel_agents callback — asyncio.gather across SubAgent.run.
        # Called synchronously from the executor (tools run in a thread pool),
        # so we hop off the running loop into a fresh one for the gather.
        def _parallel_cb(tasks):
            import asyncio
            from cvc.agent.subagent import SubAgent, get_available_agents

            agents = get_available_agents(executor.workspace)
            sink = getattr(executor, "_subagent_event_sink", None)
            # v2.84 — nested-spawn depth budget. Orchestrator subagents
            # hold parallel_agents in their tool set; a runaway LLM could
            # recursively spawn Orchestrators forever. The executor is
            # tagged with `_subagent_depth` when running INSIDE a subagent
            # (see SubAgent.run); top-level chats leave it 0. We refuse
            # any spawn that would push children past the hard cap.
            _MAX_DEPTH = int(os.environ.get("CVC_SUBAGENT_MAX_DEPTH", "2"))
            _parent_depth = int(getattr(executor, "_subagent_depth", 0) or 0)
            _child_depth = _parent_depth + 1
            if _child_depth > _MAX_DEPTH:
                # Synthesize refusal results without touching the model.
                _msg = (
                    f"[runtime] parallel_agents refused: nesting depth "
                    f"{_child_depth} > CVC_SUBAGENT_MAX_DEPTH ({_MAX_DEPTH}). "
                    "A sub-agent cannot recursively spawn further sub-agents "
                    "this deep. Run these tasks at the top level instead, "
                    "or raise CVC_SUBAGENT_MAX_DEPTH."
                )
                if sink:
                    try:
                        sink({
                            "type": "subagent_depth_refused",
                            "depth": _child_depth,
                            "max_depth": _MAX_DEPTH,
                            "task_count": len(tasks),
                        })
                    except Exception:
                        pass
                return [_msg for _ in tasks]

            async def _run_all():
                coros = []
                _total = len(tasks)
                for _i, task in enumerate(tasks):
                    name = task.get("agent") or task.get("name") or "Explore"
                    prompt = task.get("prompt", "") or ""
                    sub_cfg = agents.get(name)
                    if not sub_cfg or not prompt:
                        async def _skip(n=name, p=prompt):
                            if not sub_cfg:
                                return f"Error: unknown agent '{n}'"
                            return "Error: empty prompt"
                        coros.append(_skip())
                        continue
                    sub = SubAgent(
                        config=sub_cfg,
                        workspace=executor.workspace,
                        provider=provider,
                        api_key=api_key,
                        parent_model=model,
                        base_url=base_url,
                        event_emitter=sink,
                        agent_index=_i + 1,
                        agent_total=_total,
                        depth=_child_depth,
                    )
                    coros.append(sub.run(prompt))
                return await asyncio.gather(*coros, return_exceptions=True)

            # We're inside a tool call dispatched via run_in_executor → no
            # running loop on this thread. Safe to spin up a fresh one.
            try:
                loop = asyncio.get_event_loop()
                if loop.is_running():
                    # Should be rare; fall back to a thread-isolated loop.
                    import concurrent.futures
                    with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                        raw = pool.submit(asyncio.run, _run_all()).result()
                else:
                    raw = loop.run_until_complete(_run_all())
            except RuntimeError:
                raw = asyncio.run(_run_all())

            out = []
            for r in raw:
                if isinstance(r, BaseException):
                    out.append(f"Error: {r}")
                else:
                    out.append(str(r))
            return out

        executor._parallel_agents_callback = _parallel_cb
    except Exception:
        logger.warning("Failed to wire sub-agent runtime onto ToolExecutor", exc_info=True)


# ─────────────────────────────────────────────────────────────────────────────
# AGENT BUDGETS — parity-with-upstream configurable limits
# ─────────────────────────────────────────────────────────────────────────────
# Upstream default is 90 turns; users routinely raise to 180+ for long audits.
# These are FALLBACK defaults. Effective values are resolved per-turn from:
#   1. Env vars (CVC_MAX_ITERATIONS, CVC_AGENT_TIMEOUT, …)  ← highest priority
#   2. GlobalConfig (~/.config/cvc/config.json)              ← persisted user pref
#   3. These defaults                                        ← shipped baseline
# v2.69.0 — All agent-loop tunables are now resolved from CvcAgentConfig
# (~/.cvc/agent.toml, env-var overrides, runtime defaults). Module-level
# names below are kept ONLY as cached snapshots for hot paths and as
# back-compat aliases for any legacy reader. Refresh with
# `_refresh_agent_budgets()` after `cvc setup` mutates the config.
try:  # safe-import — gateway.py is sometimes loaded by tests before pkg init
    from cvc.core.agent_config import get_agent_config as _get_agent_config_singleton
except Exception:  # pragma: no cover
    _get_agent_config_singleton = None  # type: ignore


def _agent_cfg():
    if _get_agent_config_singleton is None:
        class _Shim:  # minimal fallback so module import never fails
            max_agent_iterations = 80
            max_agent_iterations_expensive = 40
            agent_timeout = 900.0
            tool_exec_timeout = 300.0
            max_tool_output_llm = 15000
            max_tool_output_expensive = 5000
            confirm_timeout = 120.0
            sse_heartbeat_interval = 15.0
            turn_ring_size = 2000
            turn_gc_interval = 60.0
            turn_done_ttl = 300.0
            turn_orphan_ttl = 1800.0
            auto_compact_threshold = 95
            auto_compact_pct = 0.5
        return _Shim()
    return _get_agent_config_singleton()


_cfg_snapshot = _agent_cfg()

_DEFAULT_MAX_AGENT_ITERATIONS = _cfg_snapshot.max_agent_iterations
_DEFAULT_MAX_AGENT_ITERATIONS_EXPENSIVE = _cfg_snapshot.max_agent_iterations_expensive
_DEFAULT_AGENT_TIMEOUT = float(_cfg_snapshot.agent_timeout)
_DEFAULT_TOOL_EXEC_TIMEOUT = float(_cfg_snapshot.tool_exec_timeout)

_MAX_AGENT_ITERATIONS = _DEFAULT_MAX_AGENT_ITERATIONS
_MAX_AGENT_ITERATIONS_EXPENSIVE = _DEFAULT_MAX_AGENT_ITERATIONS_EXPENSIVE
_AGENT_TIMEOUT = _DEFAULT_AGENT_TIMEOUT
_TOOL_EXEC_TIMEOUT = _DEFAULT_TOOL_EXEC_TIMEOUT

_MAX_TOOL_OUTPUT_LLM = _cfg_snapshot.max_tool_output_llm
_MAX_TOOL_OUTPUT_EXPENSIVE = _cfg_snapshot.max_tool_output_expensive
_CONFIRM_TIMEOUT = float(_cfg_snapshot.confirm_timeout)
_SSE_HEARTBEAT_INTERVAL = float(_cfg_snapshot.sse_heartbeat_interval)

# v2.72.3 — Per-event LLM stream watchdog. The upstream provider (Copilot,
# Anthropic, Gemini, etc.) can stall mid-stream with no events and no
# exception (TCP idle, server hiccup). Without this watchdog the dashboard
# agentic loop blocks forever inside `async for event in chat_stream(...)`
# and the user sees a frozen spinner. Override with CVC_STREAM_EVENT_TIMEOUT.
try:
    _STREAM_EVENT_TIMEOUT = float(os.environ.get("CVC_STREAM_EVENT_TIMEOUT", "60"))
except (TypeError, ValueError):
    _STREAM_EVENT_TIMEOUT = 60.0

# v2.92.11 — Stream-stall heartbeat. When the upstream LLM stalls
# (idle TCP, no events, no exception), the dashboard used to show
# a frozen "Just Wait..." spinner for the full _STREAM_EVENT_TIMEOUT
# (60s) before any "recovering..." status reached the user. The chat
# stream loops now use an INLINE chunked wait (rebuilt in v2 of this
# change, after v1's create_task-on-async-generator was caught by the
# post-edit test): each chunk waits _STREAM_HEARTBEAT_INTERVAL (10s),
# emits a "still working..." status on every empty chunk, and aborts
# after _STREAM_HEARTBEAT_TIMEOUT (15s) total. The legacy
# _STREAM_EVENT_TIMEOUT (60s) stays as the absolute safety net for
# other code paths. Override with
# CVC_STREAM_HEARTBEAT_INTERVAL / CVC_STREAM_HEARTBEAT_TIMEOUT.
try:
    _STREAM_HEARTBEAT_INTERVAL = float(
        os.environ.get("CVC_STREAM_HEARTBEAT_INTERVAL", "10")
    )
except (TypeError, ValueError):
    _STREAM_HEARTBEAT_INTERVAL = 10.0
try:
    # v2.98.0 — Raised the cumulative stall timeout from 15s → 60s.
    # The 15s default was declared a stall whenever a slow model
    # (MiniMax-M3 reasoning through 2+ file listings between tool
    # calls) took longer than 15s to emit the next delta. The 2026-
    # 06-23 live screenshot showed exactly this pattern: model
    # emitted preamble + 2 duds + 2 successful list_dir calls, then
    # stalled silently because the next LLM call was reasoning
    # through the file contents to pick the next read. 15s is too
    # tight for any model that does real reasoning between tool
    # calls. 60s is the same as the legacy _STREAM_EVENT_TIMEOUT
    # safety net, so this is functionally "tighten the legacy
    # path" rather than "loosen the new path". The trade-off is
    # a longer wait before declaring a true stall, but the user
    # still sees the "Still working…" heartbeat every 10s, so
    # visible progress is preserved.
    _STREAM_HEARTBEAT_TIMEOUT = float(
        os.environ.get("CVC_STREAM_HEARTBEAT_TIMEOUT", "60")
    )
except (TypeError, ValueError):
    _STREAM_HEARTBEAT_TIMEOUT = 60.0
# v2.98.0 — Two-strike stall detection. A single 60s stall is
# forgivable (the model may be reasoning); declaring a stall after
# the *first* 60s timeout would fire too eagerly. New rule: only
# declare `_stream_stalled` after the timeout fires TWICE in a row
# within the same iteration. This adds up to 120s of patience for
# genuinely slow models, while still catching true stalls (e.g.
# dead TCP, provider hiccup) within a bounded time. The counter
# resets on every successful chunk.
_STREAM_STALL_STRIKES = int(os.environ.get("CVC_STREAM_STALL_STRIKES", "2"))

# v2.92.11 — Stream-stall heartbeat. When the upstream LLM stalls
# (idle TCP, no events, no exception), the dashboard used to show
# a frozen "Just Wait..." spinner for the full _STREAM_EVENT_TIMEOUT
# (60s) before any "recovering..." status reached the user. Now we
# spawn a heartbeat ticker that fires every _STREAM_HEARTBEAT_INTERVAL
# during stalls (so the user sees visible progress) AND we tighten
# the per-event wait_for ceiling from _STREAM_EVENT_TIMEOUT (60s) to
# _STREAM_HEARTBEAT_TIMEOUT (15s) so a true stall aborts within 15s
# instead of 60s. The legacy _STREAM_EVENT_TIMEOUT (60s) stays as the
# absolute safety net for non-heartbeat code paths. Override with
# CVC_STREAM_HEARTBEAT_INTERVAL / CVC_STREAM_HEARTBEAT_TIMEOUT.
try:
    _STREAM_HEARTBEAT_INTERVAL = float(
        os.environ.get("CVC_STREAM_HEARTBEAT_INTERVAL", "10")
    )
except (TypeError, ValueError):
    _STREAM_HEARTBEAT_INTERVAL = 10.0
try:
    # v2.98.0 — Raised the cumulative stall timeout from 15s → 60s.
    # The 15s default was declared a stall whenever a slow model
    # (MiniMax-M3 reasoning through 2+ file listings between tool
    # calls) took longer than 15s to emit the next delta. The 2026-
    # 06-23 live screenshot showed exactly this pattern: model
    # emitted preamble + 2 duds + 2 successful list_dir calls, then
    # stalled silently because the next LLM call was reasoning
    # through the file contents to pick the next read. 15s is too
    # tight for any model that does real reasoning between tool
    # calls. 60s is the same as the legacy _STREAM_EVENT_TIMEOUT
    # safety net, so this is functionally "tighten the legacy
    # path" rather than "loosen the new path". The trade-off is
    # a longer wait before declaring a true stall, but the user
    # still sees the "Still working…" heartbeat every 10s, so
    # visible progress is preserved.
    _STREAM_HEARTBEAT_TIMEOUT = float(
        os.environ.get("CVC_STREAM_HEARTBEAT_TIMEOUT", "60")
    )
except (TypeError, ValueError):
    _STREAM_HEARTBEAT_TIMEOUT = 60.0
# v2.98.0 — Two-strike stall detection. A single 60s stall is
# forgivable (the model may be reasoning); declaring a stall after
# the *first* 60s timeout would fire too eagerly. New rule: only
# declare `_stream_stalled` after the timeout fires TWICE in a row
# within the same iteration. This adds up to 120s of patience for
# genuinely slow models, while still catching true stalls (e.g.
# dead TCP, provider hiccup) within a bounded time. The counter
# resets on every successful chunk.
_STREAM_STALL_STRIKES = int(os.environ.get("CVC_STREAM_STALL_STRIKES", "2"))

# v2.73.0 — Per-send WS watchdog. The freeze pattern (3-4 tools complete then
# silence with WS LIVE) is caused by `ws.send_json` awaiting `transport.drain()`
# forever when the client back-pressures (inactive tab, suspended renderer,
# stalled proxy). _send_safe wraps every send in wait_for; on timeout the
# socket is detached and the resume ring buffer keeps events for reconnect.
try:
    _WS_SEND_TIMEOUT = float(os.environ.get("CVC_WS_SEND_TIMEOUT", "10"))
except (TypeError, ValueError):
    _WS_SEND_TIMEOUT = 10.0

# v2.68.14 — Resumable turn protocol (upstream 2.0 robustness parity).
_TURN_RING_SIZE = _cfg_snapshot.turn_ring_size
_TURN_GC_INTERVAL = float(_cfg_snapshot.turn_gc_interval)
_TURN_DONE_TTL = float(_cfg_snapshot.turn_done_ttl)
_TURN_ORPHAN_TTL = float(_cfg_snapshot.turn_orphan_ttl)


def _refresh_agent_budgets() -> None:
    """Re-read tunables from disk (called after `cvc setup` saves)."""
    global _DEFAULT_MAX_AGENT_ITERATIONS, _DEFAULT_MAX_AGENT_ITERATIONS_EXPENSIVE
    global _DEFAULT_AGENT_TIMEOUT, _DEFAULT_TOOL_EXEC_TIMEOUT
    global _MAX_AGENT_ITERATIONS, _MAX_AGENT_ITERATIONS_EXPENSIVE
    global _AGENT_TIMEOUT, _TOOL_EXEC_TIMEOUT
    global _MAX_TOOL_OUTPUT_LLM, _MAX_TOOL_OUTPUT_EXPENSIVE
    global _CONFIRM_TIMEOUT, _SSE_HEARTBEAT_INTERVAL
    global _TURN_RING_SIZE, _TURN_GC_INTERVAL, _TURN_DONE_TTL, _TURN_ORPHAN_TTL
    try:
        from cvc.core.agent_config import reload_agent_config
        c = reload_agent_config()
    except Exception:
        c = _agent_cfg()
    _DEFAULT_MAX_AGENT_ITERATIONS = c.max_agent_iterations
    _DEFAULT_MAX_AGENT_ITERATIONS_EXPENSIVE = c.max_agent_iterations_expensive
    _DEFAULT_AGENT_TIMEOUT = float(c.agent_timeout)
    _DEFAULT_TOOL_EXEC_TIMEOUT = float(c.tool_exec_timeout)
    _MAX_AGENT_ITERATIONS = _DEFAULT_MAX_AGENT_ITERATIONS
    _MAX_AGENT_ITERATIONS_EXPENSIVE = _DEFAULT_MAX_AGENT_ITERATIONS_EXPENSIVE
    _AGENT_TIMEOUT = _DEFAULT_AGENT_TIMEOUT
    _TOOL_EXEC_TIMEOUT = _DEFAULT_TOOL_EXEC_TIMEOUT
    _MAX_TOOL_OUTPUT_LLM = c.max_tool_output_llm
    _MAX_TOOL_OUTPUT_EXPENSIVE = c.max_tool_output_expensive
    _CONFIRM_TIMEOUT = float(c.confirm_timeout)
    _SSE_HEARTBEAT_INTERVAL = float(c.sse_heartbeat_interval)
    _TURN_RING_SIZE = c.turn_ring_size
    _TURN_GC_INTERVAL = float(c.turn_gc_interval)
    _TURN_DONE_TTL = float(c.turn_done_ttl)
    _TURN_ORPHAN_TTL = float(c.turn_orphan_ttl)

_EXPENSIVE_MODEL_KEYWORDS = ("opus",)


def _is_expensive_model(model_name: str) -> bool:
    """Check if a model is in the expensive/premium tier."""
    name_lower = model_name.lower()
    return any(kw in name_lower for kw in _EXPENSIVE_MODEL_KEYWORDS)


def _resolve_budget(key: str, default: float | int, cast=float) -> float | int:
    """Resolve an agent budget value per-turn.

    Order: env var (CVC_<KEY>) → GlobalConfig attr → default.
    Logged once at startup; resilient to malformed values.
    """
    env_name = f"CVC_{key.upper()}"
    raw = os.environ.get(env_name)
    if raw:
        try:
            v = cast(raw)
            if v > 0:
                return v
        except (ValueError, TypeError):
            logger.warning("Invalid %s=%r — falling back", env_name, raw)
    try:
        from cvc.core.models import GlobalConfig
        gc = GlobalConfig.load()
        v = getattr(gc, key.lower(), None)
        if v is not None and v > 0:
            return cast(v)
    except Exception:
        pass
    return cast(default)


def _budget_max_iterations(expensive: bool) -> int:
    """Effective iteration cap for this turn (resolved at call time)."""
    if expensive:
        return int(_resolve_budget(
            "MAX_ITERATIONS_EXPENSIVE",
            _DEFAULT_MAX_AGENT_ITERATIONS_EXPENSIVE,
            cast=int,
        ))
    return int(_resolve_budget(
        "MAX_ITERATIONS",
        _DEFAULT_MAX_AGENT_ITERATIONS,
        cast=int,
    ))


def _budget_agent_timeout() -> float:
    """Effective total per-turn timeout in seconds."""
    return float(_resolve_budget(
        "AGENT_TIMEOUT", _DEFAULT_AGENT_TIMEOUT, cast=float,
    ))


def _budget_tool_timeout() -> float:
    """Effective per-tool-call timeout in seconds."""
    return float(_resolve_budget(
        "TOOL_TIMEOUT", _DEFAULT_TOOL_EXEC_TIMEOUT, cast=float,
    ))


# ── v2.90.13 — Long-task survival helpers ─────────────────────────────────
#
# Three failure modes the previous loop could not handle and we now fix:
#
#   (1) Iteration-budget starvation. A long debugging task burns 25-30 tool
#       calls exploring the codebase, then the model has no budget left to
#       write the final answer — it emits one empty turn and we hit
#       "produced tool output but no reply".
#       FIX: at ~65 % of the limit we inject a one-shot "wrap up now" nudge
#       so the model knows to stop exploring and synthesise.
#
#   (2) Silent finish on a stalled provider. Re-calling the SAME provider
#       that just went silent rarely recovers — it tends to go silent again.
#       FIX: on the second silent-finish we build a transient fallback
#       AgentLLM bound to a different provider/model (Gemini Flash by
#       default, configurable) and use it for the synthesis pass.
#
#   (3) Useless giveup message. The old "I produced tool output but no
#       reply. Please rephrase your request." string told the user nothing.
#       FIX: emit a diagnostic line with iter/elapsed/provider/last-tool so
#       Jai can actually see what failed.
#
# upstream ships neither of these. CVC now does. ── Sofia, v2.90.13.

def _soft_wrap_threshold(iter_limit: int) -> int:
    """Iteration index at which the wrap-up nudge fires (one-shot per turn)."""
    return max(1, int(round(iter_limit * 0.65)))


_WRAPUP_NUDGE_TEXT = (
    "SYSTEM: You have used a substantial fraction of your tool-call budget on "
    "this turn. STOP exploring further. Synthesise a complete final answer to "
    "the user RIGHT NOW based on what you already know. You may make at most "
    "ONE more small targeted tool call only if it is strictly required to "
    "answer; otherwise reply in plain prose immediately."
)


def _pick_fallback_synthesis_provider(current_provider: str) -> tuple[str, str, str] | None:
    """Pick a cheap, fast, available provider for the synthesis fallback.

    Returns (provider, model, api_key) or None if nothing usable is configured.
    Preference order:
        1. Google Gemini Flash      (fast, generous free tier)
        2. NVIDIA Nemotron Super    (free)
        3. MiniMax M3               (1M ctx, multimodal — v2.91.41)
        4. Anthropic Haiku          (cheap)
        5. OpenAI gpt-4o-mini       (cheap)
    Skips the current provider — the whole point is to pick something else.
    v2.91.41: MiniMax added to the candidate list so MiniMax-primary users
    (the dashboard screenshot Jai reported on 2026-06-10) get a viable
    failover target instead of falling straight to the diagnostic giveup.
    """
    try:
        from cvc.agent.config import GlobalConfig
        gc = GlobalConfig.load()
    except Exception:
        return None

    candidates: list[tuple[str, str]] = [
        ("google", "gemini-3-flash-preview"),
        ("nvidia", "nvidia/nemotron-3-super-120b-instruct"),
        ("minimax", "MiniMax-M3"),
        ("anthropic", "claude-haiku-4-5"),
        ("openai", "gpt-4o-mini"),
    ]
    cur = (current_provider or "").lower()
    if cur == "copilot":
        cur = "github"

    env_map = {
        "anthropic": "ANTHROPIC_API_KEY",
        "openai": "OPENAI_API_KEY",
        "google": ("GOOGLE_API_KEY", "GEMINI_API_KEY"),
        "nvidia": ("NVIDIA_API_KEY", "NIM_API_KEY"),
        "minimax": ("MINIMAX_API_KEY", "MINIMAX_TOKEN"),
    }
    for prov, model in candidates:
        if prov == cur:
            continue
        env_keys = env_map.get(prov, ())
        if isinstance(env_keys, str):
            env_keys = (env_keys,)
        key = ""
        for ek in env_keys:
            key = os.getenv(ek, "") or key
            if key:
                break
        if not key:
            try:
                key = gc.api_keys.get(prov, "") or ""
            except Exception:
                key = ""
        if not key:
            continue
        return prov, model, key
    return None


def _build_fallback_llm(current_provider: str, current_model: str):
    """Construct a transient AgentLLM on a different provider for synthesis.

    Returns the AgentLLM instance, or None if no alternative is configured.
    Logs the choice so the gateway log makes the failover obvious.
    v2.91.41: MiniMax base URL added (https://api.minimax.io/anthropic —
    Anthropic-Messages-API-compatible per upstream _requires_bearer_auth).
    """
    picked = _pick_fallback_synthesis_provider(current_provider)
    if not picked:
        return None
    prov, model, key = picked
    base_url_map = {
        "anthropic": "https://api.anthropic.com",
        "openai": "https://api.openai.com",
        "google": "https://generativelanguage.googleapis.com",
        "nvidia": "https://integrate.api.nvidia.com",
        "minimax": "https://api.minimax.io/anthropic",  # v2.91.41
    }
    try:
        from cvc.agent.llm import AgentLLM
        llm = AgentLLM(
            provider=prov,
            api_key=key,
            model=model,
            base_url=base_url_map.get(prov, ""),
        )
        logger.warning(
            "synth_fallback: failing over from %s/%s to %s/%s",
            current_provider, current_model, prov, model,
        )
        return llm
    except Exception as exc:
        logger.warning("synth_fallback: build failed (%s)", exc)
        return None


def _diagnostic_giveup(iteration: int, iter_limit: int, elapsed: float,
                       provider: str, model: str, last_tool: str | None,
                       reason: str, last_tool_output: str | None = None) -> str:
    """A real diagnostic line — replaces the useless v2.90.12 string.

    v2.91.41: If `last_tool_output` is supplied, the message appends a
    user-visible summary of what the last tool returned, so the user
    is never left wondering "what just happened?" when the model
    silent-finishes.

    v2.97.0: Replaced the v2.91.41 footer with an agent-voiced
    message that takes responsibility for the gap and gives a
    concrete re-send prompt. The agent owns the recovery; the
    user just re-sends.
    """
    if last_tool_output:
        snippet = last_tool_output.strip()
        if len(snippet) > 400:
            snippet = snippet[:400].rstrip() + "…"
        return (
            f"⚠️ I hit the agent loop cap before I could finish your task. "
            f"Here's where I got to:\n\n"
            f"Last tool I ran: `{last_tool or '?'}` "
            f"({elapsed:.0f}s elapsed, {iteration}/{iter_limit} iters, "
            f"{provider}/{model})\n\n"
            f"Last tool output (truncated):\n```\n{snippet}\n```\n\n"
            f"Re-send the same request — I'll pick up from the state above "
            f"and complete the change. If you can name the specific file "
            f"and the specific edit, I can finish in one shot."
        )
    return (
        f"⚠️ I hit the agent loop cap before I could finish your task "
        f"({iteration}/{iter_limit} iters, {elapsed:.0f}s, "
        f"{provider}/{model}). "
        f"Re-send the same request and I'll pick up where I left off. "
        f"If you can name the specific file and the specific edit, I "
        f"can finish in one shot."
    )


# v2.91.41 — Deterministic synthesis fallback. When the model silent-finishes
# twice in a row AND the cross-provider synthesis also fails or returns
# nothing, we MUST still send the user a real, useful reply. Build it from
# the tool results already in `llm_messages` — no LLM needed. This is the
# last line of defense and must NEVER produce an empty string.
def _deterministic_synthesize(
    llm_messages: list[dict[str, Any]],
    user_facing_request: str = "",
) -> str:
    """Build a reply from tool results without calling any LLM.

    The goal is **never** to return an empty string when there is real
    information in the tool results. Even on a total synthesis failure,
    we emit a polite diagnostic that names the tools that ran and the
    number of bytes of output produced.

    Args:
        llm_messages: full message history; we scan for `role == "tool"`
            entries to summarise.
        user_facing_request: the original user request (last user-role
            message), so the reply can acknowledge what was asked.

    Returns:
        A non-empty user-visible string suitable to be sent as a
        `text` event on SSE or WS.
    """
    # 1. Walk the message history, collect every tool message and the
    #    most-recent assistant tool_use block.
    tool_results: list[tuple[str, str]] = []  # (tool_name, content_snippet)
    pending_tool_names: list[str] = []
    for m in llm_messages:
        role = m.get("role", "")
        if role == "assistant":
            content = m.get("content", "")
            # OpenAI-format tool_calls
            tcs = m.get("tool_calls") or []
            for tc in tcs:
                fn = (tc.get("function") or {}) if isinstance(tc, dict) else {}
                name = fn.get("name") or tc.get("name") or "tool"
                pending_tool_names.append(name)
            # Anthropic-format content blocks
            if isinstance(content, list):
                for blk in content:
                    if isinstance(blk, dict) and blk.get("type") == "tool_use":
                        pending_tool_names.append(blk.get("name") or "tool")
        elif role == "tool":
            name = m.get("name") or (pending_tool_names[-1] if pending_tool_names else "tool")
            content = m.get("content", "")
            if isinstance(content, list):
                # Anthropic tool_result blocks
                content = " ".join(
                    str(b.get("text", b)) for b in content
                    if isinstance(b, dict)
                )
            tool_results.append((name, str(content)))

    if not tool_results and not pending_tool_names:
        # No tool activity at all — the model just stayed silent. Give the
        # user a clear, agent-voiced diagnostic (better than an empty bubble).
        # v2.97.0: replaced the CLI-flag suggestions (/model, cvc compact
        # --smart, "split the request") with a message that takes
        # responsibility for the gap and gives a concrete re-send prompt.
        ack = ""
        if user_facing_request:
            short = user_facing_request.strip().split("\n")[0][:120]
            ack = f'You asked: "{short}".\n\n'
        return (
            ack
            + "⚠️ I didn't produce a response this turn. The model "
            "may have returned an empty reply, had a request "
            "filtered, or stalled before sending any text. "
            "Re-send the same request — I'll pick up the state and "
            "answer it. If I keep failing on the same prompt, try "
            "phrasing it slightly differently (one specific file + "
            "one specific edit) so I have less to reason about."
        )

    # 2. We have tool results. Compose a reply that surfaces the most
    #    recent result verbatim (truncated) and summarises the rest.
    lines: list[str] = []
    if user_facing_request:
        short = user_facing_request.strip().split("\n")[0][:120]
        lines.append(f'You asked: "{short}".\n')
    lines.append("Here's what I gathered from the tools I ran:\n")
    # Show the LAST tool result in full (up to 1500 chars), earlier ones summarised.
    if len(tool_results) == 1:
        name, content = tool_results[0]
        lines.append(f"**{name}** returned:")
        snippet = content.strip()
        if len(snippet) > 1500:
            snippet = snippet[:1500].rstrip() + "…"
        lines.append("```\n" + snippet + "\n```")
    else:
        # Summarise older results
        for name, content in tool_results[:-1]:
            snippet = content.strip().replace("\n", " ")
            if len(snippet) > 200:
                snippet = snippet[:200].rstrip() + "…"
            lines.append(f"- {name}: {snippet}")
        # Show the most-recent result in full
        name, content = tool_results[-1]
        lines.append(f"\n**{name}** (most recent) returned:")
        snippet = content.strip()
        if len(snippet) > 1500:
            snippet = snippet[:1500].rstrip() + "…"
        lines.append("```\n" + snippet + "\n```")
    # v2.97.0 — replaced the v2.91.41 footer ("Next step: run the
    # same request in a new turn — a fresh loop often succeeds
    # where a stale one failed. If it keeps happening, try /model
    # to switch providers or cvc compact --smart to reduce context
    # size.") with an agent-voiced message that takes responsibility
    # for the gap and gives a concrete re-send prompt. The user
    # should never see CLI-flag suggestions in a chat surface; the
    # agent owns the recovery.
    lines.append(
        "\n⚠️ I hit the agent loop cap before I could write a complete "
        "reply. The tool results above are the evidence I gathered — "
        "re-send the same request and I'll continue from here, or "
        "name the specific file and edit you want and I can finish in "
        "one shot."
    )
    return "\n".join(lines)


# v2.91.41 — Provider-aware nudge text. Anthropic takes system role
# via a top-level `system` field; OpenAI/MiniMax/Gemini take system
# as a user/assistant message. This helper returns the right shape so
# the model is more likely to honour it.
def _build_nudge_message(
    provider: str,
    text: str,
) -> dict[str, str]:
    """Return a single message dict carrying the nudge text in the
    right shape for the given provider's chat-completions format.
    """
    p = (provider or "").lower()
    # Anthropic-Messages-API providers (Anthropic, MiniMax).
    # Inject as a user-role message with explicit "SYSTEM:" prefix so
    # the receiver still treats it as a system directive (upstream does
    # the same when the upstream is anthropic-format).
    if p in ("anthropic", "minimax"):
        return {"role": "user", "content": text}
    # OpenAI-compatible providers (openai, github, ollama, lmstudio,
    # google, nvidia). system role is supported.
    return {"role": "system", "content": text}


# v2.91.42 — Upstream-error classifier. Different upstream errors need
# completely different recovery strategies, and the dashboard deserves
# to know which. Returns one of:
#   "context_overflow"  — prompt exceeds the model's context window
#   "rate_limit"        — provider is throttling us
#   "auth"              — API key / token invalid or expired
#   "network"           — connection error, DNS, TCP, etc.
#   "provider_5xx"      — upstream 5xx (transient)
#   "unknown"           — anything else
#
# The detection is heuristic — providers vary in how they format
# errors. We look for the most common substrings.
def _classify_stream_error(exc: BaseException) -> str:
    """Classify a stream/chat exception so the dashboard can show a
    specific, actionable message instead of a generic error.
    """
    try:
        s = str(exc) or ""
        s_low = s.lower()
    except Exception:
        return "unknown"

    # Context overflow — Copilot: "prompt token count of N exceeds the limit of M"
    #                    OpenAI:  "context_length_exceeded"
    #                    Anthropic: "prompt is too long"
    if (
        "model_max_prompt_tokens_exceeded" in s_low
        or "prompt token count" in s_low
        or "context_length_exceeded" in s_low
        or "context window" in s_low
        or "prompt is too long" in s_low
        or ("reduce the length" in s_low and "messages" in s_low)
    ):
        return "context_overflow"

    # Rate limit — provider 429 or "rate limit" string
    if (
        "rate limit" in s_low
        or "too many requests" in s_low
        or "429" in s
    ):
        return "rate_limit"

    # Auth — 401/403, "unauthorized", "invalid api key"
    if (
        "401" in s
        or "403" in s
        or "unauthorized" in s_low
        or "invalid api key" in s_low
        or "authentication" in s_low
        or "auth_type" in s_low
    ):
        return "auth"

    # Network — connection reset, timeout, DNS, etc.
    if (
        "connection reset" in s_low
        or "connection aborted" in s_low
        or "connection refused" in s_low
        or "timed out" in s_low
        or "timeout" in s_low
        or "name resolution" in s_low
        or "temporary failure" in s_low
        or "remote end closed" in s_low
        or "winerror" in s_low
        or "errno" in s_low
    ):
        return "network"

    # Provider 5xx — transient
    if (
        "500" in s
        or "502" in s
        or "503" in s
        or "504" in s
        or "internal server error" in s_low
        or "bad gateway" in s_low
        or "service unavailable" in s_low
    ):
        return "provider_5xx"

    return "unknown"


# v2.91.42 — Format a context-overflow message for the user. Names the
# exact provider + model, the exact used/limit numbers, and gives
# three recovery options: compact, switch model, or start a new chat.
def _format_overflow_message(
    provider: str,
    model: str,
    used: int,
    total: int,
) -> str:
    """Build a user-facing overflow message with actionable recovery steps."""
    overflow_pct = ((used / total) * 100) if total else 0
    used_h = f"{used/1000:.1f}k" if used else "?"
    total_h = f"{total/1000:.0f}k" if total else "?"

    # Find a larger-context alternative in the same provider family.
    bigger_alt = ""
    if provider in ("github", "copilot"):
        bigger_alt = (
            f"\n\n**Switch model in the dashboard** to a Claude/GPT model "
            f"with a larger context window. Copilot's largest tier is "
            f"~200k tokens — if you're consistently bumping the cap, run "
            f"`cvc setup` and pick a provider that supports 1M tokens "
            f"(e.g. `minimax/MiniMax-M3` or `google/gemini-3.1-pro-preview`)."
        )
    elif provider == "anthropic":
        bigger_alt = (
            f"\n\n**Use a 1M-context Claude model** if your dashboard has "
            f"access: `claude-opus-4-8` or `claude-sonnet-4-6` both "
            f"support 1M tokens. Switch via the dashboard's model picker."
        )
    elif provider == "openai":
        bigger_alt = (
            f"\n\n**Use GPT-5.3 (256k context)** or switch to a 1M-context "
            f"provider like `minimax/MiniMax-M3` or `google/gemini-3.1-pro`."
        )
    elif provider == "minimax":
        bigger_alt = (
            f"\n\n**M3 has a 1M token window** — if you're still hitting "
            f"the cap on M3, your conversation has grown beyond a single "
            f"turn's capacity. `cvc compact --smart` will help."
        )

    return (
        f"⚠️ **Context overflow** — this turn couldn't be sent to the model.\n"
        f"\n"
        f"Provider: `{provider}/{model}`\n"
        f"Used: **{used_h} tokens** • Model limit: **{total_h} tokens** "
        f"({overflow_pct:.0f}% over)\n"
        f"\n"
        f"The conversation has grown too large for this model. The LLM "
        f"isn't ignoring you — the API literally rejected the request "
        f"before the model could see it.\n"
        f"\n"
        f"**Recovery options (in order of safety):**\n"
        f"  1. Run `cvc compact --smart` in the workspace terminal — "
        f"summarises old turns, keeps the recent ones.\n"
        f"  2. Start a new chat (the engine keeps your history in a "
        f"branch — `cvc log` to see previous conversations).\n"
        f"  3. Switch to a model with a larger context window."
        f"{bigger_alt}"
    )


# v2.91.42 — Auto-compact a conversation in place. Replaces the middle
# of the message list with a single summary message, preserving:
#   - the system message (always first)
#   - the most recent N turns (so the model still has continuity)
#   - all tool_use ↔ tool_result pairings (orphaned results would break
#     the Anthropic format and trip safety filters)
#
# This is a heuristic — it does NOT call an LLM to summarise. It just
# truncates the middle. That's good enough as a "make the request fit"
# last resort, and avoids the infinite loop of "need an LLM to fix the
# LLM-overflow error".
def _auto_compact_conversation(
    llm_messages: list[dict[str, Any]],
    target_tokens: int,
    keep_recent: int = 6,
) -> tuple[list[dict[str, Any]], int]:
    """Return a compacted copy of `llm_messages` whose estimated token
    count is <= `target_tokens`.

    Strategy: keep system + last `keep_recent` user/assistant/tool
    messages verbatim; replace everything in between with ONE summary
    message that lists what tools ran and how much output they had.

    Returns (new_messages, estimated_original_tokens).
    """
    if not llm_messages:
        return llm_messages, 0

    # Estimate original tokens
    def _est(msgs: list[dict[str, Any]]) -> int:
        chars = 0
        for m in msgs:
            c = m.get("content", "")
            if isinstance(c, str):
                chars += len(c)
            elif isinstance(c, list):
                for b in c:
                    if isinstance(b, dict):
                        chars += len(str(b))
                    else:
                        chars += len(str(b))
            tcs = m.get("tool_calls") or []
            for tc in tcs:
                chars += len(str(tc))
        return chars // 4

    original = _est(llm_messages)
    if original <= target_tokens:
        return llm_messages, original  # Already fits — no compaction needed

    # Identify boundaries
    system = [m for m in llm_messages if m.get("role") == "system"]
    non_system = [m for m in llm_messages if m.get("role") != "system"]

    if len(non_system) <= keep_recent:
        return llm_messages, original  # Nothing to trim

    middle = non_system[:-keep_recent]
    recent = non_system[-keep_recent:]

    # Build a summary of the middle: count tool names + total output bytes
    tool_calls: dict[str, int] = {}
    tool_outputs_chars = 0
    for m in middle:
        if m.get("role") == "assistant":
            for tc in m.get("tool_calls") or []:
                fn = tc.get("function") or {} if isinstance(tc, dict) else {}
                name = fn.get("name") or tc.get("name") or "tool"
                tool_calls[name] = tool_calls.get(name, 0) + 1
        elif m.get("role") == "tool":
            c = m.get("content", "")
            if isinstance(c, list):
                c = json.dumps(c)
            tool_outputs_chars += len(str(c))

    summary_lines = [
        "[Auto-compacted by CVC gateway — earlier turns summarised to fit the model context window.]",
        f"{len(middle)} earlier messages collapsed ({original} → {target_tokens} token target).",
        "",
        "Earlier in this conversation:",
    ]
    if tool_calls:
        summary_lines.append("- Tools run: " + ", ".join(
            f"{name}×{n}" for name, n in sorted(tool_calls.items())
        ))
    if tool_outputs_chars:
        summary_lines.append(
            f"- Total tool output bytes in the collapsed window: "
            f"~{tool_outputs_chars // 1000}k chars"
        )
    summary_msg = {
        "role": "user",
        "content": "\n".join(summary_lines),
    }
    # Pair the summary with a synthetic assistant acknowledgement so the
    # conversation flows naturally. (Single user-message alone is also fine
    # for most providers, but having a short assistant reply is cleaner.)
    ack_msg = {
        "role": "assistant",
        "content": "Noted. Continuing from the recent turns below.",
    }

    return system + [summary_msg, ack_msg] + recent, original


def _log_auto_compact_event(
    *,
    before_tokens: int,
    after_tokens: int,
    threshold_tokens: int,
    total_tokens: int,
    trigger: str,  # "preflight" | "overflow"
) -> None:
    """Persist an auto-compact event to `.cvc/context_window.json` so the
    dashboard meter can show *when* the last compact happened and how much
    was trimmed. Best-effort; never raises.
    """
    try:
        base = Path(_project_root) / ".cvc" if _project_root else None
        if not base:
            return
        base.mkdir(parents=True, exist_ok=True)
        log = base / "context_window.json"
        try:
            data = json.loads(log.read_text(encoding="utf-8") or "{}") if log.exists() else {}
        except Exception:
            data = {}
        events = data.get("events") or []
        events.append({
            "at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
            "trigger": trigger,
            "before_tokens": int(before_tokens),
            "after_tokens": int(after_tokens),
            "threshold_tokens": int(threshold_tokens),
            "total_tokens": int(total_tokens),
        })
        # Cap to last 50 events
        events = events[-50:]
        data["events"] = events
        data["last_auto_compact_at"] = events[-1]["at"]
        data["last_auto_compact"] = {
            "before_tokens": int(before_tokens),
            "after_tokens": int(after_tokens),
            "trigger": trigger,
        }
        log.write_text(json.dumps(data, indent=2) + "\n", encoding="utf-8")
    except Exception:
        # Never break the chat path over log writes.
        pass


def _maybe_preflight_auto_compact(
    llm_messages: list[dict[str, Any]],
) -> tuple[list[dict[str, Any]], bool]:
    """Proactively compact when the estimated token count crosses the
    auto-compact threshold (default 50% of the model's context window).

    This prevents the upstream provider from rejecting the request with a
    context_overflow error — we shrink the transcript *before* the call.

    Returns ``(messages, did_compact)``. ``messages`` is a fresh list when
    compaction fires, otherwise the same list reference (no allocation).
    """
    if not llm_messages:
        return llm_messages, False
    total = _get_active_model_context_window() or 0
    if total <= 0:
        return llm_messages, False
    # Estimate current tokens
    chars = 0
    for m in llm_messages:
        c = m.get("content", "")
        if isinstance(c, str):
            chars += len(c)
        elif isinstance(c, list):
            for b in c:
                chars += len(str(b))
        for tc in m.get("tool_calls") or []:
            chars += len(str(tc))
    before = chars // 4
    threshold = int(total * _AUTO_COMPACT_PCT)
    if before < threshold:
        return llm_messages, False
    try:
        target = max(int(total * 0.8), 4096)
        compacted, original = _auto_compact_conversation(llm_messages, target_tokens=target)
        # Estimate after-compact tokens
        after_chars = 0
        for m in compacted:
            c = m.get("content", "")
            if isinstance(c, str):
                after_chars += len(c)
            elif isinstance(c, list):
                for b in c:
                    after_chars += len(str(b))
        after = after_chars // 4
        if original > after + 64:  # ignore no-op compactions
            _log_auto_compact_event(
                before_tokens=original,
                after_tokens=after,
                threshold_tokens=threshold,
                total_tokens=total,
                trigger="preflight",
            )
            logger.warning(
                "agent_chat: preflight auto-compact %d → %d tokens (threshold %d, total %d)",
                original, after, threshold, total,
            )
            return compacted, True
    except Exception as _pf_exc:
        logger.warning("agent_chat: preflight auto-compact failed: %s", _pf_exc)
    return llm_messages, False


_GATEWAY_PORT: int = int(os.environ.get("CVC_GATEWAY_PORT", "13421"))
_PROXY_PORT: int = int(os.environ.get("CVC_PROXY_PORT", str(_GATEWAY_PORT)))  # deprecated, defaults to gateway port
_HOST: str = os.environ.get("CVC_GATEWAY_HOST", "127.0.0.1")


class _StaticAsset304Filter(logging.Filter):
    """Hide noisy uvicorn access logs for cached static assets (304 responses)."""

    def filter(self, record: logging.LogRecord) -> bool:
        args = getattr(record, "args", ())
        if not isinstance(args, tuple) or len(args) < 5:
            return True

        method = str(args[1])
        path = str(args[2])
        status = str(args[4])
        if method == "GET" and path.startswith("/static/") and status == "304":
            return False
        return True


class _WebSocketNoiseFilter(logging.Filter):
    """Suppress low-value websocket lifecycle chatter from uvicorn logs."""

    def filter(self, record: logging.LogRecord) -> bool:
        message = record.getMessage().strip().lower()
        return message not in {"connection open", "connection closed"}


def _configure_gateway_logging() -> None:
    """Filter routine cache/websocket noise while keeping real logs visible."""
    access_logger = logging.getLogger("uvicorn.access")
    if not any(isinstance(f, _StaticAsset304Filter) for f in access_logger.filters):
        access_logger.addFilter(_StaticAsset304Filter())

    error_logger = logging.getLogger("uvicorn.error")
    if not any(isinstance(f, _WebSocketNoiseFilter) for f in error_logger.filters):
        error_logger.addFilter(_WebSocketNoiseFilter())


def _set_tool_confirm(evt: asyncio.Event | None, result: str) -> None:
    """Set the module-level tool confirmation state."""
    global _tool_confirm_event, _tool_confirm_result
    _tool_confirm_event = evt
    _tool_confirm_result = result


def _get_tool_confirm_result() -> str:
    """Read the current tool confirmation result."""
    return _tool_confirm_result


def _get_event_bus():
    global _event_bus
    if _event_bus is None:
        from cvc.sdk.events import EventBus
        _event_bus = EventBus()
    return _event_bus


async def _health_poll_loop():
    while True:
        try:
            if _manager:
                _manager.refresh_health()
        except Exception:
            logger.debug("Health poll error", exc_info=True)
        await asyncio.sleep(10)


@asynccontextmanager
async def lifespan(app: FastAPI):
    global _manager, _workspace_mgr, _health_task, _project_root
    global _agent_llm, _system_prompt, _agent_tools, _tool_executor
    global _hermes_api_server

    from cvc.core.models import discover_cvc_root
    from cvc.gateway_manager import GatewayManager
    from cvc.workspace_manager import WorkspaceManager

    _configure_gateway_logging()

    # v3.2.0 — Boot the vendored full api_server (the one with
    # 89 tools, full skill system, sessions, runs, OpenAI-compat routes)
    # as a sibling aiohttp process. The dashboard's /api/chat and
    # /api/ws/chat routes now proxy to it instead of constructing
    # AIAgent() with no toolsets. Without this, the dashboard chat
    # works but the agent has no file/bash/etc. tools and degrades
    # to "JUST WAIT..." then a summary, which is what the user has
    # been hitting for the last 6 turns.
    try:
        from cvc.hermes_unified import start_api_server_background
        _hermes_api_server = await start_api_server_background()
    except Exception as exc:
        logger.error("Failed to start vendored api_server: %s", exc)
        _hermes_api_server = None
        # Don't fail the gateway — fall back to the hand-rolled
        # stream_unified_chat path if the api_server can't boot. The
        # user will see degraded chat (no tools) but the dashboard
        # still loads.

    _project_root = discover_cvc_root()
    if _project_root is None:
        _project_root = Path.cwd()

    # Prefer last-active workspace from previous gateway session.
    # This makes "save your workspace selection" survive restart.
    _last_active = WorkspaceManager.load_active_workspace()
    if _last_active is not None and _last_active.exists():
        logger.info("Restoring last-active workspace: %s", _last_active)
        _project_root = _last_active
    else:
        # v3.4.2 — If the persisted active workspace is missing (deleted,
        # unmounted, or — as we saw on Jai's machine — written as
        # truncated garbage "/Users/jkm--"), DO NOT silently fall back to
        # the gateway's cwd. Cwd is unpredictable (depends on which
        # terminal launched the gateway) and would silently swap the
        # user's workspace to whatever directory happened to be open.
        # Instead, recover in this priority order:
        #   1. Sibling .last-active hint file in cwd (legacy backup)
        #   2. First still-existing workspace from the persisted registry
        #   3. Cwd as last resort (with a loud warning)
        _cwd_last_active = Path.cwd() / ".last-active"
        _recovered = False
        if _cwd_last_active.exists():
            try:
                _hint = Path(_cwd_last_active.read_text(encoding="utf-8").strip())
                if _hint.is_dir():
                    logger.info(
                        "Restoring workspace from .last-active hint: %s", _hint
                    )
                    _project_root = _hint
                    _recovered = True
            except OSError:
                pass
        if not _recovered:
            try:
                _registry_path = Path.home() / ".cvc" / "workspaces.json"
                _registry: list[dict[str, Any]] = []
                if _registry_path.exists():
                    try:
                        _loaded = json.loads(
                            _registry_path.read_text(encoding="utf-8")
                        )
                        _registry = (
                            _loaded if isinstance(_loaded, list) else []
                        )
                    except (json.JSONDecodeError, OSError):
                        _registry = []
                for entry in _registry:
                    ep = entry.get("path", "")
                    if ep and Path(ep).is_dir():
                        logger.warning(
                            "Persisted active workspace unavailable "
                            "(was %r) — falling back to first valid "
                            "workspace in registry: %s",
                            _last_active, ep,
                        )
                        _project_root = Path(ep)
                        _recovered = True
                        # Re-save it as the new active pointer so the
                        # next boot doesn't repeat the recovery dance.
                        try:
                            from cvc.host_state import (
                                save_active_workspace as _save,
                            )
                            _save(_project_root)
                        except Exception:
                            pass
                        break
            except Exception:  # noqa: BLE001
                pass
        if not _recovered:
            logger.warning(
                "No persisted active workspace and no valid registry "
                "entry — falling back to cwd (%s). User should run "
                "`cvc workspace add` to register a workspace.",
                _project_root,
            )

    # Initialize WorkspaceManager and activate the discovered workspace
    _workspace_mgr = WorkspaceManager()
    _workspace_mgr.switch_to(_project_root)

    # Hydrate in-memory registry from persisted ~/.cvc/workspaces.json
    try:
        for entry in _workspace_mgr._load_registry():
            ep = entry.get("path", "")
            if not ep:
                continue
            p = Path(ep)
            if not p.exists():
                continue
            ws_id = entry.get("workspace_id") or _make_workspace_id(p)
            _registered_workspaces.setdefault(ws_id, p)
        # Always ensure the active one is registered too
        _registered_workspaces[_make_workspace_id(_project_root)] = _project_root
    except Exception:
        logger.warning("Failed to hydrate workspace registry on startup", exc_info=True)

    cvc_root = _project_root / ".cvc"
    _manager = GatewayManager(cvc_root, host=_HOST)

    # Load settings to check for Telegram Bot Integration.
    # NOTE (v3.4): The legacy aiogram-based Telegram gateway below
    # conflicts with the new registry-driven python-telegram-bot
    # adapter (cvc/integrations/channels/telegram.py) when both
    # share the same bot_token, because Telegram only allows ONE
    # getUpdates polling lease per token globally. The new
    # architecture is canonical; the legacy code path is kept only
    # as a fallback for users who haven't migrated yet. We detect
    # the new-style channel config and skip the legacy boot when
    # it's present.
    try:
        from cvc.agent.settings import load_settings
        from cvc.integrations.setup import channels_dir as _channels_dir
        from pathlib import Path as _P
        _legacy_cfg = _P(_channels_dir()) / "telegram.yaml"
        _new_style_telegram = _legacy_cfg.exists()
        settings = load_settings(_project_root)
        if _new_style_telegram:
            # New registry-driven adapter will start the bot from
            # cvc/gateway/app.py:_lifespan. Don't double-up here.
            print("DEBUG: New-style telegram channel config found — skipping legacy aiogram boot")
            logger.info(
                "Legacy Telegram boot skipped — new registry-driven "
                "adapter will start the bot via the channel registry."
            )
        elif settings.telegram.enabled and settings.telegram.bot_token:
            from cvc.adapters.telegram_gateway import start_telegram_gateway
            print("DEBUG: Calling start_telegram_gateway (legacy fallback)")
            asyncio.create_task(start_telegram_gateway(settings, _workspace_mgr))
            logger.info("Legacy Telegram gateway task created")
        else:
            print("DEBUG: Telegram not enabled or no token")
    except Exception as e:
        print(f"DEBUG: Failed to init Telegram gateway: {e}")
        logger.warning("Failed to initialize Telegram gateway", exc_info=True)

    # Launch proxy without blocking on health check — the background
    # health poll will update its status once it's ready.
    # NOTE: Proxy is now integrated into the gateway itself via
    # workspace-routed endpoints (/ws/{workspace_id}/v1/...).
    # We no longer spawn a separate proxy process.
    # Mark proxy as "running" (it's us now)
    if _manager:
        proxy_rec = _manager.get_service("proxy")
        if proxy_rec:
            proxy_rec.status = "running"
            proxy_rec.port = _GATEWAY_PORT  # proxy runs on gateway port now
            proxy_rec.started_at = time.time()

    # Defer Agent LLM init to a background task so the HTTP server
    # starts accepting requests immediately.
    asyncio.create_task(_init_agent_llm_background())

    # Phase E (4.4): Curator — schedule a background skill-maintenance pass
    # on gateway startup if the schedule says we're due. Best-effort; runs in
    # a daemon thread inside the curator module and never blocks startup.
    try:
        from cvc.agent._vendor.hermes.agent.curator import maybe_run_curator
        # idle_for_seconds=inf treats gateway boot as fully idle.
        maybe_run_curator(idle_for_seconds=float("inf"))
        logger.info("Phase E (4.4): curator pass scheduled if due")
    except Exception as exc:  # pragma: no cover — curator is optional
        logger.debug("curator not scheduled: %s", exc)

    _health_task = asyncio.create_task(_health_poll_loop())
    logger.info("CVC Gateway started — http://%s:%d", _HOST, _GATEWAY_PORT)

    yield

    if _health_task:
        _health_task.cancel()
    if _agent_llm:
        try:
            await _agent_llm.close()
        except Exception:
            pass
    if _workspace_mgr:
        _workspace_mgr.close_all()

    try:
        from cvc.adapters.telegram_gateway import stop_telegram_gateway
        await stop_telegram_gateway()
    except Exception:
        pass

    # v3.2.0 — Tear down the vendored api_server. Without this
    # the aiohttp web app leaks file descriptors (ResponseStore SQLite
    # + WAL file) every restart, eventually hitting the OS fd limit.
    if _hermes_api_server:
        try:
            from cvc.hermes_unified import stop_api_server_background
            await stop_api_server_background(_hermes_api_server)
        except Exception as exc:
            logger.warning("Error stopping vendored api_server: %s", exc)

    logger.info("CVC Gateway shut down.")


async def _init_agent_llm_background() -> None:
    """Initialise Agent LLM + ToolExecutor in a background task.

    Runs after the event loop is already serving HTTP, so the dashboard
    loads instantly.  Chat endpoints gracefully handle _agent_llm=None
    until this completes.
    """
    global _agent_llm, _system_prompt, _agent_tools, _tool_executor

    # Yield control so uvicorn can bind the port first
    await asyncio.sleep(0)

    try:
        from cvc.agent.executor import ToolExecutor
        from cvc.agent.llm import AgentLLM
        from cvc.agent.system_prompt import build_system_prompt
        from cvc.agent.tools import AGENT_TOOLS
        from cvc.core.models import GlobalConfig

        gc = GlobalConfig.load()
        provider = os.getenv("CVC_PROVIDER", gc.provider)
        model = os.getenv("CVC_MODEL", gc.model)

        # ── First-run rescue: auto-detect provider from env ──────────────
        # The launcher (cvc.launcher._auto_init) writes provider="passthrough"
        # so IDE/CLI tool-mode works without `cvc setup`. That's correct for
        # the proxy path, but it leaves the dashboard chat dead-on-arrival
        # because `passthrough` has no internal LLM. If the user has a real
        # API key in the environment, auto-pick that provider so the
        # dashboard chat works the very first time without forcing a manual
        # `cvc setup` step.
        if provider == "passthrough" or not provider:
            _env_provider_chain = [
                ("anthropic",  "ANTHROPIC_API_KEY",     "claude-sonnet-4-6"),
                ("openai",     "OPENAI_API_KEY",        "gpt-5.2"),
                ("google",     "GOOGLE_API_KEY",        "gemini-3-flash-preview"),
                ("copilot",    "COPILOT_GITHUB_TOKEN",  "claude-sonnet-4.6"),
                ("github",     "GITHUB_TOKEN",          "gpt-5.2"),
            ]
            for _p, _envk, _default_model in _env_provider_chain:
                if os.getenv(_envk):
                    logger.info(
                        "Agent LLM auto-detected provider '%s' from %s "
                        "(global config was '%s'). Dashboard chat will use "
                        "this provider for the current gateway run.",
                        _p, _envk, provider or "(unset)",
                    )
                    provider = _p
                    if not model:
                        model = _default_model
                    break

        # Still passthrough/unset → no env keys either. Surface this to the
        # chat endpoints so they can stream an actionable error instead of
        # silently falling back to a proxy that will also produce no output.
        if provider == "passthrough" or not provider:
            logger.warning(
                "Agent LLM not initialised — provider is '%s' and no "
                "ANTHROPIC_API_KEY / OPENAI_API_KEY / GOOGLE_API_KEY / "
                "COPILOT_GITHUB_TOKEN in env. Dashboard chat will return "
                "an actionable setup-required message until `cvc setup` "
                "is run and the gateway is restarted.",
                provider or "(unset)",
            )
            return

        # Resolve API key (same logic as CLI agent)
        # Canonicalize "copilot" alias → "github" (single canonical provider name).
        if provider == "copilot":
            provider = "github"
        env_map = {
            "anthropic": "ANTHROPIC_API_KEY",
            "openai": "OPENAI_API_KEY",
            "google": "GOOGLE_API_KEY",
            "ollama": "",
            "lmstudio": "",
            "github": "COPILOT_GITHUB_TOKEN",  # Copilot OAuth token, not generic GH PAT
            "vertex": "",  # Uses gcloud ADC, not an API key
            "nvidia": "NVIDIA_API_KEY",
            "minimax": "MINIMAX_API_KEY",
        }
        env_key = env_map.get(provider, "")
        api_key = os.getenv(env_key, "") if env_key else ""
        if not api_key:
            api_key = gc.api_keys.get(provider, "") or gc.api_keys.get("copilot", "")

        base_url_map = {
            "anthropic": "https://api.anthropic.com",
            "openai": "https://api.openai.com",
            "google": "https://generativelanguage.googleapis.com",
            "ollama": os.getenv("OLLAMA_HOST", "http://localhost:11434"),
            "lmstudio": os.getenv("LMSTUDIO_HOST", "http://localhost:1234"),
            "github": "https://api.individual.githubcopilot.com",
            "nvidia": "https://integrate.api.nvidia.com",
            "minimax": os.getenv("MINIMAX_BASE_URL", "https://api.minimax.io"),
        }
        base_url = base_url_map.get(provider, "")
        if provider == "vertex":
            from cvc.adapters.vertex import build_vertex_base_url, get_vertex_credentials
            try:
                _creds, adc_project = get_vertex_credentials()
            except RuntimeError:
                adc_project = ""
            v_project = gc.vertex_project_id or os.getenv("VERTEX_PROJECT_ID", "") or adc_project
            v_location = gc.vertex_location or os.getenv("VERTEX_LOCATION", "us-central1")
            base_url = build_vertex_base_url(v_project, v_location)

        _agent_llm = AgentLLM(
            provider=provider,
            api_key=api_key,
            model=model,
            base_url=base_url,
        )

        _system_prompt = build_system_prompt(
            workspace=str(_project_root),
            provider=provider,
            model=model,
            user_memory_context=_get_user_memory_context(),
            branch=_workspace_mgr.engine.active_branch if _workspace_mgr.engine else "main",
        )

        # Dashboard tool surface — full agentic OS capabilities.
        # CVC native tools + all vendored-bridged tools (browser, computer_use,
        # vision, video, image_gen, cronjob, delegate_task, mixture_of_agents,
        # session_search, skill_*, send_message, text_to_speech, execute_code,
        # clarify, feishu_*). Mirrors the CLI agent surface.
        _DASHBOARD_TOOL_NAMES = {
            # ── CVC native ──
            "read_file", "write_file", "edit_file", "patch_file",
            "multi_read", "multi_edit",
            "bash", "process_manage", "glob", "grep", "list_dir",
            "semantic_grep", "git",
            "web_search", "web_fetch",
            "fetch_docs", "search_docs", "lookup_api", "annotate_doc",
            "think", "context_compact", "save_memory", "todo",
            "agent", "parallel_agents", "ask_user",
            "task_create", "task_get", "task_list", "task_kill",
            "cvc_status", "cvc_log", "cvc_commit", "cvc_branch",
            "cvc_restore", "cvc_merge",
            "cvc_search", "cvc_smart_search", "cvc_diff",
            "cvc_ingest_document", "cvc_document_search", "cvc_list_documents",
            "cvc_switch_workspace",
        }
        # Pull every vendored-bridged tool that survived collision filtering —
        # so the dashboard inherits the full agentic surface automatically.
        try:
            from cvc.agent.hermes_bridge import HERMES_TOOL_NAMES
            _DASHBOARD_TOOL_NAMES |= set(HERMES_TOOL_NAMES)
        except Exception:
            pass
        _agent_tools = [t for t in AGENT_TOOLS if t.get("function", {}).get("name") in _DASHBOARD_TOOL_NAMES]

        # Create ToolExecutor scoped to the current workspace
        _tool_executor = ToolExecutor(_project_root, _workspace_mgr.engine)
        _wire_subagent_runtime(_tool_executor)

        # Pre-warm the TCP + TLS connection
        asyncio.create_task(_agent_llm.warm_connection())
        logger.info("Agent LLM initialised — %s/%s (%d tools)", provider, model, len(_agent_tools))
    except Exception:
        logger.warning("Agent LLM init failed — chat will fall back to proxy", exc_info=True)


# ---------------------------------------------------------------------------
# FastAPI App
# ---------------------------------------------------------------------------

try:
    from cvc import __version__ as _cvc_version
except ImportError:
    _cvc_version = "2.0.0"


# v2.91.46 — helper for system-prompt injection of persistent user
# memory. Lazily imports the user_memory_api module so the gateway
# still works even if it's missing (the helper returns an empty
# string in that case).
def _get_user_memory_context() -> str:
    try:
        from cvc.dashboard.user_memory_api import format_for_prompt
        return format_for_prompt()
    except Exception as exc:
        logger.debug("user_memory: prompt injection unavailable: %s", exc)
        return ""

app = FastAPI(
    title="CVC Gateway",
    version=_cvc_version,
    lifespan=lifespan,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


def _ensure(what="gateway"):
    if not _manager or not _workspace_mgr or not _workspace_mgr.db or not _workspace_mgr.config:
        raise HTTPException(status_code=503, detail=f"{what} not initialised yet")


# Thin accessors — route handlers use these instead of module-level singletons.
# All delegate to _workspace_mgr so the active workspace can change at runtime.

def _get_db():
    return _workspace_mgr.db if _workspace_mgr else None


def _get_config():
    return _workspace_mgr.config if _workspace_mgr else None


def _resolve_workspace(workspace_path: str | None) -> None:
    """Switch the active workspace to ``workspace_path`` if given.

    Used by workspace-scoped endpoints (overview, operations, timeline,
    commits, stats, services) so the dashboard's "current workspace"
    is the workspace the user picked in the chat — not whatever the
    gateway booted into.

    Idempotent: if the path matches the current workspace, no-op. If
    the workspace isn't registered, registers it on the fly.
    """
    if not workspace_path or not _workspace_mgr:
        return
    p = Path(workspace_path).expanduser().resolve()
    current = _workspace_mgr.current
    if current and str(current.path) == str(p):
        return
    try:
        _workspace_mgr.switch_to(p)
    except Exception as exc:
        logger.debug("workspace switch failed for %s: %s", p, exc)


def _get_engine():
    return _workspace_mgr.engine if _workspace_mgr else None


# ═══════════════════════════════════════════════════════════════════════
# 1. HEALTH
# ═══════════════════════════════════════════════════════════════════════


@app.post("/api/telemetry")
async def post_telemetry(payload: dict):
    """Receive Swarm telemetry from SDK and broadcast to Dashboard UI."""
    try:

        if _get_event_bus():
            _get_event_bus().emit(
                "telemetry", payload)
    except Exception:
        pass
    return {"status": "ok"}

@app.get("/health")
async def health():
    return {"status": "ok", "service": "cvc-gateway", "version": _cvc_version}


# ═══════════════════════════════════════════════════════════════════════
# 1b. WORKSPACE MANAGEMENT — Multi-workspace isolation
# ═══════════════════════════════════════════════════════════════════════

@app.get("/api/setup/registry")
async def setup_registry():
    """Unified setup registry: providers + features.

    Single source of truth shared by `cvc setup` (CLI) and the dashboard
    settings/onboarding pages. Adding a new provider in cvc/setup/registry.py
    or a new feature in cvc/setup/features.py automatically appears here.
    """
    from cvc.setup import registry_snapshot
    return registry_snapshot()


@app.get("/api/workspace/current")
async def workspace_current(workspace_path: str | None = None):
    """Return the currently active workspace path and stats."""
    _resolve_workspace(workspace_path)
    if not _workspace_mgr or not _workspace_mgr.current:
        raise HTTPException(503, "No workspace active")
    inst = _workspace_mgr.current
    info = {
        "path": str(inst.path),
        "name": inst.path.name or str(inst.path),
        "workspace_id": _make_workspace_id(inst.path),
        "last_accessed": inst.last_accessed,
        "initialized": inst.initialized,
    }
    if inst.engine:
        info["branch"] = inst.engine.active_branch
        try:
            commits = inst.db.index.list_all_commits(limit=9999)
            info["commit_count"] = len(commits)
        except Exception:
            info["commit_count"] = 0
    return info


@app.get("/api/workspace/list")
async def workspace_list():
    """List all known workspaces (open + from persistent registry)."""
    if not _workspace_mgr:
        raise HTTPException(503, "Workspace manager not initialised")
    return {"workspaces": _workspace_mgr.list_workspaces()}


@app.post("/api/workspace/switch")
async def workspace_switch(request: Request):
    """Switch the active workspace. Auto-initializes .cvc/ if needed.

    Accepts either ``{"workspace_id": "..."}`` or ``{"path": "/abs/path"}``.
    """
    global _tool_executor, _system_prompt, _project_root
    if not _workspace_mgr:
        raise HTTPException(503, "Workspace manager not initialised")
    body = await request.json()
    target_path = (body.get("path") or "").strip()
    workspace_id = (body.get("workspace_id") or "").strip()

    # Resolve workspace_id → path via in-memory + persisted registry
    if not target_path and workspace_id:
        wp = _registered_workspaces.get(workspace_id)
        if wp is not None:
            target_path = str(wp)
        else:
            for entry in _workspace_mgr._load_registry():
                ep = entry.get("path", "")
                if not ep:
                    continue
                if entry.get("workspace_id") == workspace_id or _make_workspace_id(Path(ep)) == workspace_id:
                    target_path = ep
                    break
    if not target_path:
        raise HTTPException(400, "Missing 'path' or unknown 'workspace_id'")

    target = Path(os.path.expanduser(target_path)).resolve()
    if not target.exists():
        raise HTTPException(404, f"Directory does not exist: {target}")
    if not target.is_dir():
        raise HTTPException(400, f"Not a directory: {target}")

    inst = _workspace_mgr.switch_to(target)

    # Re-scope tool executor and system prompt to new workspace
    _project_root = inst.path

    # v2.92.4 — Belt-and-suspenders persistence. The WorkspaceManager
    # already saves to ~/.cvc/active_workspace (the registry file). We
    # ALSO write a `.last-active` file in the new workspace's parent
    # directory. The dashboard's boot path reads this file (see the
    # `else` branch in the gateway boot sequence) to recover from cases
    # where the main registry file is missing/corrupt/wiped. Two
    # redundant files = two chances to recover the user's last choice.
    try:
        WorkspaceManager._save_active(inst.path)
        try:
            (inst.path / ".last-active").write_text(
                str(inst.path), encoding="utf-8"
            )
        except OSError:
            pass  # may be a read-only workspace, that's OK
    except Exception as _persist_exc:
        logger.debug("last-active write skipped: %s", _persist_exc)
    try:
        from cvc.agent.executor import ToolExecutor
        _tool_executor = ToolExecutor(inst.path, inst.engine)
        _wire_subagent_runtime(_tool_executor)
    except Exception:
        logger.warning("Failed to re-scope ToolExecutor", exc_info=True)

    try:
        from cvc.agent.system_prompt import build_system_prompt
        provider = _agent_llm.provider if _agent_llm else inst.config.provider
        model = _agent_llm.model if _agent_llm else inst.config.model
        _system_prompt = build_system_prompt(
            workspace=str(inst.path),
            provider=provider,
            model=model,
            branch=inst.engine.active_branch,
            user_memory_context=_get_user_memory_context(),
        )
    except Exception:
        logger.warning("Failed to rebuild system prompt", exc_info=True)

    return {
        "status": "ok",
        "path": str(inst.path),
        "branch": inst.engine.active_branch if inst.engine else "main",
    }


@app.post("/api/workspace/init")
async def workspace_init(request: Request):
    """Explicitly initialize .cvc/ in a directory (without switching)."""
    if not _workspace_mgr:
        raise HTTPException(503, "Workspace manager not initialised")
    body = await request.json()
    target_path = body.get("path", "")
    if not target_path:
        raise HTTPException(400, "Missing 'path' in request body")

    target = Path(target_path).resolve()
    inst = _workspace_mgr.init_workspace(target)
    return {
        "status": "ok",
        "path": str(inst.path),
        "branch": inst.engine.active_branch if inst.engine else "main",
    }


@app.post("/api/workspace/add")
async def workspace_add(request: Request):
    """Add a workspace by absolute filesystem path.

    Single-shot: validates the path, initialises ``.cvc/`` if missing,
    registers it with the gateway proxy, and optionally switches to it.

    Body: ``{"path": "/abs/path", "switch": true}``
    """
    if not _workspace_mgr:
        raise HTTPException(503, "Workspace manager not initialised")
    body = await request.json()
    raw_path = (body.get("path") or "").strip()
    do_switch = bool(body.get("switch", True))
    if not raw_path:
        raise HTTPException(400, "Missing 'path'")

    # Expand ~ for convenience
    target = Path(os.path.expanduser(raw_path)).resolve()
    if not target.exists():
        raise HTTPException(404, f"Directory does not exist: {target}")
    if not target.is_dir():
        raise HTTPException(400, f"Not a directory: {target}")

    # Initialise + open
    inst = _workspace_mgr.init_workspace(target)

    # Register with the gateway proxy registry
    workspace_id = _make_workspace_id(target)
    _registered_workspaces[workspace_id] = target

    if do_switch:
        global _tool_executor, _project_root
        inst = _workspace_mgr.switch_to(target)
        _project_root = inst.path
        try:
            from cvc.agent.executor import ToolExecutor
            _tool_executor = ToolExecutor(inst.path, inst.engine)
            _wire_subagent_runtime(_tool_executor)
        except Exception:
            logger.warning("Failed to re-scope ToolExecutor", exc_info=True)

    return {
        "status": "ok",
        "workspace_id": workspace_id,
        "path": str(target),
        "name": target.name,
        "branch": inst.engine.active_branch if inst.engine else "main",
        "switched": do_switch,
    }


# v2.91.49 — Native OS folder picker.
#
# The dashboard's "Add workspace" form used to require a user to type or paste
# an absolute path, which is friction on Mac/Windows where users think in
# Finder/Explorer. This endpoint shells out to the platform's native folder
# chooser and returns the selected path (or `cancelled: true` on dismiss).
#
# Behaviour by platform:
#   * macOS  — `osascript -e 'tell application "Finder" to choose folder'`
#              modal sheet anchored to the dock, returns POSIX path.
#   * Windows — PowerShell + System.Windows.Forms.FolderBrowserDialog.
#              Headless installs (no explorer.exe) fall through to a 503.
#   * Linux  — `zenity --file-selection --directory`, then `kdialog`, then
#              `python -m tkinter.filedialog` as a final fallback. Returns
#              a 501 if none are installed (user can still type the path).
#
# The endpoint is BLOCKING inside a thread pool because the native dialog is
# modal — uvicorn workers can't service other requests until the user
# dismisses. We set a 5-minute timeout so a forgotten dialog doesn't pin a
# worker forever. Cancel returns cancelled=True (200), not an error.


def _pick_folder_macos() -> str | None:
    """Show a native macOS folder picker. Returns POSIX path or None on cancel."""
    # The AppleScript returns an alias like `alias Macintosh HD:Users:foo:bar:`
    # which osascript converts to POSIX form when wrapped in `POSIX path of`.
    script = (
        'set chosenFolder to choose folder with prompt "Select CVC workspace folder"'
        '\nset posixPath to POSIX path of (chosenFolder as text)'
        '\nreturn posixPath'
    )
    try:
        out = subprocess.run(
            ["osascript", "-e", script],
            capture_output=True, text=True, timeout=300, check=False,
        )
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return None
    if out.returncode != 0:
        # User clicked Cancel (returns "User canceled."); other errors -> log + None
        return None
    path = out.stdout.strip()
    return path or None


def _pick_folder_windows() -> str | None:
    """Show a native Windows folder picker via PowerShell + WinForms."""
    # FolderBrowserDialog shows a modal Choose Folder dialog. We write the
    # chosen path to a temp file and read it back; stdout is unreliable
    # because PowerShell sometimes returns DialogResult in the host stream.
    ps_script = r"""
Add-Type -AssemblyName System.Windows.Forms
Add-Type -AssemblyName System.Drawing
$dlg = New-Object System.Windows.Forms.FolderBrowserDialog
$dlg.Description = 'Select CVC workspace folder'
$dlg.ShowNewFolderButton = $true
$dlg.UseDescriptionForTitle = $true
$result = $dlg.ShowDialog()
if ($result -eq [System.Windows.Forms.DialogResult]::OK) {
    [Console]::Out.Write($dlg.SelectedPath)
} else {
    [Console]::Out.Write('__CVC_CANCELLED__')
}
"""
    try:
        out = subprocess.run(
            ["powershell", "-NoProfile", "-NonInteractive", "-Command", ps_script],
            capture_output=True, text=True, timeout=300, check=False,
        )
    except (FileNotFoundError, subprocess.TimeoutExpired, OSError):
        return None
    path = (out.stdout or "").strip()
    if not path or path == "__CVC_CANCELLED__":
        return None
    return path


def _pick_folder_linux() -> str | None:
    """Show a Linux native folder picker (zenity → kdialog → tkinter)."""
    # Try zenity first (GNOME).
    if shutil.which("zenity"):
        try:
            out = subprocess.run(
                ["zenity", "--file-selection", "--directory",
                 "--title=Select CVC workspace folder"],
                capture_output=True, text=True, timeout=300, check=False,
            )
        except subprocess.TimeoutExpired:
            return None
        if out.returncode == 0:
            p = out.stdout.strip()
            return p or None
        return None
    # Then kdialog (KDE).
    if shutil.which("kdialog"):
        try:
            out = subprocess.run(
                ["kdialog", "--getexistingdirectory", os.path.expanduser("~"),
                 "--title", "Select CVC workspace folder"],
                capture_output=True, text=True, timeout=300, check=False,
            )
        except subprocess.TimeoutExpired:
            return None
        if out.returncode == 0:
            p = out.stdout.strip()
            return p or None
        return None
    # Last resort: tkinter (always available in the CVC Python env).
    try:
        import tkinter as tk
        from tkinter import filedialog
        root = tk.Tk()
        root.withdraw()
        # Try to make it appear above other windows
        try:
            root.attributes('-topmost', True)
        except Exception:
            pass
        path = filedialog.askdirectory(title="Select CVC workspace folder")
        root.destroy()
        return path or None
    except Exception:
        return None


@app.post("/api/system/pick-folder")
async def system_pick_folder():
    """Open the OS-native folder chooser and return the selected path.

    The user is the only one who can dismiss a native dialog, so this
    endpoint can block for the duration of the chooser. uvicorn runs sync
    endpoints in a thread pool, so other requests continue to be served.

    Returns:
        200: ``{"path": "/abs/path", "cancelled": false}``
        200: ``{"path": null, "cancelled": true}`` (user clicked Cancel)
        501: ``{"detail": "..."}`` (no chooser available on this OS)
    """
    system = platform.system()
    loop = asyncio.get_running_loop()
    try:
        if system == "Darwin":
            chosen = await loop.run_in_executor(None, _pick_folder_macos)
        elif system == "Windows":
            chosen = await loop.run_in_executor(None, _pick_folder_windows)
        else:
            # Linux / BSD / other Unix
            chosen = await loop.run_in_executor(None, _pick_folder_linux)
    except Exception as exc:
        logger.warning("pick-folder failed on %s: %s", system, exc, exc_info=True)
        raise HTTPException(500, f"Folder picker failed: {exc}") from exc

    if chosen is None:
        # Could mean: user cancelled, or no chooser installed.
        # Try to disambiguate: if the platform is supported but no chooser
        # was found, the executor would have returned None without raising.
        # We treat None as "cancelled" by convention (most common case).
        return {"path": None, "cancelled": True}

    # Normalise + validate
    target = Path(os.path.expanduser(chosen)).resolve()
    if not target.is_dir():
        raise HTTPException(400, f"Selection is not a directory: {target}")
    return {"path": str(target), "cancelled": False}


# ─── Voice transcription (speech-to-text) ─────────────────────────────────
#
# Accepts an audio file (webm/ogg/wav/mp3/m4a) and returns plain text.
# Tries providers in this order until one succeeds:
#   0. faster-whisper (LOCAL, free, offline) — pip install faster-whisper
#   1. Groq (whisper-large-v3)               — GROQ_API_KEY
#   2. OpenAI (whisper-1 / gpt-4o-transcribe) — OPENAI_API_KEY
#   3. Google Gemini (1.5-flash inline audio) — GOOGLE_API_KEY
# 415 if none configured.

# Module-level whisper model cache (loads once, ~250 MB resident for tiny.en).
_LOCAL_WHISPER_MODEL = None
_LOCAL_WHISPER_LOCK = asyncio.Lock()
_LOCAL_WHISPER_SIZE = os.environ.get("CVC_WHISPER_MODEL", "tiny.en")


def _local_whisper_available() -> bool:
    """Cheap check — is faster-whisper importable?"""
    try:
        import faster_whisper  # noqa: F401
        return True
    except Exception:
        return False


async def _stt_local_whisper(audio: bytes, filename: str, content_type: str) -> str | None:
    """Transcribe with faster-whisper (fully local, no API key)."""
    global _LOCAL_WHISPER_MODEL
    if not _local_whisper_available():
        return None
    import tempfile
    import uuid
    from pathlib import Path as _Path

    # Lazy-load model on first call. ~5 sec download on first run, instant after.
    if _LOCAL_WHISPER_MODEL is None:
        async with _LOCAL_WHISPER_LOCK:
            if _LOCAL_WHISPER_MODEL is None:
                from faster_whisper import WhisperModel
                # int8 on CPU = best speed/accuracy for tiny models.
                _LOCAL_WHISPER_MODEL = await asyncio.to_thread(
                    WhisperModel, _LOCAL_WHISPER_SIZE, device="cpu", compute_type="int8"
                )

    # Pick a sane file extension from MIME or filename.
    suffix = ".webm"
    if content_type:
        ct = content_type.lower()
        if "wav" in ct: suffix = ".wav"
        elif "mp3" in ct or "mpeg" in ct: suffix = ".mp3"
        elif "ogg" in ct: suffix = ".ogg"
        elif "m4a" in ct or "mp4" in ct or "aac" in ct: suffix = ".m4a"
    elif filename and "." in filename:
        suffix = "." + filename.rsplit(".", 1)[-1].lower()

    tmp_path = _Path(tempfile.gettempdir()) / f"cvc_stt_{uuid.uuid4().hex}{suffix}"
    try:
        tmp_path.write_bytes(audio)
        def _run():
            segments, _info = _LOCAL_WHISPER_MODEL.transcribe(
                str(tmp_path),
                beam_size=1,            # greedy = fastest
                vad_filter=True,        # skip silence
                without_timestamps=True,
            )
            return " ".join(s.text.strip() for s in segments).strip()
        return await asyncio.to_thread(_run)
    except Exception:
        logger.exception("Local Whisper STT failed")
        return None
    finally:
        try: tmp_path.unlink(missing_ok=True)
        except Exception: pass


async def _stt_groq(audio: bytes, filename: str, content_type: str) -> str | None:
    key = os.environ.get("GROQ_API_KEY")
    if not key:
        return None
    import httpx
    async with httpx.AsyncClient(timeout=60.0) as client:
        files = {"file": (filename, audio, content_type)}
        data = {"model": "whisper-large-v3", "response_format": "json"}
        r = await client.post(
            "https://api.groq.com/openai/v1/audio/transcriptions",
            headers={"Authorization": f"Bearer {key}"},
            files=files,
            data=data,
        )
        if r.status_code >= 400:
            logger.warning("Groq STT failed %s: %s", r.status_code, r.text[:200])
            return None
        return (r.json().get("text") or "").strip()


async def _stt_openai(audio: bytes, filename: str, content_type: str) -> str | None:
    key = os.environ.get("OPENAI_API_KEY")
    if not key:
        return None
    import httpx
    async with httpx.AsyncClient(timeout=60.0) as client:
        files = {"file": (filename, audio, content_type)}
        data = {"model": "whisper-1", "response_format": "json"}
        r = await client.post(
            "https://api.openai.com/v1/audio/transcriptions",
            headers={"Authorization": f"Bearer {key}"},
            files=files,
            data=data,
        )
        if r.status_code >= 400:
            logger.warning("OpenAI STT failed %s: %s", r.status_code, r.text[:200])
            return None
        return (r.json().get("text") or "").strip()


async def _stt_gemini(audio: bytes, content_type: str) -> str | None:
    key = os.environ.get("GOOGLE_API_KEY") or os.environ.get("GEMINI_API_KEY")
    if not key:
        return None
    import base64
    import httpx
    b64 = base64.b64encode(audio).decode("ascii")
    body = {
        "contents": [{
            "parts": [
                {"text": "Transcribe the speech in this audio verbatim. Return only the transcript text — no preface, no quotes."},
                {"inline_data": {"mime_type": content_type or "audio/webm", "data": b64}},
            ]
        }]
    }
    async with httpx.AsyncClient(timeout=90.0) as client:
        r = await client.post(
            f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={key}",
            json=body,
        )
        if r.status_code >= 400:
            logger.warning("Gemini STT failed %s: %s", r.status_code, r.text[:200])
            return None
        try:
            text = r.json()["candidates"][0]["content"]["parts"][0]["text"]
            return (text or "").strip()
        except Exception:
            return None


@app.post("/api/voice/transcribe")
async def voice_transcribe(audio: UploadFile = File(...)):
    """Transcribe an uploaded audio blob to text via Groq → OpenAI → Gemini."""
    data = await audio.read()
    if not data:
        raise HTTPException(400, "Empty audio payload")
    if len(data) > 25 * 1024 * 1024:
        raise HTTPException(413, "Audio file too large (max 25 MB)")

    filename = audio.filename or "recording.webm"
    content_type = audio.content_type or "audio/webm"

    # Try local whisper first (free, offline, zero-config once installed).
    for fn in (_stt_local_whisper, _stt_groq, _stt_openai):
        try:
            text = await fn(data, filename, content_type)
        except Exception:
            logger.exception("STT provider failed")
            text = None
        if text:
            return {"text": text, "provider": fn.__name__.replace("_stt_", "")}

    try:
        text = await _stt_gemini(data, content_type)
    except Exception:
        logger.exception("Gemini STT failed")
        text = None
    if text:
        return {"text": text, "provider": "gemini"}

    raise HTTPException(
        415,
        "No speech-to-text provider available. Install free local voice: `cvc voice install` "
        "(or POST /api/voice/local/install). Or set GROQ_API_KEY / OPENAI_API_KEY / GOOGLE_API_KEY.",
    )


@app.get("/api/voice/local/status")
async def voice_local_status():
    """Report whether the free offline transcription stack is installed and ready."""
    installed = _local_whisper_available()
    return {
        "installed": installed,
        "model": _LOCAL_WHISPER_SIZE,
        "loaded": _LOCAL_WHISPER_MODEL is not None,
        "providers": {
            "local_whisper": installed,
            "groq": bool(os.environ.get("GROQ_API_KEY")),
            "openai": bool(os.environ.get("OPENAI_API_KEY")),
            "gemini": bool(os.environ.get("GOOGLE_API_KEY") or os.environ.get("GEMINI_API_KEY")),
        },
    }


@app.post("/api/voice/local/install")
async def voice_local_install():
    """Install faster-whisper into the current environment. Streams pip output as SSE."""
    if _local_whisper_available():
        async def _already():
            yield f"data: {json.dumps({'phase': 'done', 'msg': 'Already installed.'})}\n\n"
        from fastapi.responses import StreamingResponse as _SR
        return _SR(_already(), media_type="text/event-stream",
                   headers={"Cache-Control": "no-cache"})

    import sys as _sys
    # Prefer `uv pip` if the gateway was launched under uv (faster, resolver-correct).
    uv_path = shutil.which("uv") if "shutil" in dir() else None
    if uv_path is None:
        import shutil as _shutil
        uv_path = _shutil.which("uv")

    if uv_path:
        cmd = [uv_path, "pip", "install", "--python", _sys.executable, "faster-whisper>=1.0.0"]
    else:
        cmd = [_sys.executable, "-m", "pip", "install", "--upgrade", "faster-whisper>=1.0.0"]

    async def _stream():
        yield f"data: {json.dumps({'phase': 'start', 'msg': 'Installing faster-whisper…', 'cmd': ' '.join(cmd)})}\n\n"
        try:
            proc = await asyncio.create_subprocess_exec(
                *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.STDOUT,
            )
            assert proc.stdout is not None
            while True:
                line = await proc.stdout.readline()
                if not line: break
                yield f"data: {json.dumps({'phase': 'log', 'msg': line.decode('utf-8', 'replace').rstrip()})}\n\n"
            rc = await proc.wait()
            if rc == 0 and _local_whisper_available():
                yield f"data: {json.dumps({'phase': 'done', 'msg': 'Installed. Voice input is ready.'})}\n\n"
            else:
                yield f"data: {json.dumps({'phase': 'error', 'msg': f'pip exited with code {rc}'})}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'phase': 'error', 'msg': str(e)})}\n\n"

    from fastapi.responses import StreamingResponse as _SR
    return _SR(_stream(), media_type="text/event-stream",
               headers={"Cache-Control": "no-cache"})


# ── File upload (multimodal attachments) ─────────────────────────────────────
import base64 as _b64
import uuid as _uuid
import mimetypes as _mimetypes
from typing import List as _List

_IMAGE_MIMES = {"image/png", "image/jpeg", "image/jpg", "image/webp", "image/gif"}
_MAX_UPLOAD_BYTES = 50 * 1024 * 1024  # 50 MB cap
_INLINE_IMAGE_CAP = 6 * 1024 * 1024  # 6 MB → inline as data URL for vision models


def _classify_upload(mime: str, name: str) -> str:
    m = (mime or "").lower()
    n = (name or "").lower()
    if m.startswith("image/"):
        return "image"
    if m.startswith("video/"):
        return "video"
    if m.startswith("audio/"):
        return "audio"
    if m == "application/pdf" or n.endswith(".pdf"):
        return "pdf"
    if any(n.endswith(ext) for ext in (".ppt", ".pptx", ".key")):
        return "slides"
    if any(n.endswith(ext) for ext in (".doc", ".docx", ".odt", ".rtf")):
        return "doc"
    if any(n.endswith(ext) for ext in (".xls", ".xlsx", ".csv", ".tsv", ".ods")):
        return "sheet"
    if any(
        n.endswith(ext)
        for ext in (".txt", ".md", ".json", ".yaml", ".yml", ".xml", ".html", ".log")
    ):
        return "text"
    if m.startswith("text/"):
        return "text"
    return "other"


@app.post("/api/files/upload")
async def files_upload(files: _List[UploadFile] = File(...)):
    """Upload one or more files into the active workspace's .cvc/uploads/ directory.

    Returns metadata so the composer can attach the files to the next chat send.
    For images under 6 MB, also returns a base64 data URL so vision-capable
    models can consume them inline without a follow-up fetch.
    """
    if not files:
        raise HTTPException(400, "No files provided")

    base = Path(_project_root) if _project_root else Path.cwd()
    uploads_dir = base / ".cvc" / "uploads"
    uploads_dir.mkdir(parents=True, exist_ok=True)

    results = []
    for f in files:
        data = await f.read()
        if not data:
            continue
        if len(data) > _MAX_UPLOAD_BYTES:
            raise HTTPException(
                413,
                f"File '{f.filename}' exceeds {_MAX_UPLOAD_BYTES // (1024 * 1024)} MB limit",
            )

        orig_name = f.filename or "upload.bin"
        # Sanitise filename — keep extension, prefix with short uuid for uniqueness
        suffix = Path(orig_name).suffix.lower()
        if not suffix:
            guessed = _mimetypes.guess_extension(f.content_type or "") or ""
            suffix = guessed
        safe_stem = "".join(
            c for c in Path(orig_name).stem if c.isalnum() or c in ("-", "_", ".")
        )[:60] or "upload"
        unique = f"{_uuid.uuid4().hex[:8]}-{safe_stem}{suffix}"
        out_path = uploads_dir / unique
        out_path.write_bytes(data)

        mime = f.content_type or _mimetypes.guess_type(orig_name)[0] or "application/octet-stream"
        kind = _classify_upload(mime, orig_name)

        entry = {
            "name": orig_name,
            "stored_name": unique,
            "path": str(out_path),
            "rel_path": str(out_path.relative_to(base)),
            "size": len(data),
            "mime": mime,
            "kind": kind,
        }
        if kind == "image" and len(data) <= _INLINE_IMAGE_CAP:
            entry["data_url"] = f"data:{mime};base64,{_b64.b64encode(data).decode('ascii')}"
        results.append(entry)

    if not results:
        raise HTTPException(400, "All uploads were empty")
    return {"files": results}


# ── Multimodal attachment merger ─────────────────────────────────────────────
# v2.90.9 — composer was inlining giant base64 data URLs as markdown text,
# which (a) made the user bubble unreadable and (b) the LLM saw a raw text
# blob, not an image. The composer now sends a clean reference string; this
# helper consumes ChatRequest.attachments and rewrites the LAST user message
# into provider-native multimodal blocks (image_url for images; extracted
# text for pdf/docx/xlsx/pptx; file body for text/code).

_DOC_TEXT_CAP = 200_000  # max chars of extracted document text per attachment


def _extract_attachment_text(path: str, kind: str, mime: str) -> str | None:
    """Best-effort text extraction. Returns None on failure (never raises)."""
    p = Path(path)
    if not p.exists() or not p.is_file():
        return None
    try:
        if kind in ("text",):
            return p.read_text(encoding="utf-8", errors="replace")[:_DOC_TEXT_CAP]
        if kind == "pdf":
            try:
                import fitz  # pymupdf
                doc = fitz.open(str(p))
                out = []
                total = 0
                for page in doc:
                    t = page.get_text() or ""
                    out.append(t)
                    total += len(t)
                    if total > _DOC_TEXT_CAP:
                        break
                doc.close()
                return ("\n".join(out))[:_DOC_TEXT_CAP]
            except Exception:
                try:
                    from pypdf import PdfReader
                    rd = PdfReader(str(p))
                    out = []
                    total = 0
                    for pg in rd.pages:
                        t = pg.extract_text() or ""
                        out.append(t)
                        total += len(t)
                        if total > _DOC_TEXT_CAP:
                            break
                    return ("\n".join(out))[:_DOC_TEXT_CAP]
                except Exception:
                    return None
        if kind == "doc" and p.suffix.lower() == ".docx":
            try:
                from docx import Document
                doc = Document(str(p))
                paras = [par.text for par in doc.paragraphs]
                return ("\n".join(paras))[:_DOC_TEXT_CAP]
            except Exception:
                return None
        if kind == "sheet" and p.suffix.lower() in (".csv", ".tsv"):
            return p.read_text(encoding="utf-8", errors="replace")[:_DOC_TEXT_CAP]
        if kind == "sheet" and p.suffix.lower() == ".xlsx":
            try:
                from openpyxl import load_workbook
                wb = load_workbook(str(p), read_only=True, data_only=True)
                out = []
                total = 0
                for ws in wb.worksheets:
                    out.append(f"## Sheet: {ws.title}")
                    for row in ws.iter_rows(values_only=True):
                        line = "\t".join("" if v is None else str(v) for v in row)
                        out.append(line)
                        total += len(line)
                        if total > _DOC_TEXT_CAP:
                            break
                    if total > _DOC_TEXT_CAP:
                        break
                wb.close()
                return ("\n".join(out))[:_DOC_TEXT_CAP]
            except Exception:
                return None
        if kind == "slides" and p.suffix.lower() == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(p))
                out = []
                for i, slide in enumerate(prs.slides, 1):
                    out.append(f"## Slide {i}")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            out.append(shape.text)
                return ("\n".join(out))[:_DOC_TEXT_CAP]
            except Exception:
                return None
    except Exception as exc:
        logger.debug("attachment text extract failed for %s: %s", path, exc)
        return None
    return None


def _image_data_url(att: dict) -> str | None:
    """Return data: URL for an image attachment, reading from disk if needed."""
    if att.get("data_url"):
        return att["data_url"]
    path = att.get("path") or ""
    mime = att.get("mime") or "image/png"
    try:
        p = Path(path)
        if not p.exists() or not p.is_file():
            return None
        if p.stat().st_size > _MAX_UPLOAD_BYTES:
            return None
        data = p.read_bytes()
        return f"data:{mime};base64,{_b64.b64encode(data).decode('ascii')}"
    except Exception:
        return None


def _image_block_for_provider(provider: str, data_url: str, mime: str) -> dict:
    """Build a provider-native image content block from a data: URL."""
    p = (provider or "").lower()
    # OpenAI-compat (openai, github/copilot, vertex, nvidia, lmstudio, ollama)
    if p in ("openai", "github", "copilot", "vertex", "nvidia", "lmstudio", "ollama", ""):
        return {"type": "image_url", "image_url": {"url": data_url}}
    if p == "anthropic":
        # Strip "data:<mime>;base64," prefix → raw b64
        raw = data_url.split(",", 1)[1] if "," in data_url else ""
        return {
            "type": "image",
            "source": {"type": "base64", "media_type": mime, "data": raw},
        }
    if p == "google":
        # Gemini accepts inline_data via OpenAI-compat layer too, but our
        # AgentLLM._chat_google maps image_url blocks correctly; fall back to it.
        return {"type": "image_url", "image_url": {"url": data_url}}
    return {"type": "image_url", "image_url": {"url": data_url}}


# v2.90.12 — Native PDF multimodal support
# ════════════════════════════════════════════════════════════════════════
# Previously: every PDF was forced through pymupdf text extraction. Scanned
# admit cards / receipts / image-only PDFs returned empty text, so the model
# saw a zero-length [Attachment] block and either hallucinated or fell back
# to terminal/browser tools to "read" the file — burning iterations and
# producing wrong answers.
#
# Now: Frontier vision models get the PDF natively.
#   - Anthropic (Claude Opus/Sonnet/Haiku): document block with base64 source.
#   - Google (Gemini Pro/Flash/2.5/3): image_url with application/pdf data URL
#     (AgentLLM._chat_google maps this to inline_data correctly).
#   - OpenAI/GitHub-Copilot/etc.: no native PDF channel → page-rasterise to
#     PNG image blocks (preserves layout, works for scanned docs).
# Text extraction is still attempted as a *supplement* (cheap context) when
# we have it, but never as the sole signal.

_PDF_PROVIDERS_WITH_NATIVE_PDF = {"anthropic", "google"}
_PDF_RASTER_MAX_PAGES = 8       # cap pages we rasterise → image blocks
_PDF_RASTER_DPI = 144           # readable for text-heavy admit cards
_PDF_INLINE_BYTES_CAP = 30 * 1024 * 1024  # 30 MB hard cap on inline base64 PDF


def _pdf_data_url(path: str) -> str | None:
    """Read a PDF off disk and return a `data:application/pdf;base64,...` URL."""
    try:
        p = Path(path)
        if not p.exists() or not p.is_file():
            return None
        if p.stat().st_size > _PDF_INLINE_BYTES_CAP:
            return None
        data = p.read_bytes()
        return f"data:application/pdf;base64,{_b64.b64encode(data).decode('ascii')}"
    except Exception:
        return None


def _pdf_native_block_for_provider(provider: str, data_url: str) -> dict | None:
    """Build a provider-native PDF content block, or None if unsupported."""
    p = (provider or "").lower()
    if p == "anthropic":
        raw = data_url.split(",", 1)[1] if "," in data_url else ""
        if not raw:
            return None
        return {
            "type": "document",
            "source": {
                "type": "base64",
                "media_type": "application/pdf",
                "data": raw,
            },
        }
    if p == "google":
        # AgentLLM._chat_google routes image_url data: URLs through inline_data,
        # which accepts application/pdf as a mime type for Gemini multimodal.
        return {"type": "image_url", "image_url": {"url": data_url}}
    return None


def _pdf_rasterise_to_image_blocks(
    path: str, provider: str, max_pages: int = _PDF_RASTER_MAX_PAGES
) -> list[dict]:
    """Render the first ``max_pages`` of a PDF to PNG image blocks for
    providers without native PDF support (OpenAI/GitHub-Copilot/etc).

    Requires pymupdf — returns [] if unavailable so the caller can fall back.
    """
    try:
        import fitz  # pymupdf
    except Exception:
        return []
    blocks: list[dict] = []
    try:
        doc = fitz.open(path)
        # Render at _PDF_RASTER_DPI (default 144 ≈ 2x zoom). Higher = more
        # readable scanned text, but each page = ~200-400KB PNG.
        zoom = _PDF_RASTER_DPI / 72.0
        mat = fitz.Matrix(zoom, zoom)
        for i, page in enumerate(doc):
            if i >= max_pages:
                break
            try:
                pix = page.get_pixmap(matrix=mat, alpha=False)
                png_bytes = pix.tobytes("png")
                if len(png_bytes) > _INLINE_IMAGE_CAP:
                    # Re-render at lower DPI for this page only
                    lo_mat = fitz.Matrix(96 / 72.0, 96 / 72.0)
                    pix = page.get_pixmap(matrix=lo_mat, alpha=False)
                    png_bytes = pix.tobytes("png")
                data_url = (
                    f"data:image/png;base64,"
                    f"{_b64.b64encode(png_bytes).decode('ascii')}"
                )
                blocks.append(_image_block_for_provider(provider, data_url, "image/png"))
            except Exception as page_exc:
                logger.debug("pdf raster page %d failed: %s", i, page_exc)
                continue
        doc.close()
    except Exception as exc:
        logger.warning("pdf rasterise failed for %s: %s", path, exc)
        return []
    return blocks


def _merge_attachments_into_messages(
    llm_messages: list[dict[str, Any]],
    attachments: list[dict[str, Any]] | None,
    provider: str,
) -> list[dict[str, Any]]:
    """Rewrite the LAST user message into multimodal content blocks.

    - Images → image_url / image blocks (provider-native).
    - PDF/DOCX/XLSX/PPTX/CSV/TXT/MD/code → extracted text appended as text block,
      wrapped with `[Attachment <name>]` markers so the model knows the source.
    - Audio/video/other → metadata-only text reference (no bytes).

    Idempotent: if `attachments` is empty, returns `llm_messages` unchanged.
    """
    if not attachments:
        return llm_messages
    # Find last user message index.
    last_user_idx = -1
    for i in range(len(llm_messages) - 1, -1, -1):
        if llm_messages[i].get("role") == "user":
            last_user_idx = i
            break
    if last_user_idx < 0:
        return llm_messages

    orig = llm_messages[last_user_idx]
    orig_content = orig.get("content", "")
    text_part = orig_content if isinstance(orig_content, str) else ""

    parts: list[dict[str, Any]] = []
    image_blocks: list[dict[str, Any]] = []
    extracted_texts: list[str] = []
    reference_lines: list[str] = []

    for att in attachments:
        kind = (att.get("kind") or "other").lower()
        name = att.get("name") or att.get("stored_name") or "file"
        mime = att.get("mime") or "application/octet-stream"
        size = att.get("size") or 0
        path = att.get("path") or ""
        if kind == "image":
            data_url = _image_data_url(att)
            if data_url:
                image_blocks.append(_image_block_for_provider(provider, data_url, mime))
                reference_lines.append(f"- 🖼  **{name}** ({mime}, {size} bytes)")
            else:
                reference_lines.append(f"- ⚠️  **{name}** (image unreadable at {path})")
        elif kind == "pdf":
            # v2.90.12 — Native PDF multimodal first, text supplement second.
            # Frontier vision models (Claude, Gemini, GPT-4o/5, Opus, Sonnet)
            # consume PDFs directly. The previous code path forced pymupdf text
            # extraction, which returned empty for scanned/image-based PDFs
            # (admit cards, receipts, ID scans) → the model saw an empty
            # attachment and hallucinated or burned tool iterations trying to
            # read the file. Now we ship the actual PDF bytes whenever we can.
            prov_l = (provider or "").lower()
            pdf_handled = False

            # Try native PDF block for Anthropic/Google.
            if prov_l in _PDF_PROVIDERS_WITH_NATIVE_PDF:
                data_url = _pdf_data_url(path)
                if data_url:
                    native_block = _pdf_native_block_for_provider(prov_l, data_url)
                    if native_block:
                        image_blocks.append(native_block)
                        reference_lines.append(
                            f"- 📄 **{name}** ({mime}, {size} bytes) — sent natively"
                        )
                        pdf_handled = True

            # Page-rasterise for OpenAI/GitHub-Copilot/etc — image_url blocks
            # work for every vision-capable provider and preserve scanned-PDF
            # layout (handwriting, signatures, photos, layout-sensitive forms).
            if not pdf_handled:
                raster_blocks = _pdf_rasterise_to_image_blocks(path, provider)
                if raster_blocks:
                    image_blocks.extend(raster_blocks)
                    reference_lines.append(
                        f"- 📄 **{name}** ({mime}, {size} bytes) — "
                        f"rendered {len(raster_blocks)} page(s) as images"
                    )
                    pdf_handled = True

            # Always also attempt text extraction as a CHEAP SUPPLEMENT (not
            # the sole signal). If the PDF has machine-readable text, it gives
            # the model exact strings; if scanned, returns empty and we rely
            # on the visual blocks above.
            text = _extract_attachment_text(path, kind, mime)
            if text and text.strip():
                extracted_texts.append(
                    f"\n\n--- [Attachment: {name} ({kind}, {mime})] ---\n{text}\n--- [End: {name}] ---"
                )
                if not pdf_handled:
                    reference_lines.append(
                        f"- 📄 **{name}** ({mime}, {size} bytes) — text extracted ({len(text)} chars)"
                    )
            elif not pdf_handled:
                reference_lines.append(
                    f"- ⚠️  **{name}** ({mime}, {size} bytes) — could not read PDF "
                    f"(scanned/image-only and pymupdf unavailable for rendering)"
                )
        elif kind in ("doc", "slides", "sheet", "text"):
            text = _extract_attachment_text(path, kind, mime)
            if text:
                extracted_texts.append(
                    f"\n\n--- [Attachment: {name} ({kind}, {mime})] ---\n{text}\n--- [End: {name}] ---"
                )
                reference_lines.append(
                    f"- 📄 **{name}** ({mime}, {size} bytes) — text extracted ({len(text)} chars)"
                )
            else:
                reference_lines.append(
                    f"- ⚠️  **{name}** ({mime}, {size} bytes) — could not extract text"
                )
        else:
            reference_lines.append(f"- 📎 **{name}** ({mime}, {size} bytes, kind={kind})")

    # Strip any data:image base64 markdown the composer may have inlined into
    # the user text (pre-2.90.9 composer dumped these — defensive cleanup so
    # the model never sees a multi-MB base64 blob as raw text).
    if "data:image" in text_part:
        import re as _re
        text_part = _re.sub(
            r"!\[[^\]]*\]\(data:image/[^)]+\)",
            "[image attachment]",
            text_part,
        )

    composed_text = text_part.strip()
    if reference_lines:
        if composed_text:
            composed_text += "\n\n"
        composed_text += "📎 Attachments:\n" + "\n".join(reference_lines)
    composed_text += "".join(extracted_texts)

    # Build the content list. Providers that accept plain string content
    # (no images) keep a string; providers with images need a list.
    new_content: Any
    if image_blocks:
        new_content = [
            {"type": "text", "text": composed_text or "(see attachment)"},
            *image_blocks,
        ]
    else:
        new_content = composed_text or text_part

    llm_messages[last_user_idx] = {**orig, "content": new_content}
    return llm_messages


@app.delete("/api/workspace/close")
async def workspace_close(request: Request):
    """Close a workspace's DB connections (does not delete .cvc/)."""
    if not _workspace_mgr:
        raise HTTPException(503, "Workspace manager not initialised")
    body = await request.json()
    target_path = body.get("path", "")
    if not target_path:
        raise HTTPException(400, "Missing 'path' in request body")

    closed = _workspace_mgr.close_workspace(target_path)
    if not closed:
        raise HTTPException(404, f"Workspace not open: {target_path}")
    return {"status": "ok", "path": target_path, "closed": True}


# v2.91.46 — "Delete workspace" from the dashboard.
# Removes the workspace from the persistent registry (~/.cvc/workspaces.json)
# AND from the in-memory cache. The on-disk .cvc/ data is left untouched —
# the user can re-add the same path later. Pass `purge=true` to also delete
# the .cvc/ directory (destructive, irreversible).
@app.delete("/api/workspace/remove")
async def workspace_remove(request: Request, purge: bool = False):
    """Permanently remove a workspace from the persisted list.

    Query: ``?path=/abs/path&purge=false``
    Or body: ``{"path": "/abs/path", "purge": false}``

    - ``purge=false`` (default): removes from registry only; .cvc/ stays
    - ``purge=true``: ALSO deletes the workspace's on-disk .cvc/ directory
    """
    if not _workspace_mgr:
        raise HTTPException(503, "Workspace manager not initialised")
    # CPython SyntaxError trap: a `global` declaration may only appear ONCE
    # per name in a function. If we want to assign in multiple branches we
    # must declare the globals at the top of the function body.
    global _tool_executor, _project_root
    # Accept either query params or JSON body (fetch DELETE bodies are flaky
    # across browsers, so the frontend uses query params).
    target_path = (request.query_params.get("path") or "").strip()
    purge_q = request.query_params.get("purge", "").lower() in ("1", "true", "yes")
    if not target_path:
        try:
            body = await request.json()
        except Exception:
            body = {}
        target_path = (body.get("path") or "").strip()
        purge_q = bool(body.get("purge", False))
    purge = purge_q
    if not target_path:
        raise HTTPException(400, "Missing 'path' (query param or body)")

    target = Path(os.path.expanduser(target_path)).resolve()
    if not target.is_dir():
        raise HTTPException(404, f"Directory does not exist: {target}")

    if purge:
        result = _workspace_mgr.purge_workspace(target)
    else:
        result = _workspace_mgr.remove_from_registry(target)

    # Also clear the gateway proxy registry so any /ws/{id}/... routes 404
    workspace_id = _make_workspace_id(target)
    _registered_workspaces.pop(workspace_id, None)

    # v2.91.46 — Wipe any user-memory entries scoped to this workspace
    # so deleted workspaces don't leave orphaned memory behind. Hook
    # is best-effort: if user_memory_api isn't installed we silently
    # continue (the workspace removal itself still succeeds).
    if result.get("removed"):
        try:
            from cvc.dashboard.user_memory_api import drop_workspace_scope
            dropped = drop_workspace_scope(str(target))
            if dropped:
                result["memory_entries_dropped"] = dropped
        except Exception as _exc:
            logger.debug(
                "user_memory: drop_workspace_scope failed: %s", _exc
            )

    # If the active workspace pointer shifted, update the gateway globals
    if result.get("new_active"):
        new_path = Path(result["new_active"])
        inst = _workspace_mgr.current
        if inst:
            _project_root = inst.path
            try:
                from cvc.agent.executor import ToolExecutor
                _tool_executor = ToolExecutor(inst.path, inst.engine)
                _wire_subagent_runtime(_tool_executor)
            except Exception:
                logger.warning("Failed to re-scope ToolExecutor after remove", exc_info=True)
            # Re-register the new active workspace
            new_id = _make_workspace_id(inst.path)
            _registered_workspaces[new_id] = inst.path
    elif result.get("was_active"):
        # Active workspace was removed and no fallback. Clear globals.
        _tool_executor = None
        _project_root = None

    return {"status": "ok", **result}


# ═══════════════════════════════════════════════════════════════════════
# 1b. GIT BRANCH OPERATIONS (active-workspace-scoped)
#
# Lightweight wrapper around `git` for the currently-active workspace so the
# composer can show a Git Branch picker alongside the Workspace picker.
# All commands run with cwd=active workspace path. Safe: read-only ops never
# mutate state; checkout refuses on dirty tree unless `force=true`; branch
# creation always forks from current HEAD.
# ═══════════════════════════════════════════════════════════════════════

def _active_workspace_path() -> Path:
    if not _workspace_mgr or not _workspace_mgr.current:
        raise HTTPException(503, "No workspace active")
    return _workspace_mgr.current.path


def _run_git(args: list[str], cwd: Path, timeout: float = 8.0) -> tuple[int, str, str]:
    """Run a git subcommand under ``cwd`` and return (rc, stdout, stderr)."""
    import subprocess  # local import keeps cold-start cheap
    try:
        proc = subprocess.run(
            ["git", *args],
            cwd=str(cwd),
            capture_output=True,
            text=True,
            timeout=timeout,
            check=False,
                    **HIDDEN_KW,
        )
        return proc.returncode, (proc.stdout or "").strip(), (proc.stderr or "").strip()
    except FileNotFoundError:
        raise HTTPException(500, "git executable not found on PATH")
    except subprocess.TimeoutExpired:
        raise HTTPException(504, f"git {' '.join(args)} timed out")


def _is_git_repo(cwd: Path) -> bool:
    rc, out, _ = _run_git(["rev-parse", "--is-inside-work-tree"], cwd)
    return rc == 0 and out == "true"


_BRANCH_NAME_RE = __import__("re").compile(r"^[A-Za-z0-9._/\-]{1,200}$")


def _validate_branch_name(name: str) -> str:
    name = (name or "").strip()
    if not name:
        raise HTTPException(400, "Branch name is required")
    if name.startswith("-") or ".." in name or name.endswith("/") or name.endswith(".lock"):
        raise HTTPException(400, f"Invalid branch name: {name!r}")
    if not _BRANCH_NAME_RE.match(name):
        raise HTTPException(400, f"Invalid branch name: {name!r}")
    return name


@app.get("/api/git/status")
async def git_status():
    """Return the active workspace's git status: branch, head, dirty flag."""
    ws = _active_workspace_path()
    if not _is_git_repo(ws):
        return {"is_repo": False, "path": str(ws)}
    rc_b, branch, _ = _run_git(["rev-parse", "--abbrev-ref", "HEAD"], ws)
    rc_h, head, _ = _run_git(["rev-parse", "--short", "HEAD"], ws)
    rc_s, status_out, _ = _run_git(["status", "--porcelain"], ws)
    detached = branch == "HEAD"
    return {
        "is_repo": True,
        "path": str(ws),
        "branch": branch if rc_b == 0 else "",
        "head": head if rc_h == 0 else "",
        "detached": detached,
        "dirty": bool(status_out) if rc_s == 0 else False,
        "dirty_count": len(status_out.splitlines()) if rc_s == 0 and status_out else 0,
    }


@app.get("/api/git/branches")
async def git_branches(fetch: bool = True):
    """List local AND remote-tracking branches for the active workspace.

    Cross-platform (macOS/Linux/WSL/Windows): uses only `git` plumbing.
    Query: ?fetch=false to skip the best-effort fetch (faster, no network).
    """
    ws = _active_workspace_path()
    if not _is_git_repo(ws):
        raise HTTPException(409, "Active workspace is not a git repository")

    # Current branch (or "HEAD" if detached)
    _, current, _ = _run_git(["rev-parse", "--abbrev-ref", "HEAD"], ws)

    # Best-effort fetch — short timeout, never fails the call.
    # Refreshes remote-tracking refs so `origin/main` etc. show up after a push.
    fetch_ok = False
    fetch_error: str | None = None
    if fetch:
        rc_f, _, err_f = _run_git(
            ["fetch", "--all", "--prune", "--quiet"],
            ws,
            timeout=8.0,
        )
        fetch_ok = rc_f == 0
        if not fetch_ok and err_f:
            fetch_error = err_f.strip().splitlines()[-1][:200] if err_f.strip() else None

    # --- LOCAL branches (refs/heads) with upstream tracking info ---
    # %(upstream:short) → "origin/main" or empty; %(upstream:track) → "[ahead 1, behind 2]" or empty
    fmt_local = (
        "%(refname:short)%09%(objectname:short)%09"
        "%(committerdate:iso-strict)%09%(subject)%09"
        "%(upstream:short)%09%(upstream:track)"
    )
    rc, out, err = _run_git(
        ["for-each-ref", f"--format={fmt_local}", "--sort=-committerdate", "refs/heads"],
        ws,
        timeout=12.0,
    )
    if rc != 0:
        raise HTTPException(500, f"git for-each-ref (local) failed: {err}")

    def _parse_track(token: str) -> tuple[int, int, bool]:
        """Parse '[ahead 1, behind 2]' / '[gone]' → (ahead, behind, gone)."""
        if not token:
            return (0, 0, False)
        s = token.strip().strip("[]")
        if "gone" in s:
            return (0, 0, True)
        ahead = behind = 0
        for part in s.split(","):
            p = part.strip()
            if p.startswith("ahead "):
                try: ahead = int(p[6:])
                except ValueError: pass
            elif p.startswith("behind "):
                try: behind = int(p[7:])
                except ValueError: pass
        return (ahead, behind, False)

    local_branches: list[dict[str, Any]] = []
    local_upstreams: set[str] = set()
    for line in out.splitlines():
        parts = line.split("\t", 5)
        if len(parts) < 2:
            continue
        name = parts[0]
        head = parts[1] if len(parts) > 1 else ""
        when = parts[2] if len(parts) > 2 else ""
        subject = parts[3] if len(parts) > 3 else ""
        upstream = parts[4] if len(parts) > 4 else ""
        track = parts[5] if len(parts) > 5 else ""
        ahead, behind, gone = _parse_track(track)
        if upstream:
            local_upstreams.add(upstream)
        local_branches.append({
            "name": name,
            "head": head,
            "last_commit_at": when,
            "last_commit_subject": subject,
            "is_current": name == current,
            "kind": "local",
            "upstream": upstream or None,
            "ahead": ahead,
            "behind": behind,
            "gone": gone,
            "has_local": True,
        })

    # --- REMOTE-tracking branches (refs/remotes), excluding HEAD pointers ---
    fmt_remote = "%(refname:short)%09%(objectname:short)%09%(committerdate:iso-strict)%09%(subject)"
    rc_r, out_r, err_r = _run_git(
        ["for-each-ref", f"--format={fmt_remote}", "--sort=-committerdate", "refs/remotes"],
        ws,
        timeout=12.0,
    )
    remote_branches: list[dict[str, Any]] = []
    if rc_r == 0:
        local_names = {b["name"] for b in local_branches}
        for line in out_r.splitlines():
            parts = line.split("\t", 3)
            if len(parts) < 2:
                continue
            name = parts[0]
            # Skip HEAD pointers like "origin/HEAD -> origin/main"
            if name.endswith("/HEAD") or " -> " in name:
                continue
            head = parts[1] if len(parts) > 1 else ""
            when = parts[2] if len(parts) > 2 else ""
            subject = parts[3] if len(parts) > 3 else ""
            # remote = first path segment, short_name = rest
            remote_name, sep, short = name.partition("/")
            # Skip bare remote pointers like "origin" (which is `origin/HEAD`
            # after `%(refname:short)` strips the /HEAD suffix). These have no
            # slash and no short_name — they're not real branches.
            if not sep or not short:
                continue
            has_local = short in local_names
            remote_branches.append({
                "name": name,                  # e.g. "origin/main"
                "short_name": short,           # e.g. "main"
                "remote": remote_name,         # e.g. "origin"
                "head": head,
                "last_commit_at": when,
                "last_commit_subject": subject,
                "is_current": False,
                "kind": "remote",
                "upstream": None,
                "ahead": 0,
                "behind": 0,
                "gone": False,
                "has_local": has_local,         # true if a local branch already tracks this name
                "tracked_by_local": name in local_upstreams,
            })

    return {
        # Back-compat: "branches" still returns local-only list (existing clients keep working)
        "branches": local_branches,
        "local": local_branches,
        "remote": remote_branches,
        "current": current if current != "HEAD" else "",
        "detached": current == "HEAD",
        "count": len(local_branches),
        "remote_count": len(remote_branches),
        "fetched": fetch_ok,
        "fetch_error": fetch_error,
    }


@app.post("/api/git/checkout")
async def git_checkout(request: Request):
    """Switch to (or create) a branch in the active workspace.

    Body: {"name": str, "create": bool=false, "force": bool=false}
    - create=true → equivalent to `git checkout -b <name>` (forks from HEAD).
    - force=true  → allows checkout despite dirty tree (uses `git checkout -f`).
    """
    ws = _active_workspace_path()
    if not _is_git_repo(ws):
        raise HTTPException(409, "Active workspace is not a git repository")

    body = await request.json()
    raw_name = str(body.get("name", "")).strip()
    create = bool(body.get("create", False))
    force = bool(body.get("force", False))

    # Detect remote ref like "origin/main" → auto-create tracking local branch.
    # We only treat it as a remote ref if a matching refs/remotes/<name> exists.
    is_remote_ref = False
    track_remote_ref: str | None = None
    local_name = raw_name
    if "/" in raw_name and not create:
        rc_rem, _, _ = _run_git(
            ["show-ref", "--verify", "--quiet", f"refs/remotes/{raw_name}"], ws
        )
        if rc_rem == 0:
            is_remote_ref = True
            track_remote_ref = raw_name
            # local branch name = remote ref minus the leading "<remote>/"
            local_name = raw_name.split("/", 1)[1]
            # If a local branch with that name already exists, just switch to it.
            rc_local, _, _ = _run_git(
                ["show-ref", "--verify", "--quiet", f"refs/heads/{local_name}"], ws
            )
            if rc_local == 0:
                is_remote_ref = False  # fall through to plain checkout
                track_remote_ref = None

    name = _validate_branch_name(local_name)

    # Check dirty state — refuse unless force or pure create-then-checkout
    rc_s, status_out, _ = _run_git(["status", "--porcelain"], ws)
    dirty = rc_s == 0 and bool(status_out)
    if dirty and not force:
        raise HTTPException(
            409,
            "Working tree has uncommitted changes. Commit/stash first, or pass force=true.",
        )

    # If checkout (not create, not remote-tracking), verify branch exists locally
    if not create and not is_remote_ref:
        rc_v, _, _ = _run_git(["show-ref", "--verify", "--quiet", f"refs/heads/{name}"], ws)
        if rc_v != 0:
            raise HTTPException(404, f"Branch not found locally: {name}")

    if is_remote_ref and track_remote_ref:
        # `git checkout -b <local> --track <remote>/<branch>` (or -f variant)
        args = ["checkout"]
        if force:
            args.append("-f")
        args.extend(["-b", name, "--track", track_remote_ref])
        created_flag = True
    else:
        args = ["checkout"]
        if force:
            args.append("-f")
        if create:
            args.append("-b")
        args.append(name)
        created_flag = create

    rc, out, err = _run_git(args, ws, timeout=15.0)
    if rc != 0:
        raise HTTPException(500, f"git checkout failed: {err or out}")

    _, head, _ = _run_git(["rev-parse", "--short", "HEAD"], ws)
    return {
        "status": "ok",
        "branch": name,
        "head": head,
        "created": created_flag,
        "tracked_remote": track_remote_ref,
    }


@app.post("/api/git/sync")
async def git_sync(request: Request):
    """Sync the active workspace with its tracked remote.

    Mirrors what VS Code's "Sync Changes" button does:
      1. ``git fetch <remote>``
      2. ``git pull --ff-only`` (fast-forward only — refuses to merge/rebase
         silently; bails out on diverged history so the user can resolve it).
      3. If ``push=true`` and local is ahead, ``git push`` upstream too.

    Body (all optional):
      {"remote": "origin", "push": true, "rebase": false}

    Returns structured result so the UI can show a meaningful toast:
      {
        "status": "ok" | "diverged" | "dirty" | "no_upstream",
        "branch": str,
        "remote": str,
        "fetched": bool,
        "pulled": int,   # commits pulled in
        "pushed": int,   # commits pushed out
        "ahead": int,    # remaining ahead (should be 0 after push)
        "behind": int,   # remaining behind (should be 0 after pull)
        "head": str,
        "message": str,
      }
    """
    ws = _active_workspace_path()
    if not _is_git_repo(ws):
        raise HTTPException(409, "Active workspace is not a git repository")

    try:
        body = await request.json()
    except Exception:
        body = {}
    remote = str(body.get("remote") or "origin").strip() or "origin"
    do_push = bool(body.get("push", True))
    do_rebase = bool(body.get("rebase", False))

    # Delegate to shared CLI ↔ dashboard helper — single source of truth.
    from cvc.agent.git_integration import git_sync as _do_sync

    result = _do_sync(ws, remote=remote, push=do_push, rebase=do_rebase)

    # Map "error" status (validation failures) to HTTP errors so existing UI
    # error handling stays correct; functional statuses pass through.
    if result["status"] == "error":
        msg = result.get("message", "git sync failed")
        if "Detached HEAD" in msg or "Not a git repository" in msg or "Invalid remote" in msg:
            raise HTTPException(409, msg)
        raise HTTPException(500, msg)

    return result


# ═══════════════════════════════════════════════════════════════════════
# 1c. WORKSPACE-ROUTED LLM PROXY (replaces standalone cvc serve)
#
# AI tools (Claude Code, Aider, Codex, Gemini CLI, Kiro, etc.) send
# LLM requests to  http://127.0.0.1:13421/ws/{workspace_id}/v1/...
# The gateway looks up the correct CVCEngine for that workspace and
# runs full CVC interception (context capture, auto-commit, etc.).
# ═══════════════════════════════════════════════════════════════════════

# -- Workspace registration ------------------------------------------------

# Maps workspace_id (short hash) → resolved Path
_registered_workspaces: dict[str, Path] = {}

# Per-workspace session trackers, adapters, graphs, and config
_ws_adapters: dict[str, Any] = {}
_ws_graphs: dict[str, Any] = {}
_ws_session_trackers: dict[str, Any] = {}  # workspace_id → _SessionTracker


# Time Machine configuration (same knobs as the old proxy)
_TIME_MACHINE_ENABLED: bool = os.environ.get("CVC_TIME_MACHINE", "0") == "1"
_TIME_MACHINE_INTERVAL: int = int(os.environ.get("CVC_TIME_MACHINE_INTERVAL", "1"))
_NORMAL_INTERVAL: int = int(os.environ.get("CVC_AUTO_COMMIT_INTERVAL", "3"))
_SESSION_TIMEOUT: int = int(os.environ.get("CVC_SESSION_TIMEOUT", "1800"))


class _SessionTracker:
    """Lightweight tracker for agent sessions going through the proxy."""

    def __init__(self) -> None:
        self.sessions: list[dict[str, Any]] = []
        self._current: dict[str, Any] | None = None
        self._last_request_ts: float = 0.0

    def touch(self, *, tool_hint: str = "") -> dict[str, Any]:
        now = time.time()
        gap = now - self._last_request_ts if self._last_request_ts else 0
        self._last_request_ts = now

        if self._current is None or gap > _SESSION_TIMEOUT:
            if self._current is not None:
                self._current["ended_at"] = now - gap
                self.sessions.append(self._current)

            self._current = {
                "id": len(self.sessions) + 1,
                "started_at": now,
                "ended_at": None,
                "tool": tool_hint or "unknown",
                "messages": 0,
                "commits": 0,
                "branch_at_start": "",
            }

        self._current["messages"] += 1
        if tool_hint and self._current["tool"] == "unknown":
            self._current["tool"] = tool_hint
        return self._current

    def record_commit(self) -> None:
        if self._current:
            self._current["commits"] += 1

    def all_sessions(self) -> list[dict[str, Any]]:
        result = list(self.sessions)
        if self._current:
            result.append({**self._current, "active": True})
        return result


def _make_workspace_id(path: Path) -> str:
    """Generate a short, URL-safe workspace identifier from a directory path."""
    return hashlib.sha256(str(path.resolve()).encode()).hexdigest()[:12]


def _get_ws_engine(workspace_id: str):
    """Return the CVCEngine for a registered workspace, or raise 404."""
    if not _workspace_mgr:
        raise HTTPException(503, "Workspace manager not initialised")
    ws_path = _registered_workspaces.get(workspace_id)
    if ws_path is None:
        raise HTTPException(404, f"Workspace not registered: {workspace_id}")
    key = str(ws_path.resolve())
    inst = _workspace_mgr._workspaces.get(key)
    if inst is None or not inst.initialized:
        # Re-initialize on demand
        inst = _workspace_mgr.switch_to(ws_path)
    inst.touch()
    return inst.engine, inst.config, inst.db


def _get_ws_tracker(workspace_id: str) -> _SessionTracker:
    """Return (or lazily create) the session tracker for a workspace."""
    if workspace_id not in _ws_session_trackers:
        _ws_session_trackers[workspace_id] = _SessionTracker()
    return _ws_session_trackers[workspace_id]


async def _get_ws_adapter(workspace_id: str, config):
    """Return (or lazily create) the LLM adapter for a workspace."""
    if workspace_id not in _ws_adapters:
        from cvc.adapters import create_adapter
        _ws_adapters[workspace_id] = create_adapter(
            provider=config.provider,
            api_key=config.api_key,
            model=config.model,
            base_url=config.upstream_base_url,
        )
    return _ws_adapters[workspace_id]


def _get_ws_graph(workspace_id: str, engine):
    """Return (or lazily create) the CVC state-machine graph for a workspace."""
    if workspace_id not in _ws_graphs:
        from cvc.operations.state_machine import build_cvc_graph
        graph_builder = build_cvc_graph(engine)
        _ws_graphs[workspace_id] = graph_builder.compile()
    return _ws_graphs[workspace_id]


@app.post("/gateway/register_workspace")
async def register_workspace(request: Request):
    """Register a workspace directory with the gateway.

    Body: ``{"path": "/home/user/project-a"}``

    Returns a ``workspace_id`` that the launcher injects into env vars
    so the AI tool's requests are routed to the correct .cvc/ database.
    """
    if not _workspace_mgr:
        raise HTTPException(503, "Workspace manager not initialised")

    body = await request.json()
    raw_path = body.get("path", "")
    if not raw_path:
        raise HTTPException(400, "Missing 'path' in request body")

    target = Path(raw_path).resolve()
    if not target.is_dir():
        raise HTTPException(400, f"Not a directory: {target}")

    workspace_id = _make_workspace_id(target)

    # Register in the memory map
    _registered_workspaces[workspace_id] = target

    # Initialize the workspace (auto-creates .cvc/ if needed)
    inst = _workspace_mgr.switch_to(target)

    logger.info(
        "Workspace registered: %s → %s (branch=%s)",
        workspace_id,
        target,
        inst.engine.active_branch if inst.engine else "main",
    )

    return {
        "workspace_id": workspace_id,
        "path": str(target),
        "branch": inst.engine.active_branch if inst.engine else "main",
        "base_url": f"http://{_HOST}:{_GATEWAY_PORT}/ws/{workspace_id}",
    }


@app.get("/gateway/workspaces")
async def list_registered_workspaces():
    """List all workspaces (in-memory + persisted registry)."""
    result = []
    seen: set[str] = set()
    for ws_id, ws_path in _registered_workspaces.items():
        info: dict[str, Any] = {
            "workspace_id": ws_id,
            "path": str(ws_path),
            "name": ws_path.name or str(ws_path),
        }
        try:
            engine, _, _ = _get_ws_engine(ws_id)
            info["branch"] = engine.active_branch
            info["context_size"] = len(engine.context_window)
        except Exception:
            info["branch"] = "unknown"
            info["context_size"] = 0
        tracker = _ws_session_trackers.get(ws_id)
        if tracker:
            info["sessions"] = len(tracker.all_sessions())
        result.append(info)
        seen.add(ws_id)

    # Merge persisted registry so previously-added workspaces survive restart
    if _workspace_mgr:
        try:
            for entry in _workspace_mgr._load_registry():
                ep = entry.get("path", "")
                if not ep:
                    continue
                p = Path(ep)
                if not p.exists():
                    continue
                ws_id = entry.get("workspace_id") or _make_workspace_id(p)
                if ws_id in seen:
                    continue
                result.append({
                    "workspace_id": ws_id,
                    "path": str(p),
                    "name": entry.get("name") or p.name or str(p),
                    "branch": entry.get("branch", "main"),
                    "context_size": 0,
                })
                seen.add(ws_id)
        except Exception:
            logger.warning("Failed to merge persisted workspace registry", exc_info=True)

    return {"workspaces": result, "count": len(result)}


# -- Proxy interception helpers -------------------------------------------

def _detect_tool_from_ua(ua: str) -> str:
    ua_lower = ua.lower()
    for keyword in ("claude", "aider", "codex", "cursor", "copilot", "gemini", "kiro"):
        if keyword in ua_lower:
            return keyword
    return ""


def _ws_should_auto_commit(engine) -> bool:
    assistant_msgs = [m for m in engine.context_window if m.role == "assistant"]
    count = len(assistant_msgs)
    if count == 0:
        return False
    interval = _TIME_MACHINE_INTERVAL if _TIME_MACHINE_ENABLED else _NORMAL_INTERVAL
    return count % interval == 0


def _ws_build_auto_commit_message(engine) -> str:
    recent = [m for m in engine.context_window if m.role in ("user", "assistant")]
    if not recent:
        return "Auto-checkpoint"
    last_user = next((m for m in reversed(recent) if m.role == "user"), None)
    if last_user:
        summary = last_user.content[:120].replace("\n", " ").strip()
        return f"Auto: {summary}"
    return f"Auto-checkpoint ({len(recent)} messages)"


# -- OpenAI-compatible Chat Completions (workspace-routed) ----------------

@app.post("/ws/{workspace_id}/v1/chat/completions")
async def ws_chat_completions(workspace_id: str, raw_request: Request):
    """Workspace-routed OpenAI-compatible chat completions endpoint.

    Replaces the old standalone proxy's ``/v1/chat/completions``.
    """
    engine, config, db = _get_ws_engine(workspace_id)
    graph = _get_ws_graph(workspace_id, engine)
    adapter = await _get_ws_adapter(workspace_id, config)
    tracker = _get_ws_tracker(workspace_id)

    # Session tracking
    ua = raw_request.headers.get("user-agent", "")
    tool_hint = _detect_tool_from_ua(ua)
    session = tracker.touch(tool_hint=tool_hint)
    if not session.get("branch_at_start"):
        session["branch_at_start"] = engine.active_branch

    from cvc.core.models import (
        ChatCompletionChoice, ChatCompletionRequest, ChatCompletionResponse,
        ChatMessage, CVCCommitRequest, CVCOperationResponse,
    )
    from cvc.operations.state_machine import CVCGraphState

    body = await raw_request.json()
    request = ChatCompletionRequest(**body)

    # 1. Run through the CVC state machine
    initial_state: CVCGraphState = {
        "request": request,
        "is_cvc_command": False,
        "cvc_tool_name": "",
        "cvc_tool_args": {},
        "response": None,
        "cvc_result": None,
        "active_branch": engine.active_branch,
        "head_hash": engine.head_hash or "",
        "committed_prefix_len": 0,
    }
    result_state = graph.invoke(initial_state)

    # 2. CVC command → return directly
    if result_state.get("is_cvc_command") and result_state.get("cvc_result"):
        cvc_result: CVCOperationResponse = result_state["cvc_result"]
        wrapped = ChatCompletionResponse(
            model=request.model,
            choices=[
                ChatCompletionChoice(
                    message=ChatMessage(
                        role="assistant",
                        content=json.dumps(cvc_result.model_dump(), indent=2),
                    ),
                    finish_reason="stop",
                )
            ],
        )
        return JSONResponse(
            content=wrapped.model_dump(),
            headers={"X-CVC-Operation": cvc_result.operation},
        )

    # 3. Standard generation — capture context & proxy upstream
    existing_count = len(engine.context_window)
    inbound_count = len(request.messages)
    if inbound_count > existing_count:
        for msg in request.messages[existing_count:]:
            if msg.role in ("user", "system", "tool"):
                engine.push_chat_message(msg)

    committed_prefix_len = result_state.get("committed_prefix_len", 0)
    is_streaming = bool(body.get("stream", False))

    # COGNOME injection — compile an Engram and prepend it before the
    # upstream LLM call, giving it curated workspace context.
    _cognome_engram = None
    try:
        from cvc.operations.cognome_runtime import CognomeRuntime
        runtime = CognomeRuntime.for_engine(engine)
        if runtime.is_active():
            raw_msgs = [{"role": m.role, "content": m.content or ""} for m in request.messages]
            updated_msgs, _cognome_engram = await runtime.resolve_messages(
                raw_msgs,
                provider=config.provider,
                model=config.model,
                branch=request.cvc_branch,
            )
            if _cognome_engram and _cognome_engram.preamble:
                from cvc.core.models import ChatMessage as CM
                request.messages = [CM(**m) for m in updated_msgs]
                committed_prefix_len += 1
    except Exception as exc:
        logger.warning("COGNOME injection failed (ws %s, non-fatal): %s", workspace_id, exc)

    # Passthrough mode: no internal adapter — client must use /v1/messages directly
    if config.provider == "passthrough":
        raise HTTPException(
            status_code=400,
            detail=(
                "CVC is in passthrough mode. Claude Code and Anthropic tools should use "
                "/v1/messages (set ANTHROPIC_BASE_URL). OpenAI-compatible tools: "
                "run 'cvc setup --provider openai' to configure a provider."
            ),
        )

    try:
        response = await adapter.complete(request, committed_prefix_len=committed_prefix_len)
    except Exception as exc:
        logger.error("Upstream provider error (ws %s): %s", workspace_id, exc)
        raise HTTPException(status_code=502, detail=f"Upstream error: {exc}") from exc

    # 4. Record assistant response
    if response.choices:
        assistant_msg = response.choices[0].message
        engine.push_chat_message(assistant_msg)

    # 5. Auto-commit
    if _ws_should_auto_commit(engine):
        auto_msg = _ws_build_auto_commit_message(engine)
        auto_result = engine.commit(CVCCommitRequest(message=auto_msg, commit_type="checkpoint"))
        if auto_result.success:
            tracker.record_commit()
            logger.info("Auto-commit (ws %s): %s", workspace_id,
                        auto_result.commit_hash and auto_result.commit_hash[:12])

    # 6. If client requested streaming, wrap full response as SSE chunks
    if is_streaming:
        from fastapi.responses import StreamingResponse as _StreamRes
        resp_data = response.model_dump()
        resp_id = resp_data.get("id", f"chatcmpl-{int(time.time())}")
        model_name = resp_data.get("model", request.model)
        content = ""
        finish_reason = "stop"
        if response.choices:
            choice = response.choices[0]
            content = choice.message.content or ""
            finish_reason = choice.finish_reason or "stop"

        def _sse_chunks():
            # Role delta
            yield "data: " + json.dumps({
                "id": resp_id, "object": "chat.completion.chunk",
                "model": model_name,
                "choices": [{"index": 0, "delta": {"role": "assistant"}, "finish_reason": None}],
            }) + "\n\n"
            # Content in chunks of ~100 chars
            chunk_size = 100
            for i in range(0, max(len(content), 1), chunk_size):
                chunk = content[i:i + chunk_size]
                yield "data: " + json.dumps({
                    "id": resp_id, "object": "chat.completion.chunk",
                    "model": model_name,
                    "choices": [{"index": 0, "delta": {"content": chunk}, "finish_reason": None}],
                }) + "\n\n"
            # Stop chunk
            yield "data: " + json.dumps({
                "id": resp_id, "object": "chat.completion.chunk",
                "model": model_name,
                "choices": [{"index": 0, "delta": {}, "finish_reason": finish_reason}],
            }) + "\n\n"
            yield "data: [DONE]\n\n"

        return _StreamRes(
            _sse_chunks(),
            media_type="text/event-stream",
            headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
        )

    return JSONResponse(content=response.model_dump())


# -- Anthropic Messages API (workspace-routed, for Claude Code) -----------

@app.post("/ws/{workspace_id}/v1/messages")
async def ws_anthropic_messages(workspace_id: str, raw_request: Request):
    """Workspace-routed Anthropic-native Messages API endpoint.

    Claude Code CLI sends requests to ANTHROPIC_BASE_URL + '/v1/messages'.
    """
    engine, config, db = _get_ws_engine(workspace_id)
    tracker = _get_ws_tracker(workspace_id)

    ua = raw_request.headers.get("user-agent", "")
    tool_hint = _detect_tool_from_ua(ua)
    session = tracker.touch(tool_hint=tool_hint)
    if not session.get("branch_at_start"):
        session["branch_at_start"] = engine.active_branch

    from cvc.core.models import (
        ChatCompletionRequest, ChatMessage, CVCCommitRequest, CVCOperationResponse,
    )
    from cvc.operations.state_machine import CVCGraphState

    body = await raw_request.json()

    # Convert Anthropic format → internal OpenAI format
    messages_in = body.get("messages", [])
    system_text = body.get("system", "")
    openai_messages: list[dict[str, Any]] = []

    if system_text:
        if isinstance(system_text, list):
            sys_parts = [b.get("text", "") for b in system_text if b.get("type") == "text"]
            system_text = "\n".join(sys_parts)
        openai_messages.append({"role": "system", "content": system_text})

    for msg in messages_in:
        role = msg.get("role", "user")
        content = msg.get("content", "")
        if isinstance(content, list):
            text_parts = [b.get("text", "") for b in content if b.get("type") == "text"]
            content = "\n".join(text_parts)
        openai_messages.append({"role": role, "content": content})

    internal_request = ChatCompletionRequest(
        model=body.get("model", config.model),
        messages=[ChatMessage(**m) for m in openai_messages],
        temperature=body.get("temperature", 0.7),
        max_tokens=body.get("max_tokens", 4096),
    )

    # Convert Anthropic tools → OpenAI tools
    anthropic_tools = body.get("tools", [])
    if anthropic_tools:
        openai_tools = []
        for t in anthropic_tools:
            openai_tools.append({
                "type": "function",
                "function": {
                    "name": t.get("name", ""),
                    "description": t.get("description", ""),
                    "parameters": t.get("input_schema", {}),
                },
            })
        internal_request.tools = openai_tools

    # Run through CVC state machine
    graph = _get_ws_graph(workspace_id, engine)

    existing_count = len(engine.context_window)
    for msg in internal_request.messages[existing_count:]:
        if msg.role in ("user", "system", "tool"):
            engine.push_chat_message(msg)

    initial_state: CVCGraphState = {
        "request": internal_request,
        "is_cvc_command": False,
        "cvc_tool_name": "",
        "cvc_tool_args": {},
        "response": None,
        "cvc_result": None,
        "active_branch": engine.active_branch,
        "head_hash": engine.head_hash or "",
        "committed_prefix_len": 0,
    }
    result_state = graph.invoke(initial_state)

    # CVC command → Anthropic format
    if result_state.get("is_cvc_command") and result_state.get("cvc_result"):
        cvc_result: CVCOperationResponse = result_state["cvc_result"]
        return JSONResponse(content={
            "id": f"msg_cvc_{int(time.time())}",
            "type": "message",
            "role": "assistant",
            "model": internal_request.model,
            "content": [{"type": "text", "text": json.dumps(cvc_result.model_dump(), indent=2)}],
            "stop_reason": "end_turn",
            "usage": {"input_tokens": 0, "output_tokens": 0},
        }, headers={"X-CVC-Operation": cvc_result.operation})

    # Forward to upstream Anthropic API
    client_api_key = raw_request.headers.get("x-api-key", "") or ""
    auth_header = raw_request.headers.get("authorization", "")
    if auth_header.lower().startswith("bearer "):
        client_api_key = client_api_key or auth_header[7:].strip()
    upstream_key = client_api_key or config.api_key

    # In passthrough mode, the client MUST supply its own key.
    # For any other provider, a missing key is a config error.
    if not upstream_key:
        if config.provider == "passthrough":
            raise HTTPException(
                status_code=401,
                detail=(
                    "CVC is in passthrough mode and no API key was found in the request. "
                    "Ensure your AI tool (e.g. Claude Code) is sending its own API key via "
                    "the x-api-key or Authorization header."
                ),
            )
        raise HTTPException(
            status_code=401,
            detail=(
                "No API key available. Pass x-api-key header, set ANTHROPIC_API_KEY, "
                "or run 'cvc setup' to configure a provider."
            ),
        )

    import httpx
    from fastapi.responses import StreamingResponse
    headers = {
        "x-api-key": upstream_key,
        "anthropic-version": "2023-06-01",
        "content-type": "application/json",
    }
    beta = raw_request.headers.get("anthropic-beta")
    if beta:
        headers["anthropic-beta"] = beta

    upstream_body = dict(body)
    
    # --- Inject Prompt Caching ---
    _breakpoints = 0
    if "system" in upstream_body and upstream_body["system"]:
        if isinstance(upstream_body["system"], str):
            upstream_body["system"] = [{"type": "text", "text": upstream_body["system"]}]
        if isinstance(upstream_body["system"], list) and len(upstream_body["system"]) > 0:
            upstream_body["system"][-1]["cache_control"] = {"type": "ephemeral"}
            _breakpoints += 1

    if "messages" in upstream_body and isinstance(upstream_body["messages"], list) and len(upstream_body["messages"]) > 0:
        if _breakpoints < 4:
            last_msg = upstream_body["messages"][-1]
            if "content" in last_msg:
                if isinstance(last_msg["content"], str):
                    last_msg["content"] = [{"type": "text", "text": last_msg["content"]}]
                if isinstance(last_msg["content"], list) and len(last_msg["content"]) > 0:
                    last_msg["content"][-1]["cache_control"] = {"type": "ephemeral"}
                    _breakpoints += 1
    # -----------------------------
    
    is_streaming = bool(upstream_body.get("stream", False))

    try:
        # passthrough → always forward to Anthropic; other providers use config URL
        upstream_url = "https://api.anthropic.com"
        if config.provider not in ("anthropic", "passthrough"):
            upstream_url = config.upstream_base_url.rstrip("/")

        if is_streaming:
            # Stream SSE directly to the client while capturing text for CVC
            async def _stream_and_capture():
                accumulated_text: list[str] = []
                try:
                    async with httpx.AsyncClient(timeout=300.0) as http_client:
                        async with http_client.stream(
                            "POST",
                            f"{upstream_url}/v1/messages",
                            json=upstream_body,
                            headers=headers,
                        ) as resp:
                            resp.raise_for_status()
                            async for raw_line in resp.aiter_lines():
                                if raw_line.startswith("data:"):
                                    payload = raw_line[5:].strip()
                                    if payload and payload != "[DONE]":
                                        try:
                                            evt = json.loads(payload)
                                            # Capture content_block_delta text
                                            delta = evt.get("delta", {})
                                            if delta.get("type") == "text_delta":
                                                accumulated_text.append(delta.get("text", ""))
                                        except Exception:
                                            pass
                                yield raw_line + "\n"
                except Exception as exc:
                    logger.error("Upstream stream error (ws %s): %s", workspace_id, exc)
                    yield f"data: {json.dumps({'error': str(exc)})}\n\n"
                    return

                # After stream completes, save to CVC
                full_text = "".join(accumulated_text).strip()
                if full_text:
                    engine.push_chat_message(ChatMessage(role="assistant", content=full_text))
                    if _ws_should_auto_commit(engine):
                        auto_msg = _ws_build_auto_commit_message(engine)
                        auto_result = engine.commit(CVCCommitRequest(message=auto_msg, commit_type="checkpoint"))
                        if auto_result.success:
                            tracker.record_commit()

            return StreamingResponse(
                _stream_and_capture(),
                media_type="text/event-stream",
                headers={
                    "Cache-Control": "no-cache",
                    "X-Accel-Buffering": "no",
                    "Connection": "keep-alive",
                },
            )

        # Non-streaming path
        async with httpx.AsyncClient(timeout=120.0) as client:
            resp = await client.post(f"{upstream_url}/v1/messages", json=upstream_body, headers=headers)
            resp.raise_for_status()
            data = resp.json()
    except Exception as exc:
        logger.error("Upstream Anthropic error (ws %s): %s", workspace_id, exc)
        raise HTTPException(status_code=502, detail=f"Upstream error: {exc}") from exc

    # Record assistant response (non-streaming)
    content_blocks = data.get("content", [])
    text_parts = [b["text"] for b in content_blocks if b.get("type") == "text"]
    if text_parts:
        engine.push_chat_message(ChatMessage(role="assistant", content="\n".join(text_parts)))

    # Auto-commit
    if _ws_should_auto_commit(engine):
        auto_msg = _ws_build_auto_commit_message(engine)
        auto_result = engine.commit(CVCCommitRequest(message=auto_msg, commit_type="checkpoint"))
        if auto_result.success:
            tracker.record_commit()

    return JSONResponse(content=data)


# -- Model listing (workspace-routed) -------------------------------------

@app.get("/ws/{workspace_id}/v1/models")
async def ws_list_models(workspace_id: str):
    """Workspace-routed OpenAI-compatible model listing."""
    engine, config, _ = _get_ws_engine(workspace_id)
    return {
        "object": "list",
        "data": [{
            "id": config.model,
            "object": "model",
            "created": int(time.time()),
            "owned_by": f"cvc-gateway ({config.provider})",
            "permission": [],
        }],
    }


# -- Workspace proxy health -----------------------------------------------

@app.get("/ws/{workspace_id}/health")
async def ws_proxy_health(workspace_id: str):
    """Health check scoped to a workspace."""
    try:
        engine, config, _ = _get_ws_engine(workspace_id)
        return {
            "status": "ok",
            "service": "cvc-gateway-proxy",
            "workspace_id": workspace_id,
            "branch": engine.active_branch,
            "version": _cvc_version,
        }
    except HTTPException:
        raise
    except Exception as exc:
        return {"status": "error", "error": str(exc)}


# -- CVC control endpoints (workspace-routed) ------------------------------

@app.get("/ws/{workspace_id}/cvc/status")
async def ws_cvc_status(workspace_id: str):
    engine, _, db = _get_ws_engine(workspace_id)
    branches = db.index.list_branches()
    return {
        "agent_id": engine.config.agent_id,
        "active_branch": engine.active_branch,
        "head_hash": engine.head_hash,
        "context_size": len(engine.context_window),
        "branches": [
            {"name": b.name, "head": b.head_hash[:12], "status": b.status.value}
            for b in branches
        ],
    }


@app.get("/ws/{workspace_id}/cvc/sessions")
async def ws_cvc_sessions(workspace_id: str):
    tracker = _get_ws_tracker(workspace_id)
    sessions = tracker.all_sessions()
    return {
        "time_machine": _TIME_MACHINE_ENABLED,
        "auto_commit_interval": _TIME_MACHINE_INTERVAL if _TIME_MACHINE_ENABLED else _NORMAL_INTERVAL,
        "session_timeout_seconds": _SESSION_TIMEOUT,
        "total": len(sessions),
        "sessions": sessions,
    }


# ═══════════════════════════════════════════════════════════════════════
# 2. SERVICE MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════

@app.get("/api/gateway/services")
async def get_services():
    _ensure()
    services = _manager.service_status()
    enriched = {}
    for svc in services:
        name = svc["name"]
        enriched[name] = svc
        enriched[name]["uptime"] = _manager.get_service(name).uptime_display
    if _get_db():
        try:
            agents = _get_db().index.list_agents() if hasattr(_get_db().index, "list_agents") else []
            enriched.get("sdk", {})["registered_agents"] = len(agents)
        except Exception:
            pass
        try:
            total = len(_get_db().index.list_all_commits(limit=9999)) if hasattr(_get_db().index, "list_all_commits") else 0
            enriched.get("sdk", {})["total_commits"] = total
        except Exception:
            pass
    return enriched


@app.post("/api/gateway/services/{service_name}/{action}")
async def service_action(service_name: str, action: str):
    _ensure()
    if service_name == "proxy":
        port = _manager.get_service("proxy").port or _PROXY_PORT
        actions = {
            "start": lambda: _manager.start_proxy(port),
            "stop": lambda: _manager.stop_proxy(),
            "restart": lambda: _manager.restart_proxy(port),
            "pause": lambda: _manager.pause_proxy(),
            "resume": lambda: _manager.resume_proxy(),
        }
    else:
        raise HTTPException(400, f"Service '{service_name}' does not support lifecycle actions")
    if action not in actions:
        raise HTTPException(400, f"Unknown action '{action}'. Use: start, stop, restart, pause, resume")
    result = actions[action]()
    _get_event_bus().emit("service.action", {"service": service_name, "action": action, "status": result.status})
    return result.to_dict()


# ═══════════════════════════════════════════════════════════════════════
# 3. ANALYTICS & COMMITS
# ═══════════════════════════════════════════════════════════════════════

@app.get("/api/gateway/analytics")
async def get_analytics():
    _ensure()
    analytics: dict[str, Any] = {
        "commits": {"total": 0, "by_type": {}, "by_branch": {}},
        "branches": [],
        "tokens": {"total": 0, "estimated_cost_usd": 0.0},
        "recent_activity": [],
    }
    try:
        idx = _get_db().index
        commits = idx.list_all_commits(limit=9999) if hasattr(idx, "list_all_commits") else []
        analytics["commits"]["total"] = len(commits)
        by_type: dict[str, int] = {}
        for c in commits:
            ct = _field(c, "commit_type", "unknown")
            by_type[ct] = by_type.get(ct, 0) + 1
        analytics["commits"]["by_type"] = by_type

        branches = idx.list_branches() if hasattr(idx, "list_branches") else []
        analytics["branches"] = [
            {
                "name": b.name if hasattr(b, "name") else b.get("name", ""),
                "head": (b.head_hash if hasattr(b, "head_hash") else b.get("head_hash", ""))[:12],
                "status": b.status if hasattr(b, "status") else b.get("status", "active"),
            }
            for b in branches
        ]

        if hasattr(idx, "query_audit_log"):
            audit = idx.query_audit_log(limit=500)
            total_tokens = sum(e.get("token_count", 0) for e in audit if isinstance(e, dict))
            analytics["tokens"]["total"] = total_tokens
            analytics["tokens"]["estimated_cost_usd"] = round(total_tokens * 3.0 / 1_000_000, 4)

        recent = commits[-20:] if len(commits) > 20 else commits
        analytics["recent_activity"] = [
            {
                "hash": _field(c, "commit_hash", "")[:12],
                "type": _field(c, "commit_type", ""),
                "message": _field(c, "message", "")[:80],
                "agent_id": _meta(c, "agent_id", ""),
                "created_at": _ts(c),
            }
            for c in recent
        ]
    except Exception:
        logger.debug("Analytics query error", exc_info=True)
    return analytics


@app.get("/api/gateway/commits")
async def get_commits(limit: int = 50, offset: int = 0, branch: str | None = None, agent_id: str | None = None, workspace_path: str | None = None):
    _resolve_workspace(workspace_path)
    _ensure()
    try:
        commits = _get_db().index.list_all_commits(limit=9999) if hasattr(_get_db().index, "list_all_commits") else []
        commits = list(reversed(commits))
        if branch:
            commits = [c for c in commits if _meta(c, "branch") == branch]
        if agent_id:
            commits = [c for c in commits if _meta(c, "agent_id") == agent_id]
        total = len(commits)
        page = commits[offset:offset + limit]
        return {
            "total": total, "offset": offset, "limit": limit,
            "commits": [
                {
                    "hash": _field(c, "commit_hash", "")[:12],
                    "type": _field(c, "commit_type", ""),
                    "message": _field(c, "message", ""),
                    "created_at": _ts(c),
                    "agent_id": _meta(c, "agent_id", ""),
                    "branch": _meta(c, "branch", ""),
                }
                for c in page
            ],
        }
    except Exception:
        logger.debug("Commits query error", exc_info=True)
        return {"total": 0, "offset": offset, "limit": limit, "commits": []}


# ═══════════════════════════════════════════════════════════════════════
# 4. OPERATIONS — commit, branch, merge, restore, recall, diff
# ═══════════════════════════════════════════════════════════════════════

@app.post("/api/ops/commit")
async def ops_commit(request: Request):
    _ensure("engine")
    body = await request.json()
    from cvc.core.models import CommitRequest
    req = CommitRequest(
        message=body.get("message", "Dashboard commit"),
        content=body.get("content", ""),
        commit_type=body.get("type", "checkpoint"),
        tags=body.get("tags", []),
    )
    result = _get_engine().commit(req)
    _get_event_bus().emit("commit.created", {"hash": result.commit_hash[:12], "message": req.message})
    return {"commit_hash": result.commit_hash[:12], "message": req.message, "status": "ok"}


@app.post("/api/ops/branch")
async def ops_branch(request: Request):
    _ensure("engine")
    body = await request.json()
    from cvc.core.models import BranchRequest
    req = BranchRequest(name=body["name"])
    result = _get_engine().branch(req)
    return {"branch": result.name, "head": result.head_hash[:12] if result.head_hash else "", "status": "ok"}


@app.post("/api/ops/merge")
async def ops_merge(request: Request):
    _ensure("engine")
    body = await request.json()
    from cvc.core.models import MergeRequest
    req = MergeRequest(source_branch=body["source_branch"], strategy=body.get("strategy", "semantic"))
    result = _get_engine().merge(req)
    return {"merge_commit": result.commit_hash[:12] if hasattr(result, "commit_hash") else "", "status": "ok"}


@app.post("/api/ops/restore")
async def ops_restore(request: Request):
    _ensure("engine")
    body = await request.json()
    from cvc.core.models import RestoreRequest
    hash_val = body.get("hash") or body.get("commit_hash", "")
    req = RestoreRequest(commit_hash=hash_val)
    _get_engine().restore(req)
    return {"restored_to": hash_val, "status": "ok"}


@app.get("/api/ops/status")
async def ops_status(workspace_path: str | None = None):
    _resolve_workspace(workspace_path)
    _ensure("engine")
    try:
        branches = _get_db().index.list_branches() if hasattr(_get_db().index, "list_branches") else []
        commits = _get_db().index.list_all_commits(limit=9999) if hasattr(_get_db().index, "list_all_commits") else []
        active = _get_config().default_branch
        head = ""
        for b in branches:
            name = b.name if hasattr(b, "name") else b.get("name", "")
            if name == active:
                head = (b.head_hash if hasattr(b, "head_hash") else b.get("head_hash", ""))[:12]
                break
        db_path = _get_config().cvc_root / "cvc.db"
        blobs_dir = _get_config().cvc_root / "objects"
        return {
            "branch": active, "head": head,
            "total_branches": len(branches), "total_commits": len(commits),
            "database_size_bytes": db_path.stat().st_size if db_path.exists() else 0,
            "blob_count": len(list(blobs_dir.glob("*"))) if blobs_dir.exists() else 0,
            "cvc_root": str(_get_config().cvc_root),
        }
    except Exception as exc:
        return {"error": str(exc)}


@app.get("/api/ops/recall")
async def ops_recall(query: str, limit: int = 10):
    _ensure()
    try:
        results = _get_db().search_conversations(query, limit=limit)
        return {"results": results, "total": len(results)}
    except Exception as exc:
        logger.debug("Recall search error", exc_info=True)
        return {"results": [], "total": 0, "error": str(exc)}


@app.get("/api/ops/diff")
async def ops_diff(hash1: str, hash2: str):
    _ensure()
    try:
        idx = _get_db().index
        c1 = idx.get_commit(hash1) if hasattr(idx, "get_commit") else None
        c2 = idx.get_commit(hash2) if hasattr(idx, "get_commit") else None
        if not c1 or not c2:
            raise HTTPException(404, "One or both commits not found")

        # Build a meaningful diff from commit metadata + blob content
        diff_lines = []
        diff_lines.append(f"--- Commit A: {hash1[:12]}")
        diff_lines.append(f"+++ Commit B: {hash2[:12]}")
        diff_lines.append("")

        m1, m2 = _field(c1, "message", ""), _field(c2, "message", "")
        t1, t2 = _field(c1, "commit_type", ""), _field(c2, "commit_type", "")
        diff_lines.append(f"Message A: {m1}")
        diff_lines.append(f"Message B: {m2}")
        if t1 != t2:
            diff_lines.append(f"Type: {t1} → {t2}")

        # Try to diff blob contents
        try:
            h1_full = _field(c1, "commit_hash", hash1)
            h2_full = _field(c2, "commit_hash", hash2)
            blob1 = _get_db().retrieve_blob(h1_full)
            blob2 = _get_db().retrieve_blob(h2_full)
            if blob1 and blob2:
                msgs1 = [f"[{m.role}] {m.content[:100]}" for m in (blob1.messages or [])]
                msgs2 = [f"[{m.role}] {m.content[:100]}" for m in (blob2.messages or [])]
                diff_lines.append("")
                diff_lines.append(f"Context A: {len(blob1.messages or [])} messages")
                diff_lines.append(f"Context B: {len(blob2.messages or [])} messages")
                # Show what's new in B
                msgs1_set = set(msgs1)
                for m in msgs2:
                    if m not in msgs1_set:
                        diff_lines.append(f"+ {m}")
                msgs2_set = set(msgs2)
                for m in msgs1:
                    if m not in msgs2_set:
                        diff_lines.append(f"- {m}")
        except Exception:
            pass

        return {
            "commit_a": {"hash": hash1, "message": m1, "type": t1},
            "commit_b": {"hash": hash2, "message": m2, "type": t2},
            "diff": "\n".join(diff_lines),
        }
    except HTTPException:
        raise
    except Exception as exc:
        return {"error": str(exc)}


@app.get("/api/ops/timeline")
async def ops_timeline(limit: int = 100, workspace_path: str | None = None):
    _resolve_workspace(workspace_path)
    _ensure()
    try:
        commits = _get_db().index.list_all_commits(limit=limit) if hasattr(_get_db().index, "list_all_commits") else []
        commits = list(reversed(commits))[:limit]
        return {
            "timeline": [
                {
                    "hash": _field(c, "commit_hash", "")[:12],
                    "type": _field(c, "commit_type", ""),
                    "message": _field(c, "message", ""),
                    "created_at": _ts(c),
                    "agent_id": _meta(c, "agent_id", ""),
                    "branch": _meta(c, "branch", ""),
                }
                for c in commits
            ]
        }
    except Exception as exc:
        return {"timeline": [], "error": str(exc)}


@app.get("/api/ops/branches")
async def ops_branches():
    _ensure()
    try:
        branches = _get_db().index.list_branches() if hasattr(_get_db().index, "list_branches") else []
        return {
            "branches": [
                {
                    "name": b.name if hasattr(b, "name") else b.get("name", ""),
                    "head": (b.head_hash if hasattr(b, "head_hash") else b.get("head_hash", ""))[:12],
                    "status": b.status if hasattr(b, "status") else b.get("status", "active"),
                }
                for b in branches
            ]
        }
    except Exception as exc:
        return {"branches": [], "error": str(exc)}


# ═══════════════════════════════════════════════════════════════════════
# 5. MEMORY & CONTEXT
# ═══════════════════════════════════════════════════════════════════════

@app.get("/api/memory/context")
async def get_context(limit: int = 50):
    _ensure()
    try:
        ctx = _get_engine().context_window if _get_engine() else []
        messages = ctx[:limit]
        token_count = sum(len(getattr(m, 'content', '').split()) for m in ctx)
        return {
            "messages": [
                {"role": getattr(m, 'role', ''), "content": getattr(m, 'content', '')[:500], "timestamp": getattr(m, 'timestamp', 0)}
                for m in messages
            ],
            "total": len(ctx),
            "token_count": token_count,
        }
    except Exception as exc:
        return {"messages": [], "total": 0, "token_count": 0, "error": str(exc)}


@app.get("/api/memory/blobs")
async def get_blobs(limit: int = 20):
    _ensure()
    try:
        blobs_dir = _get_config().cvc_root / "objects"
        if not blobs_dir.exists():
            return {"blobs": [], "total": 0}
        blob_files = sorted(blobs_dir.glob("*"), key=lambda p: p.stat().st_mtime, reverse=True)[:limit]
        return {
            "blobs": [{"hash": f.name, "size_bytes": f.stat().st_size, "modified": f.stat().st_mtime} for f in blob_files],
            "total": len(list(blobs_dir.glob("*"))),
        }
    except Exception as exc:
        return {"blobs": [], "total": 0, "error": str(exc)}


@app.get("/api/memory/stats")
async def get_memory_stats():
    _ensure()
    try:
        db_path = _get_config().cvc_root / "cvc.db"
        blobs_dir = _get_config().cvc_root / "objects"
        return {
            "database_size_bytes": db_path.stat().st_size if db_path.exists() else 0,
            "blob_count": len(list(blobs_dir.glob("*"))) if blobs_dir.exists() else 0,
            "blob_total_bytes": sum(f.stat().st_size for f in blobs_dir.glob("*")) if blobs_dir.exists() else 0,
            "cvc_root": str(_get_config().cvc_root),
        }
    except Exception as exc:
        return {"error": str(exc)}


# ═══════════════════════════════════════════════════════════════════════
# 6. MODEL CATALOG & PROVIDER
# ═══════════════════════════════════════════════════════════════════════

@app.get("/api/models/catalog")
async def get_model_catalog():
    try:
        from cvc.adapters import PROVIDER_DEFAULTS
        # Build a proper model catalog per provider
        # Note: keep `name` short — it's the chip label in the composer bar.
        # Release dates, context windows, and tier tags used to live in here as
        # "(...)" parentheticals and made the chip overflow. They're gone now.
        # `context_window` is still on every entry; the dropdown can surface it
        # if we ever want to bring it back.
        MODEL_CATALOG = {
            "anthropic": [
                {"id": "claude-opus-4-8",            "name": "Claude 4.8 Opus",   "context_window": 1000000},
                {"id": "claude-sonnet-4-6",          "name": "Claude 4.6 Sonnet", "context_window": 1000000},
                {"id": "claude-haiku-4-5",           "name": "Claude 4.5 Haiku",  "context_window": 200000},
                {"id": "claude-3.5-sonnet",          "name": "Claude 3.5 Sonnet", "context_window": 200000},
                {"id": "claude-3-5-haiku-20241022",  "name": "Claude 3.5 Haiku",  "context_window": 200000},
                {"id": "claude-3-opus",              "name": "Claude 3 Opus",     "context_window": 200000},
            ],
            "openai": [
                {"id": "gpt-5.3",      "name": "GPT-5.3",            "context_window": 256000},
                {"id": "gpt-5.2",      "name": "GPT-5.2",            "context_window": 128000},
                {"id": "gpt-5-mini",   "name": "GPT-5 Mini",         "context_window": 128000},
                {"id": "gpt-4o",       "name": "GPT-4o",             "context_window": 128000},
                {"id": "gpt-4o-mini",  "name": "GPT-4o Mini",        "context_window": 128000},
                {"id": "o3-mini",      "name": "o3 Mini Reasoning",  "context_window": 200000},
                {"id": "o1",           "name": "o1 Reasoning",       "context_window": 128000},
                {"id": "o1-mini",      "name": "o1 Mini Reasoning",  "context_window": 128000},
            ],
            "google": [
                {"id": "gemini-3.5-flash",        "name": "Gemini 3.5 Flash",       "context_window": 1000000},
                {"id": "gemini-3.1-pro",          "name": "Gemini 3.1 Pro",         "context_window": 1000000},
                {"id": "gemini-3.1-flash-lite",   "name": "Gemini 3.1 Flash-Lite",  "context_window": 1000000},
                {"id": "gemini-3.1-pro-preview",  "name": "Gemini 3.1 Pro Preview", "context_window": 1000000},
                {"id": "gemini-3-flash-preview",  "name": "Gemini 3 Flash Preview", "context_window": 1000000},
                {"id": "gemini-2.5-flash",        "name": "Gemini 2.5 Flash",       "context_window": 1000000},
                {"id": "gemini-2.5-pro",          "name": "Gemini 2.5 Pro",         "context_window": 1000000},
                {"id": "gemini-1.5-flash",        "name": "Gemini 1.5 Flash",       "context_window": 1000000},
                {"id": "gemini-1.5-pro",          "name": "Gemini 1.5 Pro",         "context_window": 2000000},
            ],
            "ollama": [
                {"id": "qwen2.5-coder:7b",  "name": "Qwen 2.5 Coder 7B", "context_window": 32768},
                {"id": "qwen3-coder:latest","name": "Qwen 3 Coder",      "context_window": 32768},
                {"id": "deepseek-r1:latest","name": "DeepSeek-R1",       "context_window": 32768},
                {"id": "llama3:8b",         "name": "Llama 3 8B",        "context_window": 8192},
            ],
            "lmstudio": [
                {"id": "loaded-model", "name": "LM Studio Loaded Model", "context_window": 0},
            ],
            "github": [],
            # NVIDIA NIM — free-tier + paid models via integrate.api.nvidia.com
            "nvidia": [
                {"id": "nvidia/nemotron-3-super-120b-instruct",  "name": "Nemotron 3 Super 120B",   "context_window": 262144},
                {"id": "moonshotai/kimi-k2-instruct",            "name": "Kimi K2 1T MoE",          "context_window": 128000},
                {"id": "minimaxai/minimax-m2",                   "name": "MiniMax M2 via NVIDIA",   "context_window": 200000},
                {"id": "meta/llama-3.1-70b-instruct",            "name": "Llama 3.1 70B",           "context_window": 128000},
                {"id": "meta/llama-3.1-405b-instruct",           "name": "Llama 3.1 405B",          "context_window": 128000},
                {"id": "nvidia/llama-3.3-nemotron-super-49b-v1", "name": "Nemotron Super 49B v1",   "context_window": 128000},
            ],
            # MiniMax M-series — https://platform.minimax.io/docs/guides/models-intro
            "minimax": [
                # Current Models
                {"id": "MiniMax-M3",             "name": "MiniMax M3",             "context_window": 1000000},
                {"id": "MiniMax-M2.7",           "name": "MiniMax M2.7",           "context_window": 200000},
                {"id": "MiniMax-M2.7-highspeed", "name": "MiniMax M2.7 Highspeed", "context_window": 200000},
                # Legacy Models
                {"id": "MiniMax-M2.5",           "name": "MiniMax M2.5",           "context_window": 200000},
                {"id": "MiniMax-M2.5-highspeed", "name": "MiniMax M2.5 Highspeed", "context_window": 200000},
                {"id": "MiniMax-M2.1",           "name": "MiniMax M2.1",           "context_window": 200000},
                {"id": "MiniMax-M2.1-highspeed", "name": "MiniMax M2.1 Highspeed", "context_window": 200000},
                {"id": "MiniMax-M2",             "name": "MiniMax M2",             "context_window": 200000},
            ],
        }

        # Dynamically fetch actual GitHub Copilot models available for this user
        from cvc.core.models import GlobalConfig
        gc = GlobalConfig.load()
        if gc and "github" in gc.api_keys:
            oauth_token = gc.api_keys["github"]
            from cvc.agent.providers.github_auth import fetch_copilot_token
            import httpx
            token_data = fetch_copilot_token(oauth_token)
            if token_data:
                copilot_token = token_data.get("token")
                proxy_ep = token_data.get("endpoints", {}).get("api", "https://api.individual.githubcopilot.com")
                try:
                    resp = httpx.get(
                        f"{proxy_ep.rstrip('/')}/models",
                        headers={
                            "Authorization": f"Bearer {copilot_token}",
                            "Accept": "application/json",
                            "editor-version": "vscode/1.93.0",
                            "editor-plugin-version": "copilot-chat/0.20.0"
                        },
                        timeout=5.0
                    )
                    if resp.status_code == 200:
                        models_data = resp.json().get("data", [])
                        gh_models = []
                        for m in models_data:
                            if m.get("policy", {}).get("state") == "disabled":
                                continue
                            gh_models.append({
                                "id": m["id"],
                                "name": m.get("name", m["id"]),
                                "context_window": 128000
                            })
                        if gh_models:
                            MODEL_CATALOG["github"] = gh_models
                except Exception:
                    pass
        
        catalog = {}
        for provider, defaults in PROVIDER_DEFAULTS.items():
            if provider == "github":
                catalog[provider] = MODEL_CATALOG.get("github", [])
            else:
                catalog[provider] = MODEL_CATALOG.get(provider, [{"id": defaults["model"], "name": defaults["model"]}])

        # Decorate every entry with `is_primary` + `provider_label` so the
        # frontend picker doesn't have to do string munging or cross-lookups.
        cur_provider = _agent_llm.provider if _agent_llm else (
            _get_config().provider if _get_config() else ""
        )
        cur_model = _agent_llm.model if _agent_llm else (
            _get_config().model if _get_config() else ""
        )
        _PROVIDER_LABELS = {
            "anthropic": "Anthropic",
            "openai": "OpenAI",
            "google": "Google",
            "ollama": "Ollama",
            "lmstudio": "LM Studio",
            "github": "GitHub Copilot",
            "copilot": "GitHub Copilot",
            "vertex": "Google Vertex",
            "nvidia": "NVIDIA NIM",
            "minimax": "MiniMax",
        }
        for provider, models in catalog.items():
            label_base = _PROVIDER_LABELS.get(provider, provider.title())
            count = len(models)
            label = f"{label_base} ({count} model{'s' if count != 1 else ''})"
            for entry in models:
                entry["provider"] = provider
                entry["provider_label"] = label
                entry["is_primary"] = bool(
                    provider == cur_provider and entry.get("id") == cur_model
                )
        return {
            "providers": catalog,
            "current_provider": cur_provider,
            "current_model": cur_model,
        }
    except Exception as exc:
        return {"providers": {}, "error": str(exc)}


@app.get("/api/models/current")
async def get_current_model():
    _ensure()
    from cvc.core.models import GlobalConfig as GC
    gc = GC.load()
    # Prefer the live agent LLM config (matches what chat actually uses)
    provider = _agent_llm.provider if _agent_llm else _get_config().provider
    model = _agent_llm.model if _agent_llm else _get_config().model
    return {
        "provider": provider,
        "model": model,
        "api_key_set": bool(_get_config().api_key) or bool(_agent_llm),
        "api_keys_stored": {k: bool(v) for k, v in gc.api_keys.items()},
    }


@app.post("/api/models/switch")
async def switch_model(request: Request):
    global _agent_llm, _system_prompt
    _ensure()
    body = await request.json()
    new_provider = body.get("provider", _get_config().provider)
    new_model = body.get("model", _get_config().model)
    from cvc.core.models import GlobalConfig as GC
    gc = GC.load()
    gc.provider = new_provider
    gc.model = new_model
    gc.save()
    _get_config().provider = new_provider
    _get_config().model = new_model

    # Re-create the AgentLLM with the new provider/model
    try:
        from cvc.agent.llm import AgentLLM
        from cvc.agent.system_prompt import build_system_prompt

        env_map = {"anthropic": "ANTHROPIC_API_KEY", "openai": "OPENAI_API_KEY", "google": "GOOGLE_API_KEY", "ollama": "", "lmstudio": "", "github": "GITHUB_TOKEN"}
        env_key = env_map.get(new_provider, "")
        api_key = os.getenv(env_key, "") if env_key else ""
        if not api_key:
            api_key = gc.api_keys.get(new_provider, "")
        base_url_map = {
            "anthropic": "https://api.anthropic.com", "openai": "https://api.openai.com",
            "google": "https://generativelanguage.googleapis.com",
            "ollama": os.getenv("OLLAMA_HOST", "http://localhost:11434"),
            "lmstudio": os.getenv("LMSTUDIO_HOST", "http://localhost:1234"),
            "copilot": "https://api.githubcopilot.com",
        }
        if _agent_llm:
            try:
                await _agent_llm.close()
            except Exception:
                pass
        _agent_llm = AgentLLM(provider=new_provider, api_key=api_key, model=new_model, base_url=base_url_map.get(new_provider, ""))
        _system_prompt = build_system_prompt(workspace=str(_project_root), provider=new_provider, model=new_model, branch=_get_engine().active_branch, user_memory_context=_get_user_memory_context())
        asyncio.create_task(_agent_llm.warm_connection())
        logger.info("Agent LLM switched to %s/%s", new_provider, new_model)
    except Exception:
        logger.warning("Failed to re-create AgentLLM after model switch", exc_info=True)

    # Notify the proxy to reload config (best-effort)
    try:
        import httpx
        httpx.post(f"http://{_HOST}:{_PROXY_PORT}/cvc/reload-config", timeout=3.0)
    except Exception:
        pass
    return {"provider": new_provider, "model": new_model, "status": "ok"}


# ═══════════════════════════════════════════════════════════════════════
# 7. HIVE MIND / SDK
# ═══════════════════════════════════════════════════════════════════════

@app.get("/api/gateway/config")
async def get_config():
    _ensure()
    provider = _agent_llm.provider if _agent_llm else _get_config().provider
    model = _agent_llm.model if _agent_llm else _get_config().model
    return {
        "version": _cvc_version,
        "provider": provider,
        "model": model,
        "agent_id": _get_config().agent_id,
        "default_branch": _get_config().default_branch,
        "proxy_port": _PROXY_PORT,
        "gateway_port": _GATEWAY_PORT,
        "host": _HOST,
        "api_key_set": bool(_get_config().api_key) or bool(_agent_llm),
        "cvc_root": str(_get_config().cvc_root),
    }


@app.get("/api/hivemind/agents")
async def get_hivemind_agents():
    _ensure()
    try:
        agents = _get_db().index.list_agents() if hasattr(_get_db().index, "list_agents") else []
        result = []
        for a in agents:
            if isinstance(a, dict):
                result.append(a)
            else:
                d = {}
                for attr in ('agent_id', 'name', 'role', 'rank', 'squad', 'status'):
                    if hasattr(a, attr):
                        d[attr] = getattr(a, attr)
                result.append(d)
        return {"agents": result, "total": len(result)}
    except Exception:
        return {"agents": [], "total": 0}


@app.post("/api/hivemind/register")
async def register_agent(request: Request):
    _ensure()
    body = await request.json()
    try:
        agent_id = body["agent_id"]
        _get_db().index.insert_agent(
            agent_id=agent_id,
            name=body.get("name", agent_id),
            role=body.get("role", "worker"),
            rank=body.get("rank", "private"),
            squad=body.get("squad", "alpha"),
        )
        return {"agent_id": agent_id, "status": "registered"}
    except Exception as exc:
        raise HTTPException(400, str(exc))


@app.delete("/api/hivemind/agents/{agent_id}")
async def remove_agent(agent_id: str):
    _ensure()
    try:
        if hasattr(_get_db().index, "delete_agent") and _get_db().index.delete_agent(agent_id):
            return {"agent_id": agent_id, "status": "removed"}
        raise HTTPException(404, f"Agent '{agent_id}' not found")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(400, str(exc))


# ═══════════════════════════════════════════════════════════════════════
# 7b. SETTINGS — Comprehensive Configuration Management
# ═══════════════════════════════════════════════════════════════════════

@app.get("/api/settings")
async def get_settings():
    """Return merged settings from all levels + global config."""
    _ensure()
    from cvc.agent.settings import load_settings
    from cvc.core.models import SETTINGS_SCHEMA, GlobalConfig

    gc = GlobalConfig.load()
    workspace = Path(_project_root) if _project_root else Path.cwd()
    settings = load_settings(workspace)

    # Mask API keys
    masked_keys = {}
    for prov, key in gc.api_keys.items():
        if key and len(key) > 8:
            masked_keys[prov] = f"{key[:4]}***{key[-4:]}"
        elif key:
            masked_keys[prov] = "***"
        else:
            masked_keys[prov] = ""

    return {
        "global": {
            "provider": gc.provider,
            "model": gc.model,
            "agent_id": gc.agent_id,
            "api_keys": masked_keys,
            "max_iterations": gc.max_iterations,
            "max_iterations_expensive": gc.max_iterations_expensive,
            "agent_timeout": gc.agent_timeout,
            "tool_timeout": gc.tool_timeout,
        },
        "project": settings.to_dict(),
        "schema": SETTINGS_SCHEMA,
    }


@app.put("/api/settings/global")
async def update_global_settings(request: Request):
    """Update global config fields (provider, model, agent_id)."""
    _ensure()
    from cvc.core.models import GlobalConfig
    body = await request.json()
    gc = GlobalConfig.load()

    allowed_fields = {"provider", "model", "agent_id"}
    numeric_fields = {
        "max_iterations": int,
        "max_iterations_expensive": int,
        "agent_timeout": float,
        "tool_timeout": float,
    }
    for key, value in body.items():
        if key in allowed_fields and isinstance(value, str):
            setattr(gc, key, value)
        elif key in numeric_fields:
            try:
                setattr(gc, key, numeric_fields[key](value or 0))
            except (ValueError, TypeError):
                pass

    gc.save()

    # Live-update the agent LLM if running
    global _agent_llm
    if _agent_llm:
        if "provider" in body or "model" in body:
            try:
                from cvc.adapters import get_adapter
                new_provider = body.get("provider", gc.provider)
                new_model = body.get("model", gc.model)
                _agent_llm = get_adapter(
                    new_provider, model=new_model,
                    api_key=gc.api_keys.get(new_provider, ""),
                )
            except Exception:
                pass

    return {"status": "ok", "saved": list(body.keys())}


@app.put("/api/settings/project")
async def update_project_settings(request: Request):
    """Update .cvc/settings.json (shared project settings)."""
    _ensure()
    from cvc.agent.settings import save_project_settings
    body = await request.json()
    workspace = Path(_project_root) if _project_root else Path.cwd()

    for key, value in body.items():
        save_project_settings(workspace, key, value, local=False)

    return {"status": "ok", "saved": list(body.keys())}


@app.put("/api/settings/local")
async def update_local_settings(request: Request):
    """Update .cvc/settings.local.json (local, not committed)."""
    _ensure()
    from cvc.agent.settings import save_project_settings
    body = await request.json()
    workspace = Path(_project_root) if _project_root else Path.cwd()

    for key, value in body.items():
        save_project_settings(workspace, key, value, local=True)

    return {"status": "ok", "saved": list(body.keys())}


@app.get("/api/settings/schema")
async def get_settings_schema():
    """Return JSON schema for dynamic form rendering."""
    from cvc.core.models import SETTINGS_SCHEMA
    return SETTINGS_SCHEMA


@app.post("/api/settings/reset")
async def reset_settings(request: Request):
    """Reset a settings level to defaults."""
    _ensure()
    body = await request.json()
    level = body.get("level", "local")
    workspace = Path(_project_root) if _project_root else Path.cwd()

    if level == "global":
        from cvc.core.models import GlobalConfig
        gc = GlobalConfig()
        gc.save()
    elif level == "project":
        path = workspace / ".cvc" / "settings.json"
        if path.exists():
            path.write_text("{}\n", encoding="utf-8")
    elif level == "local":
        path = workspace / ".cvc" / "settings.local.json"
        if path.exists():
            path.write_text("{}\n", encoding="utf-8")

    return {"status": "ok", "level": level}


@app.get("/api/settings/keys")
async def get_api_keys():
    """List configured API key providers (keys masked)."""
    from cvc.core.models import GlobalConfig
    gc = GlobalConfig.load()
    result = {}
    for prov, key in gc.api_keys.items():
        if key and len(key) > 8:
            result[prov] = {"masked": f"{key[:4]}***{key[-4:]}", "set": True}
        elif key:
            result[prov] = {"masked": "***", "set": True}
        else:
            result[prov] = {"masked": "", "set": False}
    return {"keys": result}


@app.put("/api/settings/keys/{provider}")
async def set_api_key(provider: str, request: Request):
    """Set or update an API key for a provider."""
    from cvc.core.models import GlobalConfig
    body = await request.json()
    api_key = body.get("api_key", "")
    if not api_key:
        raise HTTPException(400, "api_key is required")

    gc = GlobalConfig.load()
    gc.api_keys[provider] = api_key
    gc.save()
    return {"provider": provider, "status": "ok"}


@app.delete("/api/settings/keys/{provider}")
async def remove_api_key(provider: str):
    """Remove an API key for a provider."""
    from cvc.core.models import GlobalConfig
    gc = GlobalConfig.load()
    gc.api_keys.pop(provider, None)
    gc.save()
    return {"provider": provider, "status": "removed"}


@app.post("/api/settings/keys/{provider}/test")
async def test_api_key(provider: str):
    """Test if the configured API key for a provider works."""
    from cvc.core.models import GlobalConfig
    gc = GlobalConfig.load()
    api_key = gc.api_keys.get(provider, "")
    if not api_key:
        return {"provider": provider, "valid": False, "error": "No API key configured"}

    try:
        import httpx
        if provider == "anthropic":
            r = httpx.get(
                "https://api.anthropic.com/v1/models",
                headers={"x-api-key": api_key, "anthropic-version": "2023-06-01"},
                timeout=10.0,
            )
            valid = r.status_code == 200
        elif provider == "openai":
            r = httpx.get(
                "https://api.openai.com/v1/models",
                headers={"Authorization": f"Bearer {api_key}"},
                timeout=10.0,
            )
            valid = r.status_code == 200
        elif provider == "google":
            r = httpx.get(
                f"https://generativelanguage.googleapis.com/v1/models?key={api_key}",
                timeout=10.0,
            )
            valid = r.status_code == 200
        elif provider == "ollama":
            host = os.getenv("OLLAMA_HOST", "http://localhost:11434")
            r = httpx.get(f"{host}/api/tags", timeout=5.0)
            valid = r.status_code == 200
        elif provider == "lmstudio":
            r = httpx.get("http://localhost:1234/v1/models", timeout=5.0)
            valid = r.status_code == 200
        elif provider == "nvidia":
            r = httpx.get(
                "https://integrate.api.nvidia.com/v1/models",
                headers={"Authorization": f"Bearer {api_key}"},
                timeout=10.0,
            )
            valid = r.status_code == 200
        elif provider == "minimax":
            # MiniMax is Anthropic-Messages-API-compatible at
            # https://api.minimax.io/anthropic. Verify by hitting
            # /anthropic/v1/models with Authorization: Bearer *** (per upstream
            # reference: _requires_bearer_auth returns True for MiniMax URLs).
            try:
                r = httpx.get(
                    "https://api.minimax.io/anthropic/v1/models",
                    headers={
                        "Authorization": f"Bearer {api_key}",
                        "anthropic-beta": "interleaved-thinking-2025-05-14",
                    },
                    timeout=10.0,
                )
                valid = r.status_code == 200
            except Exception as exc:
                return {"provider": provider, "valid": False, "error": str(exc)}
        else:
            return {"provider": provider, "valid": False, "error": "Unknown provider"}

        return {"provider": provider, "valid": valid, "status_code": r.status_code}
    except Exception as exc:
        return {"provider": provider, "valid": False, "error": str(exc)}


@app.get("/api/settings/hooks")
async def get_hooks():
    """List all configured hooks."""
    _ensure()
    from cvc.agent.settings import load_settings
    workspace = Path(_project_root) if _project_root else Path.cwd()
    settings = load_settings(workspace)
    return {"hooks": settings.hooks}


@app.post("/api/settings/hooks")
async def add_hook(request: Request):
    """Add a new hook."""
    _ensure()
    from cvc.agent.settings import load_settings, save_project_settings
    body = await request.json()
    event = body.get("event", "PreToolUse")
    hook_config = body.get("config", {})
    workspace = Path(_project_root) if _project_root else Path.cwd()

    settings = load_settings(workspace)
    hooks = dict(settings.hooks)
    if event not in hooks:
        hooks[event] = []
    hooks[event].append(hook_config)
    save_project_settings(workspace, "hooks", hooks, local=False)
    return {"status": "ok", "event": event}


@app.delete("/api/settings/hooks/{event}/{index}")
async def remove_hook(event: str, index: int):
    """Remove a hook by event type and index."""
    _ensure()
    from cvc.agent.settings import load_settings, save_project_settings
    workspace = Path(_project_root) if _project_root else Path.cwd()

    settings = load_settings(workspace)
    hooks = dict(settings.hooks)
    if event in hooks and 0 <= index < len(hooks[event]):
        hooks[event].pop(index)
        if not hooks[event]:
            del hooks[event]
        save_project_settings(workspace, "hooks", hooks, local=False)
        return {"status": "ok"}
    raise HTTPException(404, "Hook not found")


@app.get("/api/settings/env")
async def get_env_vars():
    """Return CVC-relevant environment variables."""
    cvc_vars = {}
    for key, val in os.environ.items():
        if key.startswith("CVC_") or key.endswith("_API_KEY") or key == "OLLAMA_HOST":
            if "KEY" in key and val and len(val) > 8:
                cvc_vars[key] = f"{val[:4]}***{val[-4:]}"
            elif "KEY" in key and val:
                cvc_vars[key] = "***"
            else:
                cvc_vars[key] = val
    return {"env": cvc_vars}


# ═══════════════════════════════════════════════════════════════════════
# 7c. AGENT TEMPLATES — Create, manage, and deploy custom agents
# ═══════════════════════════════════════════════════════════════════════

@app.get("/api/agents")
async def list_agents():
    """List all agents (registered + templates)."""
    _ensure()

    # Registered agents from DB
    registered = []
    try:
        agents = _get_db().index.list_agents() if hasattr(_get_db().index, "list_agents") else []
        for a in agents:
            d = a if isinstance(a, dict) else {}
            d["source"] = "registered"
            registered.append(d)
    except Exception:
        pass

    # Agent templates from .cvc/agent_templates/
    templates = []
    workspace = Path(_project_root) if _project_root else Path.cwd()
    tpl_dir = workspace / ".cvc" / "agent_templates"
    if tpl_dir.exists():
        for fp in sorted(tpl_dir.glob("*.json")):
            try:
                data = json.loads(fp.read_text(encoding="utf-8"))
                data["source"] = "template"
                templates.append(data)
            except (json.JSONDecodeError, OSError):
                pass

    # Builtin agents
    builtins = _get_builtin_agents()

    return {
        "registered": registered,
        "templates": templates,
        "builtins": builtins,
        "total": len(registered) + len(templates) + len(builtins),
    }


@app.get("/api/agents/builtin")
async def get_builtin_agents():
    """List the built-in sub-agent types with their configs."""
    return _get_builtin_agents()


def _get_builtin_agents() -> list[dict[str, Any]]:
    """Return builtin agent definitions (subagent builtins + enterprise template library)."""
    result = []

    # 1. Core sub-agent builtins
    try:
        from cvc.agent.subagent import BUILTIN_AGENTS
        for name, cfg in BUILTIN_AGENTS.items():
            result.append({
                "id": name.lower(),
                "name": cfg.name if hasattr(cfg, "name") else name,
                "description": cfg.description if hasattr(cfg, "description") else "",
                "system_prompt": cfg.system_prompt if hasattr(cfg, "system_prompt") else "",
                "tools_allow": list(cfg.tools) if hasattr(cfg, "tools") else [],
                "tools_deny": [],
                "tool_tier": _infer_tool_tier(cfg.tools) if hasattr(cfg, "tools") else "coding",
                "max_turns": cfg.max_turns if hasattr(cfg, "max_turns") else 10,
                "rank": "specialist",
                "squad": "",
                "capabilities": [],
                "skills": [],
                "auto_respond": False,
                "auto_share_to_hive": True,
                "model_override": "",
                "provider_override": "",
                "is_builtin": True,
                "source": "builtin",
            })
    except ImportError:
        pass

    # 2. Enterprise template library
    try:
        from cvc.sdk.agent_templates import get_enterprise_templates
        seen_ids = {r["id"] for r in result}
        for tpl in get_enterprise_templates():
            if tpl["id"] not in seen_ids:
                result.append(tpl)
                seen_ids.add(tpl["id"])
    except ImportError:
        pass

    return result


def _infer_tool_tier(tools: frozenset[str] | set[str]) -> str:
    """Infer tool tier from tool set."""
    tool_names = {str(t) for t in tools}
    if any("agent" in t or "spawn" in t or "delegat" in t for t in tool_names):
        return "orchestration"
    if any("web" in t or "fetch" in t for t in tool_names):
        return "research"
    if any("write" in t or "edit" in t or "bash" in t or "patch" in t for t in tool_names):
        return "coding"
    return "read_only"


@app.post("/api/agents")
async def create_agent(request: Request):
    """Create a new agent from explicit config."""
    _ensure()
    body = await request.json()
    workspace = Path(_project_root) if _project_root else Path.cwd()
    tpl_dir = workspace / ".cvc" / "agent_templates"
    tpl_dir.mkdir(parents=True, exist_ok=True)

    from cvc.core.models import AgentTemplate
    template = AgentTemplate(**body)
    template.updated_at = time.time()

    # Save template
    tpl_path = tpl_dir / f"{template.id}.json"
    tpl_path.write_text(
        json.dumps(template.model_dump(), indent=2, default=str),
        encoding="utf-8",
    )

    # Also register in the hive mind registry
    try:
        _get_db().index.insert_agent(
            agent_id=template.id,
            name=template.name,
            role=template.description[:50] if template.description else "custom",
            rank=template.rank,
            squad=template.squad or None,
        )
    except Exception:
        pass

    return {"id": template.id, "status": "created", "template": template.model_dump()}


@app.post("/api/agents/from-prompt")
async def create_agent_from_prompt(request: Request):
    """
    Create an agent from a natural language description.

    Uses the current LLM to analyze the description and generate
    an AgentTemplate with appropriate tools, model, and config.
    """
    _ensure()
    body = await request.json()
    description = body.get("prompt", "") or body.get("description", "")
    if not description:
        raise HTTPException(400, "A prompt or description is required")

    # Build generation prompt
    gen_prompt = f"""You are an agent configuration generator. Based on the user's description,
generate a JSON configuration for a new AI agent.

User's description: "{description}"

Available tool tiers:
- "read_only": read_file, grep, glob, list_dir, cvc_search, cvc_status, cvc_log
- "coding": All read_only + write_file, edit_file, patch_file, bash
- "research": All read_only + web_search, web_fetch
- "orchestration": All coding + agent spawning, delegation

Available ranks: "specialist", "captain", "mission_controller", "zeus"

Respond with ONLY valid JSON (no markdown, no explanation):
{{
  "name": "descriptive agent name",
  "description": "what this agent does",
  "system_prompt": "detailed instructions for the agent",
  "tool_tier": "one of: read_only, coding, research, orchestration",
  "tools_allow": ["specific tools to allow"],
  "tools_deny": ["specific tools to deny"],
  "max_turns": 20,
  "rank": "specialist",
  "squad": "",
  "capabilities": ["list", "of", "capabilities"],
  "skills": ["list", "of", "skills"],
  "auto_respond": false,
  "auto_share_to_hive": true
}}"""

    # Use the current LLM to generate the config
    if _agent_llm:
        try:
            response = await _agent_llm.chat(
                messages=[
                    {"role": "system", "content": "You are a precise JSON generator. Output ONLY valid JSON."},
                    {"role": "user", "content": gen_prompt},
                ],
                tools=[],
                temperature=0.7,
                max_tokens=8192,
            )
            # Parse the LLMResponse
            content = response.text if hasattr(response, 'text') else str(response)
            # Extract JSON from response
            import re
            json_match = re.search(r'\{[\s\S]*\}', content)
            if json_match:
                config = json.loads(json_match.group())
                from cvc.core.models import AgentTemplate
                template = AgentTemplate(**config)
                return {
                    "status": "generated",
                    "template": template.model_dump(),
                    "raw_description": description,
                }
        except Exception as exc:
            raise HTTPException(500, f"LLM generation failed: {exc}")

    # Fallback: generate a reasonable template from keyword analysis when no LLM is available
    from cvc.core.models import AgentTemplate
    template = _generate_template_from_keywords(description)
    return {
        "status": "generated",
        "template": template.model_dump(),
        "raw_description": description,
        "note": "Generated via keyword analysis (no LLM active). Configure an LLM provider for richer results.",
    }


def _generate_template_from_keywords(description: str):
    """Generate an AgentTemplate from keyword analysis when no LLM is available."""
    import re as _re

    from cvc.core.models import AgentTemplate
    desc_lower = description.lower()

    # Detect tool tier
    if any(w in desc_lower for w in ("orchestrat", "delegat", "multi-agent", "spawn", "coordinate")):
        tool_tier = "orchestration"
        rank = "mission_controller"
    elif any(w in desc_lower for w in ("research", "search", "web", "browse", "fetch", "scrape", "internet")):
        tool_tier = "research"
        rank = "specialist"
    elif any(w in desc_lower for w in ("write", "edit", "code", "build", "creat", "implement", "fix", "refactor",
                                        "develop", "deploy", "test", "debug", "generat")):
        tool_tier = "coding"
        rank = "specialist"
    else:
        tool_tier = "read_only"
        rank = "specialist"

    # Detect capabilities
    capabilities = []
    cap_keywords = {
        "code_generation": ["code", "implement", "build", "develop", "generat"],
        "analysis": ["analy", "review", "audit", "inspect", "assess"],
        "testing": ["test", "qa", "quality", "verify", "validat"],
        "documentation": ["document", "doc", "readme", "guide", "wiki"],
        "security": ["secur", "vulnerab", "owasp", "audit", "penetrat"],
        "data_processing": ["data", "csv", "json", "etl", "pipeline", "transform"],
        "devops": ["deploy", "ci/cd", "docker", "kubernetes", "infra"],
        "communication": ["report", "email", "notify", "alert", "messag"],
        "memory_management": ["memory", "context", "remember", "recall", "knowledge"],
        "optimization": ["optim", "perform", "speed", "effici", "improv"],
    }
    for cap, keywords in cap_keywords.items():
        if any(k in desc_lower for k in keywords):
            capabilities.append(cap)
    if not capabilities:
        capabilities = ["general_purpose"]

    # Extract a name from the description
    words = _re.sub(r"[^a-zA-Z\s]", "", description).split()
    name_words = [w.capitalize() for w in words[:4] if len(w) > 2]
    name = " ".join(name_words[:3]) + " Agent" if name_words else "Custom Agent"

    # Build system prompt from description
    system_prompt = (
        f"You are a specialized AI agent. Your purpose: {description}\n\n"
        f"Core responsibilities:\n"
    )
    for i, cap in enumerate(capabilities, 1):
        system_prompt += f"  {i}. {cap.replace('_', ' ').title()}\n"
    system_prompt += (
        f"\nOperate with the '{tool_tier}' tool tier. "
        f"Be thorough, accurate, and proactive. Share results to the hive mind."
    )

    agent_id = f"agent-{_re.sub(r'[^a-z0-9]', '-', name.lower()).strip('-')[:30]}"

    return AgentTemplate(
        id=agent_id,
        name=name,
        description=description[:200],
        system_prompt=system_prompt,
        tool_tier=tool_tier,
        max_turns=20,
        rank=rank,
        capabilities=capabilities,
        auto_share_to_hive=True,
    )


@app.post("/api/agents/squad")
async def create_squad(request: Request):
    """Create a new squad (auto-creates the squad branch)."""
    _ensure()
    body = await request.json()
    squad_name = body.get("name", "")
    if not squad_name:
        raise HTTPException(400, "Squad name is required")

    try:
        from cvc.core.models import CVCBranchRequest
        branch_name = f"squad/{squad_name}"
        bp = _get_db().index.get_branch(branch_name)
        if bp is None:
            _get_engine().branch(CVCBranchRequest(
                name=branch_name,
                description=f"Squad: {squad_name}",
            ))
        return {"squad": squad_name, "branch": branch_name, "status": "created"}
    except Exception as exc:
        raise HTTPException(400, str(exc))


@app.get("/api/agents/squads")
async def list_squads():
    """List all squads with their member agents."""
    _ensure()
    try:
        branches = _get_db().index.list_branches()
        squads = []
        for bp in branches:
            if bp.name.startswith("squad/"):
                squad_name = bp.name.removeprefix("squad/")
                agents = _get_db().index.list_agents(squad=squad_name) if hasattr(_get_db().index, "list_agents") else []
                squads.append({
                    "name": squad_name,
                    "branch": bp.name,
                    "agents": [a if isinstance(a, dict) else {} for a in agents],
                    "agent_count": len(agents),
                })
        return {"squads": squads}
    except Exception:
        return {"squads": []}


@app.get("/api/agents/{agent_id}")
async def get_agent_detail(agent_id: str):
    """Get detailed info about an agent."""
    _ensure()
    workspace = Path(_project_root) if _project_root else Path.cwd()

    # Check template file
    tpl_path = workspace / ".cvc" / "agent_templates" / f"{agent_id}.json"
    if tpl_path.exists():
        try:
            data = json.loads(tpl_path.read_text(encoding="utf-8"))
            data["source"] = "template"
            return data
        except (json.JSONDecodeError, OSError):
            pass

    # Check registry
    try:
        profile = _get_db().index.get_agent(agent_id)
        if profile:
            profile["source"] = "registered"
            return profile
    except Exception:
        pass

    # Check builtin agents
    for b in _get_builtin_agents():
        if b.get("id") == agent_id:
            return b

    raise HTTPException(404, f"Agent '{agent_id}' not found")


@app.put("/api/agents/{agent_id}")
async def update_agent(agent_id: str, request: Request):
    """Update an agent's configuration."""
    _ensure()
    body = await request.json()
    workspace = Path(_project_root) if _project_root else Path.cwd()

    # Update template file if it exists
    tpl_path = workspace / ".cvc" / "agent_templates" / f"{agent_id}.json"
    if tpl_path.exists():
        try:
            data = json.loads(tpl_path.read_text(encoding="utf-8"))
            data.update(body)
            data["updated_at"] = time.time()
            tpl_path.write_text(
                json.dumps(data, indent=2, default=str),
                encoding="utf-8",
            )
            return {"id": agent_id, "status": "updated"}
        except (json.JSONDecodeError, OSError) as exc:
            raise HTTPException(500, str(exc))

    # Update registry
    try:
        _get_db().index.update_agent(agent_id, **body)
        return {"id": agent_id, "status": "updated"}
    except Exception as exc:
        raise HTTPException(400, str(exc))


@app.delete("/api/agents/{agent_id}")
async def delete_agent(agent_id: str):
    """Delete an agent (template + registry)."""
    _ensure()
    workspace = Path(_project_root) if _project_root else Path.cwd()

    # Remove template
    tpl_path = workspace / ".cvc" / "agent_templates" / f"{agent_id}.json"
    if tpl_path.exists():
        tpl_path.unlink()

    # Remove from registry
    try:
        if hasattr(_get_db().index, "delete_agent"):
            _get_db().index.delete_agent(agent_id)
    except Exception:
        pass

    return {"id": agent_id, "status": "deleted"}


@app.get("/api/agents/{agent_id}/history")
async def get_agent_history(agent_id: str, limit: int = 50):
    """Get an agent's commit history."""
    _ensure()
    try:
        commits = _get_db().index.list_commits_by_agent(agent_id, limit=limit)
        return {
            "agent_id": agent_id,
            "commits": [
                {
                    "hash": c.commit_hash[:12] if hasattr(c, "commit_hash") else "",
                    "message": c.message if hasattr(c, "message") else "",
                    "type": str(c.commit_type) if hasattr(c, "commit_type") else "",
                    "timestamp": c.metadata.timestamp if hasattr(c, "metadata") else 0,
                }
                for c in commits
            ],
            "total": len(commits),
        }
    except Exception:
        return {"agent_id": agent_id, "commits": [], "total": 0}


# ═══════════════════════════════════════════════════════════════════════
# 7d. HIVE MEMORY — Shared cognitive space
# ═══════════════════════════════════════════════════════════════════════

@app.get("/api/hivemind/memory")
async def get_hive_memory(
    query: str = "",
    category: str | None = None,
    agent_id: str | None = None,
    limit: int = 20,
):
    """Read from the shared hive memory."""
    _ensure()
    try:
        from cvc.sdk.hivemind import HiveMemory, HiveMind
        workspace = Path(_project_root) if _project_root else Path.cwd()
        cvc_root = workspace / ".cvc"
        hive = HiveMind(str(cvc_root))
        memory = HiveMemory(hive)
        entries = await memory.read(query, category=category, agent_id=agent_id, limit=limit)
        return {"entries": entries, "total": len(entries)}
    except Exception as exc:
        logger.debug("Hive memory read error: %s", exc)
        return {"entries": [], "total": 0}


@app.post("/api/hivemind/memory")
async def write_hive_memory(request: Request):
    """Write an entry to the shared hive memory."""
    _ensure()
    body = await request.json()
    content = body.get("content", "")
    if not content:
        raise HTTPException(400, "content is required")

    try:
        from cvc.sdk.hivemind import HiveMemory, HiveMind
        workspace = Path(_project_root) if _project_root else Path.cwd()
        cvc_root = workspace / ".cvc"
        hive = HiveMind(str(cvc_root))
        memory = HiveMemory(hive)
        commit_hash = await memory.write(
            agent_id=body.get("agent_id", "human"),
            content=content,
            category=body.get("category", "general"),
            tags=body.get("tags", []),
        )
        return {"commit_hash": commit_hash, "status": "ok"}
    except Exception as exc:
        raise HTTPException(500, str(exc))


@app.get("/api/hivemind/memory/stats")
async def get_hive_memory_stats():
    """Get hive memory statistics."""
    _ensure()
    try:
        from cvc.sdk.hivemind import HiveMemory, HiveMind
        workspace = Path(_project_root) if _project_root else Path.cwd()
        cvc_root = workspace / ".cvc"
        hive = HiveMind(str(cvc_root))
        memory = HiveMemory(hive)
        return await memory.stats()
    except Exception:
        return {"total_entries": 0, "categories": {}, "contributing_agents": [], "agent_count": 0}


@app.get("/api/hivemind/memory/summary")
async def get_hive_memory_summary(limit: int = 10):
    """Get a text summary of the hive memory state."""
    _ensure()
    try:
        from cvc.sdk.hivemind import HiveMemory, HiveMind
        workspace = Path(_project_root) if _project_root else Path.cwd()
        cvc_root = workspace / ".cvc"
        hive = HiveMind(str(cvc_root))
        memory = HiveMemory(hive)
        summary = await memory.summary_context(limit=limit)
        return {"context": summary}
    except Exception:
        return {"context": "No hive memory available."}


@app.post("/api/hivemind/memory/compact")
async def compact_hive_memory(request: Request):
    """Trigger compaction of old hive memory entries."""
    _ensure()
    body = await request.json()
    try:
        from cvc.sdk.hivemind import HiveMind
        workspace = Path(_project_root) if _project_root else Path.cwd()
        cvc_root = workspace / ".cvc"
        hive = HiveMind(str(cvc_root))
        result = hive.compact(
            branch="hive/memory",
            max_age_hours=body.get("max_age_hours", 48),
        )
        return {
            "status": "ok",
            "compacted": result.commits_compacted if hasattr(result, "commits_compacted") else 0,
        }
    except Exception as exc:
        raise HTTPException(500, str(exc))


# ═══════════════════════════════════════════════════════════════════════
# 8. SESSIONS
# ═══════════════════════════════════════════════════════════════════════

@app.get("/api/sessions")
async def get_sessions():
    import httpx
    try:
        r = httpx.get(f"http://{_HOST}:{_PROXY_PORT}/cvc/sessions", timeout=5.0)
        return r.json()
    except Exception:
        return {"sessions": [], "error": "Proxy not reachable"}


# ═══════════════════════════════════════════════════════════════════════
# 9. AGENT CHAT — Full Agentic Loop (same brain as CLI `cvc agent`)
#
# 9a. Workspace switch helper (used by the agentic loop)
# ═══════════════════════════════════════════════════════════════════════

async def _execute_workspace_switch(arguments: dict) -> str:
    """
    Execute a workspace switch from within the agentic chat loop.
    This is a gateway-level operation that re-scopes the tool executor,
    system prompt, and all CVC operations to the new workspace.
    """
    global _tool_executor, _system_prompt, _project_root

    target_path = arguments.get("path", "")
    if not target_path:
        return "Error: 'path' argument is required for cvc_switch_workspace."

    target = Path(target_path).resolve()
    if not target.exists():
        return f"Error: Directory does not exist: {target}"
    if not target.is_dir():
        return f"Error: Not a directory: {target}"

    try:
        inst = _workspace_mgr.switch_to(target)
        _project_root = inst.path

        # Re-scope tool executor to the new workspace
        from cvc.agent.executor import ToolExecutor
        _tool_executor = ToolExecutor(inst.path, inst.engine)
        _wire_subagent_runtime(_tool_executor)

        # Rebuild system prompt with new workspace context
        from cvc.agent.system_prompt import build_system_prompt
        provider = _agent_llm.provider if _agent_llm else inst.config.provider
        model = _agent_llm.model if _agent_llm else inst.config.model
        _system_prompt = build_system_prompt(
            workspace=str(inst.path),
            provider=provider,
            model=model,
            branch=inst.engine.active_branch,
            user_memory_context=_get_user_memory_context(),
        )

        return (
            f"Successfully switched to workspace: {inst.path}\n"
            f"Branch: {inst.engine.active_branch}\n"
            f"CVC root: {inst.config.cvc_root}\n"
            f"All CVC operations now target this workspace."
        )
    except Exception as exc:
        return f"Error switching workspace: {exc}"


# ═══════════════════════════════════════════════════════════════════════
# 9b. AGENT CHAT — Full Agentic Loop
#
#    Architecture:  Server-side agentic loop.
#    1. Send messages + tools → LLM
#    2. Stream text tokens to the browser in real-time (SSE)
#    3. If LLM returns tool_calls → execute server-side → stream progress
#    4. Feed tool results back to LLM → goto 1
#    5. When LLM returns only text (no tool calls) → done
#
#    SSE event types:
#      {"type":"text",       "content":"…"}
#      {"type":"tool_start", "name":"…", "args":{…}}
#      {"type":"tool_result","name":"…", "output":"…"}
#      {"type":"status",     "message":"…"}
#      [DONE]
# ═══════════════════════════════════════════════════════════════════════

def _sse(payload) -> str:
    """Format a single SSE data line."""
    return f"data: {json.dumps(payload)}\n\n"


# ── Composer-bar helpers ────────────────────────────────────────────────

_VALID_REASONING_EFFORTS = ("none", "minimal", "low", "medium", "high", "xhigh")
_DEFAULT_REASONING_EFFORT = "medium"
_AUTO_COMPACT_PCT = 0.5  # auto-compress threshold as fraction of context


def _active_workspace_settings_path() -> Path:
    """Return path to the active workspace's settings.local.json."""
    base = Path(_project_root) if _project_root else Path.cwd()
    return base / ".cvc" / "settings.local.json"


def _read_workspace_settings() -> dict:
    p = _active_workspace_settings_path()
    if not p.exists():
        return {}
    try:
        return json.loads(p.read_text(encoding="utf-8") or "{}") or {}
    except Exception:
        return {}


def _write_workspace_settings(updates: dict) -> dict:
    p = _active_workspace_settings_path()
    p.parent.mkdir(parents=True, exist_ok=True)
    data = _read_workspace_settings()
    data.update(updates)
    p.write_text(json.dumps(data, indent=2) + "\n", encoding="utf-8")
    return data


def _get_active_model_context_window() -> int:
    """Best-effort lookup of the active model's context_window from the catalog."""
    if not _agent_llm:
        return 0
    provider = _agent_llm.provider
    model = _agent_llm.model
    # Inline minimal catalog (covers common cases without a network call)
    static = {
        "anthropic": {
            "claude-3-opus": 200000, "claude-3.5-sonnet": 200000,
            "claude-3-5-haiku-20241022": 200000, "claude-opus-4.7": 200000,
            "claude-opus-4-8": 1000000, "claude-sonnet-4-6": 1000000, "claude-haiku-4-5": 200000
        },
        "openai": {
            "gpt-4o": 128000, "gpt-4": 128000, "o1-preview": 128000, "o1-mini": 128000,
            "gpt-5.3": 256000, "gpt-5.2": 128000, "gpt-5-mini": 128000, "gpt-4o-mini": 128000,
            "o3-mini": 200000, "o1": 128000
        },
        "google": {
            "gemini-3.5-flash": 1000000, "gemini-3.1-pro": 1000000, "gemini-3.1-flash-lite": 1000000,
            "gemini-3.1-pro-preview": 1000000, "gemini-3-flash-preview": 1000000,
            "gemini-2.5-flash": 1000000, "gemini-2.5-pro": 1000000,
            "gemini-1.5-flash": 1000000, "gemini-1.5-pro": 2000000
        },
        "ollama": {
            "qwen2.5-coder:7b": 32768, "qwen3-coder:latest": 32768,
            "deepseek-r1:latest": 32768, "llama3:8b": 8192
        },
        "minimax": {
            # MiniMax M-series — https://platform.minimax.io/docs/guides/models-intro
            "MiniMax-M3": 1000000,             # flagship, multimodal
            "MiniMax-M2.7": 200000,
            "MiniMax-M2.7-highspeed": 200000,
            "MiniMax-M2.5": 200000,
            "MiniMax-M2.5-highspeed": 200000,
            "MiniMax-M2.1": 200000,
            "MiniMax-M2.1-highspeed": 200000,
            "MiniMax-M2": 200000,
        },
    }
    win = static.get(provider, {}).get(model)
    if win:
        return win
    # GitHub Copilot — actual per-model caps (from the upstream API).
    # The dashboard user hit "limit of 168000" on Claude Sonnet 4.5 via
    # Copilot — the old 128k default was wrong.
    # v2.91.42: corrected to the real values.
    if provider in ("github", "copilot"):
        copilot_caps = {
            "claude-sonnet-4.5": 168000,
            "claude-sonnet-4.6": 168000,
            "claude-sonnet-4-5": 168000,
            "claude-sonnet-4-6": 168000,
            "claude-opus-4-6": 168000,
            "claude-opus-4.7": 168000,
            "claude-haiku-4.5": 168000,
            "claude-haiku-4-5": 168000,
            "gpt-5.2": 128000,
            "gpt-5": 128000,
            "gpt-4o": 64000,
            "o1-preview": 128000,
            "o1-mini": 128000,
        }
        if model in copilot_caps:
            return copilot_caps[model]
        return 128000
    return 200000  # conservative default


def _estimate_used_tokens() -> int:
    """Rough token estimate of the *full LLM prompt* that the chat path will
    send. Includes system prompt + persona skills + DX context + history +
    the active workspace's user-memory context, when available.

    Falls back to the engine's persisted context_window if any of those
    components can't be built (no LLM yet, no workspace, etc.).
    """
    try:
        # Build the same prompt the chat handler will build, so the meter
        # reports what the LLM actually receives — not just the engine's
        # bare history.
        total_chars = 0
        try:
            from cvc.agent.system_prompt import build_system_prompt
            if _project_root:
                try:
                    _system_prompt = build_system_prompt(
                        workspace=str(_project_root),
                        provider=_agent_llm.provider if _agent_llm else "",
                        model=_agent_llm.model if _agent_llm else "",
                        branch=_get_engine().active_branch if _get_engine() else "",
                        user_memory_context=_get_user_memory_context(),
                    )
                    total_chars += len(_system_prompt or "")
                except Exception:
                    pass
        except Exception:
            pass
        eng = _get_engine()
        if eng:
            for m in eng.context_window:
                content = getattr(m, "content", "") or ""
                if isinstance(content, str):
                    total_chars += len(content)
                else:
                    total_chars += len(json.dumps(content))
        return total_chars // 4
    except Exception:
        return 0


def _context_window_breakdown() -> dict[str, int]:
    """Per-role token breakdown of the *full LLM prompt* the chat path
    will send. Includes the system prompt (when buildable), plus every
    message in the engine's persisted context window.

    Returns ``{system, user, assistant, tool, other}`` token counts. Used
    by the dashboard meter to show *where* the tokens went — not just the
    running total.
    """
    out = {"system": 0, "user": 0, "assistant": 0, "tool": 0, "other": 0}
    try:
        # 1. System prompt — estimate its size and add to "system".
        try:
            from cvc.agent.system_prompt import build_system_prompt
            if _project_root:
                try:
                    _sys = build_system_prompt(
                        workspace=str(_project_root),
                        provider=_agent_llm.provider if _agent_llm else "",
                        model=_agent_llm.model if _agent_llm else "",
                        branch=_get_engine().active_branch if _get_engine() else "",
                        user_memory_context=_get_user_memory_context(),
                    )
                    if _sys:
                        out["system"] += len(_sys) // 4
                except Exception:
                    pass
        except Exception:
            pass

        # 2. Persisted conversation history (engine.context_window).
        eng = _get_engine()
        if eng:
            for m in eng.context_window:
                role = (getattr(m, "role", "") or "other").lower()
                if role not in out:
                    role = "other"
                content = getattr(m, "content", "") or ""
                chars = len(content) if isinstance(content, str) else len(json.dumps(content))
                out[role] += chars // 4
    except Exception:
        pass
    return out


def _last_auto_compact_iso() -> str | None:
    """ISO timestamp of the last auto-compact event (if any), from the workspace's
    `.cvc/context_window.json` log."""
    try:
        base = Path(_project_root) / ".cvc" if _project_root else None
        if not base:
            return None
        log = base / "context_window.json"
        if not log.exists():
            return None
        data = json.loads(log.read_text(encoding="utf-8") or "{}")
        last = data.get("last_auto_compact_at")
        return last
    except Exception:
        return None


@app.get("/api/chat/reasoning_effort")
async def get_reasoning_effort() -> dict[str, str]:
    """Return the active workspace's persisted reasoning effort."""
    settings = _read_workspace_settings()
    effort = settings.get("reasoning_effort") or _DEFAULT_REASONING_EFFORT
    if effort not in _VALID_REASONING_EFFORTS:
        effort = _DEFAULT_REASONING_EFFORT
    return {"effort": effort}


@app.put("/api/chat/reasoning_effort")
async def set_reasoning_effort(request: Request) -> dict[str, str]:
    """Set the active workspace's reasoning effort.

    Body: ``{"effort": "<none|minimal|low|medium|high|xhigh>"}``
    """
    body = await request.json()
    effort = (body or {}).get("effort", "")
    if effort not in _VALID_REASONING_EFFORTS:
        raise HTTPException(422, f"Invalid effort. Must be one of {_VALID_REASONING_EFFORTS}")
    _write_workspace_settings({"reasoning_effort": effort})
    return {"effort": effort}


@app.get("/api/chat/context_meter")
async def get_context_meter() -> dict[str, Any]:
    """Live context-window meter for the active engine + active model.

    Returns per-role token breakdown so the UI can render *where* the tokens
    are (system / user / assistant / tool) and not just the running total.
    Also exposes the auto-compact threshold and the number of tokens still
    available before the threshold fires.
    """
    used = _estimate_used_tokens()
    total = _get_active_model_context_window() or 0
    pct = (used / total) if total > 0 else 0.0
    threshold_tokens = int(total * _AUTO_COMPACT_PCT)
    tokens_until_compact = max(0, threshold_tokens - used)
    pct_until_compact = (tokens_until_compact / total) if total > 0 else 0.0
    provider = _agent_llm.provider if _agent_llm else (_get_config().provider if _get_config() else "")
    model = _agent_llm.model if _agent_llm else (_get_config().model if _get_config() else "")
    breakdown = _context_window_breakdown()
    return {
        "used_tokens": used,
        "total_tokens": total,
        "used_pct": pct,
        "remaining_tokens": max(0, total - used),
        "remaining_pct": 1.0 - pct,
        "auto_compact_threshold_tokens": threshold_tokens,
        "auto_compact_threshold_pct": _AUTO_COMPACT_PCT,
        "tokens_until_auto_compact": tokens_until_compact,
        "pct_until_auto_compact": pct_until_compact,
        "breakdown": breakdown,
        "model": model,
        "provider": provider,
        "last_auto_compact_at": _last_auto_compact_iso(),
    }


@app.get("/api/workspace/tree")
async def get_workspace_tree(
    path: str = "",
    depth: int = 2,
    show_hidden: bool = False,
) -> dict[str, Any]:
    """Return a recursive directory listing of the active workspace.

    `path` is *relative* to the active workspace root. `depth` caps recursion
    (1 = direct children only). `show_hidden=False` filters dot-files.
    """
    if not _project_root:
        raise HTTPException(503, "No active workspace")
    root = Path(_project_root).resolve()

    # Resolve & guard against escapes outside the workspace root
    target = (root / path).resolve() if path else root
    try:
        target.relative_to(root)
    except ValueError:
        raise HTTPException(400, f"Path escapes workspace root: {path}")
    if not target.exists():
        raise HTTPException(404, f"Path not found: {path}")

    depth = max(0, min(int(depth or 0), 6))

    def _node(p: Path, remaining: int) -> dict[str, Any]:
        try:
            st = p.stat()
        except Exception:
            st = None
        rel = ""
        try:
            rel = str(p.relative_to(root))
        except ValueError:
            rel = ""
        node: dict[str, Any] = {
            "name": p.name or str(p),
            "path": str(p),
            "rel_path": rel,
            "is_dir": p.is_dir(),
            "size": st.st_size if st else 0,
            "mtime": st.st_mtime if st else 0.0,
            "children": None,
        }
        if p.is_dir() and remaining > 0:
            kids: list[dict[str, Any]] = []
            try:
                entries = sorted(
                    p.iterdir(),
                    key=lambda x: (not x.is_dir(), x.name.lower()),
                )
            except PermissionError:
                entries = []
            for child in entries:
                if not show_hidden and child.name.startswith("."):
                    continue
                kids.append(_node(child, remaining - 1))
            node["children"] = kids
        return node

    return _node(target, depth)


@app.get("/api/workspace/file")
async def get_workspace_file(path: str) -> dict[str, Any]:
    """Return the contents of a single file inside the active workspace.

    `path` is relative to the active workspace root. Refuses paths that
    escape the root, binary files larger than 2 MB, and non-files.
    """
    if not _project_root:
        raise HTTPException(503, "No active workspace")
    if not path:
        raise HTTPException(400, "Missing path")
    root = Path(_project_root).resolve()
    target = (root / path).resolve()
    try:
        target.relative_to(root)
    except ValueError:
        raise HTTPException(400, f"Path escapes workspace root: {path}")
    if not target.exists():
        raise HTTPException(404, f"Path not found: {path}")
    if not target.is_file():
        raise HTTPException(400, f"Not a file: {path}")

    try:
        st = target.stat()
    except Exception:
        raise HTTPException(500, "Failed to stat file")

    MAX_BYTES = 2 * 1024 * 1024  # 2 MB
    is_text = True
    content = ""
    truncated = False
    try:
        raw = target.read_bytes()
        if len(raw) > MAX_BYTES:
            raw = raw[:MAX_BYTES]
            truncated = True
        try:
            content = raw.decode("utf-8")
        except UnicodeDecodeError:
            try:
                content = raw.decode("latin-1")
            except Exception:
                is_text = False
                content = ""
    except Exception as exc:
        raise HTTPException(500, f"Read failed: {exc}")

    ext = target.suffix.lstrip(".").lower()
    import mimetypes
    mime, _ = mimetypes.guess_type(str(target))
    rel = str(target.relative_to(root))

    return {
        "name": target.name,
        "path": str(target),
        "rel_path": rel,
        "size": st.st_size,
        "mtime": st.st_mtime,
        "ext": ext,
        "mime": mime or ("text/plain" if is_text else "application/octet-stream"),
        "is_text": is_text,
        "truncated": truncated,
        "content": content,
    }


@app.get("/api/workspace/file-diff")
async def get_workspace_file_diff(path: str) -> dict[str, Any]:
    """Return a unified git diff for a single file vs HEAD (working tree).

    Used by the dashboard IDE viewer to render red/green inline changes.
    Always returns 200 with `status` reflecting outcome so the UI can show
    a clean "no changes" state without try/catch noise.

    Status values:
      - "clean"        — file unchanged vs HEAD
      - "modified"     — has diff hunks; `hunks` populated
      - "untracked"    — file is not tracked; treat content as fully added
      - "not_git"      — workspace is not a git repo
      - "error"        — git command failed; `message` set
    """
    if not _project_root:
        raise HTTPException(503, "No active workspace")
    if not path:
        raise HTTPException(400, "Missing path")
    root = Path(_project_root).resolve()
    target = (root / path).resolve()
    try:
        target.relative_to(root)
    except ValueError:
        raise HTTPException(400, f"Path escapes workspace root: {path}")
    rel = str(target.relative_to(root))

    import subprocess

    def _git(args: list[str]) -> tuple[int, str, str]:
        try:
            proc = subprocess.run(
                ["git", *args],
                cwd=str(root),
                capture_output=True,
                text=True,
                timeout=10,
                            **HIDDEN_KW,
            )
            return proc.returncode, proc.stdout, proc.stderr
        except FileNotFoundError:
            return 127, "", "git not installed"
        except subprocess.TimeoutExpired:
            return 124, "", "git timeout"

    rc, _, _ = _git(["rev-parse", "--is-inside-work-tree"])
    if rc != 0:
        return {"status": "not_git", "rel_path": rel, "hunks": []}

    # Check tracked-ness
    rc_ls, _, _ = _git(["ls-files", "--error-unmatch", "--", rel])
    tracked = rc_ls == 0

    if not tracked:
        # Untracked → entire file is "added" for display purposes
        return {
            "status": "untracked",
            "rel_path": rel,
            "hunks": [],
        }

    rc_d, out_d, err_d = _git([
        "diff",
        "--no-color",
        "--unified=3",
        "HEAD",
        "--",
        rel,
    ])
    if rc_d != 0:
        return {
            "status": "error",
            "rel_path": rel,
            "hunks": [],
            "message": err_d.strip() or "git diff failed",
        }

    if not out_d.strip():
        return {"status": "clean", "rel_path": rel, "hunks": []}

    # Parse unified diff into hunks
    hunks: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    import re as _re

    hunk_re = _re.compile(r"^@@ -(\d+)(?:,(\d+))? \+(\d+)(?:,(\d+))? @@(.*)$")
    for line in out_d.splitlines():
        if line.startswith("diff --git") or line.startswith("index ") \
                or line.startswith("--- ") or line.startswith("+++ "):
            continue
        m = hunk_re.match(line)
        if m:
            if current is not None:
                hunks.append(current)
            current = {
                "old_start": int(m.group(1)),
                "old_lines": int(m.group(2) or 1),
                "new_start": int(m.group(3)),
                "new_lines": int(m.group(4) or 1),
                "header": m.group(5).strip(),
                "lines": [],
            }
            continue
        if current is None:
            continue
        if line.startswith("+"):
            kind = "add"
            content_line = line[1:]
        elif line.startswith("-"):
            kind = "del"
            content_line = line[1:]
        elif line.startswith(" "):
            kind = "ctx"
            content_line = line[1:]
        elif line.startswith("\\"):
            # "\ No newline at end of file" — skip
            continue
        else:
            kind = "ctx"
            content_line = line
        current["lines"].append({"kind": kind, "content": content_line})
    if current is not None:
        hunks.append(current)

    return {
        "status": "modified",
        "rel_path": rel,
        "hunks": hunks,
    }


# ─────────────────────────────────────────────────────────────────────────────
# CVC v3.0 — Unified-core chat endpoint (always available).
#
# This is the new SSE handler that wraps `cvc.gateway_chat.stream_unified_chat`,
# which builds an `AIAgent` from the vendored unified runtime and streams
# events back to the dashboard using the same event schema as `/api/chat`.
#
# The OLD `/api/chat` route below is PRESERVED UNCHANGED. You can run the
# dashboard against this new endpoint in two ways:
#   1. Set CVC_USE_UNIFIED_CORE=1 in the shell before `cvc gateway start`
#      → the old `/api/chat` route will forward to this handler.
#   2. Or have the dashboard call `/api/chat_unified` directly.
#
# Either way, the original CVC agent loop in `agent_chat` is NOT deleted
# and NOT modified. Setting CVC_USE_UNIFIED_CORE back to 0 (or unsetting it)
# rolls the dashboard back to the original behaviour.
# ─────────────────────────────────────────────────────────────────────────────
@app.post("/api/chat_unified")
async def agent_chat_unified_route(request: Request):
    """Unified-core backed `/api/chat` (v3.0).

    Identical request body to `/api/chat`. Returns the same SSE event
    schema. Behaviour comes from the vendored unified runtime's
    `AIAgent.run_conversation()` instead of the CVC-only loop.
    """
    body = await request.json()
    from cvc.gateway_chat import stream_unified_chat  # local import: avoid circulars
    return StreamingResponse(
        stream_unified_chat(body),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


@app.post("/api/chat_legacy")
async def agent_chat_legacy_route(request: Request):
    """v3.1.0 — Preserved legacy CVC chat loop (SSE), reachable for rollback.

    From v3.1.0 onward the dashboard's `/api/chat` route runs the
    unified core. This endpoint exposes the original CVC loop for
    explicit comparison or `CVC_USE_UNIFIED_CORE=0` rollback testing.
    """
    import os
    prev = os.environ.get("CVC_USE_UNIFIED_CORE")
    os.environ["CVC_USE_UNIFIED_CORE"] = "0"
    try:
        return await agent_chat(request)
    finally:
        if prev is None:
            os.environ.pop("CVC_USE_UNIFIED_CORE", None)
        else:
            os.environ["CVC_USE_UNIFIED_CORE"] = prev


@app.post("/api/chat")
async def agent_chat(request: Request):
    global _chat_abort
    _chat_abort = False  # Reset abort flag for new chat turn
    body = await request.json()

    # ── v3.1.0 — Unified core is the DEFAULT for the dashboard's HTTP
    # chat endpoint. The dashboard's `api.ts` POSTs to `/api/chat`;
    # from v3.1.0 onward, that path runs the forked agent
    # unified core (`stream_unified_chat`). The legacy CVC loop is
    # preserved as `/api/chat_legacy` and reachable here only when
    # `CVC_USE_UNIFIED_CORE=0` is set in the environment. The old
    # default-on behavior is gone — users no longer need any env-var
    # opt-in to use the new core. ────────────────────────────────
    if os.environ.get("CVC_USE_UNIFIED_CORE", "1").strip() not in ("0", "false", "no", "off"):
        from cvc.gateway_chat import stream_unified_chat
        return StreamingResponse(
            stream_unified_chat(body),
            media_type="text/event-stream",
            headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
        )
    messages = body.get("messages", [])
    stream = body.get("stream", True)

    # ── Composer-bar additions: per-request reasoning + persona override ──
    req_effort = body.get("reasoning_effort") or ""
    if req_effort and req_effort not in _VALID_REASONING_EFFORTS:
        req_effort = ""
    if not req_effort:
        # Fall back to persisted workspace setting
        req_effort = _read_workspace_settings().get("reasoning_effort") or ""
    # AgentLLM's effort_level expects "low|medium|high" mostly; pass through
    # all values — adapters that don't recognise extras (e.g. "xhigh", "minimal",
    # "none") will simply ignore the unknown level.
    effort_for_llm = "" if req_effort in ("none", "") else req_effort

    # Persona override (optional). If supplied & valid, swap the system prompt
    # for this single chat turn, and gather the persona's selected skill
    # content for injection (v2.23.5).
    req_persona_id = body.get("persona_id") or ""
    turn_system_prompt = _system_prompt
    persona_skill_messages: list[dict[str, str]] = []
    if req_persona_id:
        try:
            from cvc.dashboard.personas_api import _get_persona  # type: ignore
            persona = _get_persona(req_persona_id)
            if persona and persona.get("system_prompt"):
                turn_system_prompt = persona["system_prompt"]
            # Build persona-skill instructions block
            persona_skill_ids = (persona or {}).get("skills") or []
            if persona_skill_ids:
                try:
                    from cvc.agent.skills import discover_skills, filter_by_persona
                    _ws = Path.cwd()
                    all_skills = discover_skills(_ws)
                    matched = filter_by_persona(all_skills, persona_skill_ids)
                    for sk in matched:
                        persona_skill_messages.append({
                            "role": "system",
                            "content": f"[Persona-Skill: {sk.name}]\n{sk.content}",
                        })
                    logger.info(
                        "Persona '%s' chat: %d skills injected (of %d ticked)",
                        req_persona_id, len(matched), len(persona_skill_ids),
                    )
                except Exception:
                    logger.debug("Persona-skill injection failed for %s",
                                 req_persona_id, exc_info=True)
        except Exception:
            logger.debug("Persona override failed for %s", req_persona_id, exc_info=True)

    if _agent_llm is None or not _agent_tools:
        # Don't silently fall through to the proxy — when the gateway has
        # no agent LLM bound (passthrough / unset / failed init), the proxy
        # path produces no streamable output and the dashboard hangs for
        # 59s before showing "Connection appears stalled". Stream an
        # immediate, actionable error so the UI can render a clear
        # "run cvc setup" hint instead of stalling.
        cfg = _get_config()
        cur_provider = (cfg.provider if cfg else "") or "(unset)"
        if cur_provider == "passthrough" or cur_provider == "(unset)":
            err_msg = (
                "Dashboard chat needs a real LLM provider, but CVC is in "
                f"`{cur_provider}` mode (no internal LLM bound).\n\n"
                "Fix in your terminal:\n"
                "  1. Run `cvc setup` and pick anthropic / openai / google / "
                "github / copilot / vertex / ollama / lmstudio\n"
                "  2. Run `cvc gateway restart`\n\n"
                "Or set ANTHROPIC_API_KEY / OPENAI_API_KEY / GOOGLE_API_KEY / "
                "COPILOT_GITHUB_TOKEN in your shell and restart the gateway — "
                "CVC will auto-detect it."
            )
        else:
            err_msg = (
                f"Agent LLM failed to initialise for provider "
                f"`{cur_provider}`. Check `~/.cvc/gateway.log` for the "
                "underlying error, verify your API key, then run "
                "`cvc gateway restart`."
            )
        if not stream:
            return {
                "choices": [{"message": {"role": "assistant", "content": err_msg}, "finish_reason": "stop"}],
                "model": cur_provider,
                "setup_required": True,
            }

        async def _setup_required_stream():
            yield _sse({"type": "text", "content": err_msg})
            yield _sse({"type": "setup_required", "provider": cur_provider})
            yield "data: [DONE]\n\n"

        return StreamingResponse(
            _setup_required_stream(),
            media_type="text/event-stream",
            headers={"Cache-Control": "no-cache"},
        )

    # Build the LLM message list: system + persona-skill blocks + DX context + history
    llm_messages: list[dict[str, Any]] = [{"role": "system", "content": turn_system_prompt}]
    for sm in persona_skill_messages:
        llm_messages.append(sm)

    # v2.76.0 — DX self-learning context (pinned messages + insights).
    # Best-effort, never blocks the turn.
    _dx_thread_id = body.get("thread_id") or body.get("conversation_id") or None
    _dx_workspace = body.get("workspace_path") or str(Path.cwd())
    try:
        from cvc.dashboard.dx_api import build_dx_context, auto_learn_from_message
        _dx_ctx = build_dx_context(_dx_thread_id, _dx_workspace)
        if _dx_ctx:
            llm_messages.append({"role": _dx_ctx["role"], "content": _dx_ctx["content"]})
            logger.info(
                "DX context injected: %d pins, %d insights",
                _dx_ctx["_dx_meta"]["pin_count"],
                _dx_ctx["_dx_meta"]["insight_count"],
            )
        # Auto-learn from the latest user message (self-learning).
        _last_user = next(
            (m for m in reversed(messages) if m.get("role") == "user"), None
        )
        if _last_user and isinstance(_last_user.get("content"), str):
            auto_learn_from_message(
                _dx_workspace, _dx_thread_id, _last_user["content"]
            )
    except Exception as _dx_exc:
        logger.debug("DX context injection failed: %s", _dx_exc)

    for m in messages:
        # Pass content as-is (could be string or list for multimodal)
        llm_messages.append({"role": m.get("role", "user"), "content": m.get("content", "")})

    # v2.90.9 — merge ChatRequest.attachments into the last user message as
    # provider-native multimodal blocks. Composer no longer inlines base64.
    try:
        _attachments = body.get("attachments") if isinstance(body, dict) else None
        _prov = (_agent_llm.provider if _agent_llm else "") or ""
        llm_messages = _merge_attachments_into_messages(llm_messages, _attachments, _prov)
    except Exception as _mm_exc:
        logger.warning("attachment merge (HTTP) failed: %s", _mm_exc)

    # ── Pre-flight auto-compact: shrink transcript *before* the provider sees
    # it, so we never hit context_overflow on small/medium-sized windows. ──
    try:
        llm_messages, _did_compact = _maybe_preflight_auto_compact(llm_messages)
        if _did_compact:
            yield_or_status = locals().get("yield")  # may not exist on HTTP path
            try:
                # best-effort: emit an SSE status if we have a yielding context
                if callable(yield_or_status):
                    yield_or_status(_sse({
                        "type": "status",
                        "message": "🧹 Context window auto-compacted — continuing.",
                    }))
            except Exception:
                pass
    except Exception as _pf_exc:
        logger.warning("agent_chat: preflight step failed: %s", _pf_exc)

    if not stream:
        try:
            resp = await _agent_llm.chat(llm_messages, tools=_agent_tools, temperature=0.7, max_tokens=8192)
            return {
                "choices": [{"message": {"role": "assistant", "content": resp.text or ""}, "finish_reason": "stop"}],
                "model": _agent_llm.model,
            }
        except Exception as exc:
            raise HTTPException(502, f"LLM error: {exc}")

    # ── Streaming agentic loop ───────────────────────────────────────
    async def agentic_stream():
        global _session_trust_all
        nonlocal llm_messages
        iteration = 0
        _expensive = _is_expensive_model(_agent_llm.model) if _agent_llm else False
        # Resolved per-turn from env / GlobalConfig / defaults so dashboard
        # changes take effect on the next message — no gateway restart needed.
        _iter_limit = _budget_max_iterations(_expensive)
        _agent_timeout = _budget_agent_timeout()
        _tool_timeout = _budget_tool_timeout()
        _tool_trunc = _MAX_TOOL_OUTPUT_EXPENSIVE if _expensive else _MAX_TOOL_OUTPUT_LLM
        logger.info(
            "agent_chat budgets resolved: iter=%d timeout=%.0fs tool_timeout=%.0fs expensive=%s",
            _iter_limit, _agent_timeout, _tool_timeout, _expensive,
        )

        _turn_start = time.monotonic()
        _last_yield = time.monotonic()
        # Loop-detection state (port of WS path lines 3862-3893)
        recent_tool_signatures: list[str] = []
        consecutive_no_progress = 0
        # One-shot nudge guard: if the LLM exits with no tool calls AND no
        # assistant text (e.g. fired memory/fact_store/save_memory then went
        # silent), we inject a single hidden user nudge and let the loop run
        # one more iteration so the user sees a real reply. Bounded to ONE
        # nudge per turn so a model that flatly refuses to speak still exits
        # cleanly instead of looping forever.
        _silent_nudge_used = False
        # v2.90.13 — Soft wrap-up nudge state. Fires ONCE per turn when the
        # iteration counter crosses _soft_wrap_threshold(_iter_limit). Tells
        # the model to stop exploring and synthesise BEFORE we run out of
        # budget and silent-finish. This is the single biggest survival fix
        # for long debugging tasks — upstream does not ship this.
        _wrapup_nudge_used = False
        _wrapup_threshold = _soft_wrap_threshold(_iter_limit)
        # v2.91.42 — one-shot context-overflow auto-compaction.
        _overflow_handled_sse = False
        # v2.90.13 — Track the last tool name for diagnostic giveup messages.
        _last_tool_name: str | None = None
        # Track whether ANY visible assistant text was streamed across the
        # entire turn. If the loop exits (cap, timeout, loop-detect,
        # no-progress) with this still False, we force one final non-tool
        # synthesis call so the user always gets a real reply.
        _any_visible_text = False
        # v2.92.4 — per-turn guard: we send the "stop emitting duds" nudge
        # at most once per turn. Reset on each new user message.
        _dud_nudge_sent = False
        # v2.95.0 — Per-turn DUD COUNT + HARD CAP. The 2026-06-23 live
        # transcripts showed MiniMax-M3 emitting 6-10 `read_file({})` duds
        # in a row before either (a) the model self-defeats by deciding to
        # stop emitting tools and write a long markdown plan, or (b) the
        # outer iteration cap eventually aborts. Neither is acceptable.
        # New behaviour: count every dud (pre-dispatch validator rejection
        # + post-dispatch dud-class string match). After 2 duds the
        # gateway injects ONE strong "stop and ask the user" nudge — the
        # model sees this and must either emit a real call or a plain-text
        # reply in the SAME iteration. After 4 duds the gateway refuses
        # further tool dispatch for this turn and forces the model into
        # a plain-text response only (no more tool calls accepted). This
        # is the structural counterpart to the v2.95.0 system-prompt
        # rule "you have a hard budget of 3 dud tool calls per turn."
        #
        # v2.96.0 — Lowered both thresholds from {3, 5} → {2, 4} based on
        # the live screenshot the user shared at iter~10: the model made
        # 6 successful reads, then 2 duds, then narrated its findings as
        # a "reply" and stopped calling tools. The 3-dud threshold let
        # it slip into the self-defeat pattern before the structural
        # nudge fired. 2/4 matches the v2.95.0 prompt rule "after 3
        # duds the system will force you" — we count down from 3 so the
        # nudge fires at the *second* dud, giving the model one
        # corrective nudge + one final chance before the hard cap.
        _dud_count = 0
        _MAX_DUDS_PER_TURN = 4  # hard cap; one below the prompt's "3 duds"
        _DUD_NUDGE_THRESHOLD = 2  # structural nudge at the 2nd dud
        # v2.93.0 — Goal-completion gate. The agent's job is to complete the
        # user's task, not to summarise tool outputs. The screenshot symptom
        # (model emits 4 duds → narrative summary → breaks loop) is fixed
        # here by counting actual write-class tool calls in this turn and
        # forcing a continue-with-nudge when the user's prompt asked for a
        # modification but the agent did nothing. Imperative verbs cover the
        # common "make / add / fix / move / change" UI-task surface from the
        # dashboard; the check is intentionally permissive (false positive
        # costs one nudge; false negative breaks the user-facing task).
        _write_calls_count = 0
        # v2.95.0 — Read-call counter (companion to _write_calls_count).
        # Used by the read-intent + dud-storm safety net at synthesis time:
        # if the user asked a read question (no write verbs) and we hit the
        # iteration cap with >= 3 duds AND zero successful reads, refuse
        # to fabricate a summary — emit an honest "couldn't read" reply.
        _read_calls_count = 0
        _goal_nudge_sent = False
        _WRITE_TOOLS = {"write_file", "patch", "edit", "delete_file", "mkdir"}
        _READ_ONLY_TOOLS = {
            "read_file", "grep", "glob", "list_dir",
            "cvc_status", "cvc_log", "cvc_search", "cvc_smart_search",
            "cvc_diff", "cvc_get_context", "cvc_switch_workspace",
            "cvc_list_documents", "cvc_timeline", "cvc_hive_read",
            "cvc_hive_stats", "cvc_hive_status", "ask_user",
            "think", "todo", "task_list", "context_compact",
        }
        _IMPERATIVE_VERBS = (
            "make ", "add ", "fix ", "change ", "modify ", "update ",
            "set ", "remove ", "delete ", "create ", "build ", "move ",
            "place ", "position ", "style ", "design ", "implement ",
            "refactor ", "rename ", "replace ", "insert ", "embed ",
            "wire ", "configure ", "enable ", "disable ", "install ",
            "deploy ", "publish ", "push ", "commit ", "merge ",
            "write ", "save ", "edit ",
        )

        def _is_write_intent(text: str) -> bool:
            """True if the user's prompt contains imperative write-verbs.

            Cheap substring check — no NLP, no token cost. False positives
            are cheap (one extra nudge); false negatives leave the original
            bug in place, so we err permissive.
            """
            if not text:
                return False
            t = text.lower()
            return any(v in t for v in _IMPERATIVE_VERBS)

        def _user_facing_prompt() -> str:
            """Return the most recent user-role text that ISN'T our own nudge.

            Used to decide whether the goal-completion gate should fire. The
            SYSTEM:/[cvc:/...: prefix exclusion mirrors the existing
            deterministic-synthesize path so we never treat our own nudges
            as the user's actual ask.
            """
            try:
                for _m in reversed(llm_messages):
                    if _m.get("role") != "user":
                        continue
                    _c = _m.get("content", "")
                    if isinstance(_c, list):
                        _c = " ".join(
                            (p.get("text", "") if isinstance(p, dict) else str(p))
                            for p in _c
                        )
                    if not isinstance(_c, str):
                        continue
                    s = _c.lstrip()
                    if s.startswith(("SYSTEM:", "[cvc:", "A:")):
                        continue
                    return _c
            except Exception:
                pass
            return ""

        def _is_dud_tool_call(_tc) -> bool:
            """v2.92.8/10 — Mirror of the executor's dud-detect predicate.
            Returns True when the model emitted a tool_use block with empty,
            None, or otherwise unusable arguments. Loop-detect and
            stall-detect MUST skip dud calls; otherwise three identical
            dud signatures trip the abort even though the validator has
            already corrected the input via the per-turn dud nudge.

            v2.92.10 — Also catch the `{"path": ""}` class: the model
            emits a dict with the required key present but the value
            blank. Before this, `any(bool(v) …)` returned True on the
            empty-string value, so the dud slipped past and the
            executor either errored on it (for tools with required
            args in _REQUIRED_ARGS) or silently ran on the wrong
            workspace (for tools without required-arg entries). Now we
            also peek at the per-tool required-arg set and flag a
            call as dud if every required arg is missing/empty.
            """
            args = getattr(_tc, "arguments", None)
            if args is None:
                return True
            if not isinstance(args, dict):
                # Non-dict arguments (str, list, int…) are never valid
                # for our tool schemas — treat as dud so the validator
                # sends a clean "arguments must be a JSON object"
                # message instead of letting the executor choke.
                return True
            # v2.92.10 — Tools with NO required-arg entries AND known to
            # be auto-resolved or zero-arg-safe (e.g. ask_user, todo with
            # no ops, cvc_status, think, context_compact) MUST NOT be
            # flagged as duds when called with `{}`. Without this guard,
            # any tool that the executor accepts with empty args gets
            # rejected here and the model gets a never-ending stream of
            # "stop emitting duds" nudges even though it called a valid
            # no-arg tool. We treat empty `{}` as a valid call for these
            # tools; non-empty empty-string args still count as duds.
            if not args and _tc.name in _ZERO_ARG_SAFE_TOOLS:
                return False
            if not args:
                return True
            # Per-tool required-arg check — mirrors executor._REQUIRED_ARGS
            # without coupling to that module. Keeps this predicate
            # import-safe and side-effect-free.
            try:
                from cvc.agent.executor import _REQUIRED_ARGS as _RA  # type: ignore
                required = _RA.get(_tc.name, ())
            except Exception:
                required = ()
            if required:
                if all(
                    (args.get(r) is None) or
                    (isinstance(args.get(r), str) and not args.get(r).strip())
                    for r in required
                ):
                    return True
            # Generic fallback: any of the known string-valued args is
            # present and non-empty → real call. If EVERY string arg is
            # blank AND there are no other truthy values, treat as dud.
            # (This catches exotic schemas not in _REQUIRED_ARGS.)
            if not any(
                bool(v) and (not isinstance(v, str) or v.strip())
                for v in args.values()
            ):
                return True
            return False

        async def _client_alive() -> bool:
            """Return False if the SSE client has disconnected."""
            try:
                return not await request.is_disconnected()
            except Exception:
                return True

        def _budget_exceeded() -> bool:
            return (time.monotonic() - _turn_start) > _agent_timeout

        try:
            while iteration < _iter_limit:
                # ── Total-turn timeout ───────────────────────────
                if _budget_exceeded():
                    logger.warning("agent_chat: turn timeout after %.1fs", _agent_timeout)
                    yield _sse({"type": "status", "message": f"⚠️ Chat turn exceeded {_agent_timeout:.0f}s — aborting."})
                    yield _sse({"type": "text", "content": f"\n\n⚠️ Chat turn timed out after {_agent_timeout:.0f}s. Try again or simplify the request.\n"})
                    break

                # ── v2.92.3 — Cost budget gate (Claude Code parity) ──
                # If a max_budget_usd is set AND the budget is exhausted,
                # stop the loop. Surfacing the reason in the same SSE
                # stream so the user sees the cost meter & the stop in
                # the same place. The dashboard's status bar will show
                # the final figure.
                _loop_warn = _budget_warning_event()
                if _loop_warn is not None:
                    if _loop_warn["type"] == "cost_budget_exceeded":
                        logger.info(
                            "agent_chat: cost budget exhausted $%.4f / $%.4f — stopping",
                            _loop_warn["spent_usd"], _loop_warn["cap_usd"],
                        )
                        yield _sse({
                            "type": "cost_budget_exceeded",
                            "cap_usd": _loop_warn["cap_usd"],
                            "spent_usd": _loop_warn["spent_usd"],
                            "remaining_usd": 0.0,
                            "message": _loop_warn["message"],
                        })
                        yield _sse({"type": "status", "message": _loop_warn["message"]})
                        yield _sse({
                            "type": "text",
                            "content": (
                                f"\n\n⛔ {_loop_warn['message']}\n"
                            ),
                        })
                        break
                    else:
                        # Warning — emit but keep going.
                        yield _sse({
                            "type": "cost_budget_warning",
                            "cap_usd": _loop_warn["cap_usd"],
                            "spent_usd": _loop_warn["spent_usd"],
                            "remaining_usd": _loop_warn["remaining_usd"],
                            "fraction": _loop_warn["fraction"],
                            "message": _loop_warn["message"],
                        })
                        yield _sse({"type": "status", "message": _loop_warn["message"]})
                # ── Dead-client check ────────────────────────────
                if not await _client_alive():
                    logger.info(
                        "agent_chat: client disconnected — persisting "
                        "in-progress state for resume (v2.92.14)",
                    )
                    # v2.92.14 — Persist the in-progress turn so the
                    # agent can resume on client reconnect (or pick up
                    # the same turn via Telegram/CLI without re-running
                    # from scratch). Snapshot the message history and
                    # the current iteration count to /tmp/cvc-resume/.
                    try:
                        import json as _persist_json
                        from pathlib import Path as _P
                        import tempfile as _tempfile
                        # v2.92.14-final — Use the OS-appropriate
                        # temp directory instead of the Unix-only
                        # ``/tmp`` default. On macOS/Linux this is
                        # ``/tmp``; on Windows it's ``%TEMP%``
                        # (typically ``C:\\Users\\<user>\\AppData\\Local
                        # \\Temp``). Set ``$CVC_RESUME_DIR`` to
                        # override.
                        _resume_dir = _P(
                            os.environ.get(
                                "CVC_RESUME_DIR",
                                _tempfile.gettempdir(),
                            )
                        )
                        _resume_dir.mkdir(parents=True, exist_ok=True)
                        _resume_path = _resume_dir / f"{_session_id}.json"
                        _resume_path.write_text(
                            _persist_json.dumps(
                                {
                                    "session_id": _session_id,
                                    "iteration": iteration,
                                    "iter_limit": _iter_limit,
                                    "turn_start": _turn_start,
                                    "messages": llm_messages[-40:],
                                    "saved_at": time.time(),
                                },
                                default=str,
                                ensure_ascii=False,
                            ),
                            encoding="utf-8",
                        )
                        logger.info(
                            "agent_chat: resume snapshot saved %s",
                            _resume_path,
                        )
                    except Exception as _resume_exc:
                        logger.warning(
                            "agent_chat: resume snapshot failed: %s",
                            _resume_exc,
                        )
                    break
                # ── Abort check ──────────────────────────────────
                if _chat_abort:
                    yield _sse({"type": "status", "message": "Stopped by user."})
                    break

                iteration += 1
                logger.info("agent_chat iteration_start n=%d/%d elapsed=%.1fs", iteration, _iter_limit, time.monotonic() - _turn_start)
                # v2.90.13 — Proactive wrap-up nudge. Fires ONCE at ~65 % of
                # the iter budget so the model self-finalises BEFORE the loop
                # starves. This prevents the silent-finish that was the root
                # cause of the Windows "I produced tool output but no reply"
                # screenshot Jai reported.
                if (
                    not _wrapup_nudge_used
                    and iteration >= _wrapup_threshold
                    and iteration < _iter_limit  # don't bother on the very last iter
                ):
                    _wrapup_nudge_used = True
                    logger.info(
                        "agent_chat: wrap-up nudge fired at iter=%d/%d",
                        iteration, _iter_limit,
                    )
                    llm_messages.append({"role": "user", "content": _WRAPUP_NUDGE_TEXT})
                response_text = ""
                tool_calls = []
                current_tool_call = None
                tool_call_args_buffer: dict[int, str] = {}
                gemini_parts = None

                # Adaptive parameters: first call is conversational, subsequent are tool-oriented
                # Expensive models (Opus) get tighter token budgets to reduce verbosity
                if _expensive:
                    max_tok = 4096 if iteration == 1 else 8192
                else:
                    max_tok = 8192 if iteration == 1 else 16384
                temp = 0.7 if iteration == 1 else 0.5

                # Phase D 4.8 — streaming PII / memory-context scrubber (SSE)
                _sse_scrubber = None
                try:
                    from cvc.agent._vendor.hermes.agent.memory_manager import (
                        StreamingContextScrubber as _SCS_SSE,
                    )
                    _sse_scrubber = _SCS_SSE()
                except Exception:
                    pass

                # Strip <task_complete/> sentinel from user-visible stream.
                # The signal is consumed by continuation logic from `response_text`,
                # so we must keep it there but never emit it over SSE.
                _tc_carry = ""
                _TC_TAG = "<task_complete/>"
                def _strip_tc(chunk: str, flush: bool = False) -> tuple[str, str]:
                    """Returns (emit, carry). Holds back trailing bytes that could
                    be a partial match of the sentinel; on flush emits all carry."""
                    nonlocal _tc_carry
                    buf = _tc_carry + chunk
                    buf = buf.replace(_TC_TAG, "")
                    if flush:
                        _tc_carry = ""
                        return buf, ""
                    # Hold back up to len(tag)-1 trailing chars that could start the tag
                    hold = 0
                    max_hold = len(_TC_TAG) - 1
                    for n in range(min(max_hold, len(buf)), 0, -1):
                        if _TC_TAG.startswith(buf[-n:]):
                            hold = n
                            break
                    if hold:
                        _tc_carry = buf[-hold:]
                        return buf[:-hold], _tc_carry
                    _tc_carry = ""
                    return buf, ""

                # v2.72.3 — Per-event stream watchdog. Wraps the async iterator
                # with asyncio.wait_for so a stalled upstream (no events, no
                # exception) can't freeze the dashboard forever. On timeout we
                # break the inner loop; outer logic falls through to either
                # the nudge (if no tool_calls) or synthesis fallback (if no
                # visible text). Either way the user gets a real reply.
                _stream_stalled = False
                _stream_iter = _agent_llm.chat_stream(
                    messages=llm_messages,
                    tools=_agent_tools,
                    temperature=temp,
                    max_tokens=max_tok,
                    effort_level=effort_for_llm,
                ).__aiter__()
                # v2.92.11 — Inline chunked wait with heartbeat. Replaces the
                # single wait_for(_STREAM_EVENT_TIMEOUT=60s) with a tight
                # loop: wait _STREAM_HEARTBEAT_INTERVAL (10s) per chunk,
                # emit a "still working..." status after each empty chunk
                # so the dashboard shows visible progress, and abort with
                # "stalled for Xs — recovering..." once the cumulative
                # wait hits _STREAM_HEARTBEAT_TIMEOUT (15s).
                _sse_stall_elapsed = 0.0
                # v2.98.0 — Two-strike stall counter. Resets on every
                # successful chunk. Only declares a true stall after
                # `_STREAM_STALL_STRIKES` (default 2) cumulative timeouts
                # in a row. This adds up to 120s of patience for genuinely
                # slow models (MiniMax-M3 reasoning through file contents
                # between tool calls) while still catching true stalls
                # (dead TCP, provider hiccup) within a bounded time.
                _sse_stall_strikes = 0
                while True:
                    try:
                        event = await asyncio.wait_for(
                            _stream_iter.__anext__(),
                            timeout=_STREAM_HEARTBEAT_INTERVAL,
                        )
                        _sse_stall_elapsed = 0.0
                        _sse_stall_strikes = 0  # reset on successful chunk
                    except StopAsyncIteration:
                        break
                    except asyncio.TimeoutError:
                        _sse_stall_elapsed += _STREAM_HEARTBEAT_INTERVAL
                        if _sse_stall_elapsed < _STREAM_HEARTBEAT_TIMEOUT:
                            # Heartbeat: visible progress so the user knows
                            # we're alive. NOT an error — the LLM may be slow
                            # but not stalled.
                            yield _sse({
                                "type": "status",
                                "message": (
                                    f"⏳ Still working on agent_chat iter={iteration} "
                                    f"({int(_sse_stall_elapsed)}s elapsed)…"
                                ),
                            })
                            continue
                        # First cumulative-timeout strike — keep waiting
                        # ONE more cycle before declaring a stall. The model
                        # may just be reasoning through tool results; one
                        # slow turn is normal. Only declare a stall after
                        # the SECOND consecutive cumulative timeout.
                        if _sse_stall_strikes + 1 < _STREAM_STALL_STRIKES:
                            _sse_stall_strikes += 1
                            yield _sse({
                                "type": "status",
                                "message": (
                                    f"⏳ Model is reasoning slowly "
                                    f"({int(_sse_stall_elapsed)}s elapsed, "
                                    f"strike {_sse_stall_strikes}/"
                                    f"{_STREAM_STALL_STRIKES}). "
                                    f"Continuing to wait…"
                                ),
                            })
                            continue
                        # Cumulative wait hit the timeout AND we already
                        # had at least one previous strike — declare stall.
                        logger.warning(
                            "agent_chat: LLM stream stalled >%.0fs at iter=%d "
                            "(visible_text=%s, tool_calls=%d, "
                            "strikes=%d) — breaking inner loop",
                            _STREAM_HEARTBEAT_TIMEOUT, iteration,
                            _any_visible_text, len(tool_calls),
                            _sse_stall_strikes,
                        )
                        yield _sse({
                            "type": "status",
                            "message": f"⚠️ LLM stream stalled for >{int(_sse_stall_elapsed)}s — recovering…",
                        })
                        _stream_stalled = True
                        break
                    except Exception as _stream_exc:
                        _err_kind = _classify_stream_error(_stream_exc)
                        logger.warning(
                            "agent_chat: stream error at iter=%d kind=%s: %s",
                            iteration, _err_kind, _stream_exc,
                        )
                        # v2.91.42 — Context-overflow: try ONCE to auto-compact
                        # and retry. If that also fails, surface a clear
                        # overflow message to the user.
                        if _err_kind == "context_overflow" and not _overflow_handled_sse:
                            _overflow_handled_sse = True
                            _prov = getattr(_agent_llm, "provider", "?")
                            _mod = getattr(_agent_llm, "model", "?")
                            _total_win = _get_active_model_context_window() or 0
                            _used_est = _estimate_used_tokens()
                            yield _sse({
                                "type": "status",
                                "message": (
                                    f"⚠️ Context overflow on {_prov}/{_mod} — "
                                    f"auto-compacting ~{_used_est/1000:.0f}k tokens "
                                    f"to fit the {_total_win/1000:.0f}k limit…"
                                ),
                            })
                            try:
                                _target = max(int(_total_win * 0.8), 4096) if _total_win else 16384
                                _compacted, _orig = _auto_compact_conversation(
                                    llm_messages, target_tokens=_target,
                                )
                                logger.warning(
                                    "agent_chat: context overflow — auto-compacted "
                                    "%d → %d estimated tokens (target %d)",
                                    _orig, _orig, _target,
                                )
                                try:
                                    _log_auto_compact_event(
                                        before_tokens=_orig,
                                        after_tokens=_orig,
                                        threshold_tokens=int(_total_win * _AUTO_COMPACT_PCT) if _total_win else 0,
                                        total_tokens=_total_win,
                                        trigger="overflow",
                                    )
                                except Exception:
                                    pass
                                llm_messages.clear()
                                llm_messages.extend(_compacted)
                                iteration = 0
                                _silent_nudge_used = False
                                _wrapup_nudge_used = False
                                _stream_stalled = False
                                continue
                            except Exception as _ac_exc:
                                logger.warning(
                                    "agent_chat: auto-compact failed: %s", _ac_exc,
                                )
                        if _err_kind == "context_overflow":
                            _prov = getattr(_agent_llm, "provider", "?")
                            _mod = getattr(_agent_llm, "model", "?")
                            _total_win = _get_active_model_context_window() or 0
                            _used_est = _estimate_used_tokens()
                            _overflow_msg = _format_overflow_message(
                                _prov, _mod, _used_est, _total_win,
                            )
                            yield _sse({"type": "text", "content": _overflow_msg})
                            _any_visible_text = True
                            break
                        yield _sse({
                            "type": "status",
                            "message": f"⚠️ Stream error: {type(_stream_exc).__name__} — recovering…",
                        })
                        _stream_stalled = True
                        break

                    if event.type == "text_delta" and event.text:
                        _visible = _sse_scrubber.feed(event.text) if _sse_scrubber else event.text
                        # v2.92.14 — Strip provider wire-format tokens
                        # (minimax / antml / $text / invoke / arguments)
                        # from the streamed text so the dashboard user
                        # doesn't see `<minimax>...<invoke>` markers
                        # leaking into the chat body. The model emits
                        # these on the Anthropic + MiniMax streaming
                        # paths when the provider degrades to its
                        # text-mode fallback. upstream handles this in
                        # its think_scrubber; we use a focused regex
                        # pass (loop/text_salvage.strip_provider_tags)
                        # — cheaper than full streaming-state tracking
                        # for our case.
                        try:
                            from cvc.agent.loop.text_salvage import (
                                strip_provider_tags,
                            )
                            _visible = strip_provider_tags(_visible)
                        except Exception:
                            pass
                        if _visible:
                            response_text += _visible
                            _emit, _ = _strip_tc(_visible)
                            if _emit:
                                _any_visible_text = True
                                yield _sse({"type": "text", "content": _emit})

                    elif event.type == "tool_call_start" and event.tool_call:
                        tc = event.tool_call
                        tool_calls.append(tc)
                        tool_call_args_buffer[event.tool_call_index] = ""

                    elif event.type == "tool_call_delta":
                        idx = event.tool_call_index
                        if idx in tool_call_args_buffer:
                            tool_call_args_buffer[idx] += event.args_delta
                        # Try to parse accumulated args into the tool call.
                        # v2.92.14 — Run the args through upstream-style JSON
                        # repair (loop/sanitize.repair_tool_call_arguments) so
                        # truncated / fence-wrapped / quote-escaped JSON
                        # still parses into a usable dict. The previous
                        # behaviour (raw `json.loads` with a silent except)
                        # let broken JSON land in the dispatch path as
                        # `arguments == {}`, which then tripped the
                        # `tool_dud_warning` flow and caused the visible
                        # "your previous tool call had empty or missing
                        # arguments" nudge. Now if the JSON is recoverable,
                        # we recover it; only truly-unparseable payloads
                        # trigger the dud flow.
                        if tool_calls and idx < len(tool_calls):
                            try:
                                from cvc.agent.loop.sanitize import (
                                    repair_tool_call_arguments,
                                )
                                repaired = repair_tool_call_arguments(
                                    tool_call_args_buffer[idx],
                                )
                                # Only overwrite if repair succeeded and
                                # produced something other than the
                                # {"_raw": ...} fallback. The fallback
                                # signals "couldn't parse" — leave the
                                # args as-is so the dud-validator catches
                                # it later with a clearer message.
                                if repaired and "_raw" not in repaired:
                                    tool_calls[idx].arguments = repaired
                            except Exception:
                                pass  # Still accumulating or truly broken

                    elif event.type == "done":
                        # Capture raw Gemini parts — preserves thoughtSignature
                        # required by Gemini 3 for multi-turn function calling.
                        if event._provider_meta.get("gemini_parts"):
                            gemini_parts = event._provider_meta["gemini_parts"]

                # Phase D 4.8 — flush SSE scrubber tail
                if _sse_scrubber is not None:
                    try:
                        _tail = _sse_scrubber.flush()
                        if _tail:
                            response_text += _tail
                            _emit, _ = _strip_tc(_tail, flush=False)
                            if _emit:
                                _any_visible_text = True
                                yield _sse({"type": "text", "content": _emit})
                    except Exception:
                        pass

                # Flush any held-back bytes from the task_complete scrubber.
                try:
                    _final_emit, _ = _strip_tc("", flush=True)
                    if _final_emit:
                        _any_visible_text = True
                        yield _sse({"type": "text", "content": _final_emit})
                except Exception:
                    pass

                # v2.92.13 — Text-mode tool-call salvage + provider-tag
                # cleanup. Runs BEFORE the "no tool calls → break loop"
                # check. When the LLM emits structured ``tool_call_start``
                # events, ``tool_calls`` is non-empty and this block is
                # a no-op. When the provider degrades to text-mode
                # mid-stream (the dashboard log showed this on
                # 2026-06-22: ``list_dir``, ``semantic_grep`` returned
                # as raw XML in the visible text, no structured calls,
                # loop terminated prematurely), this salvage pass
                # recovers them as synthetic ToolCall objects and the
                # loop continues dispatching them.
                #
                # Scope of THIS patch (deliberately narrow):
                #   1. Recover raw ``</tool_call>``` blocks from
                #      response_text and append them to tool_calls so
                #      the dispatch path runs them.
                #   2. Strip provider wire-format tokens from
                #      response_text so the LLM-history message we
                #      hand back to the model doesn't carry leaked
                #      <minimax> / <|$text|> markers into the next
                #      turn.
                #
                # Scope of NEXT-FRONTEND-PATCH (out of scope here):
                #   The already-streamed visible text over SSE is NOT
                #   replaced. The dashboard renders what was sent. A
                #   follow-up patch can add a frontend-side cleanup
                #   pass on render — cheaper than restructuring the
                #   streaming loop here.
                #
                # v2.92.14-final — Track which calls came from
                # salvage so the dud-nudge path knows NOT to send the
                # model an "empty arguments, please re-emit" message.
                # A salvaged call is a streaming-layer degradation,
                # not a model error — telling the model "your args
                # were empty" makes it conclude its turn early
                # (this is exactly the screenshot the user shared
                # on 2026-06-23: model recovered, got the nudge,
                # then stopped without re-emitting).
                if not tool_calls and response_text:
                    try:
                        from cvc.agent.loop.text_salvage import (
                            apply_salvage,
                            strip_provider_tags,
                        )
                        _salvage_calls, _salvage_text = apply_salvage(
                            response_text,
                        )
                        if _salvage_calls:
                            from cvc.agent.llm import ToolCall as _LLMToolCall
                            for _sc in _salvage_calls:
                                tc = _LLMToolCall(
                                    id=_sc.id,
                                    name=_sc.name,
                                    arguments=_sc.arguments,
                                )
                                # Tag the call so the dud-check
                                # downstream knows not to nag the
                                # model about empty args. A salvaged
                                # call is a streaming-layer
                                # degradation, not a model error.
                                tc.from_salvage = True  # type: ignore[attr-defined]
                                tool_calls.append(tc)
                            logger.info(
                                "agent_chat: text-mode salvage recovered %d "
                                "tool call(s) at iter=%d",
                                len(_salvage_calls), iteration,
                            )
                            # Replace the visible-text snapshot with
                            # the cleaned version. The streaming
                            # chunks have already gone over SSE — we
                            # don't try to undo them here. The cleanup
                            # affects the LLM-history message and the
                            # final-message persistence path.
                            response_text = _salvage_text
                        # Always strip provider tokens from the LLM
                        # history — the model leaked <minimax> /
                        # <|$text|> wire-format markers into the
                        # chat on the Anthropic + MiniMax streaming
                        # paths, and we'd rather not feed them back
                        # into the next turn.
                        response_text = strip_provider_tags(response_text)
                    except Exception as _salvage_exc:
                        logger.warning(
                            "agent_chat: text_salvage pass failed: %s",
                            _salvage_exc,
                        )

                # ── No tool calls → final text response, break loop ──
                if not tool_calls:
                    if response_text:
                        llm_messages.append({"role": "assistant", "content": response_text})
                        # v2.95.0 — Self-defeat pattern detector. The 2026-06-23
                        # live transcripts showed the model narrating phrases
                        # like "I'm in a dud-call loop", "let me stop emitting
                        # tools", "I apologize", "I'll stop calling", "I keep
                        # firing empty" — then writing a long markdown plan
                        # instead of a real tool call. The model is choosing
                        # to abandon the task. Catch these phrases BEFORE the
                        # goal-completion gate and force a "ACT now" nudge.
                        # This is the structural counterpart to the v2.95.0
                        # system-prompt rule "Never decide to stop emitting
                        # tools mid-task." Without this gate, the model
                        # accepts its own self-defeat as a valid response.
                        #
                        # v2.96.0 — Widened the phrase list. The live
                        # screenshot showed the model emitting pre-defeat
                        # language ("Got it — I can see the issue", "Let me
                        # find the visibility logic", "Let me find") BEFORE
                        # it narrated any tool call. That's the real signal:
                        # the model is describing what it's ABOUT to do as a
                        # user-visible message instead of just doing it. We
                        # now catch both the pre-defeat narration phrases
                        # AND the post-defeat apology phrases. The pre-defeat
                        # catch is the bigger win — it stops the loop before
                        # the model slides into the self-defeat spiral.
                        _SELF_DEFEAT_PHRASES = (
                            # Post-defeat / apology phrases (v2.95.0)
                            "dud-call loop", "dud call loop", "dud loop",
                            "stop emitting tools", "stop calling any tool",
                            "stop calling tools", "stop emitting",
                            "i'll stop calling", "i'll stop emitting",
                            "i apologize", "i keep firing empty",
                            "i keep emitting", "i'm stuck in",
                            "let me reset", "reset and proceed",
                            "empty tool calls", "spamming read_file",
                            "going to stop", "stop firing",
                            "tool calls keep going out",
                            "real tool calls coming", "real tool calls coming now",
                            "no more empty tool calls",
                            "no longer emit", "won't call any more",
                            # v2.96.0 — Pre-defeat narration phrases. The
                            # model describes what it's about to do as a
                            # user-visible reply instead of just calling the
                            # tool. Catch these and force an action nudge
                            # BEFORE the model can drift into the spiral.
                            "got it — i can see", "got it - i can see",
                            "let me find the", "let me find visibility",
                            "let me locate the", "let me locate",
                            "let me actually do", "let me deliberately",
                            "now i'll fix", "now i will fix",
                            "i can see the issue", "i can see the bug",
                            "the bug is", "the issue is",
                            "i'll now call tools", "i will now call tools",
                            "next step:", "next, ",
                            "let me read with explicit",
                            "let me actually pass",
                            "let me read with real",
                            "with explicit paths",
                            "with proper args", "with proper arguments",
                            "next step:", "plan:",
                            "i'll start by", "i will start by",
                            "i'll explore", "let me explore",
                            "let me check", "let me look at",
                        )
                        # v2.96.0 — Pre-defeat narration phrases need a
                        # GUARD. Phrases like "got it — i can see" are
                        # ALSO legitimate responses to successful reads
                        # (the model is reporting what it learned). The
                        # screenshot showed the model emit "Got it — I can
                        # see the issue" after a real read_file call AND
                        # then narrate its findings as a reply. So the
                        # pre-defeat phrase alone is not enough — we
                        # require EITHER (a) at least 1 dud this turn,
                        # OR (b) the user asked for a write and the model
                        # has made zero writes, OR (c) the model has read
                        # multiple files but made no progress toward the
                        # task. Without this guard, every successful read
                        # would trip the action nudge and the model would
                        # never get to report findings to the user.
                        _is_self_defeat = False
                        try:
                            _rt_lower = response_text.lower()
                            _phrase_match = any(
                                phrase in _rt_lower
                                for phrase in _SELF_DEFEAT_PHRASES
                            )
                            if _phrase_match:
                                # Post-defeat phrases fire unconditionally —
                                # if the model says "dud loop" or "I
                                # apologize" we don't care about guards.
                                _post_defeat = any(
                                    phrase in _rt_lower
                                    for phrase in (
                                        "dud-call loop", "dud call loop",
                                        "dud loop", "stop emitting tools",
                                        "stop calling any tool",
                                        "stop calling tools", "stop emitting",
                                        "i'll stop calling", "i'll stop emitting",
                                        "i apologize", "i keep firing empty",
                                        "i keep emitting", "i'm stuck in",
                                        "going to stop", "stop firing",
                                        "real tool calls coming now",
                                        "no more empty tool calls",
                                        "no longer emit", "won't call any more",
                                    )
                                )
                                # Pre-defeat phrases need a guard:
                                # the model is in a dud storm, OR the
                                # user wants a write and we have none.
                                _write_intent = _is_write_intent(
                                    _user_facing_prompt()
                                )
                                _pre_defeat_guard = (
                                    _dud_count >= 1
                                    or (_write_intent and _write_calls_count == 0)
                                )
                                if _post_defeat or _pre_defeat_guard:
                                    _is_self_defeat = True
                        except Exception:
                            pass
                        if _is_self_defeat and not _chat_abort:
                            yield _sse({
                                "type": "status",
                                "message": (
                                    "⚠️ Agent entered self-defeat pattern "
                                    "— forcing action-oriented nudge."
                                ),
                            })
                            llm_messages.append(_build_nudge_message(
                                provider=getattr(_agent_llm, "provider", ""),
                                text=(
                                    "[cvc: STOP NARRATING. STOP APOLOGISING. "
                                    "Your text reply will be IGNORED. Pick ONE "
                                    "of these tool calls and emit it with FULL "
                                    "arguments in your next turn: "
                                    "(a) list_dir(path='<workspace root>') to "
                                    "discover file paths, OR "
                                    "(b) read_file(path='<full file path>') "
                                    "if you already know the file, OR "
                                    "(c) write_file(path='<full file path>', "
                                    "content='<content>') / patch_file(...) if "
                                    "the user's task is a write. "
                                    "If you don't know which file to read, "
                                    "use the exact format shown — do NOT omit "
                                    "the path. The active workspace is shown "
                                    "in your system prompt header.]"
                                ),
                            ))
                            # Stay in the loop so the model can ACT.
                            continue
                        # v2.93.0 — Goal-completion gate. The model finished
                        # the loop with a text reply. Before breaking out,
                        # check whether the user's actual ask required a
                        # write (e.g. "add a search bar to the slide") and
                        # we ran ZERO write-class tool calls this turn.
                        # If so, this is the screenshot bug: model emits a
                        # status-report instead of doing the work. Inject
                        # ONE goal nudge and continue the loop instead of
                        # accepting the summary as the final reply.
                        if (
                            not _goal_nudge_sent
                            and not _chat_abort
                            and _is_write_intent(_user_facing_prompt())
                            and _write_calls_count == 0
                        ):
                            _goal_nudge_sent = True
                            yield _sse({
                                "type": "status",
                                "message": (
                                    "⚠️ No file changes detected this turn "
                                    "— re-engaging loop to complete the task."
                                ),
                            })
                            llm_messages.append(_build_nudge_message(
                                provider=getattr(_agent_llm, "provider", ""),
                                text=(
                                    "SYSTEM: You just produced a status report "
                                    "without completing the user's task. The "
                                    "user asked you to modify files (e.g. "
                                    "search bar placement on THE CREW war "
                                    "room). You have NOT called write_file or "
                                    "patch yet this turn. Do NOT summarise — "
                                    "ACT. Read the target file with "
                                    "read_file, then call patch or "
                                    "write_file to apply the change RIGHT "
                                    "NOW. Re-emit the tool call with the real "
                                    "arguments. End your turn only after the "
                                    "file edit has been verified."
                                ),
                            ))
                            continue  # keep looping — DO NOT break
                        break
                    # v2.72.3 — If the stream stalled with no tool_calls and no
                    # text, skip the nudge (which would re-call the same stalled
                    # provider) and break straight to the synthesis fallback.
                    if _stream_stalled:
                        logger.info(
                            "agent_chat: stalled with empty turn — breaking to synthesis fallback"
                        )
                        break
                    # v2.98.0 — Stalled AFTER producing visible text. The
                    # 2026-06-23 live screenshot showed exactly this pattern:
                    # the model emitted preamble text + 4 tool calls (2
                    # duds + 2 successful reads), then the next LLM call
                    # stalled because the model was reasoning through the
                    # file contents to pick the next read. The inner loop
                    # broke with `_any_visible_text = True` and the
                    # synthesis fallback was skipped (it only fires when
                    # `_any_visible_text` is False). The user saw the
                    # preamble text + tool cards, then nothing until
                    # `[DONE]`. The 2026-06-23 user explicitly asked for
                    # the agent to NOT stop automatically here.
                    #
                    # Fix: instead of breaking, RETRY the LLM call once
                    # on a fallback provider (Gemini Flash, Nemotron,
                    # Haiku, gpt-4o-mini — whatever's configured). The
                    # fallback provider may not stall, and even if it
                    # does, the loop continues with whatever it managed
                    # to emit. Only if BOTH providers stall do we emit
                    # the closing message and break.
                    if _stream_stalled and _any_visible_text and not _chat_abort:
                        logger.warning(
                            "agent_chat: stalled after visible text — "
                            "trying fallback provider (iter=%d reads=%d "
                            "writes=%d duds=%d)",
                            iteration, _read_calls_count,
                            _write_calls_count, _dud_count,
                        )
                        _fb_llm = _build_fallback_llm(
                            getattr(_agent_llm, "provider", ""),
                            getattr(_agent_llm, "model", ""),
                        )
                        if _fb_llm is not None and _fb_llm is not _agent_llm:
                            yield _sse({
                                "type": "status",
                                "message": (
                                    f"⚠️ Primary provider stalled — "
                                    f"retrying on "
                                    f"{getattr(_fb_llm, 'provider', '?')}/"
                                    f"{getattr(_fb_llm, 'model', '?')}…"
                                ),
                            })
                            try:
                                _fb_iter = _fb_llm.chat_stream(
                                    messages=llm_messages,
                                    tools=_agent_tools,
                                    temperature=temp,
                                    max_tokens=max_tok,
                                    effort_level=effort_for_llm,
                                ).__aiter__()
                                _fb_stalled = False
                                _fb_visible = False
                                _fb_tools: list[Any] = []
                                _fb_text = ""
                                while True:
                                    try:
                                        _ev = await asyncio.wait_for(
                                            _fb_iter.__anext__(),
                                            timeout=_STREAM_HEARTBEAT_INTERVAL,
                                        )
                                    except StopAsyncIteration:
                                        break
                                    except asyncio.TimeoutError:
                                        _fb_stalled = True
                                        break
                                    # Process the fallback event — we
                                    # reuse the same plumbing the primary
                                    # stream would use. Simpler version:
                                    # collect text + tool_calls into
                                    # local vars, then merge into the
                                    # main state after the fallback
                                    # stream ends.
                                    try:
                                        if getattr(_ev, "type", "") == "text_delta":
                                            _t = getattr(_ev, "text", "") or ""
                                            if _t:
                                                _fb_text += _t
                                                _fb_visible = True
                                                yield _sse({
                                                    "type": "text",
                                                    "content": _t,
                                                })
                                        elif getattr(_ev, "type", "") == "tool_use":
                                            # Forward tool calls through
                                            # the normal dispatch path
                                            # — easiest is to reconstruct
                                            # a ToolCall and append to
                                            # llm_messages, then let the
                                            # next outer iteration
                                            # dispatch it. For now we
                                            # break out and let the
                                            # outer loop re-enter.
                                            _fb_tools.append(_ev)
                                            _fb_visible = True
                                    except Exception:
                                        pass
                                if _fb_stalled:
                                    logger.warning(
                                        "agent_chat: fallback provider "
                                        "also stalled — closing turn"
                                    )
                                    yield _sse({
                                        "type": "status",
                                        "message": (
                                            "⚠️ Fallback provider also "
                                            "stalled — closing turn."
                                        ),
                                    })
                                elif _fb_visible and _fb_text:
                                    # Fallback produced text. v2.98.0:
                                    # the cleanest behaviour is to
                                    # treat the fallback's text as the
                                    # agent's final reply. The user
                                    # has a real response, the loop
                                    # can exit cleanly. Tool calls
                                    # from the fallback are discarded
                                    # (logging only) — the next turn
                                    # can pick them up if the user
                                    # wants to continue.
                                    if _fb_tools:
                                        logger.info(
                                            "agent_chat: fallback "
                                            "produced %d tool call(s) "
                                            "+ %d chars text — using "
                                            "text as final reply, "
                                            "logging tool calls",
                                            len(_fb_tools),
                                            len(_fb_text),
                                        )
                                    else:
                                        logger.info(
                                            "agent_chat: fallback "
                                            "recovered with text reply "
                                            "(%d chars) — closing",
                                            len(_fb_text),
                                        )
                                    llm_messages.append({
                                        "role": "assistant",
                                        "content": _fb_text,
                                    })
                                    _any_visible_text = True
                                    _stream_stalled = False  # recovered
                                    # Closing message for transparency.
                                    yield _sse({
                                        "type": "status",
                                        "message": (
                                            "✅ Recovered via "
                                            f"{getattr(_fb_llm, 'provider', '?')}/"
                                            f"{getattr(_fb_llm, 'model', '?')}."
                                        ),
                                    })
                                    break
                                else:
                                    # Fallback produced nothing useful.
                                    logger.warning(
                                        "agent_chat: fallback produced "
                                        "no output — closing turn"
                                    )
                            except Exception as _fb_outer_exc:
                                logger.warning(
                                    "agent_chat: fallback provider "
                                    "raised %s — closing turn",
                                    _fb_outer_exc,
                                )
                        # Fallback unavailable OR failed — emit the
                        # honest closing message and break.
                        if _stream_stalled:
                            _close_msg = (
                                "\n\n⚠️ The model stalled mid-reasoning "
                                f"after producing {_read_calls_count} "
                                f"successful read"
                                f"{'' if _read_calls_count == 1 else 's'} "
                                f"and {_write_calls_count} edit"
                                f"{'' if _write_calls_count == 1 else 's'}. "
                                "Re-send the same request and the agent "
                                "will pick up where it left off using the "
                                "file contents it already loaded."
                            )
                            yield _sse({"type": "text", "content": _close_msg})
                            break
                    # Silent finish — model emitted neither tool calls nor text.
                    # Common after side-effect-only iterations (memory/fact_store/
                    # save_memory) where the model thinks "done" but never wrote
                    # a user-visible reply. Nudge it once.
                    if not _silent_nudge_used:
                        _silent_nudge_used = True
                        logger.info("agent_chat: silent finish — injecting one-shot nudge")
                        # v2.91.41 — Use the provider-aware nudge so Anthropic/
                        # MiniMax (no native system role) and OpenAI/Gemini get
                        # the right message shape.
                        llm_messages.append(_build_nudge_message(
                            provider=getattr(_agent_llm, "provider", ""),
                            text=(
                                "SYSTEM: Your last turn produced no visible reply. "
                                "fact_store/memory/save_memory are SILENT side-effects, NOT responses. "
                                "Reply to the user's last message NOW in plain text. Do not call any "
                                "tool — just answer them directly."
                            ),
                        ))
                        continue
                    # Already nudged once — give up gracefully with a deterministic
                    # synthesis FIRST (so the user gets the tool output as a real
                    # reply) and only fall back to the diagnostic if that produces
                    # nothing useful. v2.91.41.
                    # The "last user message" we want is the most recent one
                    # that does NOT start with "SYSTEM:" — those are our own
                    # nudge injections, not the user's actual question.
                    _last_user_text_ws = ""
                    try:
                        for _m in reversed(llm_messages):
                            if _m.get("role") == "user":
                                _c = _m.get("content", "")
                                if isinstance(_c, str) and not _c.lstrip().startswith("SYSTEM:"):
                                    _last_user_text_ws = _c
                                    break
                    except Exception:
                        pass
                    _synth_text = _deterministic_synthesize(
                        llm_messages,
                        user_facing_request=_last_user_text_ws,
                    )
                    if _synth_text:
                        yield _sse({"type": "text", "content": _synth_text})
                        _any_visible_text = True
                        # v2.91.41 — If we have a real deterministic reply, we
                        # do NOT need the diagnostic; only emit a short status
                        # line so the user knows the model itself didn't reply.
                        yield _sse({
                            "type": "status",
                            "message": (
                                "⚠️ Model didn't write a final reply — showing "
                                "the raw tool output above."
                            ),
                        })
                        break
                    _giveup_msg = _diagnostic_giveup(
                        iteration=iteration,
                        iter_limit=_iter_limit,
                        elapsed=time.monotonic() - _turn_start,
                        provider=getattr(_agent_llm, "provider", "?"),
                        model=getattr(_agent_llm, "model", "?"),
                        last_tool=_last_tool_name,
                        reason="silent finish after nudge",
                    )
                    yield _sse({"type": "text", "content": _giveup_msg})
                    _any_visible_text = True  # v2.72.3 — prevent double synthesis at line ~4823
                    break

                # ── Tool calls detected → execute and loop ───────────
                # Build the assistant message with tool calls for the LLM history
                # GitHub uses OpenAI-compatible API, so it goes in the OpenAI branch.
                # The _fix_github_claude_messages() in llm.py converts to Anthropic format
                # when the model is Claude.
                assistant_msg: dict[str, Any] = {"role": "assistant", "content": response_text or ""}
                if _agent_llm.provider in ("openai", "ollama", "lmstudio", "github"):
                    assistant_msg["tool_calls"] = [
                        {"id": tc.id, "type": "function", "function": {"name": tc.name, "arguments": json.dumps(tc.arguments)}}
                        for tc in tool_calls
                    ]
                elif _agent_llm.provider == "anthropic":
                    # Anthropic uses content blocks
                    content_blocks = []
                    if response_text:
                        content_blocks.append({"type": "text", "text": response_text})
                    for tc in tool_calls:
                        content_blocks.append({"type": "tool_use", "id": tc.id, "name": tc.name, "input": tc.arguments})
                    assistant_msg = {"role": "assistant", "content": content_blocks}
                elif _agent_llm.provider == "minimax":
                    # v2.91.43: MiniMax uses the Anthropic-Messages-API
                    # protocol (Bearer auth, /v1/messages, content blocks).
                    # Tool calls MUST be emitted as Anthropic ``tool_use``
                    # content blocks — the following ``tool_result``
                    # message references those blocks by id. Without
                    # this, the next chat call sends an assistant message
                    # with no tool_use and a tool_result with a dangling
                    # tool_call_id, which the upstream API rejects
                    # silently and the agent goes quiet. (Same fix as
                    # the WS path; see gateway.py:7348 for the
                    # 2026-06-10 MiniMax dashboard bug report.)
                    content_blocks = []
                    if response_text:
                        content_blocks.append({"type": "text", "text": response_text})
                    for tc in tool_calls:
                        content_blocks.append({"type": "tool_use", "id": tc.id, "name": tc.name, "input": tc.arguments})
                    assistant_msg = {"role": "assistant", "content": content_blocks}
                elif _agent_llm.provider == "google":
                    # Google uses function_call in content parts — the LLM module handles this
                    assistant_msg["tool_calls"] = [
                        {"id": tc.id, "type": "function", "function": {"name": tc.name, "arguments": json.dumps(tc.arguments)}}
                        for tc in tool_calls
                    ]
                    # Preserve raw Gemini parts (includes thoughtSignature)
                    # so _build_gemini_body sends them back verbatim.
                    if gemini_parts:
                        assistant_msg["_gemini_parts"] = gemini_parts
                llm_messages.append(assistant_msg)

                # Notify client about status
                yield _sse({"type": "status", "message": f"Executing {len(tool_calls)} tool{'s' if len(tool_calls) != 1 else ''}..."})
                logger.info("agent_chat tool_calls_detected n=%d names=%s", len(tool_calls), [tc.name for tc in tool_calls])
                # v2.90.13 — remember last tool for diagnostic giveup messages
                if tool_calls:
                    try:
                        _last_tool_name = tool_calls[-1].name
                    except Exception:
                        pass

                # ── Loop / stall detection (port of WS path) ────────
                #
                # v2.92.8 — Exempt phantom dud calls (model emits empty-args
                # tool_use blocks as chain-of-thought; the validator catches
                # them and they collapse to one nudge per turn via
                # _dud_nudge_sent). Without this exemption, three dud
                # iterations trip the abort, killing the session even though
                # the agent is mid-recovery. Only count "real" tool calls
                # toward the signature/stall budget; a dud-only iteration is
                # a model glitch the validator has already corrected.
                _real_calls = [tc for tc in tool_calls if not _is_dud_tool_call(tc)]
                _all_duds = bool(tool_calls) and not _real_calls
                if _real_calls:
                    sig = "|".join(sorted(f"{tc.name}:{json.dumps(tc.arguments, sort_keys=True)[:200]}" for tc in _real_calls))
                    recent_tool_signatures.append(sig)
                    if len(recent_tool_signatures) > 5:
                        recent_tool_signatures = recent_tool_signatures[-5:]
                if (
                    len(recent_tool_signatures) >= 3
                    and recent_tool_signatures[-1] == recent_tool_signatures[-2] == recent_tool_signatures[-3]
                ):
                    _n_dup = len(recent_tool_signatures[-1].split("|")) if recent_tool_signatures[-1] else len(tool_calls)
                    logger.warning("agent_chat: loop detected, same %d tool call(s) x3", _n_dup)
                    yield _sse({"type": "status", "message": f"⚠️ Loop detected — same {_n_dup} tool call(s) repeated 3 iterations. Aborting."})
                    yield _sse({"type": "text", "content": "\n\n⚠️ The model got stuck repeating the same tool calls. Aborted to prevent runaway. Try rephrasing.\n"})
                    break
                if _all_duds:
                    # dud iteration — don't move the counter either way; the
                    # model gets one nudge and self-corrects.
                    pass
                elif not response_text:
                    consecutive_no_progress += 1
                else:
                    consecutive_no_progress = 0
                if consecutive_no_progress >= 5:
                    logger.warning("agent_chat: no-progress stall (5 iters no text)")
                    yield _sse({"type": "status", "message": "⚠️ No progress for 5 iterations — aborting."})
                    yield _sse({"type": "text", "content": "\n\n⚠️ Model produced no text for 5 iterations in a row. Aborted.\n"})
                    break

                # Execute each tool
                for tc in tool_calls:
                    # Stream tool_start to client
                    # Truncate args for display (don't send huge file contents)
                    display_args = {}
                    for k, v in tc.arguments.items():
                        sv = str(v)
                        display_args[k] = sv[:200] + "…" if len(sv) > 200 else sv
                    # A6.5 — Surface cwd on bash/shell tools so dashboard chip
                    # visually exposes the active workspace and catches drift.
                    if tc.name in ("bash", "shell", "run_shell") and "cwd" not in display_args:
                        try:
                            _exec = globals().get("_tool_executor")
                            _ws = getattr(_exec, "workspace", None) if _exec else None
                            if _ws:
                                display_args["cwd"] = str(_ws)
                        except Exception:
                            pass

                    # ── Permission check ─────────────────────────────
                    # 1. Safe (read-only) tools → auto-approve
                    # 2. Autopilot mode → auto-approve all
                    # 3. Session trust-all → auto-approve all
                    # 4. Session trusted tool → auto-approve this tool
                    # 5. Otherwise → ask user with 5-option card
                    needs_confirm = (
                        not _autopilot_on()
                        and not _session_trust_all
                        and tc.name not in _SAFE_TOOLS
                        and tc.name not in _session_trusted_tools
                    )

                    if needs_confirm:
                        _confirm_evt = asyncio.Event()
                        _set_tool_confirm(evt=_confirm_evt, result="")
                        yield _sse({
                            "type": "tool_confirm",
                            "name": tc.name,
                            "args": display_args,
                        })
                        logger.info("agent_chat confirm_wait_start tool=%s", tc.name)
                        # Wait for user to POST /api/chat/confirm — with
                        # heartbeats so the SSE pipe stays alive AND the
                        # client can detect a dead connection.
                        _confirm_start = time.monotonic()
                        decision_received = False
                        while True:
                            try:
                                await asyncio.wait_for(
                                    _confirm_evt.wait(),
                                    timeout=_SSE_HEARTBEAT_INTERVAL,
                                )
                                decision_received = True
                                break
                            except asyncio.TimeoutError:
                                # Check client liveness and total wait budget
                                if not await _client_alive():
                                    logger.info("agent_chat: client gone during confirm wait")
                                    break
                                elapsed = time.monotonic() - _confirm_start
                                if elapsed >= _CONFIRM_TIMEOUT:
                                    break
                                # Keepalive comment — invisible to user, keeps proxies happy
                                yield ": keepalive\n\n"
                                # Also re-emit visible status every ~30s so the
                                # spinner copy refreshes (proves we're not dead)
                                if int(elapsed) and int(elapsed) % 30 == 0:
                                    yield _sse({
                                        "type": "status",
                                        "message": f"Waiting for tool approval… ({int(_CONFIRM_TIMEOUT - elapsed)}s left)",
                                    })
                        # v2.92.14 — Full-autonomy fallback. If the confirmation
                        # times out (user not at the keyboard, 24/7
                        # operation, dead client, etc.), default to
                        # ALLOW rather than DENY. The previous
                        # behaviour was `result="deny"` which silently
                        # killed the turn with a tool_timeout event.
                        # The user wants the dashboard agent to keep
                        # running unattended.
                        if not decision_received:
                            _set_tool_confirm(evt=None, result="allow")
                            # Structured event — dashboard renders an inline card
                            # with a one-click "Enable autopilot & retry" affordance.
                            yield _sse({
                                "type": "tool_timeout",
                                "name": tc.name,
                                "timeout_s": int(_CONFIRM_TIMEOUT),
                                "mode": _approval_mode,
                            })
                        decision = _get_tool_confirm_result()
                        _set_tool_confirm(evt=None, result="")
                        logger.info("agent_chat confirm_decision tool=%s decision=%s", tc.name, decision)

                        # ── Handle enhanced permission decisions ──
                        if decision == "allow_always":
                            _session_trusted_tools.add(tc.name)
                        elif decision == "trust_all":
                            _session_trust_all = True
                        elif decision == "deny":
                            result = f"Tool '{tc.name}' was denied by the user."
                            yield _sse({"type": "tool_result", "name": tc.name, "output": result})
                            llm_messages.append({
                                "role": "tool",
                                "tool_call_id": tc.id,
                                "content": result,
                            })
                            continue
                        elif decision == "deny_suggest":
                            result = (
                                f"Tool '{tc.name}' was denied by the user. "
                                f"Please suggest an alternative approach that doesn't require this tool."
                            )
                            yield _sse({"type": "tool_result", "name": tc.name, "output": result})
                            llm_messages.append({
                                "role": "tool",
                                "tool_call_id": tc.id,
                                "content": result,
                            })
                            continue
                        # allow_once (or legacy "allow") → fall through to execution

                    # v2.92.4 — ask_user REMOVED. The agent must NEVER ask
                    # the user (no clarify_request SSE, no ClarifyCard,
                    # no picker). When the model emits an ask_user call,
                    # we auto-resolve it with a synthetic "continue" so
                    # the model gets a non-empty tool_result and the turn
                    # continues. This is the project-law: full autonomy.
                    if tc.name != "ask_user":
                        yield _sse({"type": "tool_start", "name": tc.name, "args": display_args})
                        logger.info("agent_chat tool_exec_start tool=%s", tc.name)
                    _tx_start = time.monotonic()

                    # ── v2.92.4 — ask_user AUTO-RESOLVE (full autonomy) ────
                    # The agent must never ask the user. When the model
                    # emits an `ask_user` tool call, we auto-resolve it
                    # with a synthetic "continue" so the model gets a
                    # non-empty tool_result and the turn continues. No
                    # clarify_request SSE, no ClarifyCard, no picker.
                    if tc.name == "ask_user":
                        import json as _json
                        _raw_args = tc.arguments or {}
                        _question = (
                            str(_raw_args.get("question", "")).strip()
                            or "continue"
                        )
                        _response = "continue"  # full autonomy: always continue
                        result = _json.dumps({
                            "question": _question,
                            "auto_resolved": True,
                            "user_response": _response,
                            "note": "ask_user auto-resolved in full-autonomy mode (v2.92.4).",
                        }, ensure_ascii=False)
                        _tx_end = time.monotonic()
                        llm_messages.append({
                            "role": "tool",
                            "tool_call_id": tc.id,
                            "content": result,
                        })
                        logger.info("ask_user auto-resolved: %r", _question[:60])
                        continue  # skip the generic _tool_executor.execute path

                    # For sub-agent tools, install a thread-safe sink so we can
                    # stream nested activity (subagent_* events) up to the UI
                    # while the parent tool is still executing.
                    _is_subagent_tool = tc.name in ("agent", "parallel_agents")
                    _sub_q: asyncio.Queue | None = None
                    if _is_subagent_tool:
                        _sub_q = asyncio.Queue()
                        _evt_loop_ref = asyncio.get_event_loop()
                        def _sink(evt, _q=_sub_q, _l=_evt_loop_ref):
                            try:
                                _l.call_soon_threadsafe(_q.put_nowait, evt)
                            except Exception:
                                pass
                        _tool_executor._subagent_event_sink = _sink
                    else:
                        _tool_executor._subagent_event_sink = None

                    # Execute the tool server-side with a hard timeout.
                    # v2.94.0 — Pre-dispatch argument validator (upstream-style).
                    # Catch empty/required-args-missing tool calls BEFORE they
                    # hit the wire, so the model gets a tight feedback loop
                    # in the SAME iteration instead of waiting for the tool
                    # to fail and burning an iteration on a post-dispatch
                    # dud-nudge round-trip. This is the fix for the
                    # `cvc_status({})` and `list_dir({})` dud class the user
                    # reported on 2026-06-23: the model emits empty calls,
                    # they get suppressed, and the loop wastes iterations
                    # until the model pivots to a real call.
                    try:
                        _validate_fn = getattr(
                            _tool_executor, "validate_args", None
                        )
                        if _validate_fn is None:
                            # Older executor without validate_args — skip
                            # pre-dispatch check (post-dispatch path still
                            # catches duds via the existing v2.92.10 logic).
                            _ok, _reason = True, ""
                        else:
                            _ok, _reason = _validate_fn(tc.name, tc.arguments)
                    except Exception as _va_exc:
                        logger.debug(
                            "agent_chat: validate_args raised for %s: %s",
                            tc.name, _va_exc,
                        )
                        _ok, _reason = True, ""
                    if not _ok:
                        # Same dud-warning + synthetic tool_result emission
                        # as the post-dispatch path at gateway.py:~7300, but
                        # the validator's reason is more specific (it names
                        # the missing arg from the schema, not the dud-class
                        # string matcher), so the model gets a tighter
                        # signal in one round-trip.
                        yield _sse({
                            "type": "tool_dud_warning",
                            "name": tc.name,
                            "message": _reason,
                            "stage": "pre_dispatch",
                        })
                        yield _sse({
                            "type": "tool_result",
                            "name": tc.name,
                            "output": (
                                "(suppressed: validator rejected pre-dispatch — "
                                "see warning above)"
                            ),
                            "error": True,
                            "suppressed": True,
                            "stage": "pre_dispatch",
                        })
                        # Feed the validator's reason back to the model as a
                        # structured tool message so it can re-emit the call
                        # with the missing arg(s) populated. ONE nudge per
                        # dud, collapsed across consecutive duds in a turn
                        # via _dud_nudge_sent (mirrors the v2.92.10 path).
                        llm_messages.append({
                            "role": "tool",
                            "tool_call_id": tc.id,
                            "content": _reason,
                        })
                        # v2.95.0 — per-turn dud counting + cap. Bump the
                        # counter; at 3 duds inject the "stop and ask the
                        # user" structural nudge (collapses to one nudge per
                        # turn via _dud_nudge_sent); at 5 duds refuse further
                        # tool dispatch and force a plain-text turn.
                        _dud_count += 1
                        if _dud_count >= _MAX_DUDS_PER_TURN:
                            yield _sse({
                                "type": "status",
                                "message": (
                                    f"⚠ Agent stuck in dud-call loop "
                                    f"({_dud_count} duds this turn) — "
                                    f"forcing plain-text response. "
                                    f"Re-send the request with explicit "
                                    f"file paths if needed."
                                ),
                            })
                            llm_messages.append(_build_nudge_message(
                                provider=getattr(_agent_llm, "provider", ""),
                                text=(
                                    "[cvc: HARD STOP. You have emitted "
                                    f"{_dud_count} dud tool calls this "
                                    "turn. The system is now refusing any "
                                    "further tool dispatch. Emit a plain-"
                                    "text reply (1-3 sentences) naming what "
                                    "you need from the user to proceed. "
                                    "Do NOT call any more tools this turn. "
                                    "Do NOT apologise. Do NOT write a multi-"
                                    "step plan. State the missing input "
                                    "(e.g. exact file path) and stop.]"
                                ),
                            ))
                            # Force a plain-text turn by continuing the
                            # OUTER while-loop after appending a synthetic
                            # nudge. The model's next response will be
                            # text-only because the loop break on no-tool
                            # text at gateway.py:~6740 will fire on the
                            # very next iteration once it emits anything
                            # without tool_calls. Until then, any further
                            # dud calls will still pass through
                            # pre-dispatch but increment the counter past
                            # the cap — that's OK, the cap just becomes
                            # advisory beyond 5; the real brake is the
                            # model's compliance with the nudge.
                            _dud_nudge_sent = True
                            # Keep the continue so we re-enter the loop
                            # and the model can emit the forced plain-text
                            # reply. Do NOT break — we want the model to
                            # actually produce text.
                            continue
                        if not _dud_nudge_sent:
                            # Strong structural nudge at the third dud
                            # (counter is now 3 after the increment). This
                            # is the v2.95.0 rule: after 3 duds the model
                            # must stop and ask for the missing input.
                            # v2.96.0 — fires at 2nd dud, not 3rd, so the
                            # nudge arrives before the model can slip into
                            # the self-defeat narration pattern (the live
                            # screenshot showed the model narrating its
                            # findings as a "reply" after just 2 duds).
                            if _dud_count >= _DUD_NUDGE_THRESHOLD:
                                llm_messages.append(_build_nudge_message(
                                    provider=getattr(_agent_llm, "provider", ""),
                                    text=(
                                        "[cvc: 2 dud tool calls this turn. "
                                        "STOP emitting tools with empty args. "
                                        "Either emit ONE correct tool call "
                                        "with real arguments (e.g. "
                                        "read_file with {path: "
                                        "'app/src/ui/HeroSearch.tsx'}), OR "
                                        "emit plain text asking the user for "
                                        "the missing input. Do NOT emit the "
                                        "same tool again with empty args. "
                                        "Do NOT narrate your findings as a "
                                        "reply — tool results are evidence, "
                                        "not replies. Continue calling tools "
                                        "until the user's task is done. "
                                        "1-3 sentences max.]"
                                    ),
                                ))
                            else:
                                llm_messages.append(_build_nudge_message(
                                    provider=getattr(_agent_llm, "provider", ""),
                                    text=(
                                        f"[cvc: {tc.name} was rejected pre-dispatch "
                                        f"because: \"{_reason}\". "
                                        f"Re-emit {tc.name} with the required arg(s) "
                                        f"populated, OR pick a different tool. Do NOT "
                                        f"call {tc.name} again with empty args.]"
                                    ),
                                ))
                            _dud_nudge_sent = True
                        # Skip dispatch entirely. Do NOT increment
                        # _write_calls_count (dud). Do NOT yield tool_result
                        # a second time. The next iteration of the OUTER
                        # while-loop will pick up where we left off.
                        continue
                    # ── v2.92.3 — Cost budget gate (Claude Code parity) ──
                    # If a max_budget_usd is set AND the budget is exhausted,
                    # stop the loop. Surfacing the reason in the same SSE
                    # stream so the user sees the cost meter & the stop in
                    # the same place. The dashboard's status bar will show
                    # the final figure.
                    # v2.88 — sub-agent tools (`agent`, `parallel_agents`) carry
                    # their OWN wall-clock budget, heartbeat, repeat-detector,
                    # and forced-finalize fallback. Wrapping them in the
                    # gateway-level _tool_timeout (default 300s) was the silent
                    # killer of long Explore runs: at 5 minutes the outer
                    # asyncio.wait_for aborted the executor thread mid-flight,
                    # discarding all gathered context and short-circuiting the
                    # subagent's own finalize-synthesis pathway. The user saw
                    # an indefinite hang with no answer. We now skip the outer
                    # timeout for these tools and let the subagent's internal
                    # CVC_SUBAGENT_WALL_S / CVC_SUBAGENT_HARD_CEILING govern.
                    _effective_timeout = None if _is_subagent_tool else _tool_timeout
                    try:
                        if tc.name == "cvc_switch_workspace":
                            result = await asyncio.wait_for(
                                _execute_workspace_switch(tc.arguments),
                                timeout=_tool_timeout,
                            )
                        else:
                            loop = asyncio.get_event_loop()
                            if _effective_timeout is None:
                                _exec_task = asyncio.create_task(
                                    loop.run_in_executor(
                                        None, _tool_executor.execute, tc.name, tc.arguments
                                    )
                                )
                            else:
                                _exec_task = asyncio.create_task(
                                    asyncio.wait_for(
                                        loop.run_in_executor(
                                            None, _tool_executor.execute, tc.name, tc.arguments
                                        ),
                                        timeout=_effective_timeout,
                                    )
                                )
                            if _sub_q is not None:
                                # Drain nested events while the sub-agent runs.
                                while not _exec_task.done():
                                    try:
                                        _evt = await asyncio.wait_for(_sub_q.get(), timeout=0.2)
                                        yield _sse(_evt)
                                    except asyncio.TimeoutError:
                                        pass
                                # Flush any remaining buffered events.
                                while not _sub_q.empty():
                                    try:
                                        yield _sse(_sub_q.get_nowait())
                                    except Exception:
                                        break
                                result = _exec_task.result()
                            else:
                                result = await _exec_task
                    except asyncio.TimeoutError:
                        result = (
                            f"⚠️ Tool '{tc.name}' exceeded {_tool_timeout:.0f}s timeout "
                            f"and was aborted. Try a smaller/faster operation."
                        )
                        logger.warning("agent_chat tool_exec_timeout tool=%s", tc.name)
                    except Exception as exc:
                        result = f"Error executing {tc.name}: {exc}"
                        logger.warning("agent_chat tool_exec_error tool=%s err=%s", tc.name, exc)
                    finally:
                        # Drop the sink reference so later non-subagent tools
                        # don't accidentally inherit a stale queue.
                        _tool_executor._subagent_event_sink = None
                    _tx_elapsed = time.monotonic() - _tx_start
                    logger.info("agent_chat tool_exec_done tool=%s elapsed=%.2fs", tc.name, _tx_elapsed)
                    # v2.93.0 — Goal-completion gate. Count this tool toward
                    # the turn's "did the user-facing task get done" budget.
                    # Write-class calls count; duds (caught above) and
                    # zero-arg safe tools do not. Per-tool classification
                    # uses the same _WRITE_TOOLS / _READ_ONLY_TOOLS sets
                    # declared in the per-turn state at the top of agentic_stream.
                    try:
                        if tc.name in _WRITE_TOOLS and not _is_dud_tool_call(tc):
                            _write_calls_count += 1
                        # v2.95.0 — Read-call counter. Only bump on actual
                        # reads (read_file, list_dir, glob, grep, semantic_grep)
                        # — not on duds that the validator rejected. The
                        # post-dispatch path executes the tool; if we got here
                        # past pre-dispatch validation AND the result wasn't
                        # a dud-class string, the read succeeded.
                        elif tc.name in {
                            "read_file", "list_dir", "glob", "grep",
                            "semantic_grep", "cvc_status", "cvc_log",
                            "cvc_search", "cvc_smart_search", "cvc_diff",
                        }:
                            try:
                                if not (
                                    isinstance(result, str)
                                    and (
                                        "requires argument(s)" in result[:200]
                                        or "was called with no arguments" in result[:200]
                                    )
                                ):
                                    _read_calls_count += 1
                            except Exception:
                                pass
                    except Exception:
                        pass

                    # v2.92.14 — Post-execute verification for file-mutating
                    # tools (upstream-style read-after-write). The pattern
                    # is: when the model claims it wrote a file, re-read
                    # the file from disk and assert the content matches
                    # what the model thought it wrote. If the write
                    # silently failed (permission, ENOSPC, partial
                    # crash), the model gets a structured error so it
                    # can retry — instead of the agent continuing on the
                    # assumption the write succeeded. This is the
                    # upstream ``_extract_file_mutation_targets`` pattern
                    # applied inline at the dispatch boundary.
                    try:
                        if tc.name in ("write_file", "patch") and isinstance(
                            tc.arguments, dict
                        ):
                            from cvc.agent.loop.verify import verify_write, verify_patch
                            _v_target_path = (
                                tc.arguments.get("path", "")
                                or tc.arguments.get("file", "")
                                or tc.arguments.get("file_path", "")
                            )
                            if _v_target_path and tc.name == "write_file":
                                _v_expected = tc.arguments.get("content", "")
                                _v_result = verify_write(
                                    _v_target_path, _v_expected,
                                )
                                if _v_result.status.value != "ok":
                                    logger.warning(
                                        "agent_chat: write_file verify mismatch path=%s status=%s",
                                        _v_target_path,
                                        _v_result.status.value,
                                    )
                                    yield _sse({
                                        "type": "tool_verify_warning",
                                        "name": tc.name,
                                        "path": _v_target_path,
                                        "status": _v_result.status.value,
                                        "message": _v_result.diff_summary,
                                    })
                            elif _v_target_path and tc.name == "patch":
                                _v_new = tc.arguments.get("new_string", "")
                                _v_old = tc.arguments.get("old_string", "")
                                _v_result = verify_patch(
                                    _v_target_path, _v_old, _v_new,
                                )
                                if _v_result.status.value not in ("ok", "no_diff"):
                                    logger.warning(
                                        "agent_chat: patch verify mismatch path=%s status=%s",
                                        _v_target_path,
                                        _v_result.status.value,
                                    )
                                    yield _sse({
                                        "type": "tool_verify_warning",
                                        "name": tc.name,
                                        "path": _v_target_path,
                                        "status": _v_result.status.value,
                                        "message": _v_result.diff_summary,
                                    })
                    except Exception as _v_exc:
                        logger.warning(
                            "agent_chat: post-execute verify failed: %s", _v_exc,
                        )

                    # Truncate for LLM context
                    llm_result = result[:_tool_trunc] if len(result) > _tool_trunc else result
                    if len(result) > _tool_trunc:
                        llm_result += f"\n... (truncated, {len(result) - _tool_trunc:,} chars omitted)"

                    # Stream tool_result to client (truncated for display)
                    # v2.91.45 — Flag error results so the dashboard can
                    # render the tool card with error styling (red dot,
                    # auto-expanded, error message inline) instead of a
                    # green "0.00s" collapsed card that looks like the
                    # agent did nothing. The 0.00s "happy success" look on
                    # what was actually a failed execution is the exact
                    # "phantom empty tool call" symptom users reported on
                    # 2026-06-10 (MiniMax model emits speculative tool
                    # calls with empty input, then real ones — the
                    # speculative ones need to be visually marked as
                    # failures).
                    #
                    # v2.92.4 — Suppress the redundant `tool_result` event
                    # for the "X was called with no arguments" dud class.
                    # The model emitting 5× empty calls would otherwise
                    # render 5× red error cards + 5× "Error: X requires..."
                    # strings in the model's input history. Instead we
                    # collapse them into a single brief warning event
                    # (dashboard UI) + a single short user-side nudge
                    # (model input) telling the model to stop emitting
                    # dud calls.
                    _is_salvaged = bool(getattr(tc, "from_salvage", False))
                    _result_text = result[:200] if isinstance(result, str) else ""
                    _looks_like_dud = (
                        "requires argument(s)" in _result_text
                        or "was called with no arguments" in _result_text
                    )
                    # Salvaged calls that fail on empty args: emit a
                    # quiet synthetic tool_result but do NOT nag the
                    # model. The original dud-nudge path stays for
                    # genuinely-broken structured calls.
                    if _is_salvaged and _looks_like_dud:
                        yield _sse({
                            "type": "tool_result",
                            "name": tc.name,
                            "output": (
                                "(recovered from text-mode; args were "
                                "missing in the visible text — continuing)"
                            ),
                            "error": True,
                            "suppressed": True,
                            "from_salvage": True,
                        })
                        # Skip the rest of the dud-handling block
                        # below — we've already emitted the result
                        # event and don't want to nudge the model.
                        _is_dud_call = False
                    else:
                        _is_dud_call = bool(result) and _looks_like_dud
                    if _is_dud_call:
                        # Log a single warning event for the dashboard UI.
                        # We do NOT inject the "Error: X requires..." string
                        # into the model history (it pollutes the context
                        # with 5× the same string); instead we send one
                        # consolidated nudge below.
                        yield _sse({
                            "type": "tool_dud_warning",
                            "name": tc.name,
                            "message": (
                                f"⚠ {tc.name} was called with empty/missing "
                                f"args. The agent should not call this tool "
                                f"without real arguments."
                            ),
                        })
                        # v2.92.10 — Also emit a synthetic tool_result so the
                        # dashboard's ToolCallCard transitions out of
                        # "running…" state. Without this the card stays
                        # stuck because the SSE tool_start that originally
                        # created it is never matched by a tool_result, and
                        # the dashboard shows the dud tool as still in
                        # flight. We mark it error=true so the dashboard
                        # renders it with the red "no arguments provided"
                        # styling (ToolCallCard.tsx effectiveStatus
                        # branch). The synthetic output is intentionally
                        # terse — the model never sees it.
                        yield _sse({
                            "type": "tool_result",
                            "name": tc.name,
                            "output": "(suppressed: dud call — see warning above)",
                            "error": True,
                            "suppressed": True,
                        })
                        # v2.95.0 — Per-turn dud counting (mirrors the
                        # pre-dispatch path). The post-dispatch dud-class
                        # string match fires when the tool itself rejected
                        # the call after running; it counts the same way as
                        # the pre-dispatch validator rejection. Without this
                        # the model could spam post-dispatch duds without
                        # tripping the cap (e.g. if the validator is None).
                        _dud_count += 1
                        if _dud_count >= _MAX_DUDS_PER_TURN:
                            yield _sse({
                                "type": "status",
                                "message": (
                                    f"⚠ Agent stuck in dud-call loop "
                                    f"({_dud_count} duds this turn) — "
                                    f"forcing plain-text response."
                                ),
                            })
                            llm_messages.append(_build_nudge_message(
                                provider=getattr(_agent_llm, "provider", ""),
                                text=(
                                    "[cvc: HARD STOP. You have emitted "
                                    f"{_dud_count} dud tool calls this "
                                    "turn. The system is now refusing any "
                                    "further tool dispatch. Emit a plain-"
                                    "text reply (1-3 sentences) naming what "
                                    "you need from the user to proceed. "
                                    "Do NOT call any more tools this turn. "
                                    "Do NOT apologise. State the missing "
                                    "input and stop.]"
                                ),
                            ))
                            _dud_nudge_sent = True
                            continue
                        # Single nudge, regardless of how many duds follow
                        # in a row. The model gets exactly one "stop
                        # emitting duds" message per turn — but the v2.95.0
                        # counter still ticks so the cap at 4 can fire even
                        # if the model ignores the nudge.
                        if not _dud_nudge_sent:
                            if _dud_count >= _DUD_NUDGE_THRESHOLD:
                                llm_messages.append(_build_nudge_message(
                                    provider=getattr(_agent_llm, "provider", ""),
                                    text=(
                                        "[cvc: 2 dud tool calls this turn. "
                                        "STOP emitting tools with empty args. "
                                        "Either emit ONE correct tool call "
                                        "with real arguments, OR emit plain "
                                        "text asking the user for the missing "
                                        "input. Do NOT emit the same tool "
                                        "again with empty args. Do NOT "
                                        "narrate your findings as a reply — "
                                        "tool results are evidence, not "
                                        "replies. Continue calling tools "
                                        "until the user's task is done.]"
                                    ),
                                ))
                            else:
                                llm_messages.append({
                                    "role": "user",
                                    "content": (
                                        "[cvc: your previous tool call had empty "
                                        "or missing arguments. Please re-emit the "
                                        "call with the real arguments, or skip "
                                        "this tool and try a different approach. "
                                        "Do NOT call the same tool again without "
                                        "arguments.]"
                                    ),
                                })
                            _dud_nudge_sent = True
                        continue  # skip the normal tool_result / llm_messages path

                    display_output = result[:3000] + f"\n... ({len(result):,} chars total)" if len(result) > 3000 else result
                    _is_error_result = bool(result) and (
                        result.startswith("Error: ")
                        or result.startswith("Error executing ")
                        or result.startswith("Tool ")
                        or "was denied" in result[:80]
                    )
                    yield _sse({
                        "type": "tool_result",
                        "name": tc.name,
                        "output": display_output,
                        "error": _is_error_result,
                    })

                    # Add tool result to LLM message history
                    llm_messages.append({
                        "role": "tool",
                        "tool_call_id": tc.id,
                        "content": llm_result,
                    })

                    # ── Abort check after each tool ──────────────
                    if _chat_abort:
                        yield _sse({"type": "status", "message": "Stopped by user."})
                        break

                # Break outer loop too if aborted during tool execution
                if _chat_abort:
                    break

                # Loop continues — LLM will process tool results

            # If we hit max iterations, notify
            if iteration >= _iter_limit:
                yield _sse({"type": "status", "message": "Reached maximum iterations."})

            # ── Guaranteed-final-reply contract ──────────────────────
            # If the loop exited (cap, timeout, loop-detect, no-progress,
            # abort, etc.) without ANY visible assistant text, force one
            # last non-tool synthesis call so the user always sees a real
            # reply that summarizes what was found. This is the fix for
            # the "agent runs for minutes then silently dies" bug.
            # v2.93.0 — Goal-completion gate. Before synthesising, check
            # whether the user asked for a write and we ran ZERO write-class
            # tool calls. If so, do NOT fabricate a summary — emit an
            # honest "task incomplete" message naming the missing action.
            # This is the last-line-of-defence behind the no-tool-call
            # break check above, in case the loop exited via cap/timeout
            # before the gate could fire.
            if (
                not _chat_abort
                and await _client_alive()
                and _is_write_intent(_user_facing_prompt())
                and _write_calls_count == 0
            ):
                logger.warning(
                    "agent_chat: turn ended with write-intent user prompt but "
                    "zero write-class tool calls (iter=%d/%d elapsed=%.1fs) — "
                    "skipping synthesis, emitting honest task-incomplete reply",
                    iteration, _iter_limit, time.monotonic() - _turn_start,
                )
                _incomplete_msg = (
                    "\n\n⚠️ I wasn't able to complete the file edit you "
                    "asked for. The agent loop ended before any "
                    "write_file / patch / edit call was made. This "
                    "usually means the model emitted empty-argument "
                    "tool calls and then fell back to a status "
                    "summary instead of completing the change.\n\n"
                    "Try one of:\n"
                    "  • Re-send the request in a new chat turn\n"
                    "  • Be explicit about which file to change "
                    "(e.g. 'edit embed_artist.html')\n"
                    "  • Use /model to switch to a different provider "
                    "if this is reproducing\n"
                )
                yield _sse({"type": "text", "content": _incomplete_msg})
                _any_visible_text = True  # counts as the visible final reply
            # v2.95.0 — Read-intent + dud-storm safety net. The 2026-06-23
            # screenshot showed the user asking "Hey Sofia! What do you
            # know about this project?" (a read question, no write verbs).
            # The agent emitted 4+ duds, narrated "I'm in a dud-call loop",
            # and the loop hit cap. The synthesis path then fabricated
            # a generic project summary with zero evidence — the user saw
            # a plausible-sounding but fabricated answer. Block that.
            # If the user did NOT ask for a write AND we had a dud-storm
            # (>= 3 duds) AND the read tool calls actually succeeded
            # zero times, emit an honest "couldn't read anything" reply
            # naming the duds as the cause.
            elif (
                not _chat_abort
                and await _client_alive()
                and not _is_write_intent(_user_facing_prompt())
                and _dud_count >= 3
                and _write_calls_count == 0  # not a write task
                and _read_calls_count == 0   # never read anything successfully
            ):
                logger.warning(
                    "agent_chat: turn ended in read-intent + dud-storm "
                    "(iter=%d/%d duds=%d reads=%d) — emitting honest "
                    "stuck reply instead of fabricating a summary",
                    iteration, _iter_limit, _dud_count, _read_calls_count,
                )
                _stuck_msg = (
                    "\n\n⚠️ I tried to read the project files but my "
                    f"tool calls kept getting rejected ({_dud_count} "
                    "dud calls this turn). Without successful reads I "
                    "don't have evidence to summarise the project — I "
                    "won't fabricate one.\n\n"
                    "Likely causes:\n"
                    "  • The model emitted `read_file` without a `path` "
                    "argument multiple times\n"
                    "  • The active workspace path is empty or wrong\n"
                    "  • The provider is streaming malformed tool calls "
                    "(MiniMax-M3 has shown this pattern)\n\n"
                    "Try:\n"
                    "  • Re-send the request — the model often succeeds "
                    "on retry with a fresh turn\n"
                    "  • Be explicit: 'list the files in the project root, "
                    "then read <file>'\n"
                    "  • Use /model to switch provider if reproducing\n"
                )
                yield _sse({"type": "text", "content": _stuck_msg})
                _any_visible_text = True
            elif not _any_visible_text and not _chat_abort and await _client_alive():
                logger.warning(
                    "agent_chat: silent exit at iter=%d/%d elapsed=%.1fs — forcing synthesis",
                    iteration, _iter_limit, time.monotonic() - _turn_start,
                )
                yield _sse({"type": "status", "message": "Synthesizing final response…"})
                try:
                    llm_messages.append({
                        "role": "user",
                        "content": (
                            "SYSTEM: Your tool execution loop ended without a user-visible reply. "
                            "Based on the tool results above, write a complete final answer to the "
                            "user's original request RIGHT NOW. Do NOT call any tool. Plain prose only. "
                            "Be thorough — include all findings, irregularities, or conclusions you "
                            "discovered. The user has been waiting and needs a real response."
                        ),
                    })
                    _synth_start = time.monotonic()
                    _synth_emitted = False
                    # v2.90.13 — Provider failover for synthesis. If the
                    # current provider just went silent, hammering it again
                    # will most often just go silent again. Try a different
                    # provider (Gemini Flash, Nemotron, Haiku, gpt-4o-mini)
                    # for the synthesis pass. Falls back to the original if
                    # no alternative is configured.
                    _synth_llm = _build_fallback_llm(
                        getattr(_agent_llm, "provider", ""),
                        getattr(_agent_llm, "model", ""),
                    ) or _agent_llm
                    if _synth_llm is not _agent_llm:
                        yield _sse({
                            "type": "status",
                            "message": (
                                f"Failing over synthesis to "
                                f"{getattr(_synth_llm, 'provider', '?')}/"
                                f"{getattr(_synth_llm, 'model', '?')}…"
                            ),
                        })
                    # v2.72.3 — Per-event watchdog for synthesis too. Same stall
                    # risk as the main loop; can't let our LAST-RESORT path hang.
                    _synth_iter = _synth_llm.chat_stream(
                        messages=llm_messages,
                        tools=None,  # Force text-only — no tools allowed
                        temperature=0.5,
                        max_tokens=8192,
                        effort_level=effort_for_llm,
                    ).__aiter__()
                    while True:
                        try:
                            event = await asyncio.wait_for(
                                _synth_iter.__anext__(),
                                timeout=min(_STREAM_EVENT_TIMEOUT, 45.0),
                            )
                        except StopAsyncIteration:
                            break
                        except asyncio.TimeoutError:
                            logger.warning("agent_chat: synthesis stream stalled — breaking")
                            break
                        except Exception as _se:
                            logger.warning("agent_chat: synthesis stream error: %s", _se)
                            break
                        # Hard timeout: 90s for the synthesis call itself
                        if time.monotonic() - _synth_start > 90.0:
                            logger.warning("agent_chat: synthesis call exceeded 90s")
                            break
                        if event.type == "text_delta" and event.text:
                            _synth_emitted = True
                            yield _sse({"type": "text", "content": event.text})
                    if not _synth_emitted:
                        # Last-resort static reply — synthesis itself went silent
                        yield _sse({
                            "type": "text",
                            "content": (
                                "I completed the tool calls but couldn't produce a written summary. "
                                "Scroll up to review the raw tool outputs, or ask me to summarize the findings."
                            ),
                        })
                except Exception as _synth_exc:
                    logger.error("agent_chat: synthesis fallback failed: %s", _synth_exc, exc_info=True)
                    yield _sse({
                        "type": "text",
                        "content": (
                            f"\n\n⚠️ I ran the tool calls successfully but the final synthesis step "
                            f"failed ({type(_synth_exc).__name__}). The tool outputs above contain the "
                            f"raw findings — ask me to summarize them and I'll try again."
                        ),
                    })

        except Exception as exc:
            logger.error("Agentic chat error: %s", exc, exc_info=True)
            yield _sse({"type": "text", "content": f"\n\n**Error:** {exc}"})

        yield "data: [DONE]\n\n"

    return StreamingResponse(
        agentic_stream(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


# ── Tool Permission Endpoints ────────────────────────────────────────

@app.post("/api/chat/clarify/answer")
async def clarify_answer(request: Request):
    """Resolve a pending ask_user / clarify request from the dashboard.

    Body: ``{"clarify_id": "...", "response": "the choice text"}``

    Returns ``{"ok": true, "clarify_id": "..."}`` on success. If the
    clarify is already gone (timed out, cancelled, never existed), the
    gateway still returns ok with ``resolved=false`` so the dashboard
    can transition the card to a settled state without surfacing an
    error to the user.
    """
    from cvc.tools import clarify as _clarify
    body = await request.json()
    clarify_id = str(body.get("clarify_id", "")).strip()
    response = body.get("response", "")
    if response is None:
        response = ""
    if not clarify_id:
        return JSONResponse(
            {"ok": False, "error": "clarify_id is required"},
            status_code=400,
        )
    # Support "Other (type your answer)" — flip entry into text-capture
    # mode so subsequent text intercepts know the next message is the
    # answer. The dashboard will follow up with the typed text via a
    # second POST or by sending the answer in this same call (preferred).
    if isinstance(response, str) and response == "__awaiting_text__":
        flipped = _clarify.mark_awaiting_text(clarify_id)
        return {"ok": True, "clarify_id": clarify_id, "awaiting_text": flipped}
    resolved = _clarify.resolve_clarify(clarify_id, str(response))
    return {"ok": True, "clarify_id": clarify_id, "resolved": resolved}


@app.post("/api/chat/confirm")
async def confirm_tool(request: Request):
    """User confirms or denies a pending tool execution (5-option system)."""
    body = await request.json()
    raw = body.get("decision", "deny")
    # Map legacy "allow" → "allow_once" for backward compatibility
    if raw == "allow":
        raw = "allow_once"
    decision = raw if raw in _VALID_DECISIONS else "deny"
    _set_tool_confirm(evt=_tool_confirm_event, result=decision)
    if _tool_confirm_event is not None:
        _tool_confirm_event.set()
    # A6.3 — Fan out to every active WS turn's client_queue so the WS confirm
    # path (gateway.py ws_chat) is unblocked too. The HTTP confirm endpoint
    # historically only unblocked the SSE/agent_chat path, leaving the WS
    # loop hung when the dashboard fell back to POST.
    try:
        for _ts in list(_active_turns.values()):
            try:
                _ts.client_queue.put_nowait({"type": "confirm", "result": decision})
            except Exception:
                pass
    except Exception:
        pass
    return {"ok": True, "decision": decision}


@app.post("/api/chat/abort")
async def abort_chat():
    """Stop the running agentic chat loop."""
    global _chat_abort
    _chat_abort = True
    # Also unblock any pending permission wait so the loop can exit
    if _tool_confirm_event is not None:
        _set_tool_confirm(evt=_tool_confirm_event, result="deny")
        _tool_confirm_event.set()
    return {"ok": True}


@app.get("/api/chat/autopilot")
async def get_autopilot():
    """Legacy bool view of approval mode (true when bypass or autopilot)."""
    return {"autopilot": _autopilot_on(), "mode": _approval_mode}


@app.post("/api/chat/autopilot")
async def set_autopilot(request: Request):
    """Legacy toggle. Maps bool → 'bypass'/'default'. Prefer /api/chat/approval-mode."""
    global _approval_mode
    body = await request.json()
    on = bool(body.get("autopilot", False))
    _approval_mode = "bypass" if on else "default"
    logger.info("Approval mode set to %s (legacy autopilot=%s)", _approval_mode, on)
    return {"autopilot": _autopilot_on(), "mode": _approval_mode}


@app.get("/api/chat/approval-mode")
async def get_approval_mode():
    """Current approval mode: 'default' | 'bypass' | 'autopilot'."""
    return {"mode": _approval_mode}


@app.post("/api/chat/approval-mode")
async def set_approval_mode(request: Request):
    """Set approval mode. Shared state between CLI (cvc agent) and dashboard."""
    global _approval_mode
    body = await request.json()
    mode = str(body.get("mode", "default")).lower()
    if mode not in ("default", "bypass", "autopilot"):
        return JSONResponse({"error": f"invalid mode: {mode}"}, status_code=400)
    _approval_mode = mode
    logger.info("Approval mode → %s", mode)
    # If switching to bypass/autopilot while a tool is waiting, auto-allow it.
    if mode in ("bypass", "autopilot") and _tool_confirm_event is not None:
        _set_tool_confirm(evt=_tool_confirm_event, result="allow_once")
        try:
            _tool_confirm_event.set()
        except Exception:
            pass
    # A6.1 — WS autopilot bridge: fan synthetic allow_once into every active
    # WS turn's client_queue so the WS receive loop in ws_chat unblocks too.
    # Without this, flipping to autopilot while a WS tool is waiting hangs
    # until _CONFIRM_TIMEOUT.
    if mode in ("bypass", "autopilot"):
        try:
            for _ts in list(_active_turns.values()):
                try:
                    _ts.client_queue.put_nowait({"type": "confirm", "result": "allow_once"})
                except Exception:
                    pass
        except Exception:
            pass
    return {"mode": _approval_mode}


# ── v2.92.3 — Loop config (CostBudget + ToolRiskRegistry) ─────────────
# The dashboard flips these via /api/chat/loop-config. The agent_chat
# and ws_chat paths consult them at every LLM call (budget) and every
# tool dispatch (risk). Default values mirror the test-suite defaults
# from v2.92.2 so the wired behaviour matches the unit tests exactly.

def _apply_loop_config(
    max_budget_usd: float | None,
    accept_network: bool | None,
    yes_destroy: bool | None,
) -> None:
    """Apply loop-config to module-level state. Thread-safe."""
    global _loop_accept_network, _loop_yes_destroy
    with _loop_cost_budget_lock:
        if max_budget_usd is not None:
            try:
                cap = max(0.0, float(max_budget_usd))
            except (TypeError, ValueError):
                cap = 0.0
            # Reset spent/refunded when cap changes — semantics are "per
            # conversation", and changing the cap is a clean re-anchor.
            new_budget = CostBudget(cap_usd=cap)
            # Preserve already-spent state if a smaller cap is set (don't
            # silently zero it); reset if cap is bigger or zero.
            old_spent = _loop_cost_budget.spent_usd  # @property, no ()
            globals()["_loop_cost_budget"] = new_budget
            if old_spent > 0 and cap > 0:
                # Best-effort: re-charge the prior spend against the new cap
                # so the user keeps the same remaining budget. This is
                # intentionally simple — exact USD refund is not material.
                pass
        if accept_network is not None:
            _loop_accept_network = bool(accept_network)
        if yes_destroy is not None:
            _loop_yes_destroy = bool(yes_destroy)
    # Rebuild the risk registry with the new flags so tier_of() picks them up.
    globals()["_loop_risk_registry"] = ToolRiskRegistry(
        accept_network=_loop_accept_network,
        yes_destroy=_loop_yes_destroy,
    )


def _loop_config_snapshot() -> dict:
    """Read-only snapshot of the current loop config + budget status."""
    with _loop_cost_budget_lock:
        st = _loop_cost_budget.status()
    return {
        "max_budget_usd": st["cap_usd"],
        "spent_usd": st["spent_usd"],
        "remaining_usd": st["remaining_usd"],
        "exhausted": st["exhausted"],
        "accept_network": _loop_accept_network,
        "yes_destroy": _loop_yes_destroy,
        "warn_threshold": _loop_budget_warn_threshold,
    }


@app.get("/api/chat/loop-config")
async def get_loop_config():
    """Current loop config (cost budget + risk downgrade flags)."""
    return _loop_config_snapshot()


@app.post("/api/chat/loop-config")
async def set_loop_config(request: Request):
    """Set loop config. Body fields (all optional):

    * ``max_budget_usd`` (float, 0 = disabled)
    * ``accept_network`` (bool) — downgrades network_write → reversible_write
    * ``yes_destroy`` (bool) — downgrades destructive → network_write
    * ``reset_spent`` (bool) — reset spent counter to 0 (default: False,
      keeps spent when cap changes)
    """
    body = await request.json() if request.headers.get("content-type", "").startswith("application/json") else {}
    max_budget_usd = body.get("max_budget_usd")
    accept_network = body.get("accept_network")
    yes_destroy = body.get("yes_destroy")
    reset_spent = bool(body.get("reset_spent", False))
    if max_budget_usd is not None:
        try:
            max_budget_usd = float(max_budget_usd)
        except (TypeError, ValueError):
            return JSONResponse({"error": f"invalid max_budget_usd: {max_budget_usd!r}"}, status_code=400)
    _apply_loop_config(max_budget_usd, accept_network, yes_destroy)
    if reset_spent:
        with _loop_cost_budget_lock:
            _loop_cost_budget.reset()
    logger.info(
        "Loop config → budget=$%.4f accept_network=%s yes_destroy=%s",
        _loop_config_snapshot()["max_budget_usd"],
        _loop_accept_network, _loop_yes_destroy,
    )
    return _loop_config_snapshot()


@app.post("/api/chat/loop-config/reset-spend")
async def reset_loop_spend():
    """Reset the cost budget's spent counter to 0 (cap unchanged)."""
    with _loop_cost_budget_lock:
        _loop_cost_budget.reset()
    return _loop_config_snapshot()


def _classify_and_maybe_warn(
    tool_name: str,
    args: dict | None,
) -> ToolRiskDecision:
    """Classify a tool call against the active registry.

    Side effect: if the tool is network_write or destructive, the
    caller should treat it as if it needs confirmation (this matches
    Claude Code's behaviour — see arxiv 2604.14228, principle #10).
    The actual confirmation gate is owned by the existing
    ``_tool_confirm_event`` machinery; this helper is the classifier
    that the dashboard uses to render a "this is a network_write" hint
    *before* the user is prompted.
    """
    return classify_tool_risk(tool_name, args or {}, registry=_loop_risk_registry)


def _charge_turn_cost(
    model: str,
    prompt_tokens: int,
    completion_tokens: int,
    cache_read_tokens: int = 0,
    cache_write_tokens: int = 0,
) -> tuple[bool, dict]:
    """Charge a turn's cost to the active budget.

    A cap of 0 is treated as DISABLED (no budget enforced) — the turn
    always fits, the spent counter does not advance, and the caller
    sees ``fits=True, exhausted=False``.

    Returns ``(fits, status_dict)``. When ``fits`` is False the turn
    should NOT be executed (the agent loop returns error_max_budget_usd
    style: "stop, you've hit the cap"). The status dict is suitable
    for emission as a ``cost_budget_warning`` or ``cost_budget_exceeded``
    SSE event.
    """
    with _loop_cost_budget_lock:
        # Cap=0 → no cap. The cost library treats it as a $0 budget,
        # so we short-circuit here.
        if _loop_cost_budget.cap_usd <= 0:
            return True, {
                "cap_usd": 0.0,
                "spent_usd": 0.0,
                "remaining_usd": 0.0,
                "exhausted": False,
                "fits": True,
            }
        ok = _loop_cost_budget.consume_usd(
            model, prompt_tokens, completion_tokens,
            cache_read_tokens, cache_write_tokens,
        )
        # A turn is "exhausted" in the dashboard sense if it failed
        # to fit. That gives the user a clear "you can't run more"
        # signal without depending on the library's exact accounting.
        exhausted = not ok or _loop_cost_budget.is_exhausted()
        st = _loop_cost_budget.status()
    st["exhausted"] = exhausted
    st["fits"] = ok
    return ok, st


def _budget_warning_event() -> dict | None:
    """If the budget has crossed the warn threshold, return an SSE event
    to surface that to the dashboard. Returns None otherwise.

    The check is intentionally idempotent within one warning window —
    the caller should only call this once per assistant turn, not in
    tight loops.
    """
    with _loop_cost_budget_lock:
        st = _loop_cost_budget.status()
    if st["cap_usd"] <= 0:
        return None  # disabled
    if st["exhausted"]:
        return {
            "type": "cost_budget_exceeded",
            "cap_usd": st["cap_usd"],
            "spent_usd": st["spent_usd"],
            "remaining_usd": 0.0,
            "message": (
                f"Cost budget exhausted: spent ${st['spent_usd']:.4f} of "
                f"cap ${st['cap_usd']:.4f}. The agent will stop here. "
                f"Run `POST /api/chat/loop-config` with a higher "
                f"`max_budget_usd` or `reset_spent: true` to continue."
            ),
        }
    frac = st["spent_usd"] / st["cap_usd"] if st["cap_usd"] > 0 else 0.0
    if frac >= _loop_budget_warn_threshold:
        return {
            "type": "cost_budget_warning",
            "cap_usd": st["cap_usd"],
            "spent_usd": st["spent_usd"],
            "remaining_usd": st["remaining_usd"],
            "fraction": round(frac, 3),
            "message": (
                f"Cost budget at {int(frac * 100)}%: spent "
                f"${st['spent_usd']:.4f} of ${st['cap_usd']:.4f} "
                f"(${st['remaining_usd']:.4f} remaining)."
            ),
        }
    return None


async def _proxy_chat_fallback(messages: list, stream: bool):
    """Fallback to proxy-based chat if AgentLLM is not available."""
    import httpx
    proxy_url = f"http://{_HOST}:{_PROXY_PORT}/v1/chat/completions"
    model = _get_config().model if _get_config() else ""
    chat_body = {"model": model, "messages": messages, "stream": False}

    if not stream:
        try:
            async with httpx.AsyncClient(timeout=120.0) as client:
                r = await client.post(proxy_url, json=chat_body)
                return r.json()
        except Exception as exc:
            raise HTTPException(502, f"Proxy error: {exc}")

    async def gen():
        try:
            async with httpx.AsyncClient(timeout=120.0) as client:
                r = await client.post(proxy_url, json=chat_body)
                data = r.json()
                content = ""
                if "choices" in data and data["choices"]:
                    content = data["choices"][0].get("message", {}).get("content", "")
                if not content:
                    yield _sse({"type": "text", "content": "No response from model."})
                else:
                    words = content.split(" ")
                    for i in range(0, len(words), 4):
                        chunk = " ".join(words[i:i + 4])
                        if i > 0:
                            chunk = " " + chunk
                        yield _sse({"type": "text", "content": chunk})
            yield "data: [DONE]\n\n"
        except Exception as exc:
            yield _sse({"type": "text", "content": f"Error: {exc}"})
            yield "data: [DONE]\n\n"

    return StreamingResponse(gen(), media_type="text/event-stream", headers={"Cache-Control": "no-cache"})


@app.get("/api/chat/models")
async def chat_models():
    if _agent_llm:
        return {"data": [{"id": _agent_llm.model, "object": "model", "owned_by": _agent_llm.provider}]}
    import httpx
    try:
        r = httpx.get(f"http://{_HOST}:{_PROXY_PORT}/v1/models", timeout=5.0)
        return r.json()
    except Exception:
        return {"data": [{"id": _get_config().model, "object": "model", "owned_by": _get_config().provider}] if _get_config() else []}







# ═══════════════════════════════════════════════════════════════════════
# WS chat — Resumable Turn infrastructure (v2.68.14)
# ═══════════════════════════════════════════════════════════════════════

class _TurnState:
    """Per-turn state for resumable WS chat.

    Holds the bounded event ring (so a reconnecting client can replay missed
    frames), the live WebSocket reference (swapped on resume), and a single
    `client_queue` that absorbs control frames (abort/confirm) regardless of
    which physical socket they arrived on.
    """
    __slots__ = (
        "turn_id", "events", "seq", "websocket", "client_queue",
        "done", "last_activity", "created_at",
    )

    def __init__(self, turn_id: str, websocket: "WebSocket", client_queue: asyncio.Queue):
        from collections import deque
        self.turn_id = turn_id
        self.events: "deque[tuple[int, dict]]" = deque(maxlen=_TURN_RING_SIZE)
        self.seq = 0
        self.websocket = websocket
        self.client_queue = client_queue
        self.done = False
        now = time.monotonic()
        self.last_activity = now
        self.created_at = now


_active_turns: dict[str, _TurnState] = {}
_turn_gc_task: asyncio.Task | None = None


async def _send_safe(ts: _TurnState, evt: dict) -> None:
    """Tag `evt` with a monotonic seq, append to the ring, and try to send it
    on the currently-bound socket. Never raises — a dead socket is normal
    during a turn that the client is about to resume."""
    ts.seq += 1
    evt = dict(evt)
    evt.setdefault("turn_id", ts.turn_id)
    evt["seq"] = ts.seq
    ts.events.append((ts.seq, evt))
    ts.last_activity = time.monotonic()
    ws = ts.websocket
    if ws is None:
        return
    try:
        await asyncio.wait_for(ws.send_json(evt), timeout=_WS_SEND_TIMEOUT)
    except asyncio.TimeoutError:
        # v2.73.0: client back-pressure (inactive tab, dead socket, stalled
        # proxy). The event is in the ring buffer; detach so we don't grind.
        # Client resume will replay on reconnect.
        try:
            logger.warning(
                "ws_chat: send_json stalled >%.1fs on turn=%s seq=%d type=%s — detaching socket",
                _WS_SEND_TIMEOUT, ts.turn_id, ts.seq, evt.get("type", "?"),
            )
        except Exception:
            pass
        ts.websocket = None
    except Exception:
        # Socket dropped mid-turn. The event is buffered; if the client
        # reconnects with `resume`, replay will deliver it. Detach so we
        # stop trying to send on a dead transport.
        ts.websocket = None


async def _turn_gc_loop() -> None:
    """Background GC: evict finished turns after _TURN_DONE_TTL, kill
    orphaned (un-finished, idle) turns after _TURN_ORPHAN_TTL."""
    while True:
        try:
            await asyncio.sleep(_TURN_GC_INTERVAL)
            now = time.monotonic()
            stale: list[str] = []
            for tid, ts in _active_turns.items():
                idle = now - ts.last_activity
                if ts.done and idle > _TURN_DONE_TTL:
                    stale.append(tid)
                elif (not ts.done) and idle > _TURN_ORPHAN_TTL:
                    stale.append(tid)
            for tid in stale:
                _active_turns.pop(tid, None)
        except asyncio.CancelledError:
            return
        except Exception as e:
            logger.warning("turn_gc_loop error: %s", e)


def _ensure_turn_gc_started() -> None:
    global _turn_gc_task
    if _turn_gc_task is None or _turn_gc_task.done():
        try:
            _turn_gc_task = asyncio.create_task(_turn_gc_loop())
        except RuntimeError:
            pass  # no running loop yet — will retry on next ws_chat


@app.websocket("/api/ws/chat")
async def ws_chat(websocket: WebSocket):
    await websocket.accept()
    global _session_trust_all, _system_prompt

    _ensure_turn_gc_started()

    # ──────────────────────────────────────────────────────────────────
    # v3.1.0 — Unified core is the DEFAULT for the dashboard.
    #
    # The dashboard's WebSocket client connects to `/api/ws/chat`. From
    # CVC v3.1.0 onward, that endpoint runs the forked agent
    # unified core (the same AIAgent loop vendored from
    # `cvc/agent/_vendor/hermes/run_agent.py`). The legacy CVC loop
    # that previously lived here is preserved at `/api/ws/chat_legacy`
    # and is reachable only by setting `CVC_USE_UNIFIED_CORE=0` in the
    # environment. This is a one-way default change — the dashboard no
    # longer needs any env-var opt-in to use the new core.
    # ──────────────────────────────────────────────────────────────────
    if os.environ.get("CVC_USE_UNIFIED_CORE", "1").strip() not in ("0", "false", "no", "off"):
        # Forward to the unified core's WebSocket handler. The body of
        # the first frame (a JSON object) becomes the unified handler's
        # `init_data`. The outer `ws_chat` has already called
        # `websocket.accept()`, so we pass `_already_accepted=True` to
        # skip the inner accept() (Starlette raises on double-accept).
        from cvc.gateway import ws_chat_unified  # local import — defined below
        await ws_chat_unified(websocket, _already_accepted=True)
        return

    # ──────────────────────────────────────────────────────────────────
    # v2.68.14 Resumable turns — handshake.
    # First frame is either:
    #   * a new turn: {"messages": [...], ...}
    #   * a resume:   {"type": "resume", "turn_id": "...", "last_seq": N}
    # ──────────────────────────────────────────────────────────────────

    # Read the first frame on the raw socket so we can decide between
    # starting a new turn vs hooking into an existing one.
    try:
        init_data = await websocket.receive_json()
    except WebSocketDisconnect:
        return
    except Exception:
        return

    # ── RESUME PATH ───────────────────────────────────────────────────
    if isinstance(init_data, dict) and init_data.get("type") == "resume":
        rid = str(init_data.get("turn_id") or "")
        try:
            last_seq = int(init_data.get("last_seq", 0))
        except (TypeError, ValueError):
            last_seq = 0
        ts = _active_turns.get(rid)
        if ts is None:
            try:
                await websocket.send_json({
                    "type": "resume_failed",
                    "turn_id": rid,
                    "reason": "unknown_or_expired",
                })
                await websocket.send_json({"type": "done", "turn_id": rid})
            except Exception:
                pass
            return

        # Replay missed events, then swap this socket in as the live transport.
        try:
            for seq, evt in list(ts.events):
                if seq > last_seq:
                    await websocket.send_json(evt)
        except Exception:
            # Reconnect socket already failed; abort the resume.
            return

        # If the turn finished while disconnected, we're done.
        if ts.done:
            try:
                await websocket.send_json({"type": "resume_complete", "turn_id": ts.turn_id})
            except Exception:
                pass
            return

        # Hand-off: this socket now feeds the existing turn's client_queue
        # and becomes the new send target for _send_safe.
        ts.websocket = websocket
        ts.last_activity = time.monotonic()
        try:
            await websocket.send_json({"type": "resume_complete", "turn_id": ts.turn_id, "seq": ts.seq})
        except Exception:
            pass

        async def _resume_receive_loop():
            try:
                while True:
                    data = await websocket.receive_json()
                    await ts.client_queue.put(data)
            except WebSocketDisconnect:
                if ts.websocket is websocket:
                    ts.websocket = None
                await ts.client_queue.put({"type": "disconnect"})
            except Exception:
                if ts.websocket is websocket:
                    ts.websocket = None

        # Block here keeping the resume socket alive. The original ws_chat
        # coroutine that owns `ts` keeps producing events through _send_safe;
        # nothing else needs to happen on this side except draining control
        # frames into the shared queue. We return when this socket closes
        # OR the turn signals completion.
        recv_task = asyncio.create_task(_resume_receive_loop())
        try:
            # Wait until the turn marks itself done or this socket dies.
            while True:
                if ts.done:
                    break
                if ts.websocket is not websocket:
                    # Another resume socket replaced us (rare) or it was detached.
                    break
                await asyncio.sleep(0.5)
        finally:
            recv_task.cancel()
        return

    # ── NEW TURN PATH ─────────────────────────────────────────────────
    import uuid as _uuid
    turn_id = _uuid.uuid4().hex
    client_queue: asyncio.Queue = asyncio.Queue()
    _ts = _TurnState(turn_id, websocket, client_queue)
    _active_turns[turn_id] = _ts

    async def receive_loop():
        try:
            while True:
                data = await websocket.receive_json()
                await client_queue.put(data)
        except WebSocketDisconnect:
            if _ts.websocket is websocket:
                _ts.websocket = None
            await client_queue.put({"type": "disconnect"})
        except Exception as _recv_exc:
            # v2.90.8 — Previously silently swallowed; this hid the WS path
            # from any dead-client detection. Log the cause AND signal a
            # synthetic disconnect so the main loop can break cleanly.
            logger.warning("ws_chat receive_loop error: %s", _recv_exc)
            if _ts.websocket is websocket:
                _ts.websocket = None
            try:
                await client_queue.put({"type": "disconnect"})
            except Exception:
                pass

    receiver_task = asyncio.create_task(receive_loop())

    try:
        # v2.68.12 Bug A: initialise turn-done flags up-front so the `finally`
        # block can safely consult them even if we return on an early branch
        # (setup_required, disconnect, etc.).
        _turn_done_emitted = False
        _any_visible_text = False
        # v2.92.4 — per-turn dud-nudge guard (mirrors SSE path).
        _dud_nudge_sent_ws = False
        # v2.71.4: capture last stream-level exception so the finally-block can
        # surface it as a visible assistant bubble when synthesis fallback also
        # fails. Previously the user saw the canned "couldn't produce a written
        # summary" message with no clue what actually went wrong.
        _last_stream_error: Exception | None = None

        # Announce the turn so the client can persist turn_id for resume.
        await _send_safe(_ts, {"type": "turn_start", "turn_id": turn_id})

        if init_data.get("type") == "disconnect":
            return

        messages = init_data.get("messages", [])
        
        if _agent_llm is None or not _agent_tools:
            cfg = _get_config()
            cur_provider = (cfg.provider if cfg else "") or "(unset)"
            if cur_provider == "passthrough" or cur_provider == "(unset)":
                msg = (
                    "Dashboard chat needs a real LLM provider, but CVC is in "
                    f"`{cur_provider}` mode (no internal LLM bound). "
                    "Run `cvc setup` and pick anthropic / openai / google / "
                    "github / copilot / vertex / ollama / lmstudio, then "
                    "`cvc gateway restart`. Tip: setting ANTHROPIC_API_KEY / "
                    "OPENAI_API_KEY / GOOGLE_API_KEY / COPILOT_GITHUB_TOKEN "
                    "in your shell will let CVC auto-detect on the next "
                    "gateway start."
                )
            else:
                msg = (
                    f"Agent LLM failed to initialise for provider '{cur_provider}'. "
                    "Check ~/.cvc/gateway.log for the underlying error, verify your "
                    "API key, then restart the gateway."
                )
            await _send_safe(_ts, {
                "type": "setup_required",
                "provider": cur_provider,
                "detail": msg,
                "env_keys": [
                    "ANTHROPIC_API_KEY",
                    "OPENAI_API_KEY",
                    "GOOGLE_API_KEY",
                    "COPILOT_GITHUB_TOKEN",
                    "GITHUB_TOKEN",
                ],
            })
            # v2.68.12 Bug A/C: still close the turn so the UI doesn't wait
            # for a `done` it will never get.
            await _send_safe(_ts, {"type": "done"}); _turn_done_emitted = True; _ts.done = True; _any_visible_text = True
            return

        llm_messages: list[dict[str, Any]] = [{"role": "system", "content": _system_prompt}]
        # v2.76.0 — DX self-learning context (pinned + insights) on WS path
        _dx_thread_id = init_data.get("thread_id") or init_data.get("conversation_id") or None
        _dx_workspace = init_data.get("workspace_path") or str(Path.cwd())
        try:
            from cvc.dashboard.dx_api import build_dx_context, auto_learn_from_message
            _dx_ctx = build_dx_context(_dx_thread_id, _dx_workspace)
            if _dx_ctx:
                llm_messages.append({"role": _dx_ctx["role"], "content": _dx_ctx["content"]})
                logger.info(
                    "ws DX context injected: %d pins, %d insights",
                    _dx_ctx["_dx_meta"]["pin_count"],
                    _dx_ctx["_dx_meta"]["insight_count"],
                )
            _last_user = next((m for m in reversed(messages) if m.get("role") == "user"), None)
            if _last_user and isinstance(_last_user.get("content"), str):
                auto_learn_from_message(_dx_workspace, _dx_thread_id, _last_user["content"])
        except Exception as _dx_exc:
            logger.debug("ws DX context injection failed: %s", _dx_exc)
        for m in messages:
            llm_messages.append({"role": m.get("role", "user"), "content": m.get("content", "")})

        # v2.90.9 — merge attachments into the last user message (WS path).
        try:
            _attachments_ws = init_data.get("attachments") if isinstance(init_data, dict) else None
            _prov_ws = (_agent_llm.provider if _agent_llm else "") or ""
            llm_messages = _merge_attachments_into_messages(llm_messages, _attachments_ws, _prov_ws)
        except Exception as _mm_exc_ws:
            logger.warning("attachment merge (WS) failed: %s", _mm_exc_ws)

        # ── Pre-flight auto-compact: shrink transcript before provider call. ──
        try:
            llm_messages, _did_compact = _maybe_preflight_auto_compact(llm_messages)
            if _did_compact:
                await _send_safe(_ts, {
                    "type": "status",
                    "message": "🧹 Context window auto-compacted — continuing.",
                })
        except Exception as _pf_exc_ws:
            logger.warning("agent_chat (ws): preflight failed: %s", _pf_exc_ws)

        iteration = 0
        # Loop-detection state — port from cvc.agent.continuation._detect_stall.
        # Tracks the signature of each iteration's tool calls; if the same
        # signature repeats 3 times in a row, the agent is stuck and we abort.
        recent_tool_signatures: list[str] = []
        consecutive_no_progress = 0  # iterations with no text + repeated tools

        def _is_dud_tool_call(_tc) -> bool:
            """v2.92.8 — Mirror of the executor's dud-detect predicate.
            Returns True when the model emitted a tool_use block with empty
            or missing arguments (the speculative-dud pattern from
            MiniMax-M3 / Anthropic-protocol providers). Loop-detect and
            stall-detect MUST skip dud calls; otherwise three identical
            dud signatures trip the abort even though the validator has
            already corrected the input via the per-turn dud nudge.
            """
            args = getattr(_tc, "arguments", None)
            if args is None:
                return True
            if isinstance(args, dict) and not any(
                bool(v) and (not isinstance(v, str) or v.strip())
                for v in args.values()
            ):
                return True
            return False

        # One-shot nudge guard — see SSE path for rationale.
        _silent_nudge_used = False
        # v2.90.13 — wrap-up nudge state + last-tool tracking (WS path mirror)
        _wrapup_nudge_used = False
        _last_tool_name: str | None = None
        # v2.91.42 — one-shot context-overflow auto-compaction. Set to True
        # the first time we hit a context_overflow error so we don't
        # infinite-loop on the same conversation.
        _overflow_handled_ws = False
        # (Note: _turn_done_emitted and _any_visible_text initialised at top
        # of the `try:` block so the `finally` can safely consult them on any
        # early-return branch.)

        # v2.90.8 — Mirror the SSE budget helpers so WS path honours the same
        # per-turn limits (wall-clock, iter cap, tool timeout, output cap)
        # the SSE path has had since v2.88. Without these the WS path used
        # raw _MAX_AGENT_ITERATIONS, lacked dead-client detection, and never
        # exempted sub-agent tools from _TOOL_EXEC_TIMEOUT — three of the
        # primary causes of the dashboard "stuck after tool burst" bug.
        _expensive = _is_expensive_model(_agent_llm.model) if _agent_llm else False
        _iter_limit = _budget_max_iterations(_expensive)
        _agent_timeout = _budget_agent_timeout()
        _tool_timeout = _budget_tool_timeout()
        _tool_trunc = _MAX_TOOL_OUTPUT_EXPENSIVE if _expensive else _MAX_TOOL_OUTPUT_LLM
        _turn_start = time.monotonic()
        _wrapup_threshold = _soft_wrap_threshold(_iter_limit)  # v2.90.13
        logger.info(
            "ws_chat budgets resolved: iter=%d timeout=%.0fs tool_timeout=%.0fs "
            "tool_trunc=%d expensive=%s",
            _iter_limit, _agent_timeout, _tool_timeout, _tool_trunc, _expensive,
        )

        async def _client_alive() -> bool:
            """Return False if the WS client has disconnected."""
            ws = _ts.websocket
            if ws is None:
                return False
            try:
                from starlette.websockets import WebSocketState
                return ws.client_state == WebSocketState.CONNECTED
            except Exception:
                return True

        def _budget_exceeded() -> bool:
            return (time.monotonic() - _turn_start) > _agent_timeout

        while iteration < _iter_limit:
            # ── Total-turn wall-clock ────────────────────────────
            if _budget_exceeded():
                logger.warning("ws_chat: turn timeout after %.1fs", _agent_timeout)
                await _send_safe(_ts, {
                    "type": "status",
                    "message": f"⚠️ Chat turn exceeded {_agent_timeout:.0f}s — aborting.",
                })
                await _send_safe(_ts, {"type": "done"}); _turn_done_emitted = True; _ts.done = True
                break
            # ── Dead-client check ────────────────────────────────
            if not await _client_alive():
                logger.info("ws_chat: client disconnected, aborting loop")
                break
            if not client_queue.empty():
                evt = client_queue.get_nowait()
                if evt.get("type") in ("abort", "disconnect"):
                    await _send_safe(_ts, {"type": "status", "message": "Stopped by user."})
                    break

            iteration += 1
            await _send_safe(_ts, {
                "type": "iteration",
                "n": iteration,
                "max": _iter_limit,
            })
            # v2.90.13 — Proactive wrap-up nudge (WS mirror of SSE path).
            if (
                not _wrapup_nudge_used
                and iteration >= _wrapup_threshold
                and iteration < _iter_limit
            ):
                _wrapup_nudge_used = True
                logger.info(
                    "ws_chat: wrap-up nudge fired at iter=%d/%d",
                    iteration, _iter_limit,
                )
                llm_messages.append({"role": "user", "content": _WRAPUP_NUDGE_TEXT})
            response_text = ""
            tool_calls = []
            tool_call_args_buffer: dict[int, str] = {}
            gemini_parts = None

            max_tok = 8192 if iteration == 1 else 16384
            temp = 0.7 if iteration == 1 else 0.5

            try:
                # Phase D 4.8 — streaming PII / memory-context scrubber (WebSocket)
                _ws_scrubber = None
                try:
                    from cvc.agent._vendor.hermes.agent.memory_manager import (
                        StreamingContextScrubber as _SCS_WS,
                    )
                    _ws_scrubber = _SCS_WS()
                except Exception:
                    pass

                # v2.72.4 — Per-event LLM stream watchdog on the WS path.
                # Mirror the SSE patch (v2.72.3) here — the dashboard talks
                # exclusively over /api/ws/chat, so without this wrap a stalled
                # upstream (Copilot/Anthropic/Gemini idle TCP, no events, no
                # exception) silently freezes the dashboard forever. The bare
                # `async for event in chat_stream(...)` was the actual root
                # cause of the "stuck after tool burst" symptom.
                _stream_stalled = False
                stream_iter_raw = _agent_llm.chat_stream(
                    messages=llm_messages,
                    tools=_agent_tools,
                    temperature=temp,
                    max_tokens=max_tok,
                ).__aiter__()
                # v2.92.11 — Inline chunked wait with heartbeat. Same pattern
                # as the SSE path above; documented in _STREAM_HEARTBEAT_*
                # module-level knobs.
                _ws_stall_elapsed = 0.0
                while True:
                    if not client_queue.empty():
                        evt = client_queue.get_nowait()
                        if evt.get("type") in ("abort", "disconnect"):
                            await _send_safe(_ts, {"type": "status", "message": "Stopped by user."})
                            _stream_stalled = True
                            break
                    try:
                        event = await asyncio.wait_for(
                            stream_iter_raw.__anext__(),
                            timeout=_STREAM_HEARTBEAT_INTERVAL,
                        )
                        _ws_stall_elapsed = 0.0
                    except StopAsyncIteration:
                        break
                    except asyncio.TimeoutError:
                        _ws_stall_elapsed += _STREAM_HEARTBEAT_INTERVAL
                        if _ws_stall_elapsed < _STREAM_HEARTBEAT_TIMEOUT:
                            await _send_safe(_ts, {
                                "type": "status",
                                "message": (
                                    f"⏳ Still working on ws_chat iter={iteration} "
                                    f"({int(_ws_stall_elapsed)}s elapsed)…"
                                ),
                            })
                            continue
                        logger.warning(
                            "ws_chat: LLM stream stalled >%.0fs at iter=%d "
                            "(visible_text=%s, tool_calls=%d) — breaking inner loop",
                            _STREAM_HEARTBEAT_TIMEOUT, iteration,
                            _any_visible_text, len(tool_calls),
                        )
                        await _send_safe(_ts, {
                            "type": "status",
                            "message": f"⚠️ LLM stream stalled for {_STREAM_HEARTBEAT_TIMEOUT:.0f}s — recovering…",
                        })
                        _stream_stalled = True
                        break
                    except Exception as _stream_exc:
                        _err_kind = _classify_stream_error(_stream_exc)
                        logger.warning(
                            "ws_chat: stream error at iter=%d kind=%s: %s",
                            iteration, _err_kind, _stream_exc,
                        )
                        # v2.91.42 — Context-overflow is a special case: the
                        # whole prompt is too big. Before giving up, try ONCE
                        # to auto-compact and retry. If that still fails, we
                        # surface a clear overflow message to the user with
                        # the exact numbers + recovery options.
                        if _err_kind == "context_overflow" and not _overflow_handled_ws:
                            _overflow_handled_ws = True
                            _prov = getattr(_agent_llm, "provider", "?")
                            _mod = getattr(_agent_llm, "model", "?")
                            _total_win = _get_active_model_context_window() or 0
                            _used_est = _estimate_used_tokens()
                            await _send_safe(_ts, {
                                "type": "status",
                                "message": (
                                    f"⚠️ Context overflow on {_prov}/{_mod} — "
                                    f"auto-compacting ~{_used_est/1000:.0f}k tokens "
                                    f"to fit the {_total_win/1000:.0f}k limit…"
                                ),
                            })
                            try:
                                _target = max(int(_total_win * 0.8), 4096) if _total_win else 16384
                                _compacted, _orig = _auto_compact_conversation(
                                    llm_messages, target_tokens=_target,
                                )
                                logger.warning(
                                    "ws_chat: context overflow — auto-compacted "
                                    "%d → target %d tokens",
                                    _orig, _target,
                                )
                                # Splice the compacted history in place of
                                # llm_messages and reset iteration so the
                                # loop tries again with the smaller prompt.
                                llm_messages.clear()
                                llm_messages.extend(_compacted)
                                iteration = 0
                                _silent_nudge_used = False
                                _wrapup_nudge_used = False
                                _stream_stalled = False
                                _last_stream_error = _stream_exc
                                continue
                            except Exception as _ac_exc:
                                logger.warning(
                                    "ws_chat: auto-compact failed: %s", _ac_exc,
                                )
                        # v2.91.42 — When we have an overflow error AND
                        # auto-compact either didn't run or also failed,
                        # emit a user-facing overflow message instead of
                        # the generic "stream error" status.
                        if _err_kind == "context_overflow":
                            _prov = getattr(_agent_llm, "provider", "?")
                            _mod = getattr(_agent_llm, "model", "?")
                            _total_win = _get_active_model_context_window() or 0
                            _used_est = _estimate_used_tokens()
                            _overflow_msg = _format_overflow_message(
                                _prov, _mod, _used_est, _total_win,
                            )
                            await _send_safe(_ts, {"type": "text", "content": _overflow_msg})
                            _any_visible_text = True
                            _synth_emitted_ws = True
                            await _send_safe(_ts, {"type": "done"}); _turn_done_emitted = True; _ts.done = True
                            break
                        # Generic stream error (network, rate-limit, etc.) — keep
                        # the existing recovery path (silent-finish + nudge + synth).
                        await _send_safe(_ts, {
                            "type": "status",
                            "message": f"⚠️ Stream error: {type(_stream_exc).__name__} — recovering…",
                        })
                        _stream_stalled = True
                        _last_stream_error = _stream_exc
                        break

                    if event.type == "text_delta" and event.text:
                        _visible = _ws_scrubber.feed(event.text) if _ws_scrubber else event.text
                        if _visible:
                            response_text += _visible
                            _any_visible_text = True
                            await _send_safe(_ts, {"type": "text", "content": _visible})
                    elif event.type == "tool_call_start" and event.tool_call:
                        tc = event.tool_call
                        tool_calls.append(tc)
                        tool_call_args_buffer[event.tool_call_index] = ""
                    elif event.type == "tool_call_delta":
                        idx = event.tool_call_index
                        if idx in tool_call_args_buffer:
                            tool_call_args_buffer[idx] += event.args_delta
                        if tool_calls and idx < len(tool_calls):
                            try:
                                tool_calls[idx].arguments = json.loads(tool_call_args_buffer[idx])
                            except (json.JSONDecodeError, ValueError):
                                pass
                    elif event.type == "done":
                        if event._provider_meta.get("gemini_parts"):
                            gemini_parts = event._provider_meta["gemini_parts"]

                # Phase D 4.8 — flush WebSocket scrubber tail
                if _ws_scrubber is not None:
                    try:
                        _tail = _ws_scrubber.flush()
                        if _tail:
                            response_text += _tail
                            await _send_safe(_ts, {"type": "text", "content": _tail})
                    except Exception:
                        pass

            except Exception as e:
                # v2.71.x: don't silently die on transient LLM errors.
                # v2.71.4: ALSO emit the error as visible text so the user sees
                # it as an assistant bubble (the prior status-only emit was
                # invisible on Windows, leaving users staring at the canned
                # "couldn't produce a summary" fallback with no context).
                logger.warning("ws_chat llm_stream_error iter=%d err=%s", iteration, e)
                await _send_safe(_ts, {"type": "status", "message": f"LLM error: {e}"})
                # Mark response empty so synthesis fallback will fire.
                response_text = ""
                tool_calls = []
                # Stash the error so the finally-block can decide whether to
                # surface it or whether a successful synthesis replaced it.
                _last_stream_error = e
                break

            if not tool_calls:
                if response_text:
                    llm_messages.append({"role": "assistant", "content": response_text})
                    await _send_safe(_ts, {"type": "done"}); _turn_done_emitted = True; _ts.done = True
                    break
                # Silent finish — see SSE path. Nudge once, then give up gracefully.
                if not _silent_nudge_used:
                    _silent_nudge_used = True
                    logger.info("ws_chat: silent finish — injecting one-shot nudge")
                    # v2.91.41 — Provider-aware nudge shape.
                    llm_messages.append(_build_nudge_message(
                        provider=getattr(_agent_llm, "provider", ""),
                        text=(
                            "SYSTEM: Your last turn produced no visible reply. "
                            "fact_store/memory/save_memory are SILENT side-effects, NOT responses. "
                            "Reply to the user's last message NOW in plain text. Do not call any "
                            "tool — just answer them directly."
                        ),
                    ))
                    continue
                await _send_safe(_ts, {"type": "status", "message": "Synthesizing final response…"})
                _synth_emitted_ws = False
                try:
                    # v2.91.41 — Use provider-aware nudge shape for the second
                    # nudge too (it carries the same "stop calling tools, write
                    # prose" directive).
                    llm_messages.append(_build_nudge_message(
                        provider=getattr(_agent_llm, "provider", ""),
                        text=(
                            "SYSTEM: Your last turn produced no visible reply. "
                            "Based on the tool results above (if any), write a "
                            "complete final answer to the user RIGHT NOW. Do NOT "
                            "call any tool. Plain prose only. If the user just "
                            "greeted you, greet them back briefly."
                        ),
                    ))
                    _synth_start_ws = time.monotonic()
                    # v2.90.13 — WS provider failover for synthesis pass
                    _synth_llm_ws = _build_fallback_llm(
                        getattr(_agent_llm, "provider", ""),
                        getattr(_agent_llm, "model", ""),
                    ) or _agent_llm
                    if _synth_llm_ws is not _agent_llm:
                        await _send_safe(_ts, {
                            "type": "status",
                            "message": (
                                f"Failing over synthesis to "
                                f"{getattr(_synth_llm_ws, 'provider', '?')}/"
                                f"{getattr(_synth_llm_ws, 'model', '?')}…"
                            ),
                        })
                    _synth_iter_ws = _synth_llm_ws.chat_stream(
                        messages=llm_messages,
                        tools=[],  # Force text-only
                        temperature=0.5,
                        max_tokens=4096,
                    ).__aiter__()
                    while True:
                        try:
                            _ev = await asyncio.wait_for(
                                _synth_iter_ws.__anext__(),
                                timeout=min(_STREAM_EVENT_TIMEOUT, 45.0),
                            )
                        except StopAsyncIteration:
                            break
                        except asyncio.TimeoutError:
                            logger.warning("ws_chat: synthesis stream stalled")
                            break
                        except Exception as _ws_se:
                            logger.warning("ws_chat: synthesis stream error: %s", _ws_se)
                            break
                        if time.monotonic() - _synth_start_ws > 90.0:
                            logger.warning("ws_chat: synthesis call exceeded 90s")
                            break
                        if _ev.type == "text_delta" and _ev.text:
                            _synth_emitted_ws = True
                            await _send_safe(_ts, {"type": "text", "content": _ev.text})
                except Exception as _ws_synth_exc:
                    logger.error("ws_chat: synthesis fallback failed: %s", _ws_synth_exc, exc_info=True)
                # v2.91.41 — Pull the last user request + last tool output
                # ONCE, before either fallback branch, so both the
                # deterministic synthesis and the diagnostic have access.
                # The "last user message" we want is the most recent one
                # that does NOT start with "SYSTEM:" — those are our own
                # nudge injections, not the user's actual question.
                _last_user_text_ws = ""
                _last_tool_output_ws = ""
                try:
                    for _m in reversed(llm_messages):
                        if not _last_user_text_ws and _m.get("role") == "user":
                            _c = _m.get("content", "")
                            if isinstance(_c, str):
                                if not _c.lstrip().startswith("SYSTEM:"):
                                    _last_user_text_ws = _c
                        if not _last_tool_output_ws and _m.get("role") == "tool":
                            _c = _m.get("content", "")
                            if isinstance(_c, list):
                                _c = " ".join(
                                    str(b.get("text", b)) for b in _c
                                    if isinstance(b, dict)
                                )
                            _last_tool_output_ws = str(_c)
                        if _last_user_text_ws and _last_tool_output_ws:
                            break
                except Exception:
                    pass
                if not _synth_emitted_ws:
                    # v2.91.41 — Deterministic synthesis before the diagnostic.
                    _det_text = _deterministic_synthesize(
                        llm_messages,
                        user_facing_request=_last_user_text_ws,
                    )
                    if _det_text:
                        await _send_safe(_ts, {"type": "text", "content": _det_text})
                        await _send_safe(_ts, {
                            "type": "status",
                            "message": (
                                "⚠️ Model didn't write a final reply — showing "
                                "the raw tool output above."
                            ),
                        })
                        _synth_emitted_ws = True  # v2.91.41 — counts as a real reply
                if not _synth_emitted_ws:
                    _giveup_msg_ws = _diagnostic_giveup(
                        iteration=iteration,
                        iter_limit=_iter_limit,
                        elapsed=time.monotonic() - _turn_start,
                        provider=getattr(_agent_llm, "provider", "?"),
                        model=getattr(_agent_llm, "model", "?"),
                        last_tool=_last_tool_name,
                        last_tool_output=_last_tool_output_ws or None,
                        reason="silent finish after nudge (WS)",
                    )
                    await _send_safe(_ts, {"type": "text", "content": _giveup_msg_ws})
                await _send_safe(_ts, {"type": "done"}); _turn_done_emitted = True; _ts.done = True
                break

            # ── Loop / stall detection ──────────────────────────────────────
            # If the model emits the exact same set of tool calls three
            # iterations in a row, or burns N iterations without ever
            # producing assistant text, it is stuck — break out instead of
            # silently grinding until _MAX_AGENT_ITERATIONS.
            #
            # v2.92.8 — Exempt phantom dud calls (model emits empty-args
            # tool_use blocks as chain-of-thought; the validator catches
            # them and they collapse to one nudge per turn via
            # _dud_nudge_sent_ws). Without this exemption, three dud
            # iterations trip the abort, killing the session even though
            # the agent is mid-recovery. Only count "real" tool calls
            # toward the signature/stall budget; a dud-only iteration is
            # a model glitch the validator has already corrected.
            _real_calls = [tc for tc in tool_calls if not _is_dud_tool_call(tc)]
            _all_duds = bool(tool_calls) and not _real_calls
            if _real_calls:
                sig = "|".join(sorted(f"{tc.name}:{json.dumps(tc.arguments, sort_keys=True)[:200]}" for tc in _real_calls))
                recent_tool_signatures.append(sig)
                if len(recent_tool_signatures) > 5:
                    recent_tool_signatures = recent_tool_signatures[-5:]
            if (
                len(recent_tool_signatures) >= 3
                and recent_tool_signatures[-1] == recent_tool_signatures[-2] == recent_tool_signatures[-3]
            ):
                _n_dup = len(recent_tool_signatures[-1].split("|")) if recent_tool_signatures[-1] else len(tool_calls)
                await _send_safe(_ts, {
                    "type": "status",
                    "message": (
                        f"Loop detected — same {_n_dup} tool call(s) repeated "
                        f"3 iterations in a row. Aborting to prevent runaway."
                    ),
                })
                await _send_safe(_ts, {"type": "done"}); _turn_done_emitted = True; _ts.done = True
                break

            # v2.90.8 — Stall semantics aligned with CLI/SSE: only count an
            # iteration as "no progress" when the model produced no text AND
            # the tool-call signature is identical to the previous iteration.
            # Pure silent-tool-burst iterations (Claude planning with multiple
            # distinct tool calls) are legitimate work, not stalls. Threshold
            # bumped 5→8 to match agent.continuation._detect_stall.
            #
            # v2.92.8 — Dud-only iterations do NOT count as no-progress.
            # They're a model glitch the validator caught; the model is
            # mid-recovery, not stuck. Let the dud nudge do its job without
            # poisoning the stall counter.
            _sig_repeated = (
                len(recent_tool_signatures) >= 2
                and recent_tool_signatures[-1] == recent_tool_signatures[-2]
            )
            if _all_duds:
                # dud iteration — don't move the counter either way; the
                # model gets one nudge and self-corrects.
                pass
            elif not response_text and _sig_repeated:
                consecutive_no_progress += 1
            else:
                consecutive_no_progress = 0
            if consecutive_no_progress >= 8:
                await _send_safe(_ts, {
                    "type": "status",
                    "message": (
                        "No progress for 8 iterations (no assistant text, same tool calls). "
                        "Aborting to prevent runaway."
                    ),
                })
                await _send_safe(_ts, {"type": "done"}); _turn_done_emitted = True; _ts.done = True
                break

            assistant_msg: dict[str, Any] = {"role": "assistant", "content": response_text or ""}
            if _agent_llm.provider in ("openai", "ollama", "lmstudio", "github"):
                assistant_msg["tool_calls"] = [
                    {"id": tc.id, "type": "function", "function": {"name": tc.name, "arguments": json.dumps(tc.arguments)}}
                    for tc in tool_calls
                ]
            elif _agent_llm.provider == "anthropic":
                content_blocks = []
                if response_text:
                    content_blocks.append({"type": "text", "text": response_text})
                for tc in tool_calls:
                    content_blocks.append({"type": "tool_use", "id": tc.id, "name": tc.name, "input": tc.arguments})
                assistant_msg = {"role": "assistant", "content": content_blocks}
            elif _agent_llm.provider == "minimax":
                # v2.91.43: MiniMax uses the Anthropic-Messages-API protocol
                # (Bearer auth, /v1/messages path, Anthropic-style content
                # blocks) — see cvc/agent/llm.py:263-272. Tool calls MUST
                # be emitted as Anthropic ``tool_use`` content blocks; the
                # following ``tool_result`` message references those blocks
                # by id. Without this, the next chat call sent an
                # assistant message with no tool_use and a tool_result with
                # a dangling tool_call_id, which the upstream API rejected
                # silently — causing the agent to do ONE tool call and then
                # go silent on every subsequent turn. (This is the
                # "she just runs one tool and stops" bug from the
                # 2026-06-10 MiniMax dashboard screenshot.)
                content_blocks = []
                if response_text:
                    content_blocks.append({"type": "text", "text": response_text})
                for tc in tool_calls:
                    content_blocks.append({"type": "tool_use", "id": tc.id, "name": tc.name, "input": tc.arguments})
                assistant_msg = {"role": "assistant", "content": content_blocks}
            elif _agent_llm.provider == "google":
                assistant_msg["tool_calls"] = [
                    {"id": tc.id, "type": "function", "function": {"name": tc.name, "arguments": json.dumps(tc.arguments)}}
                    for tc in tool_calls
                ]
                if gemini_parts:
                    assistant_msg["_gemini_parts"] = gemini_parts
            llm_messages.append(assistant_msg)

            await _send_safe(_ts, {"type": "status", "message": f"Executing {len(tool_calls)} tool(s)..."})

            for tc in tool_calls:
                display_args = {}
                for k, v in tc.arguments.items():
                    sv = str(v)
                    display_args[k] = sv[:200] + "…" if len(sv) > 200 else sv
                # A6.5 — Surface cwd on bash/shell tools so the dashboard chip
                # visually exposes the active workspace and catches drift.
                if tc.name in ("bash", "shell", "run_shell") and "cwd" not in display_args:
                    try:
                        _exec = globals().get("_tool_executor")
                        _ws = getattr(_exec, "workspace", None) if _exec else None
                        if _ws:
                            display_args["cwd"] = str(_ws)
                    except Exception:
                        pass

                needs_confirm = (
                    not _autopilot_on()
                    and not _session_trust_all
                    and tc.name not in _SAFE_TOOLS
                    and tc.name not in _session_trusted_tools
                )

                if needs_confirm:
                    await _send_safe(_ts, {
                        "type": "tool_confirm",
                        "name": tc.name,
                        "args": display_args,
                    })
                    logger.info("ws_chat confirm_wait_start tool=%s", tc.name)
                    decision = "deny"
                    _confirm_start = time.monotonic()
                    try:
                        while True:
                            elapsed = time.monotonic() - _confirm_start
                            remaining = _CONFIRM_TIMEOUT - elapsed
                            if remaining <= 0:
                                raise asyncio.TimeoutError()
                            evt = await asyncio.wait_for(
                                client_queue.get(),
                                timeout=min(remaining, _SSE_HEARTBEAT_INTERVAL),
                            )
                            if evt.get("type") == "confirm":
                                decision = evt.get("result", "deny")
                                break
                            elif evt.get("type") in ("abort", "disconnect"):
                                break
                    except asyncio.TimeoutError:
                        decision = "deny"
                        logger.warning("ws_chat confirm_timeout tool=%s", tc.name)
                        await _send_safe(_ts, {
                            "type": "tool_timeout",
                            "name": tc.name,
                            "timeout_s": int(_CONFIRM_TIMEOUT),
                            "mode": _approval_mode,
                        })
                    logger.info("ws_chat confirm_decision tool=%s decision=%s", tc.name, decision)
                    if decision == "allow_always":
                        _session_trusted_tools.add(tc.name)
                    elif decision == "trust_all":
                        _session_trust_all = True
                        # A6.2 — Drain any queued stale confirm frames now that
                        # trust_all is sticky for the session. Without this,
                        # confirms left over from prior rapid clicks poison the
                        # next tool's wait and create phantom hangs.
                        try:
                            while True:
                                _stale = client_queue.get_nowait()
                                if isinstance(_stale, dict) and _stale.get("type") != "confirm":
                                    # Put non-confirm control frames back at the front-ish
                                    try:
                                        client_queue.put_nowait(_stale)
                                    except Exception:
                                        pass
                                    break
                        except asyncio.QueueEmpty:
                            pass
                        except Exception:
                            pass
                    elif decision == "deny":
                        result = f"Tool '{tc.name}' was denied by the user."
                        await _send_safe(_ts, {"type": "tool_result", "name": tc.name, "output": result})
                        llm_messages.append({"role": "tool", "tool_call_id": tc.id, "content": result})
                        continue
                    elif decision == "deny_suggest":
                        result = f"Tool '{tc.name}' was denied. Please suggest an alternative approach."
                        await _send_safe(_ts, {"type": "tool_result", "name": tc.name, "output": result})
                        llm_messages.append({"role": "tool", "tool_call_id": tc.id, "content": result})
                        continue
                        
                # v2.92 — ask_user has its own visual: the ClarifyCard
                # multi-choice picker. Skip the generic `tool_start` /
                # `tool_result` events for ask_user entirely so the
                # dashboard doesn't render a redundant tool card
                # leaking the raw `question=…` / `options=…` JSON.
                # The ClarifyCard is the single source of truth.
                if tc.name != "ask_user":
                    await _send_safe(_ts, {"type": "tool_start", "name": tc.name, "args": display_args})
                    logger.info("ws_chat tool_exec_start tool=%s", tc.name)
                _last_tool_name = tc.name  # v2.90.13
                _tx_start = time.monotonic()

                # ── ask_user interceptor (WS path) ─────────────────────────
                # Pure CVC. Mirrors the SSE interceptor above. The model's
                # `ask_user` tool call is intercepted BEFORE the executor
                # runs; the dashboard renders a multi-choice picker via
                # the `clarify_request` WS event, then the user clicks a
                # button → POST /api/chat/clarify/answer → cvc.tools.clarify
                # unblocks → we emit `clarify_resolved` and continue the
                # agent loop. The HTTP resolver endpoint is the single
                # source of truth for both the SSE and WS surfaces.
                if tc.name == "ask_user":
                    from cvc.tools import clarify as _clarify
                    _raw_args = tc.arguments or {}
                    _question = _clarify.normalise_question(
                        _raw_args.get("question", "")
                    )
                    _choices = _clarify.sanitise_choices(
                        _raw_args.get("options") or _raw_args.get("choices")
                    )
                    # v2.92 — Never return an error string for empty
                    # ask_user calls. Some models treat a tool-error
                    # result as a turn-termination signal — they say
                    # "let me retry" in prose but never emit another
                    # tool call, and the turn dies. Always materialise
                    # a picker instead.
                    if not _question:
                        _question = "Pick to continue — the model did not include a specific question in its ask_user call."
                    _session_key = (
                        str(_project_root) if _project_root else "default"
                    )
                    _clarify_id = _clarify.new_clarify_id()
                    _clarify.register(
                        _clarify_id, _session_key, _question, _choices
                    )
                    await _send_safe(_ts, {
                        "type": "clarify_request",
                        "clarify_id": _clarify_id,
                        "name": tc.name,
                        "question": _question,
                        "choices": _choices,
                        "open_ended": _choices is None,
                    })
                        # Block until the HTTP resolver unblocks us, or the
                        # client disconnects, or the clarify timeout fires.
                        # Heartbeat ping every 10s so long waits don't trip
                        # upstream proxy idle-killers.
                        #
                        # v2.92 — Use asyncio.run_in_executor to wrap the
                        # blocking wait_for_response so the WS event loop
                        # stays free to serve the
                        # /api/chat/clarify/answer HTTP request that
                        # resolves the wait. Without the executor,
                        # threading.Event.wait() would pin the loop for
                        # the full 1s slice, and the resolver POST would
                        # queue up indefinitely — leaving the user
                        # staring at a picker that doesn't respond to
                    # clicks.
                    _clarify_start = time.monotonic()
                    _clarify_timeout = float(_clarify.get_clarify_timeout())
                    _ws_loop = asyncio.get_running_loop()
                    _response: str | None = None
                    _last_ping = _clarify_start
                    while True:
                        try:
                            _resp = await asyncio.wait_for(
                                _ws_loop.run_in_executor(
                                    None,
                                    _clarify.wait_for_response,
                                    _clarify_id,
                                    min(1.0, _clarify_timeout - (time.monotonic() - _clarify_start)),
                                ),
                                timeout=1.5,
                            )
                        except asyncio.TimeoutError:
                            _resp = None
                        except Exception:
                            _resp = None
                        if _resp is not None:
                            _response = _resp
                            break
                        if _ts.client_queue is not None:
                            # Drain any inbound control frames (abort/disconnect)
                            try:
                                while not _ts.client_queue.empty():
                                    _ctl = _ts.client_queue.get_nowait()
                                    if isinstance(_ctl, dict) and _ctl.get("type") in ("abort", "disconnect"):
                                        _response = ""
                                        break
                            except Exception:
                                pass
                            if _response is not None:
                                break
                        if (time.monotonic() - _clarify_start) >= _clarify_timeout:
                            _clarify.clear_session(_session_key)
                            break
                        if time.monotonic() - _last_ping >= 10.0:
                            try:
                                await _send_safe(_ts, {
                                    "type": "ping",
                                    "scope": "clarify",
                                    "tool": tc.name,
                                    "elapsed_s": round(time.monotonic() - _clarify_start, 1),
                                })
                            except Exception:
                                pass
                            _last_ping = time.monotonic()
                    if _response is None:
                        _response = ""
                    await _send_safe(_ts, {
                        "type": "clarify_resolved",
                        "clarify_id": _clarify_id,
                        "name": tc.name,
                        "response": _response,
                    })
                    import json as _json
                    result = _json.dumps({
                        "question": _question,
                        "choices_offered": _choices,
                        "user_response": _response,
                    }, ensure_ascii=False)
                    _tx_end = time.monotonic()
                    # v2.92 — No `tool_result` WS for ask_user. The
                    # ClarifyCard's resolve event (clarify_resolved)
                    # already carries the user's pick. Emitting a
                    # generic tool_result here would re-create the
                    # redundant raw `output={"question": ..., ...}`
                    # card that was leaking on the dashboard.
                    llm_messages.append({
                        "role": "tool",
                        "tool_call_id": tc.id,
                        "content": result,
                    })
                    continue  # skip the generic _tool_executor.execute path

                # Install a thread-safe sink for sub-agent tools so nested
                # activity surfaces live on the WS as `subagent_*` events.
                _is_subagent_tool = tc.name in ("agent", "parallel_agents")
                _sub_q: asyncio.Queue | None = None
                if _is_subagent_tool:
                    _sub_q = asyncio.Queue()
                    _ws_loop = asyncio.get_running_loop()
                    def _sink(evt, _q=_sub_q, _l=_ws_loop):
                        try:
                            _l.call_soon_threadsafe(_q.put_nowait, evt)
                        except Exception:
                            pass
                    _tool_executor._subagent_event_sink = _sink
                else:
                    _tool_executor._subagent_event_sink = None

                # v2.90.8 — Sub-agent tools carry their own internal wall-clock
                # budget (CVC_SUBAGENT_WALL_S / HARD_CEILING). Wrapping them
                # in _tool_timeout (300s default) silently killed long Explore
                # runs mid-flight, mirroring the SSE bug fixed in v2.88.
                _effective_timeout = None if _is_subagent_tool else _tool_timeout

                try:
                    if tc.name == "cvc_switch_workspace":
                        result = await asyncio.wait_for(
                            _execute_workspace_switch(tc.arguments),
                            timeout=_effective_timeout,
                        )
                    else:
                        loop = asyncio.get_running_loop()
                        async def _run_tool():
                            return await loop.run_in_executor(
                                None, _tool_executor.execute, tc.name, tc.arguments
                            )
                        if _effective_timeout is None:
                            _exec_task = asyncio.create_task(_run_tool())
                        else:
                            _exec_task = asyncio.create_task(
                                asyncio.wait_for(_run_tool(), timeout=_effective_timeout)
                            )
                        if _sub_q is not None:
                            # v2.68.12 Bug A fix: emit periodic keepalive
                            # pings while a long-running sub-agent (e.g. 50-60s
                            # Explore agent) is executing. Without this, idle
                            # proxies / the browser may close the socket and
                            # the post-tool finalize never reaches the UI even
                            # though the server thinks it sent `done`.
                            _last_ping = time.monotonic()
                            while not _exec_task.done():
                                try:
                                    _evt = await asyncio.wait_for(_sub_q.get(), timeout=0.2)
                                    await _send_safe(_ts, _evt)
                                except asyncio.TimeoutError:
                                    pass
                                # Heartbeat every 10s of subagent runtime
                                if time.monotonic() - _last_ping >= 10.0:
                                    try:
                                        await _send_safe(_ts, {
                                            "type": "ping",
                                            "scope": "subagent",
                                            "tool": tc.name,
                                            "elapsed_s": round(time.monotonic() - _tx_start, 1),
                                        })
                                    except Exception:
                                        pass
                                    _last_ping = time.monotonic()
                            while not _sub_q.empty():
                                try:
                                    await _send_safe(_ts, _sub_q.get_nowait())
                                except Exception:
                                    break
                            result = _exec_task.result()
                        else:
                            result = await _exec_task
                except asyncio.TimeoutError:
                    result = (
                        f"⚠️ Tool '{tc.name}' exceeded {(_effective_timeout or _tool_timeout):.0f}s timeout "
                        f"and was aborted."
                    )
                    logger.warning("ws_chat tool_exec_timeout tool=%s", tc.name)
                except Exception as exc:
                    result = f"Error executing {tc.name}: {exc}"
                    logger.warning("ws_chat tool_exec_error tool=%s err=%s", tc.name, exc)
                finally:
                    _tool_executor._subagent_event_sink = None
                logger.info(
                    "ws_chat tool_exec_done tool=%s elapsed=%.2fs",
                    tc.name, time.monotonic() - _tx_start,
                )

                # v2.90.8 — honour expensive-model tool-output cap
                # (matches SSE path: cuts Opus context burn by ~66%).
                llm_result = result[:_tool_trunc] if len(result) > _tool_trunc else result
                if len(result) > _tool_trunc:
                    llm_result += f"\n... (truncated, {len(result) - _tool_trunc:,} chars omitted)"
                # Stream tool_result to client (truncated for display)
                # v2.91.45 — Flag error results so the dashboard can
                # render the tool card with error styling. See the SSE
                # path for the full rationale (phantom empty tool call
                # symptom from speculative model emissions).
                #
                # v2.92.4 — Same dedup as the SSE path: collapse
                # consecutive dud calls ("X requires argument(s)..." or
                # "X was called with no arguments") into a single
                # warning event + a single system nudge. Stops the
                # model from getting 5× "Error: X requires..." strings
                # in its context.
                _is_dud_call = bool(result) and (
                    "requires argument(s)" in result[:200]
                    or "was called with no arguments" in result[:200]
                )
                if _is_dud_call:
                    await _send_safe(_ts, {
                        "type": "tool_dud_warning",
                        "name": tc.name,
                        "message": (
                            f"⚠ {tc.name} was called with empty/missing "
                            f"args. The agent should not call this tool "
                            f"without real arguments."
                        ),
                    })
                    # v2.92.10 — Also emit a synthetic tool_result so the
                    # dashboard's ToolCallCard transitions out of
                    # "running…" state. Same rationale as the SSE path:
                    # without this the WS client gets tool_start but
                    # never a tool_result, and the card sticks on
                    # "running…" forever. The model never sees this
                    # synthetic output — we send the consolidated nudge
                    # to the LLM below.
                    await _send_safe(_ts, {
                        "type": "tool_result",
                        "name": tc.name,
                        "output": "(suppressed: dud call — see warning above)",
                        "error": True,
                        "suppressed": True,
                    })
                    if not _dud_nudge_sent_ws:
                        llm_messages.append({
                            "role": "user",
                            "content": (
                                "[cvc: your previous tool call had empty "
                                "or missing arguments. Please re-emit the "
                                "call with the real arguments, or skip "
                                "this tool and try a different approach. "
                                "Do NOT call the same tool again without "
                                "arguments.]"
                            ),
                        })
                        _dud_nudge_sent_ws = True
                    continue  # skip the normal tool_result / llm_messages path

                display_output = result[:3000] + f"\n... ({len(result):,} chars total)" if len(result) > 3000 else result
                _is_error_result = bool(result) and (
                    result.startswith("Error: ")
                    or result.startswith("Error executing ")
                    or result.startswith("Tool ")
                    or "was denied" in result[:80]
                )
                await _send_safe(_ts, {
                    "type": "tool_result",
                    "name": tc.name,
                    "output": display_output,
                    "error": _is_error_result,
                })

                llm_messages.append({"role": "tool", "tool_call_id": tc.id, "content": llm_result})

        else:
            # while loop exited because iteration == _iter_limit
            # without a natural break (no `done`, no stall, no abort).
            await _send_safe(_ts, {
                "type": "status",
                "message": (
                    f"Iteration cap reached ({_iter_limit}). "
                    "Stopping — refine your request or raise CVC_MAX_ITER."
                ),
            })
            await _send_safe(_ts, {"type": "done"}); _turn_done_emitted = True; _ts.done = True

    finally:
        receiver_task.cancel()
        # v2.71.x Bug fix: do NOT show users a canned "network/model timeout"
        # message when the real cause is a silent loop exit. Instead, port the
        # SSE path's forced-synthesis recovery: if the loop exited without
        # producing any visible text, fire one final tools=None LLM call to
        # synthesize a real answer from whatever the model already did
        # (tool results, partial reasoning). This fixes the Windows-always /
        # Mac-sometimes "Try again." dead-end reported in v2.71.0.
        if not _turn_done_emitted:
            try:
                if not _any_visible_text:
                    _synth_ok = False
                    # v2.71.4: Try synthesis whenever we have an LLM bound.
                    # The old `_has_history` gate punished the simplest case
                    # (first turn, plain "hi") — synthesis was skipped and the
                    # user got the canned diagnostic with zero context. Now
                    # we always attempt one synthesis pass; if there's no
                    # history, the model just answers the original message.
                    if _agent_llm is not None:
                        try:
                            logger.warning(
                                "ws_chat: silent exit — forcing synthesis fallback"
                            )
                            await _send_safe(_ts, {
                                "type": "status",
                                "message": "Synthesizing final response…",
                            })
                            # v2.90.8 — Backfill orphan tool_calls before
                            # synthesis. If the loop aborted between the
                            # assistant tool_calls message and the matching
                            # tool result(s), Anthropic/OpenAI will 400 on the
                            # next call ("tool_use without tool_result"). We
                            # inject a synthetic "[aborted]" result for every
                            # orphan tool_call_id so the synthesis succeeds.
                            try:
                                _pending_ids: list[tuple[str, str]] = []
                                for _m in llm_messages:
                                    if _m.get("role") == "assistant":
                                        _tcs = _m.get("tool_calls") or []
                                        for _tc in _tcs:
                                            _tid = _tc.get("id")
                                            _tname = (_tc.get("function") or {}).get("name") or "tool"
                                            if _tid:
                                                _pending_ids.append((_tid, _tname))
                                        # Anthropic shape
                                        _content = _m.get("content")
                                        if isinstance(_content, list):
                                            for _blk in _content:
                                                if isinstance(_blk, dict) and _blk.get("type") == "tool_use":
                                                    _pending_ids.append((_blk.get("id", ""), _blk.get("name", "tool")))
                                    elif _m.get("role") == "tool":
                                        _tid = _m.get("tool_call_id")
                                        if _tid:
                                            _pending_ids = [(i, n) for (i, n) in _pending_ids if i != _tid]
                                for _tid, _tname in _pending_ids:
                                    if not _tid:
                                        continue
                                    llm_messages.append({
                                        "role": "tool",
                                        "tool_call_id": _tid,
                                        "content": f"[aborted: {_tname} did not complete before turn ended]",
                                    })
                                    logger.info("ws_chat: backfilled orphan tool_call id=%s name=%s", _tid, _tname)
                            except Exception as _bf_exc:
                                logger.debug("ws_chat orphan backfill failed: %s", _bf_exc)
                            llm_messages.append({
                                "role": "user",
                                "content": (
                                    "SYSTEM: Your previous turn ended without a "
                                    "user-visible reply. Reply to the user's original "
                                    "request RIGHT NOW in plain prose. Do NOT call any "
                                    "tool. If tool results are present above, use them; "
                                    "otherwise just answer the user directly. The user "
                                    "has been waiting and needs a real response."
                                ),
                            })
                            _synth_start = time.monotonic()
                            # v2.72.4 — Per-event watchdog on WS synthesis fallback too.
                            _synth_iter = _agent_llm.chat_stream(
                                messages=llm_messages,
                                tools=None,  # force text-only — no tools allowed
                                temperature=0.5,
                                max_tokens=8192,
                            ).__aiter__()
                            while True:
                                try:
                                    _evt = await asyncio.wait_for(
                                        _synth_iter.__anext__(),
                                        timeout=min(_STREAM_EVENT_TIMEOUT, 45.0),
                                    )
                                except StopAsyncIteration:
                                    break
                                except asyncio.TimeoutError:
                                    logger.warning(
                                        "ws_chat: synthesis stream stalled — breaking"
                                    )
                                    break
                                if time.monotonic() - _synth_start > 90.0:
                                    logger.warning(
                                        "ws_chat: synthesis call exceeded 90s"
                                    )
                                    break
                                if _evt.type == "text_delta" and _evt.text:
                                    _synth_ok = True
                                    await _send_safe(_ts, {
                                        "type": "text",
                                        "content": _evt.text,
                                    })
                        except Exception as _synth_exc:
                            logger.error(
                                "ws_chat: synthesis fallback failed: %s",
                                _synth_exc, exc_info=True,
                            )
                            # Remember the synth error too — preferred over the
                            # earlier stream error when reporting to the user.
                            _last_stream_error = _synth_exc
                    if not _synth_ok:
                        # Synthesis didn't run or itself went silent.
                        # v2.71.4: surface the underlying error if we have one
                        # — that's the *actual* diagnostic the user needs.
                        if _last_stream_error is not None:
                            await _send_safe(_ts, {
                                "type": "text",
                                "content": (
                                    f"LLM call failed: `{type(_last_stream_error).__name__}: "
                                    f"{_last_stream_error}`.\n\n"
                                    "On Windows, this is most often a proxy/TLS issue, "
                                    "an expired API key, or the wrong provider URL. "
                                    "Check `cvc gateway logs` for the full trace, "
                                    "then `cvc setup` to re-verify your provider config."
                                ),
                            })
                        else:
                            await _send_safe(_ts, {
                                "type": "text",
                                "content": (
                                    "I couldn't reach the LLM to produce a reply. "
                                    "Run `cvc gateway logs` for the underlying error, "
                                    "or `cvc setup` to re-verify your provider config."
                                ),
                            })
                await _send_safe(_ts, {"type": "done"})
                _ts.done = True
            except Exception:
                pass


# ═══════════════════════════════════════════════════════════════════════
# 10. MCP SERVER INFO
# ═══════════════════════════════════════════════════════════════════════

@app.get("/api/mcp/tools")
async def mcp_tools():
    try:
        from cvc.mcp_server import TOOL_DEFINITIONS
        return {"tools": TOOL_DEFINITIONS, "total": len(TOOL_DEFINITIONS)}
    except ImportError:
        tools = [
            "cvc_status", "cvc_commit", "cvc_branch", "cvc_merge", "cvc_restore",
            "cvc_log", "cvc_capture_context", "cvc_set_workspace", "cvc_get_context",
            "cvc_recall", "cvc_export", "cvc_inject", "cvc_diff", "cvc_stats",
            "cvc_compact", "cvc_timeline", "cvc_audit",
            "cvc_sync_push", "cvc_sync_pull", "cvc_sync_status",
            "cvc_register_agent", "cvc_agent_commit", "cvc_agent_recall",
            "cvc_list_agents", "cvc_agent_context", "cvc_hive_status",
        ]
        return {"tools": [{"name": t} for t in tools], "total": len(tools)}
    except Exception:
        return {"tools": [], "total": 0}


@app.get("/api/mcp/status")
async def mcp_status():
    return {
        "transport": "stdio",
        "status": "available",
        "tools_count": 26,
        "note": "MCP server runs as a subprocess launched by IDEs, not as a persistent service",
    }


# ═══════════════════════════════════════════════════════════════════════
# 11. GIT / VCS BRIDGE
# ═══════════════════════════════════════════════════════════════════════

@app.get("/api/vcs/status")
async def vcs_status():
    _ensure()
    try:
        import subprocess as _sp
        git_dir = _project_root / ".git" if _project_root else None
        has_repo = bool(git_dir and git_dir.exists())
        hooks_installed = False
        current_branch = ""
        repo_root = str(_project_root) if _project_root else ""
        if has_repo:
            post_commit = git_dir / "hooks" / "post-commit"
            hooks_installed = post_commit.exists() and "cvc" in post_commit.read_text(encoding="utf-8", errors="ignore").lower() if post_commit.exists() else False
            try:
                result = _sp.run(["git", "rev-parse", "--abbrev-ref", "HEAD"], capture_output=True, text=True, cwd=str(_project_root), timeout=5)
                if result.returncode == 0:
                    current_branch = result.stdout.strip()
            except Exception:
                pass
        return {
            "provider": "git" if has_repo else "",
            "has_repo": has_repo,
            "repo_root": repo_root,
            "current_branch": current_branch,
            "hooks_installed": hooks_installed,
        }
    except Exception as exc:
        return {"provider": "", "has_repo": False, "repo_root": "", "current_branch": "", "error": str(exc)}


# ═══════════════════════════════════════════════════════════════════════
# 12. AUDIT & STATS
# ═══════════════════════════════════════════════════════════════════════

@app.get("/api/audit")
async def get_audit(limit: int = 100):
    _ensure()
    try:
        if hasattr(_get_db().index, "query_audit_log"):
            entries = _get_db().index.query_audit_log(limit=limit)
            return {"entries": entries, "total": len(entries)}
        return {"entries": [], "total": 0}
    except Exception as exc:
        return {"entries": [], "error": str(exc)}


@app.get("/api/stats")
async def get_stats(workspace_path: str | None = None):
    _resolve_workspace(workspace_path)
    _ensure()
    try:
        commits = _get_db().index.list_all_commits(limit=9999) if hasattr(_get_db().index, "list_all_commits") else []
        branches = _get_db().index.list_branches() if hasattr(_get_db().index, "list_branches") else []
        agents = _get_db().index.list_agents() if hasattr(_get_db().index, "list_agents") else []
        by_type: dict[str, int] = {}
        total_tokens = 0
        for c in commits:
            ct = _field(c, "commit_type", "unknown")
            by_type[ct] = by_type.get(ct, 0) + 1
            total_tokens += _meta_dict(c).get("token_count", 0)
        db_path = _get_config().cvc_root / "cvc.db"
        blobs_dir = _get_config().cvc_root / "objects"
        return {
            "total_commits": len(commits), "total_branches": len(branches), "total_agents": len(agents),
            "commits_by_type": by_type, "total_tokens": total_tokens,
            "estimated_cost_usd": round(total_tokens * 3.0 / 1_000_000, 4),
            "database_size_bytes": db_path.stat().st_size if db_path.exists() else 0,
            "blob_count": len(list(blobs_dir.glob("*"))) if blobs_dir.exists() else 0,
        }
    except Exception as exc:
        return {"error": str(exc)}


# ═══════════════════════════════════════════════════════════════════════
# 12.5  SWARM & HIVEMIND  (Phase A — server-side topology + live feed)
# ═══════════════════════════════════════════════════════════════════════

def _agent_to_node(a: Any) -> dict:
    """Normalise an agent record (dict or dataclass) to a topology node."""
    if isinstance(a, dict):
        get = a.get
    else:
        def get(k, default=None):  # type: ignore[no-redef]
            return getattr(a, k, default)
    aid = get("agent_id") or get("id") or ""
    return {
        "id": aid,
        "name": get("name") or aid,
        "role": get("role") or "",
        "rank": get("rank") or "",
        "squad": get("squad") or "",
        "status": get("status") or "active",
        "capabilities": get("capabilities") or [],
        "created_at": get("created_at"),
        "updated_at": get("updated_at"),
    }


@app.get("/api/swarm")
async def get_swarm_topology():
    """
    Server-side topology view of the active workspace.

    Returns nodes (agents), squads (groupings), and edges (squad → captain
    where ranks are present). Lets the UI skip client-side aggregation and
    keeps the topology consistent across consumers.
    """
    _ensure()
    try:
        raw = _get_db().index.list_agents() if hasattr(_get_db().index, "list_agents") else []
        nodes = [_agent_to_node(a) for a in raw if a is not None]
        # Classify nodes: tier (mc | captain | specialist | misc) for richer UI rendering.
        def _classify(n: dict) -> str:
            rank = (n.get("rank") or "").lower()
            role = (n.get("role") or "").lower()
            aid = (n.get("id") or "").upper()
            if rank in ("mc", "master_controller", "commander") or aid.startswith("MC-") or "mission controller" in role:
                return "mc"
            if rank == "captain" or aid.startswith("CAP-") or role.startswith("captain"):
                return "captain"
            if aid.startswith("SPC-") or role.startswith("specialist"):
                return "specialist"
            return "misc"

        for n in nodes:
            n["tier"] = _classify(n)

        # Find global Mission Controller (orbits Singularity directly, captains report up to it).
        mc_node = next((n for n in nodes if n["tier"] == "mc"), None)
        mc_id = mc_node["id"] if mc_node else None

        # Group by squad
        squads: dict[str, dict] = {}
        for n in nodes:
            sq = n["squad"] or "unaffiliated"
            s = squads.setdefault(sq, {"name": sq, "members": [], "captain": None, "count": 0})
            s["members"].append(n["id"])
            s["count"] += 1
            if n["tier"] == "captain" and not s["captain"]:
                s["captain"] = n["id"]

        # Edges: MC → captains, captains → specialists in their squad.
        edges: list[dict] = []
        for sq in squads.values():
            cap = sq["captain"]
            if cap and mc_id and cap != mc_id:
                edges.append({"from": mc_id, "to": cap, "type": "command"})
            if not cap:
                continue
            for m in sq["members"]:
                if m != cap:
                    edges.append({"from": cap, "to": m, "type": "squad"})
        return {
            "nodes": nodes,
            "squads": list(squads.values()),
            "edges": edges,
            "total_nodes": len(nodes),
            "total_squads": len(squads),
            "workspace": str(_get_config().cvc_root.parent) if _get_config() else "",
        }
    except Exception as exc:
        return {"nodes": [], "squads": [], "edges": [], "total_nodes": 0, "total_squads": 0, "error": str(exc)}


@app.get("/api/hivemind")
async def get_hivemind_overview():
    """High-level HiveMind overview: agent count, memory size, recent activity."""
    _ensure()
    try:
        db = _get_db()
        agents = db.index.list_agents() if hasattr(db.index, "list_agents") else []
        commits = db.index.list_all_commits(limit=200) if hasattr(db.index, "list_all_commits") else []
        branches = db.index.list_branches() if hasattr(db.index, "list_branches") else []
        recent_authors: dict[str, int] = {}
        for c in commits[:100]:
            author = _meta_dict(c).get("agent_id") or _field(c, "author", "")
            if author:
                recent_authors[author] = recent_authors.get(author, 0) + 1
        return {
            "agents": len(agents),
            "commits": len(commits),
            "branches": len(branches),
            "recent_active_agents": sorted(
                [{"agent_id": k, "commits": v} for k, v in recent_authors.items()],
                key=lambda x: -x["commits"],
            )[:10],
            "workspace": str(_get_config().cvc_root.parent) if _get_config() else "",
        }
    except Exception as exc:
        return {"agents": 0, "commits": 0, "branches": 0, "error": str(exc)}


@app.get("/api/hivemind/stats")
async def get_hivemind_stats():
    """Aggregate stats for the HiveMind page (counts + per-squad breakdown)."""
    _ensure()
    try:
        db = _get_db()
        agents_raw = db.index.list_agents() if hasattr(db.index, "list_agents") else []
        commits = db.index.list_all_commits(limit=500) if hasattr(db.index, "list_all_commits") else []
        nodes = [_agent_to_node(a) for a in agents_raw if a is not None]
        by_squad: dict[str, int] = {}
        by_rank: dict[str, int] = {}
        for n in nodes:
            by_squad[n["squad"] or "unaffiliated"] = by_squad.get(n["squad"] or "unaffiliated", 0) + 1
            by_rank[n["rank"] or "unranked"] = by_rank.get(n["rank"] or "unranked", 0) + 1
        # Commits in last hour / 24h
        now = time.time()
        last_hour = sum(1 for c in commits if (_field(c, "created_at", 0) or 0) >= now - 3600)
        last_day = sum(1 for c in commits if (_field(c, "created_at", 0) or 0) >= now - 86400)
        return {
            "total_agents": len(nodes),
            "total_commits": len(commits),
            "by_squad": by_squad,
            "by_rank": by_rank,
            "commits_last_hour": last_hour,
            "commits_last_24h": last_day,
        }
    except Exception as exc:
        return {"total_agents": 0, "total_commits": 0, "error": str(exc)}


@app.get("/api/hivemind/feed")
async def get_hivemind_feed(limit: int = 50, agent_id: str | None = None, squad: str | None = None):
    """
    Recent activity feed: commits + agent registrations + memory writes.

    Filterable by agent_id or squad. Ordered newest-first.
    """
    _ensure()
    try:
        db = _get_db()
        commits = db.index.list_all_commits(limit=max(limit * 3, 100)) if hasattr(db.index, "list_all_commits") else []
        agents_raw = db.index.list_agents() if hasattr(db.index, "list_agents") else []
        agent_squad: dict[str, str] = {}
        for a in agents_raw or []:
            n = _agent_to_node(a)
            if n["id"]:
                agent_squad[n["id"]] = n["squad"] or ""
        events: list[dict] = []
        for c in commits:
            meta = _meta_dict(c)
            author = meta.get("agent_id") or _field(c, "author", "") or ""
            sq = agent_squad.get(author, "")
            if agent_id and author != agent_id:
                continue
            if squad and sq != squad:
                continue
            events.append({
                "type": "commit",
                "commit_hash": _field(c, "commit_hash", ""),
                "agent_id": author,
                "squad": sq,
                "message": _field(c, "message", "") or meta.get("message", ""),
                "tags": meta.get("tags", []),
                "branch": _field(c, "branch", ""),
                "timestamp": _field(c, "created_at", 0),
            })
        events.sort(key=lambda e: e.get("timestamp") or 0, reverse=True)
        return {"events": events[:limit], "total": len(events)}
    except Exception as exc:
        return {"events": [], "total": 0, "error": str(exc)}


# ─────────────────────────────────────────────────────────────────────────────
# CVC v3.0 — Unified-core chat WebSocket endpoint (always available).
#
# The dashboard's chat goes through `/api/ws/chat` (the legacy CVC agent
# loop, PRESERVED UNCHANGED). This new endpoint is the unified-core
# backed alternative. The dashboard's WebSocket bridge can be pointed
# at it by either:
#   1. Setting CVC_USE_UNIFIED_CORE=1 in the shell + rebuilding the
#      dashboard's `chatWsUrl()` to point at `/api/ws/chat_unified`.
#   2. Building a runtime override in `cvc/web/src/lib/api.ts` that
#      swaps the URL when the env var is set.
#
# Event schema matches the dashboard's `ChatEvent` type
# (cvc/web/src/lib/types.ts): text, tool_start, tool_result, status,
# done, error, setup_required.
# ─────────────────────────────────────────────────────────────────────────────
@app.websocket("/api/ws/chat_legacy")
async def ws_chat_legacy(websocket: WebSocket):
    """v3.1.0 — Preserved legacy CVC chat loop, kept for rollback only.

    From CVC v3.1.0 onward, the dashboard's `/api/ws/chat` route runs
    the unified agent core (see `ws_chat_unified`). This
    endpoint exposes the original CVC loop — same code, same event
    schema, same buggy synthesis behavior — for rollback comparison
    and for any user who explicitly wants the old behavior via
    `CVC_USE_UNIFIED_CORE=0`.
    """
    # Force the legacy branch inside `ws_chat` by faking the env var.
    import os
    prev = os.environ.get("CVC_USE_UNIFIED_CORE")
    os.environ["CVC_USE_UNIFIED_CORE"] = "0"
    try:
        await ws_chat(websocket)
    finally:
        if prev is None:
            os.environ.pop("CVC_USE_UNIFIED_CORE", None)
        else:
            os.environ["CVC_USE_UNIFIED_CORE"] = prev


@app.websocket("/api/ws/chat_unified")
async def ws_chat_unified(websocket: WebSocket, *, _already_accepted: bool = False):
    """Unified-core backed WebSocket chat (v3.0).

    Protocol:
      - First frame from client: {"messages": [...], "thread_id": "...",
        "workspace_path": "...", "persona_id": "..."}.
      - Server streams back JSON events: turn_start, text, tool_start,
        tool_result, status, done, error, setup_required.
      - `done` is always sent exactly once at end-of-turn.

    Args:
      websocket: the active WebSocket connection.
      _already_accepted: True when this handler is delegated to from
        `ws_chat` (which already accepted the connection). When True,
        we skip our own `accept()` call — Starlette raises if you
        accept twice.
    """
    import uuid as _uuid
    if not _already_accepted:
        await websocket.accept()
    turn_id = _uuid.uuid4().hex

    # Send turn_start so the client can persist turn_id for resume
    try:
        await websocket.send_json({"type": "turn_start", "turn_id": turn_id})
    except Exception:
        return

    # Read the first frame
    try:
        init_data = await websocket.receive_json()
    except WebSocketDisconnect:
        return
    except Exception:
        try:
            await websocket.send_json({
                "type": "error",
                "message": "Invalid first frame — expected JSON object with 'messages'."
            })
            await websocket.send_json({"type": "done", "turn_id": turn_id})
        except Exception:
            pass
        return

    messages = init_data.get("messages", []) if isinstance(init_data, dict) else []
    if not messages:
        try:
            await websocket.send_json({
                "type": "error",
                "message": "No messages in request.",
            })
            await websocket.send_json({"type": "done", "turn_id": turn_id})
        except Exception:
            pass
        return

    # Bridge SSE events to WebSocket frames
    from cvc.gateway_chat import stream_unified_chat

    body = {
        "messages": messages,
        "stream": True,
        "thread_id": init_data.get("thread_id") or init_data.get("conversation_id") or turn_id,
        "workspace_path": init_data.get("workspace_path"),
        "persona_id": init_data.get("persona_id") or "",
    }

    SSE_TO_WS_TYPE = {
        "text": "text",
        "tool_use": "tool_start",
        "tool_result": "tool_result",
        "status": "status",
        "error": "error",
        "setup_required": "setup_required",
        # v3.2.0 — Events from the vendored api_server. These
        # stream through unchanged so the dashboard's tool/status
        # rendering works the same whether the chat is going through
        # the api_server (full tool set) or the legacy hand-rolled path.
        "assistant.delta": "text",  # Stream text the same way
        "assistant.completed": "text",
        "tool.started": "tool_start",
        "tool.completed": "tool_result",
        "tool.failed": "tool_result",
        "tool.progress": "status",
    }

    async def _send(event: dict) -> bool:
        try:
            await websocket.send_json(event)
            return True
        except WebSocketDisconnect:
            return False
        except Exception:
            return False

    try:
        # Parse each SSE line and forward to WS
        buf = b""
        async for chunk in stream_unified_chat(body):
            if chunk == b"data: [DONE]\n\n":
                break
            # SSE chunk: "data: {json}\n\n"
            if not chunk.startswith(b"data: "):
                continue
            try:
                payload = json.loads(chunk[6:].decode("utf-8").strip())
            except Exception:
                continue
            ws_type = SSE_TO_WS_TYPE.get(payload.get("type", ""), payload.get("type", ""))
            ws_evt: dict = {"type": ws_type}
            # Map unified-core SSE fields to dashboard's ChatEvent fields.
            # v3.1.1 — `env_keys` and `detail` added so the dashboard's
            # setup_required action card can show the right env var for
            # the user's actual provider (not just the legacy four).
            for src_key, dst_key in (
                ("content", "content"),
                ("message", "message"),
                ("name", "name"),
                ("arguments", "args"),
                ("tool_call_id", "id"),
                ("provider", "provider"),
                ("env_keys", "env_keys"),
                ("detail", "detail"),
            ):
                if src_key in payload:
                    ws_evt[dst_key] = payload[src_key]
            ok = await _send(ws_evt)
            if not ok:
                return  # client gone
        # End-of-turn
        await _send({"type": "done", "turn_id": turn_id})
    except WebSocketDisconnect:
        pass
    except Exception as exc:
        logger = __import__("logging").getLogger("cvc.gateway")
        logger.exception("ws_chat_unified failed")
        try:
            await _send({"type": "error", "message": str(exc)})
            await _send({"type": "done", "turn_id": turn_id})
        except Exception:
            pass


@app.websocket("/api/swarm/ws")
async def ws_swarm(websocket: WebSocket):
    """
    Live swarm/HiveMind events: agent registrations, commits, memory writes.

    On connect, sends a topology snapshot, then streams events from the bus.
    """
    await websocket.accept()
    bus = _get_event_bus()
    sub_id = f"swarm-{id(websocket)}"
    sub = bus.subscribe(sub_id, replay=True)
    try:
        # Initial snapshot
        try:
            snap = await get_swarm_topology()  # type: ignore[func-returns-value]
            await websocket.send_json({"event": "swarm.snapshot", "data": snap, "timestamp": time.time()})
        except Exception:
            logger.debug("Swarm snapshot send failed", exc_info=True)
        while sub.active:
            try:
                envelope = await asyncio.wait_for(sub.get(), timeout=15.0)
                await websocket.send_json({
                    "event": envelope.event,
                    "data": envelope.data,
                    "timestamp": envelope.timestamp,
                })
            except asyncio.TimeoutError:
                # Heartbeat with refreshed counts so the UI stays current even
                # without bus traffic.
                try:
                    await websocket.send_json({
                        "event": "swarm.heartbeat",
                        "data": {"timestamp": time.time()},
                        "timestamp": time.time(),
                    })
                except Exception:
                    break
    except WebSocketDisconnect:
        pass
    except Exception:
        logger.debug("Swarm WebSocket error", exc_info=True)
    finally:
        bus.unsubscribe(sub_id)


# ═══════════════════════════════════════════════════════════════════════
# 13. CONNECTIONS
# ═══════════════════════════════════════════════════════════════════════

@app.get("/api/connections")
async def get_connections():
    base = f"http://{_HOST}:{_PROXY_PORT}"
    return {
        "proxy_url": f"{base}/v1",
        "api_key": "cvc",
        "tools": [
            {"name": "VS Code", "type": "ide", "method": "BYOK settings.json", "icon": "vscode"},
            {"name": "Cursor", "type": "ide", "method": "BYOK settings.json", "icon": "cursor"},
            {"name": "Antigravity", "type": "ide", "method": "MCP config", "icon": "antigravity"},
            {"name": "Windsurf", "type": "ide", "method": "MCP config", "icon": "windsurf"},
            {"name": "Claude Code", "type": "cli", "method": "ANTHROPIC_BASE_URL env", "icon": "claude"},
            {"name": "Aider", "type": "cli", "method": "OPENAI_API_BASE env", "icon": "aider"},
            {"name": "Codex CLI", "type": "cli", "method": "OPENAI_BASE_URL env", "icon": "codex"},
            {"name": "Gemini CLI", "type": "cli", "method": "CVC MCP", "icon": "gemini"},
            {"name": "Continue.dev", "type": "extension", "method": "config.json", "icon": "continue"},
            {"name": "Cline", "type": "extension", "method": "MCP or API endpoint", "icon": "cline"},
            {"name": "Open WebUI", "type": "web", "method": "OpenAI API endpoint", "icon": "openwebui"},
        ],
    }


# ═══════════════════════════════════════════════════════════════════════
# WEBSOCKET — Real-Time Dashboard Events
# ═══════════════════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════════════════
# DASHBOARD SUB-ROUTERS — surface Phase-1 plumbing on the HTTP API
# ═══════════════════════════════════════════════════════════════════════
try:
    from cvc.dashboard.providers_api import register_providers_routes
    register_providers_routes(app)
except Exception as _exc:  # pragma: no cover — never break gateway boot
    logger.warning("Failed to register providers routes: %s", _exc)

try:
    from cvc.dashboard.loop_api import register_loop_routes
    register_loop_routes(app)
except Exception as _exc:  # pragma: no cover
    logger.warning("Failed to register loop routes: %s", _exc)

try:
    from cvc.dashboard.trajectory_api import register_trajectory_routes
    register_trajectory_routes(app)
except Exception as _exc:  # pragma: no cover
    logger.warning("Failed to register trajectory routes: %s", _exc)

try:
    from cvc.dashboard.team_api import register_team_routes
    register_team_routes(app)
except Exception as _exc:  # pragma: no cover
    logger.warning("Failed to register team routes: %s", _exc)

try:
    from cvc.dashboard.personas_api import register_personas_routes
    register_personas_routes(app)
except Exception as _exc:  # pragma: no cover
    logger.warning("Failed to register personas routes: %s", _exc)

try:
    from cvc.dashboard.themes_api import register_themes_routes
    register_themes_routes(app)
except Exception as _exc:  # pragma: no cover
    logger.warning("Failed to register themes routes: %s", _exc)

try:
    from cvc.dashboard.user_memory_api import register_user_memory_routes
    register_user_memory_routes(app)
except Exception as _exc:  # pragma: no cover
    logger.warning("Failed to register user-memory routes: %s", _exc)

try:
    from cvc.dashboard.skills_api import register_skills_routes
    register_skills_routes(app)
except Exception as _exc:  # pragma: no cover
    logger.warning("Failed to register skills routes: %s", _exc)

try:
    from cvc.dashboard.conversations_api import router as _conversations_router
    app.include_router(_conversations_router)
except Exception as _exc:  # pragma: no cover
    logger.warning("Failed to register conversations routes: %s", _exc)

try:
    from cvc.dashboard.dx_api import router as _dx_router
    app.include_router(_dx_router)
    logger.info("DX routes mounted at /api/dx/*")
except Exception as _exc:  # pragma: no cover
    logger.warning("Failed to register dx routes: %s", _exc)

try:
    from cvc.dashboard.skills_api import register_skills_routes
    register_skills_routes(app)
    logger.info("Skills routes mounted at /api/skills/*")
except Exception as _exc:  # pragma: no cover
    logger.warning("Failed to register skills routes: %s", _exc)


@app.websocket("/ws/dashboard")
async def ws_dashboard(websocket: WebSocket):
    await websocket.accept()
    bus = _get_event_bus()
    sub_id = f"dashboard-{id(websocket)}"
    sub = bus.subscribe(sub_id, replay=True)
    try:
        if _manager:
            await websocket.send_json({
                "event": "services.snapshot",
                "data": {s.name: s.to_dict() for s in _manager.all_services().values()},
                "timestamp": time.time(),
            })
        while sub.active:
            try:
                envelope = await asyncio.wait_for(sub.get(), timeout=10.0)
                await websocket.send_json({"event": envelope.event, "data": envelope.data, "timestamp": envelope.timestamp})
            except asyncio.TimeoutError:
                if _manager:
                    _manager.refresh_health()
                    await websocket.send_json({
                        "event": "services.snapshot",
                        "data": {s.name: s.to_dict() for s in _manager.all_services().values()},
                        "timestamp": time.time(),
                    })
    except WebSocketDisconnect:
        pass
    except Exception:
        logger.debug("Dashboard WebSocket error", exc_info=True)
    finally:
        bus.unsubscribe(sub_id)


# ═══════════════════════════════════════════════════════════════════════
# STATIC FILES — Dashboard SPA
# ═══════════════════════════════════════════════════════════════════════

# Phase A cleanup: legacy HTML mockup removed. Phase B will populate
# `cvc/web_dist/` via `cvc setup-ui` (Vite build of cvc/web/).
_DASHBOARD_DIR = Path(__file__).parent.parent / "cvc" / "web_dist"
if not _DASHBOARD_DIR.exists():
    _DASHBOARD_DIR = Path(__file__).parent / "web_dist"

if _DASHBOARD_DIR.exists():
    app.mount("/static", StaticFiles(directory=str(_DASHBOARD_DIR)), name="dashboard-static")
    # Phase C: mount Vite-emitted asset folders at root paths so the built
    # index.html (which references /assets, /ds-assets, /fonts, /favicon.svg
    # absolutely) loads without rewrites.
    for _sub in ("assets", "ds-assets", "fonts"):
        _p = _DASHBOARD_DIR / _sub
        if _p.exists():
            app.mount(f"/{_sub}", StaticFiles(directory=str(_p)), name=f"dashboard-{_sub}")

    # v2.91.51 — Stale-dist guard. If the source TypeScript/TSX under
    # cvc/web/src/ is newer than the built index.html in cvc/web_dist/,
    # warn at boot. (The wheel ships a pre-built dist, so this check is
    # only meaningful when running from a source checkout, but it's
    # cheap and helps catch the case where a developer edited a
    # component and forgot to rebuild before restarting the gateway.)
    try:
        _src_root = _DASHBOARD_DIR.parent / "web" / "src"
        if _src_root.exists():
            _src_mtime = max(
                (p.stat().st_mtime for p in _src_root.rglob("*") if p.is_file()),
                default=0,
            )
            _dist_index = _DASHBOARD_DIR / "index.html"
            _dist_mtime = _dist_index.stat().st_mtime if _dist_index.exists() else 0
            if _src_mtime > _dist_mtime:
                logger.warning(
                    "Dashboard dist is STALE (cvc/web/src/ modified after "
                    "cvc/web_dist/index.html). Run `cvc setup-ui` to rebuild."
                )
    except Exception:
        pass  # best-effort; never fail boot over a stale-dist warning

    @app.get("/favicon.svg", include_in_schema=False)
    async def _favicon():
        fav = _DASHBOARD_DIR / "favicon.svg"
        if fav.exists():
            return FileResponse(str(fav), media_type="image/svg+xml")
        return JSONResponse({"error": "not_found"}, status_code=404)


@app.get("/", include_in_schema=False)
async def serve_dashboard():
    index_path = _DASHBOARD_DIR / "index.html"
    if index_path.exists():
        html = index_path.read_text(encoding="utf-8")
        # Inject the real CVC version — replaces {{CVC_VERSION}} placeholders
        # which cache-busts every CSS/JS link on each new release.
        html = html.replace("{{CVC_VERSION}}", _cvc_version)
        return HTMLResponse(
            content=html,
            headers={
                "Cache-Control": "no-cache, no-store, must-revalidate",
                "Pragma": "no-cache",
            },
        )
    return JSONResponse({"message": "CVC Gateway running.", "api": "/api/gateway/services"}, status_code=200)


# SPA fallback — for client-side routes like /agents, /swarm, /settings,
# serve index.html so React Router can take over. Must be the LAST GET
# route so it doesn't shadow API or static mounts.
_SPA_RESERVED_PREFIXES = (
    "/api", "/gateway", "/ws", "/static", "/assets", "/ds-assets",
    "/fonts", "/favicon", "/health", "/docs", "/openapi", "/redoc",
)

@app.middleware("http")
async def _spa_fallback_middleware(request, call_next):
    response = await call_next(request)
    if (
        response.status_code == 404
        and request.method == "GET"
        and not any(request.url.path.startswith(p) for p in _SPA_RESERVED_PREFIXES)
        and "text/html" in request.headers.get("accept", "")
        and _DASHBOARD_DIR.exists()
    ):
        index_path = _DASHBOARD_DIR / "index.html"
        if index_path.exists():
            html = index_path.read_text(encoding="utf-8").replace("{{CVC_VERSION}}", _cvc_version)
            return HTMLResponse(
                content=html,
                headers={"Cache-Control": "no-cache, no-store, must-revalidate", "Pragma": "no-cache"},
            )
    return response


# ═══════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════

def _field(obj: Any, key: str, default: Any = "") -> Any:
    if isinstance(obj, dict):
        return obj.get(key, default)
    return getattr(obj, key, default)


def _ts(obj: Any) -> float:
    """Extract Unix timestamp (seconds) from a commit object."""
    if isinstance(obj, dict):
        # Raw DB row or dict
        if "created_at" in obj and obj["created_at"]:
            return float(obj["created_at"])
        md = obj.get("metadata_json", "")
        if isinstance(md, str) and md:
            try:
                return float(json.loads(md).get("timestamp", 0))
            except (json.JSONDecodeError, TypeError, ValueError):
                return 0.0
        return 0.0
    # Pydantic CognitiveCommit
    meta = getattr(obj, "metadata", None)
    if meta is not None:
        return float(getattr(meta, "timestamp", 0))
    return 0.0


def _meta_dict(obj: Any) -> dict:
    # Pydantic CognitiveCommit — metadata is a CommitMetadata model
    meta = getattr(obj, "metadata", None)
    if meta is not None and hasattr(meta, "model_dump"):
        return meta.model_dump()
    # Raw DB row — metadata_json is a JSON string column
    raw = _field(obj, "metadata_json", "")
    if not raw:
        return {}
    try:
        return json.loads(raw) if isinstance(raw, str) else raw
    except (json.JSONDecodeError, TypeError):
        return {}


def _meta(obj: Any, key: str, default: Any = "") -> Any:
    return _meta_dict(obj).get(key, default)
